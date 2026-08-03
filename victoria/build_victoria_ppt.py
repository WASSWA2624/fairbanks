#!/usr/bin/env python3
"""
Victoria University partnership proposal deck.
Content from victoria/raw.md only.
FairBanks branding from assets. Clear spacing — no overlaps.
No personal names, phone numbers, or personal contact details.
"""

from pathlib import Path

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
ASSETS = REPO / "assets"
LOGO = ASSETS / "fairbanks_logo.jpeg"
OUT = Path(__file__).resolve().parent / "documents"
OUT.mkdir(parents=True, exist_ok=True)
OUT_PPT = OUT / "victoria_ppt.pptx"

# FairBanks brand palette
NAVY = "0A1F2E"
TEAL = "0D6E6E"
GREEN = "2D7A55"
ORANGE = "C45C26"
GOLD = "D99A2B"
CREAM = "F7F5F0"
PALE_TEAL = "E8F3F2"
WHITE = "FFFFFF"
SLATE = "1E2F38"
MUTED = "52636C"
LINE = "CED9D8"
LIGHT = "D4E8DC"

ORG = "FairBanks Medical Centre"
SLOGAN = "Your health, our mission."
TAGLINE = "Health for All"
FOOTER = f"{ORG}  ·  {SLOGAN}"

PHOTOS = {
    "cover": ASSETS / "facility_exterior_branded_entrance_01.jpeg",
    "facility": ASSETS / "facility_exterior_entrance_01.jpg",
    "facility_wide": ASSETS / "facility_exterior_branded_entrance_02.jpeg",
    "outreach": ASSETS / "outreach_bp_screening.jpeg",
    "camp": ASSETS / "outreach_medical_camp_01.jpg",
    "mission": ASSETS / "reception_mission_wall.jpeg",
    "audience": ASSETS / "outreach_audience_full_group_01.jpg",
}

# Layout grid (inches) — generous clear zones
ML, MR = 0.70, 0.65
SW, SH = 13.333, 7.5
CW = SW - ML - MR
TOP_BAR = 0.06
LOGO_H = 0.36
LOGO_W = LOGO_H * (269 / 101)
LOGO_Y = 0.40
HEADER_BOTTOM_WITH_SUB = 1.72
HEADER_BOTTOM = 1.42
FOOTER_Y = 7.10
CONTENT_BOTTOM = 6.65


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run(run, text, size, color, bold=False, font="Calibri"):
    # Presentation bump: small type grows more; large titles a little
    if size >= 32:
        size = size + 2
    elif size >= 20:
        size = size + 3
    else:
        size = size + 4
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_fade_transition(slide):
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    transition = etree.Element(qn("p:transition"), {"spd": "med"})
    etree.SubElement(transition, qn("p:fade"))
    cSld = sld.find(qn("p:cSld"))
    if cSld is not None:
        cSld.addnext(transition)
    else:
        sld.insert(0, transition)


def build():
    missing = [p for p in [LOGO, *PHOTOS.values()] if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing brand assets:\n" + "\n".join(str(p) for p in missing))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
        else:
            shape.line.fill.background()
        if rounded:
            try:
                shape.adjustments[0] = 0.1
            except Exception:
                pass
        # Kill auto shadows that make cards look like they overlap
        try:
            spPr = shape._element.spPr
            for child in list(spPr):
                if "effectLst" in child.tag or "effectDag" in child.tag:
                    spPr.remove(child)
            from lxml import etree as ET
            effect = ET.SubElement(spPr, qn("a:effectLst"))
        except Exception:
            pass
        return shape

    def textbox(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=18,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        font="Calibri",
        line_spacing=None,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Emu(20000)
        tf.margin_right = Emu(20000)
        tf.margin_top = Emu(8000)
        tf.margin_bottom = Emu(8000)
        tf.vertical_anchor = valign
        lines = value.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            p.space_before = Pt(0)
            p.space_after = Pt(6) if len(lines) > 1 else Pt(0)
            r = p.add_run()
            set_run(r, line, size, color, bold, font)
        return box

    def bullets(slide, items, x, y, w, h, size=16, color=SLATE, space=9):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Emu(25000)
        tf.margin_right = Emu(25000)
        tf.margin_top = Emu(8000)
        tf.margin_bottom = Emu(8000)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(space)
            r = p.add_run()
            set_run(r, "•  " + item, size, color)
        return box

    def add_logo(slide, x, y, height=LOGO_H, plate=True):
        """Logo on a white plate fully below the top brand bar."""
        if plate:
            pad = 0.08
            # Keep plate below TOP_BAR
            py = max(y - pad, TOP_BAR + 0.06)
            rect(slide, x - pad, py, LOGO_W + 2 * pad, height + (y - py) + pad, WHITE, LINE, rounded=True)
        return slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), height=Inches(height))

    def crop_photo(slide, path, x, y, w, h):
        with PILImage.open(path) as im:
            iw, ih = im.size
        pic = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
        )
        frame_ratio, image_ratio = w / h, iw / ih
        if image_ratio > frame_ratio:
            amount = (1 - frame_ratio / image_ratio) / 2
            pic.crop_left = pic.crop_right = amount
        else:
            amount = (1 - image_ratio / frame_ratio) / 2
            pic.crop_top = pic.crop_bottom = amount
        return pic

    def new_slide(bg=CREAM):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, SW, SH, bg)
        add_fade_transition(s)
        return s

    def header(s, kicker, title, subtitle=""):
        """Brand chrome + logo. Returns safe content-start Y."""
        rect(s, 0, 0, SW, TOP_BAR, TEAL)
        rect(s, 0, TOP_BAR, 0.08, SH - TOP_BAR, ORANGE)

        logo_x = SW - MR - LOGO_W
        add_logo(s, logo_x, LOGO_Y, LOGO_H, plate=True)

        title_w = logo_x - ML - 0.55
        if kicker.strip():
            textbox(s, kicker.upper(), ML, 0.28, title_w, 0.28, 15, ORANGE, True)
            textbox(s, title, ML, 0.58, title_w, 0.55, 26, NAVY, True)
            title_bottom = 1.18
        else:
            textbox(s, title, ML, 0.32, title_w, 0.58, 28, NAVY, True)
            title_bottom = 1.05

        if subtitle:
            textbox(s, subtitle, ML, title_bottom + 0.08, CW, 0.36, 16, MUTED)
            return HEADER_BOTTOM_WITH_SUB
        return max(HEADER_BOTTOM, title_bottom + 0.22)

    def footer(s, number, total):
        rect(s, ML, FOOTER_Y, CW, 0.012, LINE)
        textbox(s, FOOTER, ML, FOOTER_Y + 0.08, 9.5, 0.28, 12, MUTED)
        textbox(
            s,
            f"{number}  /  {total}",
            SW - MR - 1.4,
            FOOTER_Y + 0.08,
            1.4,
            0.28,
            12,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )

    def card(s, x, y, w, h, title, items, accent=TEAL, title_size=16, item_size=15):
        """Card with inset accent bar — text never sits on the bar."""
        rect(s, x, y, w, h, WHITE, LINE, rounded=True)
        rect(s, x + 0.12, y + 0.18, 0.08, h - 0.36, accent)
        textbox(s, title, x + 0.35, y + 0.20, w - 0.55, 0.42, title_size, NAVY, True)
        bullets(s, items, x + 0.30, y + 0.70, w - 0.50, h - 0.95, item_size, SLATE, space=9)

    slides_meta = []

    # ------------------------------------------------------------------
    # 1 Cover — large, clear type for live presentation
    # ------------------------------------------------------------------
    s = new_slide(NAVY)
    panel_w = 7.90
    crop_photo(s, PHOTOS["cover"], 0, 0, SW, SH)
    rect(s, 0, 0, panel_w, SH, NAVY)
    rect(s, panel_w, 0, 0.12, SH, GOLD)

    add_logo(s, 0.55, 0.25, 0.52, plate=True)

    textbox(s, "STRATEGIC INSTITUTIONAL PROPOSAL", 0.55, 1.00, 7.1, 0.38, 16, GOLD, True)
    textbox(s, ORG, 0.55, 1.45, 7.1, 0.48, 24, TEAL, True)

    # Two-line hero title — wide bands so descenders never collide
    textbox(s, "University–Community Health", 0.55, 2.25, 7.1, 0.78, 34, WHITE, True)
    textbox(s, "Partnership Proposal", 0.55, 3.20, 7.1, 0.78, 34, GOLD, True)

    rect(s, 0.55, 4.15, 3.20, 0.10, GOLD)

    textbox(s, "Building the FairBanks Community Health Improvement Programme", 0.55, 4.40, 7.1, 0.48, 18, LIGHT)
    textbox(s, "(FCHIP) with Victoria University", 0.55, 5.05, 7.1, 0.45, 22, GOLD, True)

    textbox(s, SLOGAN, 0.55, 5.65, 7.1, 0.42, 26, GOLD, True)

    rect(s, 0, 6.40, panel_w, 1.10, TEAL)
    textbox(
        s,
        "Prepared for Victoria University leadership",
        0.55,
        6.55,
        7.1,
        0.40,
        20,
        WHITE,
        True,
    )
    textbox(s, "Kyebando–Kisalosalo, Kampala  ·  Institutional briefing", 0.55, 7.00, 7.1, 0.35, 16, LIGHT)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 2 Slide contents
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "Slide Contents")
    agenda = [
        ("01", "Executive summary"),
        ("02", "About FairBanks Medical Centre"),
        ("03", "What FCHIP is"),
        ("04", "How we work together"),
        ("05", "Why Victoria University"),
        ("06", "Partnership network and outcomes"),
        ("07", "Roadmap and long-term vision"),
        ("08", "Invitation to partner"),
    ]
    gap_x, gap_y = 0.35, 0.28
    card_w = (CW - gap_x) / 2
    avail = CONTENT_BOTTOM - (y0 + 0.15)
    card_h = (avail - 3 * gap_y) / 4
    for i, (num, label) in enumerate(agenda):
        col, row = i % 2, i // 2
        x = ML + col * (card_w + gap_x)
        y = y0 + 0.15 + row * (card_h + gap_y)
        rect(s, x, y, card_w, card_h, WHITE, LINE, rounded=True)
        rect(s, x + 0.18, y + 0.22, 0.08, card_h - 0.44, ORANGE if col == 0 else TEAL)
        textbox(s, num, x + 0.40, y + (card_h - 0.40) / 2, 0.70, 0.40, 20, ORANGE, True, valign=MSO_ANCHOR.MIDDLE)
        textbox(
            s,
            label,
            x + 1.20,
            y + (card_h - 0.40) / 2,
            card_w - 1.45,
            0.40,
            17,
            NAVY,
            True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 3 Executive summary
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Executive summary", "A strategic partnership for health, learning, and impact")
    textbox(
        s,
        "FairBanks Medical Centre proposes a strategic partnership with Victoria University "
        "to integrate healthcare service delivery, academic excellence, research, innovation, "
        "and community engagement.",
        ML,
        y0,
        CW,
        0.75,
        17,
        SLATE,
    )
    points = [
        (
            "FCHIP at the centre",
            "The FairBanks Community Health Improvement Programme is the flagship platform for "
            "preventive care, outreach, research, skills development, and partnerships.",
            TEAL,
        ),
        (
            "Shared strengths",
            "Victoria University brings education, innovation, and research. FairBanks brings "
            "a practical healthcare platform and community presence.",
            ORANGE,
        ),
        (
            "Shared goal",
            "Build a replicable University–Community Health Partnership Model that advances "
            "healthcare, education, innovation, and sustainable development in Uganda.",
            GREEN,
        ),
    ]
    start = y0 + 0.90
    gap = 0.22
    row_h = (CONTENT_BOTTOM - start - 2 * gap) / 3
    for i, (t, body, accent) in enumerate(points):
        y = start + i * (row_h + gap)
        rect(s, ML, y, CW, row_h, WHITE, LINE, rounded=True)
        rect(s, ML + 0.18, y + 0.20, 0.08, row_h - 0.40, accent)
        textbox(s, t, ML + 0.45, y + 0.22, CW - 0.75, 0.35, 18, NAVY, True)
        textbox(s, body, ML + 0.45, y + 0.62, CW - 0.75, row_h - 0.80, 15, MUTED)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 4 About FairBanks
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Who we are", "About FairBanks Medical Centre")
    left_w = 6.90
    gap = 0.30
    right_w = CW - left_w - gap
    box_h = CONTENT_BOTTOM - y0

    rect(s, ML, y0, left_w, box_h, WHITE, LINE, rounded=True)
    textbox(s, "Community-based primary healthcare", ML + 0.35, y0 + 0.30, left_w - 0.70, 0.35, 16, NAVY, True)
    textbox(s, SLOGAN, ML + 0.35, y0 + 0.70, left_w - 0.70, 0.28, 13, ORANGE, True)
    bullets(
        s,
        [
            "Located in Kyebando–Kisalosalo, Kampala",
            "Accessible, affordable, quality healthcare",
            "Strong focus on preventive health and community well-being",
            "Committed to social determinants of health through partnerships",
            "Active in innovation, research, and community engagement",
        ],
        ML + 0.35,
        y0 + 1.15,
        left_w - 0.70,
        box_h - 1.45,
        14,
        SLATE,
        space=12,
    )

    rx = ML + left_w + gap
    photo_h = box_h * 0.58
    crop_photo(s, PHOTOS["facility"], rx, y0, right_w, photo_h)
    cap_y = y0 + photo_h + 0.18
    cap_h = CONTENT_BOTTOM - cap_y
    rect(s, rx, cap_y, right_w, cap_h, TEAL, rounded=True)
    textbox(s, ORG, rx + 0.25, cap_y + 0.25, right_w - 0.50, 0.32, 13, GOLD, True)
    textbox(
        s,
        "A trusted community health home — care close to where families live.",
        rx + 0.25,
        cap_y + 0.65,
        right_w - 0.50,
        cap_h - 0.85,
        13,
        WHITE,
    )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 5 What is FCHIP
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "What is FCHIP?")

    # Clear definition card — taller for larger presentation type
    def_h = 2.55
    rect(s, ML, y0, CW, def_h, TEAL, rounded=True)
    textbox(s, "FCHIP stands for", ML + 0.40, y0 + 0.22, CW - 0.80, 0.35, 16, GOLD, True)
    textbox(
        s,
        "FairBanks Community Health Improvement Programme",
        ML + 0.40,
        y0 + 0.55,
        CW - 0.80,
        0.48,
        22,
        WHITE,
        True,
    )
    textbox(
        s,
        "It is FairBanks Medical Centre's flagship community programme — run through the "
        "FairBanks Social Enterprise. FCHIP brings healthcare closer to families by combining "
        "preventive care, outreach, education, research, skills development, and partnerships "
        "with universities, government, and community partners.",
        ML + 0.40,
        y0 + 1.15,
        CW - 0.80,
        1.20,
        15,
        LIGHT,
    )

    textbox(s, "What FCHIP focuses on", ML, y0 + def_h + 0.18, CW, 0.35, 17, NAVY, True)

    focus = [
        "Preventive healthcare",
        "Community outreach",
        "Health education",
        "Early disease detection",
        "Maternal and child health",
        "NCD prevention",
        "Elderly care & healthy ageing",
        "Community-based research",
        "Digital health innovation",
        "Skills & student learning",
    ]
    cols, rows = 5, 2
    gap_x, gap_y = 0.22, 0.18
    grid_top = y0 + def_h + 0.55
    avail_h = CONTENT_BOTTOM - grid_top
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (avail_h - (rows - 1) * gap_y) / rows
    accents = [TEAL, ORANGE, GREEN, GOLD, TEAL]
    for i, item in enumerate(focus):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = grid_top + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        rect(s, x + 0.14, y + 0.14, cell_w - 0.28, 0.07, accents[col])
        textbox(
            s,
            item,
            x + 0.10,
            y + 0.28,
            cell_w - 0.20,
            cell_h - 0.40,
            13,
            NAVY,
            True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 6 How we work together (after What, before Why)
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "How we work together")
    textbox(
        s,
        "Five strategic areas of collaboration between Victoria University and FairBanks Medical Centre.",
        ML,
        y0,
        CW,
        0.40,
        15,
        MUTED,
    )
    areas = [
        ("A", "Student Experiential Learning", "Structured placements across health and related fields", TEAL),
        ("B", "Research and Innovation", "Joint studies, publications, and evidence for policy", ORANGE),
        ("C", "Community Engagement", "Camps, schools, screening, and volunteer programmes", GREEN),
        ("D", "Digital Health & Innovation", "AI, EMRs, telemedicine, and predictive analytics", GOLD),
        ("E", "Resource Mobilisation", "Shared pursuit of grants and development partnerships", TEAL),
    ]
    badge = 0.52
    inset = 0.24
    gap_y = 0.22
    areas_top = y0 + 0.50
    row_h = (CONTENT_BOTTOM - areas_top - (len(areas) - 1) * gap_y) / len(areas)
    for i, (letter, title, desc, accent) in enumerate(areas):
        y = areas_top + i * (row_h + gap_y)
        rect(s, ML, y, CW, row_h, WHITE, LINE, rounded=True)
        by = y + (row_h - badge) / 2
        rect(s, ML + inset, by, badge, badge, accent, rounded=True)
        textbox(
            s,
            letter,
            ML + inset,
            by,
            badge,
            badge,
            18,
            WHITE,
            True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        tx = ML + inset + badge + 0.28
        tw = CW - (inset + badge + 0.28) - 0.30
        textbox(s, title, tx, y + 0.14, tw, 0.30, 15, NAVY, True)
        textbox(s, desc, tx, y + 0.44, tw, 0.30, 12, MUTED)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 7 Why Victoria (after How)
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "Why Victoria University?")
    textbox(
        s,
        "Victoria University is committed to practical learning, innovation, entrepreneurship, "
        "and community transformation — values that closely align with FairBanks Medical Centre.",
        ML,
        y0,
        CW,
        0.60,
        14,
        SLATE,
    )
    roles = [
        (
            "What the University contributes",
            ["Academic expertise", "Research capacity", "Student engagement", "Innovation culture"],
            TEAL,
        ),
        (
            "What the University gains",
            [
                "Real-world experiential learning",
                "Applied research environment",
                "Community impact pathways",
                "Visible institutional leadership",
            ],
            ORANGE,
        ),
        (
            "Proposed role",
            [
                "Founding academic partner for FCHIP",
                "Co-shaper of a national partnership model",
                "Bridge between campus and community",
                "Partner in grants and innovation",
            ],
            GREEN,
        ),
    ]
    gap = 0.28
    card_w = (CW - 2 * gap) / 3
    card_top = y0 + 0.80
    card_h = CONTENT_BOTTOM - card_top
    for i, (title, items, accent) in enumerate(roles):
        x = ML + i * (card_w + gap)
        card(s, x, card_top, card_w, card_h, title, items, accent=accent)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 8 Student learning
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration A",
        "Student experiential learning",
        "Students gain practical experience while contributing to community healthcare.",
    )
    disciplines = [
        "Nursing",
        "Public Health",
        "Human Resource Management",
        "Business Administration",
        "Information Technology",
        "Social Work",
        "Community Development",
        "Marketing and Communications",
    ]
    cols, rows = 4, 2
    gap_x, gap_y = 0.28, 0.30
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - y0 - (rows - 1) * gap_y) / rows
    for i, d in enumerate(disciplines):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        accent_w = 0.65
        rect(s, x + (cell_w - accent_w) / 2, y + 0.45, accent_w, 0.07, ORANGE)
        textbox(
            s,
            d,
            x + 0.18,
            y + 0.70,
            cell_w - 0.36,
            cell_h - 1.00,
            14,
            NAVY,
            True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 9 Research
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration B",
        "Research and innovation",
        "Joint programmes, publications, conferences, and evidence-based policy development.",
    )
    themes = [
        "Community Health",
        "Health Systems Strengthening",
        "Artificial Intelligence in Healthcare",
        "Digital Health",
        "Occupational Burnout",
        "Maternal and Child Health",
        "Non-Communicable Diseases",
        "Preventive Healthcare",
        "Healthcare Quality Improvement",
    ]
    cols, rows = 3, 3
    gap_x, gap_y = 0.28, 0.24
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - y0 - (rows - 1) * gap_y) / rows
    row_accent = [TEAL, ORANGE, GREEN]
    for i, t in enumerate(themes):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        rect(s, x + 0.18, y + 0.20, 0.08, cell_h - 0.40, row_accent[row])
        textbox(
            s,
            t,
            x + 0.40,
            y + 0.20,
            cell_w - 0.60,
            cell_h - 0.40,
            14,
            NAVY,
            True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 10 Community engagement
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration C",
        "Community engagement through FCHIP",
        "Measurable social impact with meaningful practical learning.",
    )
    photo_h = 1.55
    gap_p = 0.28
    pw = (CW - 2 * gap_p) / 3
    crop_photo(s, PHOTOS["outreach"], ML, y0, pw, photo_h)
    crop_photo(s, PHOTOS["camp"], ML + pw + gap_p, y0, pw, photo_h)
    crop_photo(s, PHOTOS["audience"], ML + 2 * (pw + gap_p), y0, pw, photo_h)

    activities = [
        "Community medical camps",
        "School health programmes",
        "Corporate wellness initiatives",
        "Community disease screening",
        "Health promotion campaigns",
        "Maternal and child health outreach",
        "Elderly care programmes",
        "Student volunteer programmes",
        "Community-based research",
    ]
    cols, rows = 3, 3
    gap_x, gap_y = 0.28, 0.22
    grid_top = y0 + photo_h + 0.32
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - grid_top - (rows - 1) * gap_y) / rows
    for i, a in enumerate(activities):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = grid_top + row * (cell_h + gap_y)
        fill = PALE_TEAL if row % 2 == 0 else WHITE
        rect(s, x, y, cell_w, cell_h, fill, LINE, rounded=True)
        textbox(
            s,
            a,
            x + 0.20,
            y + 0.10,
            cell_w - 0.40,
            cell_h - 0.20,
            12,
            NAVY,
            True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 11 Digital health
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration D",
        "Digital health and innovation",
        "Developing practical healthcare solutions together.",
    )
    digi = [
        ("Artificial Intelligence", "Applications that support smarter clinical and community decisions"),
        ("Electronic Medical Records", "Stronger continuity of care and better clinical information"),
        ("Telemedicine", "Care and advice that reach people beyond the facility walls"),
        ("Community Health Information Systems", "Better visibility of community health activity and needs"),
        ("Predictive analytics", "Earlier insight for planning and prevention"),
        ("Digital patient engagement", "Clearer communication and follow-up with patients and families"),
    ]
    cols, rows = 3, 2
    gap_x, gap_y = 0.28, 0.28
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - y0 - (rows - 1) * gap_y) / rows
    accents = [TEAL, ORANGE, GREEN, GOLD, TEAL, ORANGE]
    for i, (t, d) in enumerate(digi):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        rect(s, x + 0.22, y + 0.22, cell_w - 0.44, 0.08, accents[i])
        textbox(s, t, x + 0.25, y + 0.45, cell_w - 0.50, 0.50, 14, NAVY, True)
        textbox(s, d, x + 0.25, y + 1.05, cell_w - 0.50, cell_h - 1.25, 12, MUTED)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 12 Resource mobilisation
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration E",
        "Resource mobilisation",
        "Collaborative proposals improve competitiveness for local and international funding.",
    )
    funds = [
        "Research grants",
        "Development funding",
        "Innovation grants",
        "Corporate Social Responsibility partnerships",
        "International development partnerships",
        "Capacity-building programmes",
    ]
    cols, rows = 3, 2
    gap_x, gap_y = 0.28, 0.28
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - y0 - (rows - 1) * gap_y) / rows
    for i, f in enumerate(funds):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        textbox(s, f"{i + 1:02}", x + 0.30, y + 0.40, cell_w - 0.60, 0.38, 18, ORANGE, True)
        textbox(s, f, x + 0.30, y + 0.95, cell_w - 0.60, cell_h - 1.20, 14, NAVY, True)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 13 Network
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Ecosystem",
        "Strategic partnerships and collaborative network",
        "Victoria University as the flagship academic and research partner.",
    )
    partners = [
        "Government ministries and public health institutions",
        "Local government leadership and Village Health Teams",
        "Health insurance providers",
        "Specialist healthcare professionals",
        "Schools and educational institutions",
        "Faith-based and community organisations",
        "Corporate organisations",
        "Media partners",
        "National and international development partners",
        "Philanthropic foundations and research institutions",
    ]
    cols, rows = 2, 5
    gap_x, gap_y = 0.28, 0.16
    cell_w = (CW - gap_x) / cols
    cell_h = (CONTENT_BOTTOM - y0 - (rows - 1) * gap_y) / rows
    for i, p in enumerate(partners):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        rect(s, x + 0.16, y + 0.16, 0.08, cell_h - 0.32, TEAL if col == 0 else ORANGE)
        textbox(
            s,
            p,
            x + 0.38,
            y + 0.10,
            cell_w - 0.55,
            cell_h - 0.20,
            12,
            SLATE,
            True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 14 Outcomes
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Impact", "Expected outcomes of the partnership")
    outcomes = [
        "Improve access to quality primary healthcare",
        "Strengthen preventive healthcare programmes",
        "Expand experiential learning opportunities",
        "Increase community-based research",
        "Promote healthcare innovation",
        "Strengthen graduate employability",
        "Enhance institutional visibility",
        "Mobilise sustainable development partnerships",
        "Improve health outcomes in surrounding communities",
    ]
    cols, rows = 3, 3
    gap_x, gap_y = 0.28, 0.22
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - y0 - (rows - 1) * gap_y) / rows
    for i, o in enumerate(outcomes):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, rounded=True)
        textbox(s, f"{i + 1:02}", x + 0.28, y + 0.25, cell_w - 0.56, 0.30, 13, ORANGE, True)
        textbox(s, o, x + 0.28, y + 0.62, cell_w - 0.56, cell_h - 0.80, 13, NAVY, True)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 15 Roadmap
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Implementation", "Proposed roadmap")
    phases = [
        (
            "Phase I",
            "Foundation",
            ["Introductory planning meeting", "Joint Technical Working Group"],
            TEAL,
        ),
        (
            "Phase II",
            "Formalisation",
            ["Sign Memorandum of Understanding", "Identify priority pilot projects"],
            ORANGE,
        ),
        (
            "Phase III",
            "Delivery",
            [
                "Student placements",
                "Joint research and outreach",
                "Innovation and grant development",
            ],
            GREEN,
        ),
        (
            "Phase IV",
            "Learning and scale",
            ["Monitoring and evaluation", "Document lessons", "Expand and replicate"],
            NAVY,
        ),
    ]
    gap = 0.25
    n = len(phases)
    card_w = (CW - (n - 1) * gap) / n
    card_h = CONTENT_BOTTOM - y0
    for i, (phase, label, items, accent) in enumerate(phases):
        x = ML + i * (card_w + gap)
        rect(s, x, y0, card_w, card_h, WHITE, LINE, rounded=True)
        rect(s, x, y0, card_w, 1.00, accent)
        textbox(s, phase, x + 0.20, y0 + 0.18, card_w - 0.40, 0.28, 13, WHITE, True)
        textbox(s, label, x + 0.20, y0 + 0.50, card_w - 0.40, 0.35, 15, WHITE, True)
        bullets(s, items, x + 0.18, y0 + 1.25, card_w - 0.36, card_h - 1.45, 12, SLATE, space=10)
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 16 Vision
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Looking ahead", "Long-term vision")
    hero_h = 2.20
    left_w = CW * 0.58
    gap = 0.28
    right_w = CW - left_w - gap
    rect(s, ML, y0, left_w, hero_h, TEAL, rounded=True)
    textbox(
        s,
        "Establish FCHIP as Uganda's leading University–Community Health Partnership Model — "
        "showing how academia, healthcare providers, government, development partners, and "
        "communities can work together to improve lives.",
        ML + 0.40,
        y0 + 0.35,
        left_w - 0.80,
        hero_h - 0.70,
        15,
        WHITE,
        True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    crop_photo(s, PHOTOS["mission"], ML + left_w + gap, y0, right_w, hero_h)

    half_gap = 0.28
    half = (CW - half_gap) / 2
    pair_y = y0 + hero_h + 0.28
    pair_h = CONTENT_BOTTOM - pair_y
    rect(s, ML, pair_y, half, pair_h, WHITE, LINE, rounded=True)
    textbox(s, "For FairBanks Medical Centre", ML + 0.30, pair_y + 0.30, half - 0.55, 0.32, 13, ORANGE, True)
    textbox(
        s,
        "A living laboratory for education, research, innovation, and community service.",
        ML + 0.30,
        pair_y + 0.75,
        half - 0.55,
        pair_h - 1.00,
        14,
        NAVY,
        True,
    )
    rx = ML + half + half_gap
    rect(s, rx, pair_y, half, pair_h, WHITE, LINE, rounded=True)
    textbox(s, "For Victoria University", rx + 0.30, pair_y + 0.30, half - 0.55, 0.32, 13, ORANGE, True)
    textbox(
        s,
        "Stronger leadership in experiential learning, applied research, and social impact — "
        "with lessons that can be replicated across Uganda.",
        rx + 0.30,
        pair_y + 0.75,
        half - 0.55,
        pair_h - 1.00,
        14,
        NAVY,
        True,
    )
    slides_meta.append(s)

    # ------------------------------------------------------------------
    # 17 Closing — large readable type
    # ------------------------------------------------------------------
    s = new_slide(NAVY)
    panel_w = 7.85
    crop_photo(s, PHOTOS["facility_wide"], 0, 0, SW, SH)
    rect(s, 0, 0, panel_w, SH, NAVY)
    rect(s, panel_w, 0, 0.12, SH, GOLD)

    add_logo(s, 0.55, 0.35, 0.52, plate=True)

    textbox(s, "INVITATION TO PARTNER", 0.55, 1.15, 7.0, 0.35, 14, GOLD, True)
    textbox(s, ORG, 0.55, 1.55, 7.0, 0.40, 20, TEAL, True)
    textbox(s, "Together, we can build", 0.55, 2.15, 7.0, 0.50, 30, WHITE, True)
    textbox(s, "healthier communities", 0.55, 2.65, 7.0, 0.50, 30, WHITE, True)
    textbox(s, "and stronger professionals.", 0.55, 3.15, 7.0, 0.50, 30, GOLD, True)

    rect(s, 0.55, 3.75, 2.40, 0.07, GOLD)

    textbox(
        s,
        "FairBanks Medical Centre invites Victoria University\n"
        "to join as a strategic partner in advancing FCHIP.",
        0.55,
        4.00,
        7.0,
        0.75,
        16,
        LIGHT,
        line_spacing=1.2,
    )
    textbox(s, SLOGAN, 0.55, 4.85, 7.0, 0.40, 20, GOLD, True)

    rect(s, 0.55, 5.45, 6.90, 1.55, TEAL, rounded=True)
    textbox(s, "Institutional contact", 0.80, 5.65, 6.4, 0.30, 14, GOLD, True)
    textbox(
        s,
        "Kyebando–Kisalosalo, Tirupati Road, Kampala\n"
        "info@fairbanksmedicalcentre.org\n"
        "fairbanksmedicalcentre.org",
        0.80,
        6.05,
        6.4,
        0.80,
        15,
        WHITE,
        line_spacing=1.2,
    )
    slides_meta.append(s)

    total = len(slides_meta)
    for i, s in enumerate(slides_meta):
        if i == 0 or i == total - 1:
            continue
        footer(s, i + 1, total)

    prs.save(OUT_PPT)
    print(f"Wrote {OUT_PPT} ({total} slides)")


if __name__ == "__main__":
    build()
