#!/usr/bin/env python3
"""
MUST partnership pitch deck — FairBanks Medical Centre & FCHIP.
Prepared for Mbarara University of Science and Technology.
Large, readable fonts; simple language; picture-led layouts.
"""

from pathlib import Path

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
ASSETS = REPO / "assets"
CURSOR = REPO / ".cursor"
LOGO = ASSETS / "fairbanks_logo.jpeg"
OUT = Path(__file__).resolve().parent / "documents"
OUT.mkdir(parents=True, exist_ok=True)
OUT_PPT = OUT / "must_ppt.pptx"

NAVY, TEAL, GREEN = "0A1F2E", "0D6E6E", "2D7A55"
ORANGE, GOLD, CREAM = "C45C26", "D99A2B", "F7F5F0"
PALE_TEAL, WHITE = "E8F3F2", "FFFFFF"
SLATE, MUTED, LINE, LIGHT = "1E2F38", "52636C", "CED9D8", "D4E8DC"

ORG = "FairBanks Medical Centre"
SLOGAN = "Your health, our mission."
FOOTER = f"{ORG}  ·  {SLOGAN}"
FCHIP = "FairBanks Community Health Intelligence Platform"
MUST = "Mbarara University of Science and Technology"

PHOTOS = {
    "cover": ASSETS / "facility_exterior_branded_entrance_01.jpeg",
    "facility": ASSETS / "facility_exterior_entrance_01.jpg",
    "facility_wide": ASSETS / "facility_exterior_branded_entrance_02.jpeg",
    "outreach": ASSETS / "outreach_bp_screening.jpeg",
    "camp": ASSETS / "outreach_medical_camp_01.jpg",
    "mission": ASSETS / "reception_mission_wall.jpeg",
    "audience": ASSETS / "outreach_audience_full_group_01.jpg",
    "consult": ASSETS / "clinic_consult_compassion_01.jpg",
    "paeds": ASSETS / "clinic_paediatric_consult.jpeg",
    "maternal": ASSETS / "bloom_maternal_health_participant_01.jpg",
    "dashboard": ASSETS / "dashboard_demo.png",
    "mobile": ASSETS / "outreach_mobile_phone_demo_01.jpg",
    "gis": ASSETS / "gis_hotspots.png",
    "pharmacy": ASSETS / "pharmacy_interior_01.jpg",
    "concept": CURSOR / "concept_improved.jpeg",
}

ML, MR = 0.55, 0.55
SW, SH = 13.333, 7.5
CW = SW - ML - MR
TOP_BAR = 0.06
LOGO_H = 0.38
LOGO_W = LOGO_H * (269 / 101)
LOGO_Y = 0.26
FOOTER_Y = 7.05
CONTENT_BOTTOM = 6.68
HEADER_Y = 1.15
HEADER_Y_SUB = 1.52


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run(run, text, size, color, bold=False, font="Calibri"):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_fade(slide):
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    transition = etree.Element(qn("p:transition"), {"spd": "med"})
    etree.SubElement(transition, qn("p:fade"))
    cSld = sld.find(qn("p:cSld"))
    if cSld is not None:
        cSld.addnext(transition)
    else:
        sld.insert(0, transition)


def build():
    missing = [p for p in [LOGO, *PHOTOS.values()] if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing assets:\n" + "\n".join(map(str, missing)))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
        else:
            shape.line.fill.background()
        if rounded:
            try:
                shape.adjustments[0] = 0.09
            except Exception:
                pass
        try:
            spPr = shape._element.spPr
            for child in list(spPr):
                if "effectLst" in child.tag or "effectDag" in child.tag:
                    spPr.remove(child)
            etree.SubElement(spPr, qn("a:effectLst"))
        except Exception:
            pass
        return shape

    def textbox(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=18,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(18000)
        tf.margin_top = tf.margin_bottom = Emu(6000)
        tf.vertical_anchor = valign
        lines = value.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(4) if len(lines) > 1 else Pt(0)
            r = p.add_run()
            set_run(r, line, size, color, bold)
        return box

    def bullets(slide, items, x, y, w, h, size=17, color=SLATE, space=10):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(18000)
        tf.margin_top = tf.margin_bottom = Emu(4000)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(space)
            r = p.add_run()
            set_run(r, "•  " + item, size, color)
        return box

    def add_logo(slide, x, y, height=LOGO_H):
        pad = 0.08
        py = max(y - pad, TOP_BAR + 0.05)
        rect(slide, x - pad, py, LOGO_W + 2 * pad, height + (y - py) + pad, WHITE, LINE, True)
        return slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), height=Inches(height))

    def crop_photo(slide, path, x, y, w, h):
        with PILImage.open(path) as im:
            iw, ih = im.size
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        fr, ir = w / h, iw / ih
        if ir > fr:
            a = (1 - fr / ir) / 2
            pic.crop_left = pic.crop_right = a
        else:
            a = (1 - ir / fr) / 2
            pic.crop_top = pic.crop_bottom = a
        return pic

    def new_slide(bg=CREAM):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, SW, SH, bg)
        add_fade(s)
        return s

    def header(s, kicker, title, subtitle=""):
        rect(s, 0, 0, SW, TOP_BAR, TEAL)
        rect(s, 0, TOP_BAR, 0.08, SH - TOP_BAR, ORANGE)
        logo_x = SW - MR - LOGO_W
        add_logo(s, logo_x, LOGO_Y)
        tw = logo_x - ML - 0.40
        if kicker.strip():
            textbox(s, kicker.upper(), ML, 0.20, tw, 0.28, 14, ORANGE, True)
            textbox(s, title, ML, 0.46, tw, 0.52, 26, NAVY, True)
            y = 1.02
        else:
            textbox(s, title, ML, 0.28, tw, 0.58, 28, NAVY, True)
            y = 0.92
        if subtitle:
            textbox(s, subtitle, ML, y + 0.02, CW, 0.38, 17, MUTED)
            return HEADER_Y_SUB
        return HEADER_Y

    def footer(s, number, total):
        rect(s, ML, FOOTER_Y, CW, 0.012, LINE)
        textbox(s, FOOTER, ML, FOOTER_Y + 0.08, 9.8, 0.28, 12, MUTED)
        textbox(
            s,
            f"{number}  /  {total}",
            SW - MR - 1.3,
            FOOTER_Y + 0.08,
            1.3,
            0.28,
            12,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )

    slides = []

    # ------------------------------------------------------------------
    # 1 Cover
    # ------------------------------------------------------------------
    s = new_slide(NAVY)
    panel = 7.85
    crop_photo(s, PHOTOS["cover"], 0, 0, SW, SH)
    rect(s, 0, 0, panel, SH, NAVY)
    rect(s, panel, 0, 0.12, SH, GOLD)
    add_logo(s, 0.55, 0.28, 0.48)

    textbox(s, "UNIVERSITY PARTNERSHIP PROPOSAL", 0.55, 1.00, 7.0, 0.34, 16, GOLD, True)
    textbox(s, ORG, 0.55, 1.40, 7.0, 0.40, 20, TEAL, True)
    textbox(s, "Working with MUST", 0.55, 2.00, 7.0, 0.58, 34, WHITE, True)
    textbox(s, "for community health,", 0.55, 2.60, 7.0, 0.58, 34, WHITE, True)
    textbox(s, "research, and grants", 0.55, 3.20, 7.0, 0.58, 34, GOLD, True)
    rect(s, 0.55, 3.90, 2.80, 0.09, GOLD)

    textbox(s, f"Introducing the {FCHIP}", 0.55, 4.20, 7.0, 0.40, 17, LIGHT)
    textbox(s, "(FCHIP)", 0.55, 4.60, 7.0, 0.40, 22, GOLD, True)
    textbox(s, SLOGAN, 0.55, 5.15, 7.0, 0.40, 22, GOLD, True)

    rect(s, 0, 6.25, panel, 1.25, TEAL)
    textbox(s, f"Prepared for {MUST}", 0.55, 6.42, 7.0, 0.38, 18, WHITE, True)
    textbox(s, "Kyebando–Kisalosalo, Kampala  ·  fairbanksmedicalcentre.org", 0.55, 6.88, 7.0, 0.35, 15, LIGHT)
    slides.append(s)

    # ------------------------------------------------------------------
    # 2 Agenda
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "What we will cover")
    items = [
        ("01", "Why we are here"),
        ("02", "About FairBanks Medical Centre"),
        ("03", "How Community Reach works"),
        ("04", "What FCHIP is"),
        ("05", "Why MUST is a strong fit"),
        ("06", "How we can work together"),
        ("07", "Joint research and grants"),
        ("08", "Next steps and invitation"),
    ]
    gap_x, gap_y = 0.28, 0.20
    card_w = (CW - gap_x) / 2
    avail = CONTENT_BOTTOM - y0 - 0.08
    card_h = (avail - 3 * gap_y) / 4
    for i, (num, label) in enumerate(items):
        col, row = i % 2, i // 2
        x = ML + col * (card_w + gap_x)
        y = y0 + 0.08 + row * (card_h + gap_y)
        rect(s, x, y, card_w, card_h, WHITE, LINE, True)
        rect(s, x + 0.20, y + 0.22, 0.08, card_h - 0.44, ORANGE if col == 0 else TEAL)
        mid = y + (card_h - 0.46) / 2
        textbox(s, num, x + 0.42, mid, 0.70, 0.46, 22, ORANGE, True, valign=MSO_ANCHOR.MIDDLE)
        textbox(s, label, x + 1.20, mid, card_w - 1.45, 0.46, 19, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 3 Why we are here
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Purpose", "Why we are here")
    textbox(
        s,
        "FairBanks Medical Centre seeks a practical partnership with MUST — "
        "linking a working clinic and community programmes with university teaching, "
        "research, innovation, and grant writing.",
        ML,
        y0,
        CW,
        0.78,
        18,
        SLATE,
    )
    points = [
        (
            "Share a living field site",
            "A licensed medical centre and active community work in Kampala peri-urban areas.",
            TEAL,
        ),
        (
            "Grow FCHIP together",
            "Build and test the FairBanks Community Health Intelligence Platform with MUST scholars and students.",
            ORANGE,
        ),
        (
            "Win stronger grants",
            "Joint proposals that combine academic depth with real community and facility data.",
            GREEN,
        ),
    ]
    start = y0 + 0.95
    gap = 0.16
    row_h = (CONTENT_BOTTOM - start - 2 * gap) / 3
    for i, (t, body, accent) in enumerate(points):
        y = start + i * (row_h + gap)
        rect(s, ML, y, CW, row_h, WHITE, LINE, True)
        rect(s, ML + 0.18, y + 0.18, 0.08, row_h - 0.36, accent)
        textbox(s, t, ML + 0.42, y + 0.18, CW - 0.70, 0.36, 19, NAVY, True)
        textbox(s, body, ML + 0.42, y + 0.55, CW - 0.70, row_h - 0.70, 17, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 4 About FairBanks
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Who we are", "About FairBanks Medical Centre")
    left_w, gap = 7.10, 0.26
    right_w = CW - left_w - gap
    box_h = CONTENT_BOTTOM - y0
    rect(s, ML, y0, left_w, box_h, WHITE, LINE, True)
    textbox(s, "A community medical centre in Kampala", ML + 0.35, y0 + 0.24, left_w - 0.70, 0.40, 20, NAVY, True)
    textbox(s, SLOGAN, ML + 0.35, y0 + 0.68, left_w - 0.70, 0.32, 17, ORANGE, True)
    bullets(
        s,
        [
            "Located in Kyebando–Kisalosalo, Kampala",
            "Outpatient and inpatient care, emergency, diagnostics, and pharmacy",
            "Maternal and child health, chronic care, and preventive services",
            "Community Reach with CHWs / VHTs in nearby communities",
            "Building FCHIP as our community health intelligence layer",
        ],
        ML + 0.35,
        y0 + 1.15,
        left_w - 0.70,
        box_h - 1.40,
        17,
        SLATE,
        12,
    )
    rx = ML + left_w + gap
    photo_h = box_h * 0.58
    crop_photo(s, PHOTOS["facility"], rx, y0, right_w, photo_h)
    cap_y = y0 + photo_h + 0.14
    cap_h = CONTENT_BOTTOM - cap_y
    rect(s, rx, cap_y, right_w, cap_h, TEAL, rounded=True)
    textbox(s, ORG, rx + 0.22, cap_y + 0.22, right_w - 0.44, 0.34, 16, GOLD, True)
    textbox(
        s,
        "Care close to where families live — and a base for learning and research.",
        rx + 0.22,
        cap_y + 0.62,
        right_w - 0.44,
        cap_h - 0.85,
        15,
        WHITE,
    )
    slides.append(s)

    # ------------------------------------------------------------------
    # 5 Clinical services
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Clinical anchor",
        "What the Medical Centre offers",
        "Services listed on fairbanksmedicalcentre.org",
    )
    services = [
        "General outpatient & inpatient care",
        "24/7 emergency services",
        "Advanced diagnostics & laboratory",
        "Internal medicine",
        "Paediatrics & child health",
        "Obstetrics & gynaecology",
        "ENT & ophthalmology",
        "Urology & men's health",
        "Physiotherapy & rehab",
        "Retail pharmacy",
        "Vaccination programmes",
        "Community outreach & screening",
    ]
    cols, rows = 3, 4
    gap_x, gap_y = 0.22, 0.14
    cell_w = (CW - 2 * gap_x) / 3
    cell_h = (CONTENT_BOTTOM - y0 - 3 * gap_y) / 4
    accents = [TEAL, ORANGE, GREEN]
    for i, item in enumerate(services):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x + 0.16, y + 0.18, 0.08, cell_h - 0.36, accents[col])
        textbox(s, item, x + 0.36, y + 0.12, cell_w - 0.52, cell_h - 0.24, 17, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 6 Community Reach
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Operating model",
        "FairBanks Community Reach",
        "How care, learning, and empowerment connect — FCHIP sits on Data & Feedback.",
    )
    left_w = CW * 0.52
    gap = 0.24
    right_w = CW - left_w - gap
    photo_h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["concept"], ML, y0, left_w, photo_h)
    rx = ML + left_w + gap
    steps = [
        "Community members",
        "CHWs / VHTs — the bridge",
        "Community Reach programmes",
        "FairBanks Medical Centre",
        "Research · partnerships · skills",
        "Economic empowerment & CHIS",
    ]
    row_h = (photo_h - 5 * 0.10) / 6
    for i, step in enumerate(steps):
        y = y0 + i * (row_h + 0.10)
        rect(s, rx, y, right_w, row_h, WHITE, LINE, True)
        rect(s, rx + 0.14, y + 0.14, 0.08, row_h - 0.28, TEAL if i % 2 == 0 else ORANGE)
        textbox(s, f"{i + 1}.  {step}", rx + 0.34, y + 0.08, right_w - 0.48, row_h - 0.16, 16, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 7 What is FCHIP
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "What is FCHIP?")
    def_h = 2.35
    rect(s, ML, y0, CW, def_h, TEAL, rounded=True)
    textbox(s, "FCHIP stands for", ML + 0.40, y0 + 0.20, CW - 0.80, 0.32, 16, GOLD, True)
    textbox(s, FCHIP, ML + 0.40, y0 + 0.55, CW - 0.80, 0.45, 24, WHITE, True)
    textbox(
        s,
        "FCHIP is FairBanks’ deep-tech platform. It brings together community and facility "
        "health data, then uses AI, GIS maps, and climate signals to spot risks early — "
        "so CHWs, clinics, and partners can act before problems grow.",
        ML + 0.40,
        y0 + 1.15,
        CW - 0.80,
        1.00,
        17,
        LIGHT,
    )
    textbox(s, "What FCHIP helps with", ML, y0 + def_h + 0.16, CW, 0.34, 18, NAVY, True)
    focus = [
        "Disease early warning",
        "Maternal health risk",
        "NCD hotspots",
        "Child health signals",
        "Medicine demand",
        "CHW / VHT mobile tools",
        "GIS risk mapping",
        "Climate–health links",
        "EMR / HMS data APIs",
        "Facility dashboards",
    ]
    cols, rows = 5, 2
    gap_x, gap_y = 0.18, 0.14
    grid_top = y0 + def_h + 0.55
    cell_w = (CW - (cols - 1) * gap_x) / cols
    cell_h = (CONTENT_BOTTOM - grid_top - (rows - 1) * gap_y) / rows
    accents = [TEAL, ORANGE, GREEN, GOLD, TEAL]
    for i, item in enumerate(focus):
        col, row = i % cols, i // cols
        x = ML + col * (cell_w + gap_x)
        y = grid_top + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x + 0.12, y + 0.12, cell_w - 0.24, 0.07, accents[col])
        textbox(
            s,
            item,
            x + 0.08,
            y + 0.28,
            cell_w - 0.16,
            cell_h - 0.40,
            15,
            NAVY,
            True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides.append(s)

    # ------------------------------------------------------------------
    # 8 How FCHIP works (visual)
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Technology",
        "From field data to early action",
        "Capture → intelligence → alerts and dashboards for people who can act.",
    )
    gap = 0.22
    left_w = (CW - gap) * 0.48
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["dashboard"], ML, y0, left_w, h * 0.55)
    crop_photo(s, PHOTOS["gis"], ML, y0 + h * 0.55 + 0.14, left_w, h * 0.45 - 0.14)
    rx = ML + left_w + gap
    rect(s, rx, y0, right_w, h, WHITE, LINE, True)
    bullets(
        s,
        [
            "CHWs and VHTs collect data on mobile tools — including offline.",
            "Clinics and hospitals can share records through secure data APIs.",
            "Climate and GIS layers join health signals for earlier warning.",
            "Dashboards help facilities and partners plan outreach and stock.",
            "MUST can help validate models, study impact, and train users.",
        ],
        rx + 0.28,
        y0 + 0.35,
        right_w - 0.50,
        h - 0.55,
        17,
        SLATE,
        14,
    )
    slides.append(s)

    # ------------------------------------------------------------------
    # 9 Why MUST
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Fit with MUST",
        "Why MUST is a strong partner",
        "Drawn from MUST’s published faculties, programmes, and innovation units.",
    )
    reasons = [
        (
            "Faculty of Health Sciences",
            "Training doctors, nurses, pharmacists, lab scientists, and physiotherapists since 1989 — with COBERS community learning.",
            TEAL,
        ),
        (
            "Department of Community Health",
            "Public health training (including MPH), epidemiology, maternal and child health, NCDs, and health systems research.",
            ORANGE,
        ),
        (
            "Computing & innovation",
            "Faculty of Computing and Informatics, CITT, and CAMTech Uganda — natural homes for digital health and FCHIP.",
            GREEN,
        ),
        (
            "Grant and research culture",
            "Long record of partnered research with MoH, districts, and global universities — useful for joint proposals.",
            GOLD,
        ),
    ]
    gap = 0.16
    cell_w = (CW - gap) / 2
    cell_h = (CONTENT_BOTTOM - y0 - gap) / 2
    for i, (t, body, accent) in enumerate(reasons):
        col, row = i % 2, i // 2
        x = ML + col * (cell_w + gap)
        y = y0 + row * (cell_h + gap)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x + 0.18, y + 0.22, 0.08, cell_h - 0.44, accent)
        textbox(s, t, x + 0.40, y + 0.28, cell_w - 0.60, 0.50, 18, NAVY, True)
        textbox(s, body, x + 0.40, y + 0.85, cell_w - 0.60, cell_h - 1.10, 16, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 10 Collaboration areas
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "", "How we can work together")
    textbox(
        s,
        "Five clear areas for collaboration between MUST and FairBanks Medical Centre.",
        ML,
        y0,
        CW,
        0.36,
        17,
        MUTED,
    )
    areas = [
        ("A", "Teaching & placements", "COBERS-style learning, clinical attachments, and supervised field work", TEAL),
        ("B", "Joint research", "Community health, digital health, climate–health, MCH, and NCDs", ORANGE),
        ("C", "Community programmes", "Camps, school health, screening, and VHT-linked outreach", GREEN),
        ("D", "FCHIP co-development", "AI, GIS, mobile tools, and dashboards with Computing & CITT", GOLD),
        ("E", "Joint grants", "Shared concept notes, MoUs, and competitive funding bids", TEAL),
    ]
    badge, inset, gap_y = 0.55, 0.20, 0.12
    top = y0 + 0.45
    row_h = (CONTENT_BOTTOM - top - 4 * gap_y) / 5
    for i, (letter, title, desc, accent) in enumerate(areas):
        y = top + i * (row_h + gap_y)
        rect(s, ML, y, CW, row_h, WHITE, LINE, True)
        by = y + (row_h - badge) / 2
        rect(s, ML + inset, by, badge, badge, accent, rounded=True)
        textbox(s, letter, ML + inset, by, badge, badge, 18, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        tx = ML + inset + badge + 0.26
        tw = CW - (inset + badge + 0.26) - 0.26
        textbox(s, title, tx, y + (row_h - 0.66) / 2, tw, 0.32, 18, NAVY, True)
        textbox(s, desc, tx, y + (row_h - 0.66) / 2 + 0.32, tw, 0.32, 16, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 11 Student learning
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration A",
        "Learning that matches COBERS",
        "Students learn by doing — while communities receive useful care and follow-up.",
    )
    left_w = CW * 0.42
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["outreach"], ML, y0, left_w, h * 0.48)
    crop_photo(s, PHOTOS["consult"], ML, y0 + h * 0.48 + 0.14, left_w, h * 0.52 - 0.14)
    rx = ML + left_w + gap
    disciplines = [
        "Medicine & family medicine",
        "Nursing & midwifery",
        "Public health / MPH",
        "Pharmacy & laboratory",
        "Physiotherapy",
        "Computing & informatics",
        "Business & management",
        "Interdisciplinary studies",
    ]
    row_h = (h - 7 * 0.08) / 8
    for i, d in enumerate(disciplines):
        y = y0 + i * (row_h + 0.08)
        rect(s, rx, y, right_w, row_h, WHITE, LINE, True)
        rect(s, rx + 0.14, y + 0.12, 0.08, row_h - 0.24, ORANGE if i % 2 else TEAL)
        textbox(s, d, rx + 0.34, y + 0.04, right_w - 0.48, row_h - 0.08, 16, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 12 Research themes
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration B",
        "Research themes we can share",
        "Aligned with Community Health, maternal–child health, and digital innovation at MUST.",
    )
    themes = [
        "Community health systems",
        "Maternal & child health",
        "Non-communicable diseases",
        "Disease surveillance",
        "Climate and health",
        "CHW / VHT effectiveness",
        "Digital health & AI",
        "GIS & risk mapping",
        "Medicine demand forecasting",
        "Health equity & urban poor",
        "CHIS & financial protection",
        "Quality improvement",
    ]
    cols, rows = 3, 4
    gap_x, gap_y = 0.20, 0.12
    cell_w = (CW - 2 * gap_x) / 3
    cell_h = (CONTENT_BOTTOM - y0 - 3 * gap_y) / 4
    accents = [TEAL, ORANGE, GREEN]
    for i, t in enumerate(themes):
        col, row = i % 3, i // 3
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x + 0.16, y + 0.16, 0.08, cell_h - 0.32, accents[col])
        textbox(s, t, x + 0.36, y + 0.10, cell_w - 0.52, cell_h - 0.20, 17, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 13 Joint grants
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Collaboration E",
        "Joint grants and partnerships",
        "Together we write stronger proposals — clinic + community + university science.",
    )
    left_w = (CW - 0.24) * 0.55
    right_w = CW - left_w - 0.24
    h = CONTENT_BOTTOM - y0
    rect(s, ML, y0, left_w, h, WHITE, LINE, True)
    textbox(s, "What we can pursue together", ML + 0.30, y0 + 0.28, left_w - 0.55, 0.40, 18, NAVY, True)
    bullets(
        s,
        [
            "Research grants on community health intelligence",
            "Climate–health and early warning pilots",
            "Maternal, child, and NCD programme studies",
            "Digital health / AI innovation awards",
            "NGO and development partner M&E contracts",
            "Capacity-building and training grants",
        ],
        ML + 0.30,
        y0 + 0.85,
        left_w - 0.55,
        h - 1.10,
        17,
        SLATE,
        12,
    )
    rx = ML + left_w + 0.24
    rect(s, rx, y0, right_w, h, TEAL, rounded=True)
    textbox(s, "Why funders listen", rx + 0.28, y0 + 0.35, right_w - 0.56, 0.40, 18, GOLD, True)
    textbox(
        s,
        "MUST brings research quality, ethics, and student power.\n\n"
        "FairBanks brings a licensed facility, Community Reach, and a working FCHIP MVP.\n\n"
        "That mix lowers risk and raises impact for grant reviewers.",
        rx + 0.28,
        y0 + 0.95,
        right_w - 0.56,
        h - 1.30,
        16,
        WHITE,
    )
    slides.append(s)

    # ------------------------------------------------------------------
    # 14 Community impact photos
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "In the field",
        "Community work already underway",
        "Outreach, screening, and education in Kampala peri-urban communities.",
    )
    gap = 0.20
    pw = (CW - 2 * gap) / 3
    ph = (CONTENT_BOTTOM - y0) * 0.62
    crop_photo(s, PHOTOS["camp"], ML, y0, pw, ph)
    crop_photo(s, PHOTOS["audience"], ML + pw + gap, y0, pw, ph)
    crop_photo(s, PHOTOS["maternal"], ML + 2 * (pw + gap), y0, pw, ph)
    labels = [
        "Medical camps & screening",
        "Community education sessions",
        "Maternal & child health focus",
    ]
    ly = y0 + ph + 0.16
    lh = CONTENT_BOTTOM - ly
    for i, label in enumerate(labels):
        x = ML + i * (pw + gap)
        rect(s, x, ly, pw, lh, WHITE, LINE, True)
        textbox(s, label, x + 0.16, ly + 0.10, pw - 0.32, lh - 0.20, 16, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 15 Roadmap
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Implementation", "Proposed roadmap")
    phases = [
        (
            "Phase I",
            "Meet & plan",
            ["Introductory briefing", "Map shared priorities", "Form a small working group"],
            TEAL,
        ),
        (
            "Phase II",
            "Formalise",
            ["Draft and sign MoU", "Pick 2–3 pilot projects", "Agree ethics & data rules"],
            ORANGE,
        ),
        (
            "Phase III",
            "Deliver",
            ["Student placements", "Joint research starts", "First grant submissions"],
            GREEN,
        ),
        (
            "Phase IV",
            "Learn & grow",
            ["Review results together", "Publish and share lessons", "Scale what works"],
            NAVY,
        ),
    ]
    gap = 0.20
    card_w = (CW - 3 * gap) / 4
    card_h = CONTENT_BOTTOM - y0
    for i, (phase, label, items, accent) in enumerate(phases):
        x = ML + i * (card_w + gap)
        rect(s, x, y0, card_w, card_h, WHITE, LINE, True)
        rect(s, x, y0, card_w, 1.15, accent)
        textbox(s, phase, x + 0.16, y0 + 0.20, card_w - 0.32, 0.32, 15, WHITE, True)
        textbox(s, label, x + 0.16, y0 + 0.55, card_w - 0.32, 0.42, 18, WHITE, True)
        bullets(s, items, x + 0.12, y0 + 1.40, card_w - 0.28, card_h - 1.60, 15, SLATE, 10)
    slides.append(s)

    # ------------------------------------------------------------------
    # 16 Ask
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "The ask", "What we invite MUST to do")
    asks = [
        ("01", "Host a planning meeting", "Leadership, Faculty of Health Sciences, Community Health, Computing / CITT."),
        ("02", "Agree a partnership MoU", "Clear roles for teaching, research, community work, and FCHIP."),
        ("03", "Pick first pilots", "One teaching track, one research track, one grant concept."),
        ("04", "Nominate focal persons", "A small joint team to keep work moving."),
    ]
    gap = 0.14
    row_h = (CONTENT_BOTTOM - y0 - 3 * gap) / 4
    for i, (num, title, body) in enumerate(asks):
        y = y0 + i * (row_h + gap)
        rect(s, ML, y, CW, row_h, WHITE, LINE, True)
        rect(s, ML, y, 1.10, row_h, TEAL if i % 2 == 0 else ORANGE)
        textbox(s, num, ML, y, 1.10, row_h, 22, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        textbox(s, title, ML + 1.30, y + 0.18, CW - 1.55, 0.36, 19, NAVY, True)
        textbox(s, body, ML + 1.30, y + 0.55, CW - 1.55, row_h - 0.70, 16, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 17 Closing
    # ------------------------------------------------------------------
    s = new_slide(NAVY)
    panel = 7.85
    crop_photo(s, PHOTOS["facility_wide"], 0, 0, SW, SH)
    rect(s, 0, 0, panel, SH, NAVY)
    rect(s, panel, 0, 0.12, SH, GOLD)
    add_logo(s, 0.55, 0.28, 0.48)

    textbox(s, "INVITATION TO PARTNER", 0.55, 1.00, 7.0, 0.34, 16, GOLD, True)
    textbox(s, ORG, 0.55, 1.40, 7.0, 0.38, 18, TEAL, True)
    textbox(s, "Let us build healthier", 0.55, 2.00, 7.0, 0.50, 30, WHITE, True)
    textbox(s, "communities — and stronger", 0.55, 2.55, 7.0, 0.50, 30, WHITE, True)
    textbox(s, "science — together.", 0.55, 3.10, 7.0, 0.50, 30, GOLD, True)
    rect(s, 0.55, 3.70, 2.80, 0.09, GOLD)

    textbox(
        s,
        f"FairBanks Medical Centre invites {MUST}\n"
        "to partner on Community Reach, FCHIP, and joint grants.",
        0.55,
        4.00,
        7.0,
        0.70,
        16,
        LIGHT,
    )
    textbox(s, SLOGAN, 0.55, 4.80, 7.0, 0.38, 20, GOLD, True)

    rect(s, 0.55, 5.35, 6.90, 1.65, TEAL, rounded=True)
    textbox(s, "Institutional contact", 0.80, 5.52, 6.4, 0.32, 15, GOLD, True)
    textbox(
        s,
        "Kyebando–Kisalosalo, Tirupati Road, Kampala\n"
        "info@fairbanksmedicalcentre.org  ·  +256 748 319 052\n"
        "fairbanksmedicalcentre.org",
        0.80,
        5.92,
        6.4,
        0.90,
        16,
        WHITE,
    )
    slides.append(s)

    total = len(slides)
    for i, s in enumerate(slides):
        if i not in (0, total - 1):
            footer(s, i + 1, total)

    prs.save(OUT_PPT)
    print(f"Wrote {OUT_PPT} ({total} slides)")


if __name__ == "__main__":
    build()
