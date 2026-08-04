#!/usr/bin/env python3
"""
Institution partnership pitch decks — same layout as must_ppt.pptx.
Targets: Makerere University (MUK), IDI, UVRI.
Each deck is 10 slides with institution-specific content.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
ASSETS = REPO / "assets"
LOGO = ASSETS / "fairbanks_logo.jpeg"
OUT = Path(__file__).resolve().parent / "documents"
OUT.mkdir(parents=True, exist_ok=True)

NAVY, TEAL, GREEN = "0A1F2E", "0D6E6E", "2D7A55"
ORANGE, GOLD, CREAM = "C45C26", "D99A2B", "F7F5F0"
WHITE = "FFFFFF"
SLATE, MUTED, LINE, LIGHT = "1E2F38", "52636C", "CED9D8", "D4E8DC"

ORG = "FairBanks Medical Centre"
SLOGAN = "Your health, our mission."
FOOTER = f"{ORG}  ·  {SLOGAN}"
FCHIP = "FairBanks Community Health Intelligence Platform"

PHOTOS = {
    "cover": ASSETS / "facility_exterior_branded_entrance_01.jpeg",
    "facility": ASSETS / "facility_exterior_entrance_01.jpg",
    "facility_wide": ASSETS / "facility_exterior_branded_entrance_02.jpeg",
    "dashboard": ASSETS / "dashboard_demo.png",
}

ML, MR = 0.55, 0.55
SW, SH = 13.333, 7.5
CW = SW - ML - MR
TOP_BAR = 0.06
LOGO_H = 0.40
LOGO_W = LOGO_H * (269 / 101)
LOGO_Y = 0.24
FOOTER_Y = 7.02
CONTENT_BOTTOM = 6.62
HEADER_Y = 1.05
HEADER_Y_SUB = 1.42


@dataclass
class PartnerDeck:
    outfile: str
    partner_short: str
    partner_full: str
    cover_label: str  # e.g. UNIVERSITY / INSTITUTE
    cover_lines: list[str]  # 3 lines under Working with X
    purpose_intro: list[str]
    purpose_points: list[tuple[str, str, str]]
    why_kicker: str
    why_title: str
    why_subtitle: str
    why_reasons: list[tuple[str, str, str]]
    collab_areas: list[tuple[str, str, str, str]]
    grants_subtitle: str
    grants_bullets: list[str]
    grants_why: list[str]
    ask_title: str
    asks: list[str]
    close_invite: list[str]
    fchip_focus: list[str] = field(
        default_factory=lambda: [
            "Disease early warning",
            "Maternal health risk",
            "NCD & child health",
            "CHW mobile tools",
            "GIS & climate links",
            "Facility dashboards",
        ]
    )


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run(run, text, size, color, bold=False, font="Calibri"):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_fade(slide):
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    transition = etree.Element(qn("p:transition"), {"spd": "med"})
    etree.SubElement(transition, qn("p:fade"))
    cSld = sld.find(qn("p:cSld"))
    if cSld is not None:
        cSld.addnext(transition)
    else:
        sld.insert(0, transition)


def build_deck(cfg: PartnerDeck) -> Path:
    missing = [p for p in [LOGO, *PHOTOS.values()] if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing assets:\n" + "\n".join(map(str, missing)))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
        else:
            shape.line.fill.background()
        if rounded:
            try:
                shape.adjustments[0] = 0.09
            except Exception:
                pass
        try:
            spPr = shape._element.spPr
            for child in list(spPr):
                if "effectLst" in child.tag or "effectDag" in child.tag:
                    spPr.remove(child)
            etree.SubElement(spPr, qn("a:effectLst"))
        except Exception:
            pass
        return shape

    def textbox(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=22,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        line_spacing=1.5,
        para_gap=12,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(22000)
        tf.margin_top = tf.margin_bottom = Emu(8000)
        tf.vertical_anchor = valign
        lines = value.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            if len(lines) > 1:
                p.space_after = Pt(22 if line.strip() == "" else para_gap)
            else:
                p.space_after = Pt(0)
            p.line_spacing = line_spacing
            r = p.add_run()
            set_run(r, line, size, color, bold)
        return box

    def bullets(slide, items, x, y, w, h, size=22, color=SLATE, space=18):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(22000)
        tf.margin_top = tf.margin_bottom = Emu(6000)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4)
            p.space_after = Pt(space)
            p.line_spacing = 1.5
            r = p.add_run()
            set_run(r, "•  " + item, size, color)
        return box

    def add_logo(slide, x, y, height=LOGO_H):
        pad = 0.08
        py = max(y - pad, TOP_BAR + 0.05)
        rect(slide, x - pad, py, LOGO_W + 2 * pad, height + (y - py) + pad, WHITE, LINE, True)
        return slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), height=Inches(height))

    def crop_photo(slide, path, x, y, w, h):
        with PILImage.open(path) as im:
            iw, ih = im.size
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        fr, ir = w / h, iw / ih
        if ir > fr:
            a = (1 - fr / ir) / 2
            pic.crop_left = pic.crop_right = a
        else:
            a = (1 - ir / fr) / 2
            pic.crop_top = pic.crop_bottom = a
        return pic

    def new_slide(bg=CREAM):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, SW, SH, bg)
        add_fade(s)
        return s

    def header(s, kicker, title, subtitle=""):
        rect(s, 0, 0, SW, TOP_BAR, TEAL)
        rect(s, 0, TOP_BAR, 0.08, SH - TOP_BAR, ORANGE)
        logo_x = SW - MR - LOGO_W
        add_logo(s, logo_x, LOGO_Y)
        tw = logo_x - ML - 0.40
        if kicker.strip():
            box = s.shapes.add_textbox(Inches(ML), Inches(0.28), Inches(tw), Inches(0.62))
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(18000)
            tf.margin_top = tf.margin_bottom = Emu(4000)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.15
            r1 = p.add_run()
            set_run(r1, kicker.strip() + ": ", 28, ORANGE, True)
            r2 = p.add_run()
            set_run(r2, title, 28, NAVY, True)
            y = 1.00
        else:
            textbox(s, title, ML, 0.24, tw, 0.62, 34, NAVY, True)
            y = 0.98
        if subtitle:
            textbox(s, subtitle, ML, y, CW, 0.42, 20, MUTED)
            return HEADER_Y_SUB
        return HEADER_Y

    def footer(s, number, total):
        rect(s, ML, FOOTER_Y, CW, 0.012, LINE)
        textbox(s, FOOTER, ML, FOOTER_Y + 0.08, 9.8, 0.28, 14, MUTED)
        textbox(
            s,
            f"{number}  /  {total}",
            SW - MR - 1.3,
            FOOTER_Y + 0.08,
            1.3,
            0.28,
            14,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )

    slides = []

    # 1 Cover
    s = new_slide(NAVY)
    panel = 7.85
    crop_photo(s, PHOTOS["cover"], 0, 0, SW, SH)
    rect(s, 0, 0, panel, SH, NAVY)
    rect(s, panel, 0, 0.12, SH, GOLD)
    add_logo(s, 0.55, 0.28, 0.50)
    textbox(s, f"{cfg.cover_label} PARTNERSHIP PROPOSAL", 0.55, 1.00, 7.0, 0.36, 18, GOLD, True)
    textbox(s, ORG, 0.55, 1.42, 7.0, 0.42, 24, TEAL, True)
    textbox(s, f"Working with {cfg.partner_short}", 0.55, 2.05, 7.0, 0.58, 36, WHITE, True)
    textbox(s, cfg.cover_lines[0], 0.55, 2.65, 7.0, 0.58, 34, WHITE, True)
    textbox(s, cfg.cover_lines[1], 0.55, 3.25, 7.0, 0.58, 34, GOLD, True)
    rect(s, 0.55, 3.95, 2.80, 0.09, GOLD)
    textbox(s, "Introducing FCHIP", 0.55, 4.25, 7.0, 0.42, 24, LIGHT)
    textbox(s, FCHIP, 0.55, 4.72, 7.0, 0.55, 19, GOLD, True)
    textbox(s, SLOGAN, 0.55, 5.35, 7.0, 0.40, 24, GOLD, True)
    rect(s, 0, 6.25, panel, 1.25, TEAL)
    textbox(s, f"Prepared for {cfg.partner_full}", 0.55, 6.42, 7.0, 0.38, 18, WHITE, True)
    textbox(s, "Kyebando–Kisalosalo, Kampala  ·  fairbanksmedicalcentre.org", 0.55, 6.88, 7.0, 0.35, 17, LIGHT)
    slides.append(s)

    # 2 Purpose
    s = new_slide()
    y0 = header(s, "Purpose", "Why we are here")
    textbox(s, "\n".join(cfg.purpose_intro), ML, y0, CW, 1.20, 22, SLATE, para_gap=10)
    start = y0 + 1.35
    gap = 0.14
    row_h = (CONTENT_BOTTOM - start - 2 * gap) / 3
    for i, (t, body, accent) in enumerate(cfg.purpose_points):
        y = start + i * (row_h + gap)
        rect(s, ML, y, CW, row_h, WHITE, LINE, True)
        rect(s, ML + 0.18, y + 0.14, 0.10, row_h - 0.28, accent)
        textbox(s, t, ML + 0.45, y + 0.12, CW - 0.75, 0.42, 24, NAVY, True)
        textbox(s, body, ML + 0.45, y + 0.58, CW - 0.75, row_h - 0.72, 21, MUTED)
    slides.append(s)

    # 3 About FairBanks
    s = new_slide()
    y0 = header(s, "Who we are", "About FairBanks Medical Centre")
    left_w, gap = 7.15, 0.26
    right_w = CW - left_w - gap
    box_h = CONTENT_BOTTOM - y0
    rect(s, ML, y0, left_w, box_h, WHITE, LINE, True)
    textbox(s, "A community medical centre in Kampala", ML + 0.35, y0 + 0.28, left_w - 0.70, 0.42, 23, NAVY, True)
    textbox(s, SLOGAN, ML + 0.35, y0 + 0.75, left_w - 0.70, 0.34, 20, ORANGE, True)
    bullets(
        s,
        [
            "Located in Kyebando–Kisalosalo, Kampala",
            "Outpatient & inpatient care, emergency, diagnostics, pharmacy",
            "Maternal and child health, chronic care, prevention",
            "Community Reach with CHWs / VHTs nearby",
            "Building FCHIP as our health intelligence layer",
        ],
        ML + 0.35,
        y0 + 1.25,
        left_w - 0.70,
        box_h - 1.50,
        20,
        SLATE,
        18,
    )
    rx = ML + left_w + gap
    photo_h = box_h * 0.58
    crop_photo(s, PHOTOS["facility"], rx, y0, right_w, photo_h)
    cap_y = y0 + photo_h + 0.14
    cap_h = CONTENT_BOTTOM - cap_y
    rect(s, rx, cap_y, right_w, cap_h, TEAL, rounded=True)
    textbox(s, ORG, rx + 0.22, cap_y + 0.20, right_w - 0.44, 0.36, 18, GOLD, True)
    textbox(
        s,
        "Care close to where families live —\nand a base for learning and research.",
        rx + 0.22,
        cap_y + 0.60,
        right_w - 0.44,
        cap_h - 0.80,
        18,
        WHITE,
        para_gap=8,
    )
    slides.append(s)

    # 4 Community Reach
    s = new_slide()
    y0 = header(
        s,
        "Operating model",
        "FairBanks Community Reach",
        "How care, learning, and empowerment connect. FCHIP supports Data & Feedback.",
    )
    steps = [
        ("1", "Community members", "Communities name needs and own solutions"),
        ("2", "CHWs / VHTs", "The bridge — outreach, referrals, and data"),
        ("3", "Community programmes", "Screening, education, maternal & child care"),
        ("4", "Medical Centre", "Clinical care, diagnostics, pharmacy, follow-up"),
        ("5", "Research & skills", "Evidence, partners, training, and learning"),
        ("6", "Empowerment & CHIS", "Livelihoods and affordable shared protection"),
    ]
    gap_x, gap_y = 0.24, 0.18
    cell_w = (CW - gap_x) / 2
    cell_h = (CONTENT_BOTTOM - y0 - 2 * gap_y) / 3
    accents = [TEAL, ORANGE, GREEN, TEAL, ORANGE, GREEN]
    for i, (num, title, desc) in enumerate(steps):
        col, row = i % 2, i // 2
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x, y, 0.85, cell_h, accents[i])
        textbox(s, num, x, y, 0.85, cell_h, 30, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        textbox(s, title, x + 1.05, y + 0.26, cell_w - 1.25, 0.48, 23, NAVY, True)
        textbox(s, desc, x + 1.05, y + 0.80, cell_w - 1.25, cell_h - 1.00, 21, MUTED)
    slides.append(s)

    # 5 FCHIP
    s = new_slide()
    y0 = header(s, "", "What is FCHIP?")
    left_w = CW * 0.48
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    rect(s, ML, y0, left_w, h, TEAL, rounded=True)
    textbox(s, "FCHIP stands for", ML + 0.35, y0 + 0.28, left_w - 0.70, 0.38, 18, GOLD, True)
    textbox(s, FCHIP, ML + 0.35, y0 + 0.72, left_w - 0.70, 0.90, 26, WHITE, True)
    textbox(
        s,
        "FCHIP brings community and facility health data together.\n"
        "It uses AI, GIS maps, and climate signals to spot risks early —\n"
        "so CHWs, clinics, and partners can act sooner.",
        ML + 0.35,
        y0 + 1.80,
        left_w - 0.70,
        2.40,
        21,
        LIGHT,
        para_gap=12,
    )
    crop_photo(s, PHOTOS["dashboard"], ML + left_w + gap, y0, right_w, h * 0.55)
    fx = ML + left_w + gap
    fy = y0 + h * 0.55 + 0.16
    fh = CONTENT_BOTTOM - fy
    gap_x, gap_y = 0.14, 0.12
    cell_w = (right_w - gap_x) / 2
    cell_h = (fh - 2 * gap_y) / 3
    for i, item in enumerate(cfg.fchip_focus):
        col, row = i % 2, i // 2
        x = fx + col * (cell_w + gap_x)
        y = fy + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        textbox(s, item, x + 0.12, y + 0.08, cell_w - 0.24, cell_h - 0.16, 19, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    slides.append(s)

    # 6 Why partner
    s = new_slide()
    y0 = header(s, cfg.why_kicker, cfg.why_title, cfg.why_subtitle)
    gap = 0.18
    cell_w = (CW - gap) / 2
    cell_h = (CONTENT_BOTTOM - y0 - gap) / 2
    for i, (t, body, accent) in enumerate(cfg.why_reasons):
        col, row = i % 2, i // 2
        x = ML + col * (cell_w + gap)
        y = y0 + row * (cell_h + gap)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x + 0.18, y + 0.22, 0.10, cell_h - 0.44, accent)
        textbox(s, t, x + 0.45, y + 0.28, cell_w - 0.70, 0.55, 22, NAVY, True)
        textbox(s, body, x + 0.45, y + 0.95, cell_w - 0.70, cell_h - 1.25, 20, MUTED)
    slides.append(s)

    # 7 Collaboration
    s = new_slide()
    y0 = header(s, "", "How we can work together")
    badge, inset, gap_y = 0.62, 0.22, 0.14
    top = y0 + 0.08
    row_h = (CONTENT_BOTTOM - top - 4 * gap_y) / 5
    for i, (letter, title, desc, accent) in enumerate(cfg.collab_areas):
        y = top + i * (row_h + gap_y)
        rect(s, ML, y, CW, row_h, WHITE, LINE, True)
        by = y + (row_h - badge) / 2
        rect(s, ML + inset, by, badge, badge, accent, rounded=True)
        textbox(s, letter, ML + inset, by, badge, badge, 22, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        tx = ML + inset + badge + 0.28
        tw = CW - (inset + badge + 0.28) - 0.28
        textbox(s, title, tx, y + (row_h - 0.78) / 2, tw, 0.38, 22, NAVY, True)
        textbox(s, desc, tx, y + (row_h - 0.78) / 2 + 0.38, tw, 0.38, 20, MUTED)
    slides.append(s)

    # 8 Grants
    s = new_slide()
    y0 = header(s, "Funding & science", "Joint research and grants", cfg.grants_subtitle)
    left_w = (CW - 0.24) * 0.52
    right_w = CW - left_w - 0.24
    h = CONTENT_BOTTOM - y0
    rect(s, ML, y0, left_w, h, WHITE, LINE, True)
    textbox(s, "What we can pursue", ML + 0.32, y0 + 0.28, left_w - 0.60, 0.42, 22, NAVY, True)
    bullets(s, cfg.grants_bullets, ML + 0.32, y0 + 0.90, left_w - 0.60, h - 1.15, 21, SLATE, 20)
    rx = ML + left_w + 0.24
    rect(s, rx, y0, right_w, h, TEAL, rounded=True)
    textbox(s, "Why funders listen", rx + 0.32, y0 + 0.35, right_w - 0.64, 0.42, 22, GOLD, True)
    textbox(s, "\n\n".join(cfg.grants_why), rx + 0.32, y0 + 1.00, right_w - 0.64, h - 1.40, 20, WHITE, para_gap=10)
    slides.append(s)

    # 9 Roadmap
    s = new_slide()
    y0 = header(s, "Next steps", "Roadmap and invitation")
    phases = [
        ("I", "Meet & plan", "Briefing · shared priorities · working group", TEAL),
        ("II", "Formalise", "MoU · 2–3 pilots · ethics & data rules", ORANGE),
        ("III", "Deliver", "Joint work starts · first grant bids", GREEN),
        ("IV", "Grow", "Review · publish · scale what works", NAVY),
    ]
    gap = 0.18
    card_w = (CW - 3 * gap) / 4
    card_h = 2.35
    for i, (num, label, detail, accent) in enumerate(phases):
        x = ML + i * (card_w + gap)
        rect(s, x, y0, card_w, card_h, WHITE, LINE, True)
        rect(s, x, y0, card_w, 0.95, accent)
        textbox(s, f"Phase {num}", x + 0.14, y0 + 0.18, card_w - 0.28, 0.32, 18, WHITE, True)
        textbox(s, label, x + 0.14, y0 + 0.50, card_w - 0.28, 0.36, 22, WHITE, True)
        textbox(s, detail, x + 0.14, y0 + 1.15, card_w - 0.28, 1.00, 18, SLATE)
    asks_y = y0 + card_h + 0.22
    asks_h = CONTENT_BOTTOM - asks_y
    rect(s, ML, asks_y, CW, asks_h, WHITE, LINE, True)
    textbox(s, cfg.ask_title, ML + 0.35, asks_y + 0.16, CW - 0.70, 0.40, 22, NAVY, True)
    bullets(s, cfg.asks, ML + 0.35, asks_y + 0.60, CW - 0.70, asks_h - 0.75, 21, SLATE, 16)
    slides.append(s)

    # 10 Close
    s = new_slide(NAVY)
    panel = 7.85
    crop_photo(s, PHOTOS["facility_wide"], 0, 0, SW, SH)
    rect(s, 0, 0, panel, SH, NAVY)
    rect(s, panel, 0, 0.12, SH, GOLD)
    add_logo(s, 0.55, 0.28, 0.50)
    textbox(s, "INVITATION TO PARTNER", 0.55, 1.00, 7.0, 0.36, 18, GOLD, True)
    textbox(s, ORG, 0.55, 1.42, 7.0, 0.40, 22, TEAL, True)
    textbox(s, "Let us build healthier", 0.55, 2.05, 7.0, 0.52, 34, WHITE, True)
    textbox(s, "communities — and stronger", 0.55, 2.60, 7.0, 0.52, 34, WHITE, True)
    textbox(s, "science — together.", 0.55, 3.15, 7.0, 0.52, 34, GOLD, True)
    rect(s, 0.55, 3.75, 2.80, 0.09, GOLD)
    textbox(s, "\n".join(cfg.close_invite), 0.55, 4.05, 7.0, 0.80, 19, LIGHT, para_gap=8)
    textbox(s, SLOGAN, 0.55, 4.95, 7.0, 0.40, 24, GOLD, True)
    rect(s, 0.55, 5.45, 6.90, 1.55, TEAL, rounded=True)
    textbox(s, "Institutional contact", 0.80, 5.58, 6.4, 0.34, 17, GOLD, True)
    textbox(
        s,
        "Kyebando–Kisalosalo, Tirupati Road, Kampala\n"
        "info@fairbanksmedicalcentre.org  ·  +256 748 319 052\n"
        "fairbanksmedicalcentre.org",
        0.80,
        5.98,
        6.4,
        0.90,
        18,
        WHITE,
        para_gap=8,
    )
    slides.append(s)

    total = len(slides)
    assert total == 10, f"Expected 10 slides, got {total}"
    for i, s in enumerate(slides):
        if i not in (0, total - 1):
            footer(s, i + 1, total)

    out_path = OUT / cfg.outfile
    prs.save(out_path)
    print(f"Wrote {out_path} ({total} slides)")
    return out_path


# ---------------------------------------------------------------------------
# Partner content (researched from official public sources)
# ---------------------------------------------------------------------------

MUK = PartnerDeck(
    outfile="muk_ppt.pptx",
    partner_short="Makerere",
    partner_full="Makerere University",
    cover_label="UNIVERSITY",
    cover_lines=["for community health,", "research, and grants"],
    purpose_intro=[
        "FairBanks Medical Centre seeks a practical partnership with Makerere University —",
        "linking a Kampala clinic and Community Reach work with MakCHS teaching,",
        "MakSPH research, and joint grant writing.",
    ],
    purpose_points=[
        ("Share a Kampala field site", "A licensed medical centre and peri-urban Community Reach near Makerere’s catchment.", TEAL),
        ("Grow FCHIP with Mak scholars", "Validate community health intelligence with MakSPH, MakCHS, and computing teams.", ORANGE),
        ("Win stronger grants together", "Combine Makerere’s research depth with real facility and community data.", GREEN),
    ],
    why_kicker="Fit with Makerere",
    why_title="Why Makerere is a strong partner",
    why_subtitle="Based on Makerere College of Health Sciences, MakSPH, and published research programmes.",
    why_reasons=[
        ("College of Health Sciences", "Medicine, nursing, pharmacy, and related programmes — with Mulago teaching links and strong clinical training.", TEAL),
        ("School of Public Health", "MakSPH leads epidemiology, health systems, HIV surveys, METS, and climate–health research such as Cities of Youth.", ORANGE),
        ("Digital & data capacity", "Computing, health informatics, and partner platforms that fit FCHIP’s AI, GIS, and dashboard work.", GREEN),
        ("Grant track record", "Long history of MoH, CDC, DANIDA, and global university partnerships — ideal for joint proposals.", GOLD),
    ],
    collab_areas=[
        ("A", "Student placements", "Community and clinic attachments for MakCHS and MakSPH learners", TEAL),
        ("B", "Joint research", "Urban community health, climate–health, MCH, NCDs, and digital health", ORANGE),
        ("C", "Community programmes", "Outreach, school health, screening, and VHT-linked work in Kampala", GREEN),
        ("D", "FCHIP co-development", "AI, GIS, mobile tools, and dashboards with Makerere researchers", GOLD),
        ("E", "Joint grants", "Shared concept notes, ethics, MoUs, and competitive funding bids", TEAL),
    ],
    grants_subtitle="Clinic + Community Reach + Makerere science — a strong package for funders.",
    grants_bullets=[
        "Urban community health intelligence studies",
        "Climate–health and early warning pilots",
        "Maternal, child, and NCD community research",
        "Digital health / AI innovation awards",
        "Health systems M&E and capacity-building grants",
    ],
    grants_why=[
        "Makerere brings research quality, ethics, and student power.",
        "FairBanks brings a licensed facility, Community Reach, and a working FCHIP MVP.",
        "That mix lowers risk and raises impact for grant reviewers.",
    ],
    ask_title="What we invite Makerere to do",
    asks=[
        "Host a planning meeting with MakCHS, MakSPH, and computing / informatics leads",
        "Agree a partnership MoU for teaching, research, community work, and FCHIP",
        "Pick first pilots: one teaching track, one research track, one grant concept",
    ],
    close_invite=[
        "FairBanks Medical Centre invites Makerere University",
        "to partner on Community Reach, FCHIP, and joint grants.",
    ],
)

IDI = PartnerDeck(
    outfile="idi_ppt.pptx",
    partner_short="IDI",
    partner_full="Infectious Diseases Institute",
    cover_label="INSTITUTIONAL",
    cover_lines=["for infectious disease", "intelligence and impact"],
    purpose_intro=[
        "FairBanks Medical Centre seeks a practical partnership with IDI —",
        "linking community and clinic signals to IDI’s research, training,",
        "health systems, and global health security work.",
    ],
    purpose_points=[
        ("Add a community field node", "CHW/VHT and peri-urban clinic data that complements IDI’s research platforms.", TEAL),
        ("Strengthen early warning", "Use FCHIP with IDI expertise to spot infectious-disease risk earlier.", ORANGE),
        ("Co-write stronger grants", "Pair IDI’s national reach with FairBanks’ live Community Reach site.", GREEN),
    ],
    why_kicker="Fit with IDI",
    why_title="Why IDI is a strong partner",
    why_subtitle="Drawn from IDI’s published programmes, district reach, and Makerere ownership.",
    why_reasons=[
        ("Six core programmes", "Prevention & treatment, training, research, labs, health systems strengthening, and global health security.", TEAL),
        ("National HIV & ID footprint", "IDI supports a large share of Uganda’s ART effort and works across most districts.", ORANGE),
        ("Innovation & data science", "Academy for Health Innovation and African Centre of Excellence in Bioinformatics & Big Data.", GREEN),
        ("Training at scale", "Tens of thousands of health workers trained — a natural path for FCHIP skills transfer.", GOLD),
    ],
    collab_areas=[
        ("A", "Research platforms", "Community and clinic cohorts linked to IDI infectious-disease studies", TEAL),
        ("B", "Global health security", "Community early warning feeds for outbreak preparedness and response", ORANGE),
        ("C", "Training & capacity", "Placements and short courses for CHWs, clinicians, and data users", GREEN),
        ("D", "FCHIP + data science", "Join FCHIP with IDI bioinformatics, dashboards, and digital health work", GOLD),
        ("E", "Joint grants", "Shared proposals for HIV, TB, emerging infections, and health systems", TEAL),
    ],
    grants_subtitle="IDI science + FairBanks community site — stronger bids for infectious disease funders.",
    grants_bullets=[
        "Community early warning for infectious diseases",
        "HIV / TB / emerging infection field studies",
        "Digital surveillance and dashboard innovation",
        "Health systems strengthening and M&E grants",
        "Training and capacity-building awards",
    ],
    grants_why=[
        "IDI brings national reach, research depth, and MoH partnership.",
        "FairBanks brings a licensed facility, Community Reach, and a working FCHIP MVP.",
        "Together we offer funders both community last-mile data and institutional science.",
    ],
    ask_title="What we invite IDI to do",
    asks=[
        "Host a planning meeting with Research, Global Health Security, and Academy leads",
        "Agree an MoU covering research, training, community surveillance, and FCHIP",
        "Pick first pilots: one research track, one GHS early-warning track, one grant concept",
    ],
    close_invite=[
        "FairBanks Medical Centre invites the Infectious Diseases Institute",
        "to partner on Community Reach, FCHIP, and joint grants.",
    ],
    fchip_focus=[
        "Infectious disease alerts",
        "Outbreak clustering",
        "CHW symptom signals",
        "Facility EMR APIs",
        "GIS risk mapping",
        "Partner dashboards",
    ],
)

UVRI = PartnerDeck(
    outfile="uvri_ppt.pptx",
    partner_short="UVRI",
    partner_full="Uganda Virus Research Institute",
    cover_label="INSTITUTIONAL",
    cover_lines=["for viral research,", "surveillance, and grants"],
    purpose_intro=[
        "FairBanks Medical Centre seeks a practical partnership with UVRI —",
        "linking community and clinic health signals to viral surveillance,",
        "epidemiology, and climate-aware early detection.",
    ],
    purpose_points=[
        ("Connect last-mile signals", "CHW/VHT and clinic reports that can complement UVRI surveillance networks.", TEAL),
        ("Fuse climate and health", "Use FCHIP GIS and climate feeds with UVRI arbovirus and outbreak science.", ORANGE),
        ("Co-develop grant concepts", "Pair Entebbe laboratory excellence with a live Kampala community site.", GREEN),
    ],
    why_kicker="Fit with UVRI",
    why_title="Why UVRI is a strong partner",
    why_subtitle="Based on UVRI’s public mandate, departments, and MoH / WHO-linked roles.",
    why_reasons=[
        ("Viral research excellence", "Government biomedical institute in Entebbe focused on viral diseases and public health labs.", TEAL),
        ("Surveillance & epidemiology", "Outbreak investigation, sentinel surveillance, risk assessment, and public health analytics.", ORANGE),
        ("Arbovirology & climate links", "Strong fit for climate-sensitive viral risk mapped through FCHIP’s GIS and climate APIs.", GREEN),
        ("Policy & capacity role", "WHO-linked reference work, training, and advice that help turn evidence into practice.", GOLD),
    ],
    collab_areas=[
        ("A", "Surveillance linkage", "Community fever and symptom clusters feeding viral risk assessment", TEAL),
        ("B", "Joint research", "Arbovirus, outbreak, and climate–health studies with field validation", ORANGE),
        ("C", "Diagnostics pathway", "Clear referral and sample pathways from clinic/outreach to UVRI expertise", GREEN),
        ("D", "FCHIP early warning", "GIS, climate, and AI risk scores aligned to UVRI surveillance needs", GOLD),
        ("E", "Joint grants", "Shared proposals for viral surveillance, capacity building, and innovation", TEAL),
    ],
    grants_subtitle="UVRI viral science + FairBanks community capture — ready for surveillance funders.",
    grants_bullets=[
        "Community-linked viral early warning pilots",
        "Climate-sensitive arbovirus risk mapping",
        "Outbreak preparedness and response studies",
        "Surveillance data systems and analytics grants",
        "Training and research capacity awards",
    ],
    grants_why=[
        "UVRI brings laboratory excellence, epidemiology, and national viral mandate.",
        "FairBanks brings a licensed facility, Community Reach, and a working FCHIP MVP.",
        "Funders get both deep viral science and real community signal capture.",
    ],
    ask_title="What we invite UVRI to do",
    asks=[
        "Host a planning meeting with Epidemiology, Arbovirology, and partnerships leads",
        "Agree an MoU for surveillance linkage, joint research, and FCHIP early warning",
        "Pick first pilots: one surveillance track, one climate–arbovirus track, one grant concept",
    ],
    close_invite=[
        "FairBanks Medical Centre invites Uganda Virus Research Institute",
        "to partner on Community Reach, FCHIP, and joint grants.",
    ],
    fchip_focus=[
        "Fever cluster alerts",
        "Arbovirus risk maps",
        "Climate–virus links",
        "CHW field signals",
        "Outbreak dashboards",
        "Partner reporting",
    ],
)


def main():
    for cfg in (MUK, IDI, UVRI):
        build_deck(cfg)


if __name__ == "__main__":
    main()
