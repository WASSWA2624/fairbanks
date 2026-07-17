#!/usr/bin/env python3
"""
Gadfly Project Custom Web Application Grant 2026 — engineering scope proposal.

Generates synced Word, PDF, and PowerPoint for FairBanks / FCHIP.
Run: python applications/gadfly/build_gadfly_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT = PROJECT / "documents"
SLUG = "gadfly"

OUT_DOC = OUT / f"{SLUG}_word.docx"
OUT_PDF = OUT / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT / f"{SLUG}_ppt.pptx"

CALL_URL = (
    "https://opportunitiesforyouth.org/2026/07/16/"
    "the-gadfly-project-custom-web-application-grant-2026-apply-for-in-kind-software-"
    "development-support-worth-up-to-us100000/"
)

NAVY, TEAL, ACCENT = "0A1F2E", "0D6E6E", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "mobile": "hero_chw_mobile.png",
    "capture": "mobile_capture.png",
    "dashboard": "dashboard_demo.png",
    "architecture": "data_flow_iso_labeled.png",
    "outreach": "outreach_mobile_phone_demo_01.jpg",
    "training": "indoor_training_staff_presenting_01.jpg",
    "gis": "gis_hotspots.png",
    "team": "staff_outreach_conversation_01.jpg",
}

PROGRAMME = "The Gadfly Project — Custom Web Application Grant 2026"
DOC_TITLE = "FCHIP Pilot Build — Community Health Data Platform"
SUBTITLE = "In-kind engineering scope for CHW capture, sync, and dashboards"

OPENING = [
    "FairBanks runs a medical centre and active community outreach across Kampala-area "
    "communities. Community Health Workers and Village Health Teams collect vital signs, "
    "screening results, and referral notes every week — mostly on paper. That data rarely "
    "reaches a shared system in time to prevent outbreaks, flag high-risk pregnancies, or "
    "plan medicine stock.",
    "FairBanks Community Health Intelligence Platform (FCHIP) is the digital layer on our "
    "community health cascade: offline mobile capture in the field, secure sync to the cloud, "
    "and dashboards that turn community signals into early warnings. We are not asking "
    "Gadfly for cash. We are asking for professional in-kind development to build the pilot "
    "product FairBanks can run, own, and hand to partners.",
]

WIN_WIN = [
    ("For FairBanks", "A production-ready pilot: CHW capture app, sync pipeline, and "
     "operations dashboard — validated in live outreach, not a slide deck."),
    ("For Gadfly", "A mission-aligned build with real users (CHWs, clinic staff, programme "
     "managers), measurable community impact, and a portfolio case in African primary care."),
]

MODULES = [
    ("CHW Field Capture", "Offline-first mobile web forms for household visits: vitals, "
     "maternal flags, immunisation, NCD screening, referrals. Works on low-end Android "
     "phones; queues submissions when connectivity drops."),
    ("Sync & Data Layer", "Secure API, role-based access, conflict handling, and audit trail "
     "from device to cloud. Structured records ready for analytics and export."),
    ("Operations Dashboard", "Live view of submissions by area and programme; alert queue for "
     "high-risk cases; basic GIS map of activity hotspots; export for reporting."),
]

USERS = [
    "40+ Community Health Workers and VHTs in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya",
    "FairBanks clinic nurses and outreach coordinators who triage referrals",
    "Programme managers tracking maternal, child, geriatric (Gericare), and NCD outreach",
    "Future district and NGO partners who need interoperable community data feeds",
]

ARCH_ROWS = [
    ["Presentation", "Responsive mobile web (PWA-style) for CHW devices"],
    ["Offline", "Local queue + retry sync when network returns"],
    ["Backend", "REST API, PostgreSQL, role-based auth"],
    ["Dashboard", "Web admin for alerts, maps, and CSV export"],
    ["Hosting", "Documented deploy; FairBanks-operated cloud account"],
]

PHASES = [
    ["Discovery & UX", "Weeks 1–3", "Field shadowing with CHWs; form design; acceptance criteria"],
    ["Core build", "Weeks 4–10", "Capture app, API, sync, dashboard MVP"],
    ["Pilot hardening", "Weeks 11–14", "Bug fixes, performance, training materials"],
    ["Handover", "Week 15+", "Source repo, docs, admin guide, knowledge transfer sessions"],
]

HANDOVER = [
    "Full source code in FairBanks-owned repository with clear README and environment setup",
    "Architecture diagram and API documentation for future developers",
    "Admin and CHW user guides (simple language, screenshots)",
    "Two live knowledge-transfer sessions with FairBanks technical lead",
    "30-day post-handover support window for critical defects",
]

METRICS = [
    "Active CHW users submitting at least weekly during pilot",
    "Median sync time from field submission to dashboard visibility",
    "Referrals triggered from flagged high-risk captures",
    "Community members reached through structured digital records (not paper-only)",
]


def asset(name: str) -> Path:
    p = ASSETS / name
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def photo(key: str) -> Path:
    return asset(PHOTOS[key])


def embed(key_or_path, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key_or_path) if key_or_path in PHOTOS else Path(key_or_path)
    cache = REPO / "tmp" / f"{SLUG}_assets"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from PIL import Image as PILImage

    def rgb(c):
        return RGBColor.from_string(c)

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = rgb(color)

    def shade(cell, color):
        cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))

    def border(cell, color=LINE):
        cell._tc.get_or_add_tcPr().append(parse_xml(
            f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:left w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:bottom w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:right w:val="single" w:sz="8" w:color="{color}"/></w:tcBorders>'))

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes = {1: 20, 2: 14, 3: 12}
        colors = {1: NAVY, 2: TEAL, 3: SLATE}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14 if level == 1 else 8)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            p.clear()
            font(p.add_run(it))

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            shade(c, TEAL)
            border(c, TEAL)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=10)
                if ri % 2:
                    shade(c, CREAM)
                border(c)
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    def image(key, width_in=6.0, caption=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width_in, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(PROGRAMME, size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(DOC_TITLE, size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(SUBTITLE, size=13, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=6)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=12)
    image("cover", caption="FairBanks Community Reach — field operations that need a digital backbone")
    doc.add_page_break()

    heading("1. Engineering overview")
    for p in OPENING:
        para(p)
    table(["Party", "Value"], [[a, b] for a, b in WIN_WIN], widths=[1.5, 4.9])

    heading("2. Product scope — three modules")
    image("architecture", caption="Community signals → capture → sync → dashboard → action")
    for title, body in MODULES:
        heading(title, 2)
        para(body)

    heading("3. Users served in the pilot")
    image("outreach", width_in=5.8, caption="CHWs already collect data in the field — we need reliable digital capture")
    bullets(USERS)

    heading("4. Technical architecture")
    table(["Layer", "Approach"], ARCH_ROWS, widths=[1.8, 4.6])
    image("dashboard", width_in=5.6, caption="Operations dashboard — alerts, maps, exports")

    heading("5. Delivery phases")
    table(["Phase", "Timing", "Deliverables"], PHASES, widths=[1.6, 1.2, 3.6])

    heading("6. Handover & ownership")
    para("FairBanks must own and operate this system after Gadfly's build. Handover includes:")
    bullets(HANDOVER)
    image("training", width_in=5.4, caption="Training CHWs and staff on the finished product")

    heading("7. Impact metrics Gadfly can reference")
    bullets(METRICS)
    para("Grant details and cycle dates: see the official Gadfly programme page.", size=9, color=MUTED, italic=True)
    para(CALL_URL, size=8, color=TEAL, after=4)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak, KeepTogether
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    for name, kw in [
        ("CoverTitle", dict(fontName="Helvetica-Bold", fontSize=20, leading=24, textColor=navy, alignment=TA_CENTER, spaceAfter=6)),
        ("H1", dict(fontName="Helvetica-Bold", fontSize=14, leading=18, textColor=navy, spaceBefore=12, spaceAfter=6)),
        ("H2", dict(fontName="Helvetica-Bold", fontSize=11, leading=14, textColor=teal, spaceBefore=8, spaceAfter=4)),
        ("Body", dict(fontName="Helvetica", fontSize=10, leading=14, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=6)),
        ("Meta", dict(fontName="Helvetica", fontSize=9, leading=12, textColor=muted, alignment=TA_CENTER)),
        ("FBullet", dict(fontName="Helvetica", fontSize=10, leading=13, textColor=slate, leftIndent=12, spaceAfter=3)),
        ("CellHead", dict(fontName="Helvetica-Bold", fontSize=8.5, leading=11, textColor=white)),
        ("CellBody", dict(fontName="Helvetica", fontSize=8.5, leading=11, textColor=slate)),
    ]:
        st.add(ParagraphStyle(name, **kw))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None, max_h=2.8 * inch):
        path = embed(key)
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / iw
        h = w * aspect
        if h > max_h:
            h, w = max_h, max_h / aspect
        block = [Image(path, width=w, height=h)]
        if cap:
            block.append(Paragraph(cap, st["Meta"]))
        story.append(KeepTogether(block))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        data += [[Paragraph(str(c), st["CellBody"]) for c in row] for row in rows]
        t = Table(data, colWidths=widths or [pw / len(headers)] * len(headers), repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal), ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("VALIGN", (0, 0), (-1, -1), "TOP"), ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
            ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Paragraph(PROGRAMME, st["Meta"]))
    story.append(Paragraph(DOC_TITLE, st["CoverTitle"]))
    story.append(Paragraph(f'<i><font color="#{ACCENT}">{SUBTITLE}</font></i>', st["Meta"]))
    story.append(Paragraph(f'<b><i><font color="#{ACCENT}">{SLOGAN}</font></i></b>', st["Meta"]))
    story.append(Spacer(1, 10))
    story.append(img("cover", cap="FairBanks Community Reach — field operations that need a digital backbone"))
    story.append(PageBreak())

    story.append(Paragraph("1. Engineering overview", st["H1"]))
    for p in OPENING:
        story.append(Paragraph(p, st["Body"]))
    tbl(["Party", "Value"], WIN_WIN, [pw * 0.22, pw * 0.78])

    story.append(Paragraph("2. Product scope — three modules", st["H1"]))
    story.append(img("architecture", cap="Community signals → capture → sync → dashboard → action"))
    for title, body in MODULES:
        story.append(Paragraph(title, st["H2"]))
        story.append(Paragraph(body, st["Body"]))

    story.append(Paragraph("3. Users served in the pilot", st["H1"]))
    story.append(img("outreach", w=pw * 0.82, cap="CHWs collect data in the field today"))
    for u in USERS:
        story.append(Paragraph(f"• {u}", st["FBullet"]))

    story.append(PageBreak())
    story.append(Paragraph("4. Technical architecture", st["H1"]))
    tbl(["Layer", "Approach"], ARCH_ROWS, [pw * 0.28, pw * 0.72])
    story.append(img("dashboard", w=pw * 0.85, cap="Operations dashboard"))

    story.append(Paragraph("5. Delivery phases", st["H1"]))
    tbl(["Phase", "Timing", "Deliverables"], PHASES, [pw * 0.24, pw * 0.16, pw * 0.60])

    story.append(Paragraph("6. Handover & ownership", st["H1"]))
    story.append(Paragraph("FairBanks must own and operate this system after Gadfly's build.", st["Body"]))
    for h in HANDOVER:
        story.append(Paragraph(f"• {h}", st["FBullet"]))

    story.append(Paragraph("7. Impact metrics", st["H1"]))
    for m in METRICS:
        story.append(Paragraph(f"• {m}", st["FBullet"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f'<b><i><font color="#{ACCENT}">{SLOGAN}</font></i></b>', st["Meta"]))

    OUT.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                            topMargin=0.7 * inch, bottomMargin=0.7 * inch)
    doc.build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH, blank = prs.slide_width, prs.slide_height, prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        for i, ln in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.28), Inches(11), Inches(0.25),
           f"FairBanks FCHIP  |  Gadfly Grant 2026  |  {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.9), SH - Inches(0.28), Inches(0.6), Inches(0.25), str(n), size=9, color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Slide 1 — cover
    s = prs.slides.add_slide(blank)
    pic(s, "mobile", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.2), SW, Inches(3.2), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.95), Inches(12), Inches(0.35), PROGRAMME, size=13, bold=True, color=TEAL)
    tb(s, Inches(0.6), SH - Inches(2.4), Inches(12), Inches(0.7), DOC_TITLE, size=26, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.55), Inches(12), Inches(0.4), SUBTITLE, size=14, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(0.95), Inches(12), Inches(0.35), SLOGAN, size=13, bold=True, color="FFFFFF")

    # Slide 2 — the ask
    s = prs.slides.add_slide(blank)
    band(s, "The ask", "In-kind engineering, not a cash grant")
    tb(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5),
       "Gadfly builds the pilot product.\nFairBanks runs it in live CHW outreach.\nBoth sides get a measurable community-health case study.",
       size=20, color=SLATE)
    pic(s, "capture", Inches(6.8), Inches(1.2), Inches(6), Inches(5.5))
    footer(s, 2)

    # Slide 3 — modules
    s = prs.slides.add_slide(blank)
    band(s, "Product scope", "Three modules that complete the loop")
    for i, (t, b) in enumerate(MODULES):
        x = Inches(0.4 + i * 4.35)
        rect(s, x, Inches(1.3), Inches(4.1), Inches(5.3), "FFFFFF", LINE)
        rect(s, x, Inches(1.3), Inches(4.1), Inches(0.12), TEAL)
        tb(s, x + Inches(0.2), Inches(1.55), Inches(3.7), Inches(0.5), t, size=18, bold=True, color=TEAL)
        tb(s, x + Inches(0.2), Inches(2.2), Inches(3.7), Inches(4), b, size=14, color=MUTED)
    footer(s, 3)

    # Slide 4 — users
    s = prs.slides.add_slide(blank)
    band(s, "Users served", "Real people, real weekly workflows")
    pic(s, "team", Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.6))
    y = Inches(1.35)
    for u in USERS:
        tb(s, Inches(6.5), y, Inches(6.3), Inches(1.2), f"•  {u}", size=15, color=SLATE)
        y += Inches(1.25)
    footer(s, 4)

    # Slide 5 — architecture
    s = prs.slides.add_slide(blank)
    band(s, "Stack", "Pragmatic, maintainable, FairBanks-owned")
    pic(s, "architecture", Inches(0.4), Inches(1.15), Inches(7), Inches(5.7))
    y = Inches(1.3)
    for layer, approach in ARCH_ROWS:
        tb(s, Inches(7.7), y, Inches(5.2), Inches(0.35), layer, size=14, bold=True, color=TEAL)
        tb(s, Inches(7.7), y + Inches(0.38), Inches(5.2), Inches(0.55), approach, size=12, color=MUTED)
        y += Inches(1.05)
    footer(s, 5)

    # Slide 6 — delivery
    s = prs.slides.add_slide(blank)
    band(s, "Delivery", "Discovery → build → hardening → handover")
    for i, (ph, tm, deliv) in enumerate(PHASES):
        y = Inches(1.25 + i * 1.35)
        rect(s, Inches(0.5), y, Inches(12.3), Inches(1.15), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.12), Inches(2.2), Inches(0.4), ph, size=16, bold=True, color=TEAL)
        tb(s, Inches(3.0), y + Inches(0.12), Inches(1.5), Inches(0.4), tm, size=13, color=ACCENT)
        tb(s, Inches(4.6), y + Inches(0.12), Inches(7.8), Inches(0.9), deliv, size=14, color=SLATE)
    footer(s, 6)

    # Slide 7 — handover
    s = prs.slides.add_slide(blank)
    band(s, "Handover", "FairBanks owns the code and the operations")
    pic(s, "training", Inches(0.4), Inches(1.15), Inches(5.5), Inches(5.7))
    y = Inches(1.3)
    for h in HANDOVER[:4]:
        tb(s, Inches(6.2), y, Inches(6.6), Inches(1.1), f"•  {h}", size=14, color=SLATE)
        y += Inches(1.25)
    footer(s, 7)

    # Slide 8 — win-win close
    s = prs.slides.add_slide(blank)
    band(s, "Win-win", "Engineering impact both sides can measure")
    for i, (party, val) in enumerate(WIN_WIN):
        x = Inches(0.5 + i * 6.4)
        rect(s, x, Inches(1.4), Inches(6.1), Inches(4.8), "FFFFFF", LINE)
        tb(s, x + Inches(0.25), Inches(1.65), Inches(5.6), Inches(0.5), party, size=20, bold=True, color=ACCENT)
        tb(s, x + Inches(0.25), Inches(2.3), Inches(5.6), Inches(3.5), val, size=16, color=SLATE)
    tb(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4), SLOGAN, size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER, italic=True)
    footer(s, 8)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done.", OUT)
