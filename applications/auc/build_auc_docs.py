#!/usr/bin/env python3
"""
AUC Venture Lab Women Innovation Fellowship 2026 — FairBanks / FCHIP application set.

Generates:
  applications/auc/documents/auc_word.docx
  applications/auc/documents/auc_pdf.pdf
  applications/auc/documents/auc_ppt.pptx

Run: python applications/auc/build_auc_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "auc"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/07/16/"
    "auc-venture-lab-women-innovation-fellowship-empowering-women-entrepreneurs-in-egypt-and-beyond/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."
TAGLINE = "Health for All — Obulamu eri Bonna · Afya kwa Wote · Oburamu bwa Boona"

META = {
    "programme": "AUC Venture Lab Women Innovation Fellowship 2026",
    "doc_title": "Fellowship Application — Women-Led HealthTech Leadership",
    "subtitle": "From Kampala communities to Cairo — scaling FCHIP with mentorship and courage",
    "fellowship_dates": "Cairo, Egypt · 1 September – 31 October 2026",
    "deadline": "Application deadline: 25 July 2026",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "outreach": "outreach_facilitator_canopy_01.jpg",
    "training": "indoor_training_staff_presenting_01.jpg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "dashboard": "dashboard_demo.png",
    "team": "staff_team_reception.jpeg",
    "conclusion": "outreach_audience_full_group_01.jpg",
}


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


SECTIONS = [
    ("1. Why Cairo, why now", [
        "FairBanks is not only a medical centre — it is a community health ecosystem with active outreach, "
        "Community Health Workers (CHWs), and a growing intelligence layer called FCHIP. I am applying to the "
        "AUC Venture Lab Women Innovation Fellowship because HealthTech that truly serves Africa needs "
        "leadership training, cross-regional networks, and investor-ready discipline — not just code.",
        "Cairo sits at a crossroads between North and East Africa. The fellowship’s global exposure (including "
        "engagement with leaders in Egypt and Denmark) matches my goal: learn from one of Africa’s top "
        "accelerators while bringing field-tested community health innovation from Uganda.",
    ]),
    ("2. The leader I am becoming", [
        "Women-led ventures in HealthTech carry a double responsibility: build credible technology and earn "
        "trust in communities where care has often been late, costly, or out of reach. My leadership journey "
        "combines clinical operations, community outreach, and product vision for FCHIP.",
        "Through FairBanks Community Reach I work alongside CHWs and Village Health Teams in Kampala-area "
        "communities — maternal and child health, adolescent health, chronic screening, and home visits. "
        "That ground truth shapes how I make decisions: communities first, evidence always, scale with integrity.",
    ]),
    ("3. FCHIP — women-led HealthTech from Uganda", [
        "FairBanks Community Health Intelligence Platform (FCHIP) is the intelligence layer on our cascade: "
        "community members → CHWs/VHTs → outreach programmes → medical centre → research and partnerships → "
        "economic empowerment. FCHIP turns community-generated data into early warnings — outbreaks, high-risk "
        "pregnancies, chronic-disease hotspots, child-health threats, and medicine demand — before crises hit.",
        "This is deep technology: offline mobile capture, cloud sync, AI/ML risk scoring, GIS mapping, and "
        "dashboards for clinics, districts, and partners. It is built to move primary care from reactive "
        "treatment to proactive prevention.",
    ]),
    ("4. What the fellowship unlocks for FairBanks", [
        "AUC Venture Lab has supported 440+ startups and raised strong investment across Africa. For FairBanks, "
        "the fellowship offers leadership development, mentor access, and ecosystem engagement that accelerate "
        "FCHIP from live pilot to investment-ready venture.",
        "I seek structured training in strategic decision-making, storytelling to investors, and resilient "
        "founder leadership — skills that strengthen how we govern data ethically and scale across East Africa.",
    ]),
    ("5. What I bring to the cohort", [
        "A functioning medical centre and pharmacy generating real patient and community signals.",
        "Active outreach in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and surrounding areas.",
        "Trusted CHW/VHT relationships and programmes on maternal health, adolescents, and chronic care.",
        "A live HealthTech product direction (FCHIP) validated in the field, not only on slides.",
        "A bridge between grassroots health realities and continental HealthTech discourse.",
    ]),
    ("6. Cairo fellowship goals (September–October 2026)", [
        "Refine FCHIP’s investor narrative and go-to-market for clinics, districts, and NGO partners.",
        "Build mentor relationships in HealthTech, AI, and impact investing across Egypt and broader Africa.",
        "Design a cross-regional learning exchange — how community intelligence models transfer between countries.",
        "Return to Uganda with sharper leadership, a stronger network, and a clearer scale plan for FCHIP.",
    ]),
]

CAIRO_GOALS = [
    "Leadership labs: decision-making under uncertainty in community health markets",
    "Mentor hours with HealthTech founders and investors in AUC’s ecosystem",
    "Cross-cultural collaboration with fellows from Egypt, Denmark links, and beyond",
    "Pitch refinement for FCHIP as a women-led, Africa-rooted intelligence company",
]

MUTUAL_VALUE = [
    ("For FairBanks", "World-class accelerator training, continental visibility, and mentor networks that "
     "help FCHIP scale beyond Uganda."),
    ("For AUC Venture Lab", "A fellow bringing live HealthTech field validation, CHW data pipelines, and "
     "a compelling women-innovation story from East Africa."),
]


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption=TAGLINE)
    para(META["fellowship_dates"], size=10, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    for title, paras in SECTIONS:
        heading(title)
        if "FCHIP" in title:
            image("architecture", caption="FCHIP — intelligence on the FairBanks community health cascade")
        for t in paras:
            para(t)
        if title.startswith("5."):
            bullets([
                "A functioning medical centre and pharmacy generating real patient and community signals.",
                "Active outreach in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and surrounding areas.",
                "Trusted CHW/VHT relationships and programmes on maternal health, adolescents, and chronic care.",
                "A live HealthTech product direction (FCHIP) validated in the field, not only on slides.",
                "A bridge between grassroots health realities and continental HealthTech discourse.",
            ])
        if title.startswith("6."):
            bullets(CAIRO_GOALS)

    heading("7. A partnership that lifts both sides")
    para("This fellowship works when the cohort gains from real ventures — and ventures gain from structured "
         "leadership growth. FairBanks offers live community health operations; AUC Venture Lab offers the "
         "mentorship and ecosystem that help FCHIP become investment-ready across Africa.")
    for label, body in MUTUAL_VALUE:
        para(f"{label}: {body}", bold=True)

    image("conclusion", caption="Community-first HealthTech — Your health, our mission.")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = w * ih / iw
        if h > 2.8 * inch:
            h, w = 2.8 * inch, 2.8 * inch * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 8))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>FCHIP — {META["subtitle"]}</i></b></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    story.append(Spacer(1, 8))
    img("cover", cap=TAGLINE)
    story.append(Paragraph(META["fellowship_dates"], st["Meta"]))
    story.append(PageBreak())

    for title, paras in SECTIONS:
        story.append(Paragraph(title, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
        if "FCHIP" in title:
            img("architecture", w=pw * 0.85, cap="FCHIP on the community health cascade")
        for t in paras:
            story.append(Paragraph(t, st["Body"]))
        if title.startswith("5."):
            for b in [
                "Functioning medical centre and pharmacy with live data signals",
                "Outreach across Kampala communities with CHWs and VHTs",
                "Maternal, adolescent, and chronic-care programmes on the ground",
                "FCHIP validated through field operations, not theory alone",
                "Bridge between grassroots health and continental HealthTech",
            ]:
                story.append(Paragraph("• " + b, st["FBullet"]))
        if title.startswith("6."):
            for g in CAIRO_GOALS:
                story.append(Paragraph("• " + g, st["FBullet"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph("7. A partnership that lifts both sides", st["H1"]))
    story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
    for label, body in MUTUAL_VALUE:
        story.append(Paragraph(f"<b>{label}:</b> {body}", st["Body"]))
    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                            topMargin=0.7 * inch, bottomMargin=0.7 * inch,
                            title="FairBanks AUC Fellowship Application")
    doc.build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = "Calibri"
        r.font.color.rgb = C(color)
        return box

    def pic_cover(sl, key):
        path = embed(key)
        sl.shapes.add_picture(path, 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks FCHIP | AUC Women Innovation Fellowship | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    slides = [
        ("title", None),
        ("Why Cairo", "Cross-regional leadership for women in HealthTech"),
        ("Leadership journey", "Communities first — building trust before scale"),
        ("FCHIP", "Women-led intelligence on a live community health cascade"),
        ("Fellowship unlocks", "Mentorship, ecosystem, investor readiness"),
        ("What I bring", "Field validation Egypt–Africa cohorts can learn from"),
        ("Cairo goals", "September–October 2026 — concrete outcomes"),
        ("Mutual lift", "Stronger founder, stronger continental HealthTech story"),
        ("Close", None),
    ]

    # 1 Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.5), SW, Inches(3.5), NAVY)
    tb(s, Inches(0.6), SH - Inches(3.2), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.7), Inches(12), Inches(0.65), META["doc_title"], size=26, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.85), Inches(12), Inches(0.4), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.2), Inches(12), Inches(0.35), META["fellowship_dates"], size=12, color="D0E8E8")

    # 2 Why Cairo
    s = prs.slides.add_slide(blank)
    band(s, "Fellowship fit", slides[1][1])
    tb(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.5),
       "HealthTech that serves Africa needs leadership training and cross-regional networks — not code alone.\n\n"
       "AUC Venture Lab is Africa’s top accelerator heritage (440+ startups). Cairo connects North and East Africa.\n\n"
       "I bring live community health operations from Uganda; I seek mentor access and investor-ready discipline.",
       size=16, color=MUTED)
    pic_fit(s, "team", Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5))
    footer(s, 2)

    # 3 Leadership
    s = prs.slides.add_slide(blank)
    band(s, "Leadership", slides[2][1])
    pic_fit(s, "outreach", Inches(0.45), Inches(1.15), Inches(6.0), Inches(5.6))
    tb(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
       "• CHW/VHT partnerships in Kampala communities\n"
       "• Maternal, adolescent, and chronic-care programmes\n"
       "• Clinical + community + product leadership combined\n"
       "• Decisions shaped by people who need care first",
       size=17, color=SLATE)
    footer(s, 3)

    # 4 FCHIP
    s = prs.slides.add_slide(blank)
    band(s, "Innovation", slides[3][1])
    pic_fit(s, "architecture", Inches(0.4), Inches(1.1), Inches(7.2), Inches(5.7))
    tb(s, Inches(7.9), Inches(1.2), Inches(5.0), Inches(5.5),
       "Offline capture → cloud intelligence → alerts to CHWs, clinics, districts.\n\n"
       "Predict outbreaks, maternal risk, NCD hotspots, child health gaps, stock-outs.",
       size=15, color=MUTED)
    footer(s, 4)

    # 5 Fellowship unlocks
    s = prs.slides.add_slide(blank)
    band(s, "Growth", slides[4][1])
    tb(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
       "Leadership development · Mentor hours · Ecosystem engagement · Cross-cultural collaboration\n\n"
       "From live pilot toward investment-ready venture across East Africa.",
       size=18, color=SLATE)
    pic_fit(s, "training", Inches(8.5), Inches(1.3), Inches(4.5), Inches(4.8))
    footer(s, 5)

    # 6 What I bring
    s = prs.slides.add_slide(blank)
    band(s, "Cohort value", slides[5][1])
    pic_fit(s, "mobile", Inches(0.45), Inches(1.15), Inches(5.8), Inches(5.6))
    tb(s, Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5),
       "• Functioning medical centre + pharmacy\n• Live outreach data pipelines\n• Trusted community relationships\n"
       "• FCHIP field validation\n• Grassroots ↔ continental bridge",
       size=17, color=SLATE)
    footer(s, 6)

    # 7 Cairo goals
    s = prs.slides.add_slide(blank)
    band(s, "Sep–Oct 2026", slides[6][1])
    for i, g in enumerate(CAIRO_GOALS):
        y = Inches(1.25) + i * Inches(1.15)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.0), "FFFFFF", LINE)
        rect(s, Inches(0.5), y, Inches(0.12), Inches(1.0), TEAL)
        tb(s, Inches(0.75), y + Inches(0.15), Inches(11.5), Inches(0.7), g, size=15, color=SLATE)
    footer(s, 7)

    # 8 Mutual lift
    s = prs.slides.add_slide(blank)
    band(s, "Shared outcome", slides[7][1])
    rect(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "FairBanks gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[0][1], size=14, color=MUTED)
    rect(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4), "AUC Venture Lab gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[1][1], size=14, color=MUTED)
    footer(s, 8)

    # 9 Close
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.8), SW, Inches(2.8), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.4), Inches(12), Inches(0.55),
       "Ready for Cairo — women-led HealthTech with community roots.", size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.5), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(0.9), Inches(12), Inches(0.3), META["deadline"], size=12, color="D0E8E8")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. AUC document set in", OUT_DIR)
