#!/usr/bin/env python3
"""
Africa CDC AES Epidemiology Fellowship 2026–2028 — FairBanks / FCHIP application set.

Generates:
  applications/africa-cdc/documents/africa-cdc_word.docx
  applications/africa-cdc/documents/africa-cdc_pdf.pdf
  applications/africa-cdc/documents/africa-cdc_ppt.pptx

Run: python applications/africa-cdc/build_africa-cdc_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "africa-cdc"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/07/16/"
    "africa-cdc-african-epidemic-services-aes-epidemiology-fellowship-2026-2028-"
    "fully-funded-apply-by-26-august-2026/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."
TAGLINE = "Health for All — Obulamu eri Bonna · Afya kwa Wote · Oburamu bwa Boona"

META = {
    "programme": "Africa CDC African Epidemic Services (AES) Epidemiology Fellowship 2026–2028",
    "doc_title": "Fellowship Nomination — Surveillance Science for Community Health",
    "subtitle": "Professional development that returns outbreak intelligence to FCHIP and Uganda's communities",
    "timeline": "3 months training in Addis Ababa · 21 months field placement · fully funded",
    "deadline": "Apply by 26 August 2026",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "gis": "gis_hotspots.png",
    "training": "indoor_training_staff_presenting_01.jpg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "outreach": "outreach_facilitator_canopy_01.jpg",
    "dashboard": "dashboard_demo.png",
    "conclusion": "outreach_audience_full_group_01.jpg",
}

SECTIONS = [
    ("1. Fellowship purpose", [
        "Africa CDC's AES fellowship builds outbreak investigation, surveillance, and data-analysis skills "
        "over two years — including a Public Health Informatics track. FairBanks needs this depth to make "
        "FCHIP's surveillance claims scientifically sound and operationally useful in real communities.",
        "We nominate an eligible FairBanks public-health professional (African Union citizen, under 35, "
        "employed in public health with degree and three years' experience) for the Epidemiology or "
        "Informatics track. Africa CDC gains a fellow already embedded in community health work; "
        "FairBanks gains skills that return directly to CHWs, clinics, and partner districts.",
    ]),
    ("2. Why surveillance science needs community roots", [
        "Predictive community platforms fail without trained people who understand surveillance science — "
        "not only software. Outbreak signals from Community Health Workers and Village Health Teams only "
        "become action when someone can investigate cases, analyse trends, and design informatics pipelines "
        "with discipline.",
        "Africa needs more young epidemiologists and informaticians grounded in primary care realities. "
        "FairBanks sits inside that reality: live outreach, maternal and child health, chronic screening, "
        "and a growing intelligence layer on the community health cascade.",
    ]),
    ("3. Skills that return to FCHIP", [
        "A FairBanks fellow will complete AES training and bring methods home: standardised community "
        "signal → investigation → response workflows; surveillance analysis for fever clusters, maternal "
        "risk, and medicine stock-outs; and informatics design that improves CHW data quality.",
        "Employer support ensures learning returns to Community Reach programmes — not a certificate "
        "on a shelf. Supervisors and CHW leads will receive cascade training so Africa CDC's investment "
        "multiplies across Kampala catchments and partner sites.",
    ]),
    ("4. Priority application areas", [
        "Outbreak workflows: community fever reports trigger structured investigation and pharmacy "
        "pre-stocking before surges spread.",
        "Informatics: validation rules and data-quality pipelines for offline CHW capture feeding FCHIP.",
        "Training cascade: teach CHW supervisors better data use and early-warning interpretation.",
    ]),
    ("5. Employer commitment", [
        "FairBanks Medical Centre and FairBanks Community Reach will release the fellow for Addis Ababa "
        "training, provide a structured field placement linked to outreach and facility data flows, "
        "and assign internal mentors for surveillance quality improvement.",
        "We cover counterpart supervision time, device access for field work, and integration of "
        "fellowship outputs into FCHIP modules at no extra cost to Africa CDC beyond the fellowship package.",
    ]),
]

FELLOWSHIP_PATH = [
    ("Months 1–3", "Core epidemiology and informatics training in Addis Ababa"),
    ("Months 4–12", "Field placement — outbreak workflows and CHW data quality at FairBanks sites"),
    ("Months 13–24", "Advanced surveillance modules; district partner engagement; FCHIP alert validation"),
    ("Throughout", "Monthly learning reports; Africa CDC network participation; community feedback loops"),
]

MUTUAL_VALUE = [
    ("FairBanks / FCHIP gains", "World-class epidemiology training; stronger surveillance modules; "
     "stipend, travel, insurance, hardware, and software covered; continental AES network access."),
    ("Africa CDC gains", "Fellow employed in live community health work; field placement linked to "
     "real CHW and facility data; Uganda contribution to Africa's epidemic readiness; informatics "
     "track aligned with digital community health tools."),
]

FIT_ROWS = [
    ("AU citizen under 35", "Ugandan staff candidate within eligibility range"),
    ("Employed in public health", "FairBanks clinical / community health employment"),
    ("Degree + 3 years' experience", "Nominee will meet stated requirements"),
    ("Epidemiology / Informatics tracks", "Direct fit to FCHIP surveillance and data-quality goals"),
]


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{TEAL}"/>'))
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=10, color=SLATE)
                if ri % 2:
                    c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CREAM}"/>'))
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(META["subtitle"], size=12, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption=TAGLINE)
    para(META["timeline"], size=10, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    for title, paras in SECTIONS:
        heading(title)
        if "FCHIP" in title:
            image("architecture", caption="Community signals → capture → FCHIP intelligence → action")
        if "surveillance" in title.lower():
            image("gis", caption="GIS hotspot mapping — surveillance outputs communities can use")
        for t in paras:
            if title.startswith("4."):
                bullets([t])
            else:
                para(t)

    heading("6. Two-year fellowship pathway")
    table(["Phase", "Activities"], FELLOWSHIP_PATH, widths=[1.8, 4.6])

    heading("7. Fit to AES fellowship criteria")
    table(["Criterion", "How FairBanks responds"], FIT_ROWS, widths=[2.4, 4.0])

    heading("8. Mutual value — fellow, Africa CDC, and communities")
    for label, body in MUTUAL_VALUE:
        para(f"{label}: {body}", bold=True)

    heading("9. Selection ask")
    para("Epidemics are fought with skilled people. FairBanks will back our fellow completely so Africa CDC's "
         "investment returns to communities across Uganda.")
    para("We ask for selection into the 2026–2028 AES Epidemiology or Public Health Informatics cohort.")
    image("conclusion", caption="Skills that return to CHWs, clinics, and partner districts")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable, Table, TableStyle
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))
    st.add(ParagraphStyle("CellHead", fontName="Helvetica-Bold", fontSize=8.5, textColor=white))
    st.add(ParagraphStyle("CellBody", fontName="Helvetica", fontSize=8.5, textColor=slate))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = w * ih / iw
        if h > 2.8 * inch:
            h, w = 2.8 * inch, 2.8 * inch * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 8))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        for row in rows:
            data.append([Paragraph(str(c), st["CellBody"]) for c in row])
        widths = widths or [pw / len(headers)] * len(headers)
        t = Table(data, colWidths=widths, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal),
            ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(
        f'<font color="#{ACCENT}"><b><i>FairBanks Community Health Intelligence Platform (FCHIP)</i></b></font>',
        st["Meta"]))
    story.append(Paragraph(f'<font color="#{MUTED}"><i>{META["subtitle"]}</i></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    story.append(Spacer(1, 8))
    img("cover", cap=TAGLINE)
    story.append(Paragraph(META["timeline"], st["Meta"]))
    story.append(PageBreak())

    for title, paras in SECTIONS:
        story.append(Paragraph(title, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
        if "FCHIP" in title:
            img("architecture", w=pw * 0.85, cap="FCHIP on the community health cascade")
        for t in paras:
            if title.startswith("4."):
                story.append(Paragraph("• " + t, st["FBullet"]))
            else:
                story.append(Paragraph(t, st["Body"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph("6. Two-year fellowship pathway", st["H1"]))
    story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
    tbl(["Phase", "Activities"], FELLOWSHIP_PATH, [pw * 0.28, pw * 0.72])

    story.append(Paragraph("7. Fit to AES fellowship criteria", st["H1"]))
    tbl(["Criterion", "How FairBanks responds"], FIT_ROWS, [pw * 0.35, pw * 0.65])

    story.append(Paragraph("8. Mutual value", st["H1"]))
    for label, body in MUTUAL_VALUE:
        story.append(Paragraph(f"<b>{label}:</b> {body}", st["Body"]))

    story.append(Paragraph("9. Selection ask", st["H1"]))
    story.append(Paragraph(
        "We ask for selection into the 2026–2028 AES Epidemiology or Public Health Informatics cohort.", st["Body"]))
    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                            topMargin=0.7 * inch, bottomMargin=0.7 * inch,
                            title="FairBanks Africa CDC AES Fellowship")
    doc.build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    def _add_entrance_anims(slide):
        """Light fade-in on pictures and key text (documents.mdc motion)."""
        from lxml import etree
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spids = []
        for shape in slide.shapes:
            try:
                st = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                st = -1
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
            if st == 13 or has_text:
                spids.append(str(shape.shape_id))
            if len(spids) >= 3:
                break
        if not spids:
            return
        sld = slide._element
        for old in sld.findall(f"{{{NS_P}}}timing"):
            sld.remove(old)
        children = []
        nid = 3
        for i, spid in enumerate(spids):
            delay = 0 if i == 0 else 200
            children.append(
                f'<p:par xmlns:p="{NS_P}">'
                f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
                f'fill="hold" nodeType="withEffect">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
            )
            nid += 2
        xml = (
            f'<p:timing xmlns:p="{NS_P}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        )
        sld.append(etree.fromstring(xml))

    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic_cover(sl, key):
        sl.shapes.add_picture(embed(key), 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks FCHIP | Africa CDC AES Fellowship | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    # 1 Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.5), SW, Inches(3.5), NAVY)
    tb(s, Inches(0.6), SH - Inches(3.2), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.7), Inches(12), Inches(0.55), META["doc_title"], size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(2.15), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(1.75), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.2), Inches(12), Inches(0.35), META["timeline"], size=12, color="D0E8E8")

    # 2 Fellowship at a glance
    s = prs.slides.add_slide(blank)
    band(s, "AES fellowship", "Two years — training, placement, fully funded")
    tb(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.5),
       "• 3 months core training in Addis Ababa\n"
       "• 21 months field placement with employer support\n"
       "• Epidemiology or Public Health Informatics track\n"
       "• Stipend, travel, insurance, hardware, software\n"
       "• African Union citizens under 35 in public health",
       size=17, color=SLATE)
    pic_fit(s, "training", Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5))
    footer(s, 2)

    # 3 Why surveillance for FCHIP
    s = prs.slides.add_slide(blank)
    band(s, "Surveillance science", "Platforms need trained epidemiologists — not only code")
    pic_fit(s, "gis", Inches(0.45), Inches(1.15), Inches(6.0), Inches(5.6))
    tb(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
       "Community signals from CHWs only become action with investigation discipline, "
       "trend analysis, and informatics pipelines.\n\n"
       "FairBanks lives inside primary care — the right place to grow Africa's epidemic readiness.",
       size=20, color=MUTED)
    footer(s, 3)

    # 4 Skills return
    s = prs.slides.add_slide(blank)
    band(s, "Return pathway", "Training → field → FCHIP → communities")
    pic_fit(s, "architecture", Inches(0.4), Inches(1.1), Inches(7.0), Inches(5.7))
    tb(s, Inches(7.7), Inches(1.2), Inches(5.2), Inches(5.5),
       "Standardise outbreak workflows\nImprove CHW data quality\nCascade training to supervisors\n"
       "Validate FCHIP alerts with surveillance methods",
       size=19, color=MUTED)
    footer(s, 4)

    # 5 Win-win
    s = prs.slides.add_slide(blank)
    band(s, "Shared value", "Africa CDC ↔ FairBanks ↔ communities")
    rect(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "FairBanks gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[0][1], size=18, color=MUTED)
    rect(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4), "Africa CDC gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[1][1], size=18, color=MUTED)
    footer(s, 5)

    # 6 Use cases
    s = prs.slides.add_slide(blank)
    band(s, "Application areas", "Where fellowship skills meet FCHIP")
    for i, item in enumerate(SECTIONS[3][1]):
        y = Inches(1.25) + i * Inches(1.15)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.0), "FFFFFF", LINE)
        rect(s, Inches(0.5), y, Inches(0.12), Inches(1.0), TEAL)
        tb(s, Inches(0.75), y + Inches(0.15), Inches(11.5), Inches(0.7), item, size=18, color=SLATE)
    footer(s, 6)

    # 7 Timeline
    s = prs.slides.add_slide(blank)
    band(s, "24-month pathway", "Addis Ababa training → Uganda field placement")
    for i, (phase, detail) in enumerate(FELLOWSHIP_PATH):
        y = Inches(1.2) + i * Inches(1.2)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.05), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.12), Inches(2.5), Inches(0.35), phase, size=14, bold=True, color=TEAL)
        tb(s, Inches(3.3), y + Inches(0.12), Inches(9.0), Inches(0.8), detail, size=13, color=MUTED)
    footer(s, 7)

    # 8 Fit
    s = prs.slides.add_slide(blank)
    band(s, "Eligibility fit", "Nominee meets AES criteria through FairBanks employment")
    for i, (crit, resp) in enumerate(FIT_ROWS):
        y = Inches(1.2) + i * Inches(1.2)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.05), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.12), Inches(3.5), Inches(0.8), crit, size=13, bold=True, color=TEAL)
        tb(s, Inches(4.3), y + Inches(0.12), Inches(8.0), Inches(0.8), resp, size=13, color=MUTED)
    footer(s, 8)

    # 9 Traction
    s = prs.slides.add_slide(blank)
    band(s, "Field foundation", "Live community health ecosystem in Uganda")
    pic_fit(s, "outreach", Inches(0.45), Inches(1.15), Inches(5.8), Inches(5.6))
    tb(s, Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5),
       "• FairBanks Medical Centre + Community Reach\n"
       "• CHW/VHT networks in Kampala catchments\n"
       "• Maternal, child, chronic, and outreach programmes\n"
       "• FCHIP intelligence layer in development\n"
       "• Employer ready to host field placement",
       size=16, color=SLATE)
    footer(s, 9)

    # 10 Ask
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.8), SW, Inches(2.8), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.4), Inches(12), Inches(0.55),
       "Select our fellow — skills that return to FCHIP and Uganda's communities.", size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.5), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(0.9), Inches(12), Inches(0.3), META["deadline"], size=12, color="D0E8E8")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for _sl in prs.slides:
        _add_entrance_anims(_sl)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. africa-cdc document set in", OUT_DIR)
