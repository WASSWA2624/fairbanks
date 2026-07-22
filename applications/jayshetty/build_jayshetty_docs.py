#!/usr/bin/env python3
"""Build the Jay Shetty strategic partnership pack for FairBanks.

Creates one synchronized set:
  documents/jayshetty_word.docx
  documents/jayshetty_pdf.pdf
  documents/jayshetty_ppt.pptx

Send-ready email lives in partnership_email.md (not in Word/PDF).
Narrative: strategic partnership brief. Three pillars + FCHIP. Human tone; no AI buzzwords.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path


HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
ASSETS = REPO / "assets"
OUT = HERE / "documents"
DOCX = OUT / "jayshetty_word.docx"
PDF = OUT / "jayshetty_pdf.pdf"
PPTX = OUT / "jayshetty_ppt.pptx"
TMP = REPO / "tmp" / "jayshetty"

NAVY = "0A1F2E"
TEAL = "0D6E6E"
GREEN = "2D7A55"
ORANGE = "C45C26"
GOLD = "D99A2B"
CREAM = "F7F5F0"
PALE_TEAL = "E8F3F2"
PALE_ORANGE = "FBEDE6"
PALE_GREEN = "E9F2EC"
SLATE = "1E2F38"
MUTED = "52636C"
LINE = "CED9D8"
WHITE = "FFFFFF"

SLOGAN = "Your health, our mission."
PROGRAMME = "Strategic Partnership Brief"
TITLE = "Partnership proposal"
SUBTITLE = "Prepared for the Jay Shetty partnerships and impact team"
ORG = "FairBanks Medical Centre"
CONTACT_NAME = "Racheal Nabukeera"
CONTACT_TITLE = "Managing Director"
WEBSITE = "https://fairbanksmedicalcentre.org"
EMAIL = "info@fairbanksmedicalcentre.org"
PHONE = "+256 777 462 398"
LOCATION = "Kampala, Uganda"
EMAIL_MD = HERE / "partnership_email.md"
COVER_HEADLINE = "Community health rooted in Uganda"
COVER_SUMMARY = (
    "FairBanks Medical Centre is a community health organisation in Kampala. "
    "We combine clinical care, Community Reach, and FCHIP - our community health intelligence "
    "platform. We are developing an FCHIP MVP to help families get care earlier and closer to "
    "home. This brief invites a conversation with the Jay Shetty ecosystem about shared goals "
    "around wellbeing, compassion, and practical impact."
)
LOGO = ASSETS / "fairbanks_logo.jpeg"

PHOTOS = {
    "cover": "outreach_facilitator_canopy_01.jpg",
    "logo": "fairbanks_logo.jpeg",
    "facility": "facility_exterior_entrance_01.jpg",
    "facility_branded": "facility_exterior_branded_entrance_01.jpeg",
    "facility_street": "facility_exterior_street_view_01.jpg",
    "pharmacy": "pharmacy_storefront_01.jpeg",
    "pharmacy_alt": "pharmacy_exterior_01.jpg",
    "outreach": "outreach_bp_screening.jpeg",
    "outreach_camp": "outreach_medical_camp_01.jpg",
    "community": "outreach_audience_full_group_01.jpg",
    "audience": "outreach_audience_attentive_01.jpg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "dashboard": "dashboard_demo.png",
    "architecture": "data_flow_iso_labeled.png",
    "gis": "gis_hotspots.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "gericare": "gericare_wheelchair_assist_02.jpg",
    "doctor_hands": "clinic_doctor_patient_hands_01.jpg",
    "lab": "clinic_lab_fingerprick_01.jpg",
    "compassion": "clinic_consult_compassion_01.jpg",
    "reception": "reception_wheelchair_checkin_01.jpg",
    "reception_staff": "reception_staff_documents_01.jpg",
    "team": "staff_team_reception.jpeg",
    "training": "indoor_training_audience_02.jpg",
    "mission": "reception_mission_wall.jpeg",
    "mothers": "waiting_room_mothers_01.jpeg",
    "canopy": "outreach_facilitator_group_01.jpg",
    "close": "cover_hero_cinematic.jpg",
}

CONCEPT = REPO / ".cursor" / "concept_simple.jpeg"  # summarized Community Reach journey
CONCEPT_FULL = REPO / ".cursor" / "concept_improved.jpeg"  # full partner model (archive / optional)


def photo(key: str) -> Path:
    path = ASSETS / PHOTOS[key]
    if not path.exists():
        raise FileNotFoundError(path)
    return path


def _pil_size(path: Path):
    from PIL import Image as PILImage

    with PILImage.open(path) as im:
        return im.size


def fit_width(path: Path, max_width_in: float, max_height_in: float = 3.4) -> float:
    iw, ih = _pil_size(path)
    aspect = ih / float(iw)
    w = max_width_in
    h = w * aspect
    if h > max_height_in:
        w = max_height_in / aspect
    return w


def concept_print_width(path: Path) -> Path:
    """Crop empty side gutters so the Community Reach model fills printable width."""
    from PIL import Image as PILImage

    cache = TMP / "balanced"
    cache.mkdir(parents=True, exist_ok=True)
    # Bump suffix when crop logic changes.
    out = cache / f"{path.stem}_print_width_v4.jpg"
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out

    with PILImage.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        px = im.load()

        def is_diagram(rgb) -> bool:
            r, g, b = rgb
            chroma = max(r, g, b) - min(r, g, b)
            # Coloured step cards, icons, arrows - not cream paper / white padding.
            return chroma > 28 and (r < 245 or g < 245 or b < 240)

        lefts, rights = [], []
        # Focus on the cascade body (skip extreme header/footer noise).
        for y in range(int(h * 0.14), int(h * 0.82), 8):
            xs = [x for x in range(0, w, 2) if is_diagram(px[x, y])]
            if len(xs) >= 8:
                lefts.append(xs[0])
                rights.append(xs[-1])
        if not lefts:
            im.save(out, format="JPEG", quality=92, optimize=True)
            return out

        lefts.sort()
        rights.sort()
        # Percentiles ignore rare edge flares that would leave wide cream gutters.
        left = max(0, lefts[len(lefts) // 5] - 18)
        right = min(w, rights[(4 * len(rights)) // 5] + 18)

        # Trim the solid black footer strip baked into concept_simple.jpeg.
        top = 0
        bottom = h
        for y in range(h - 1, max(h // 2, 0), -1):
            row_black = 0
            samples = 0
            for x in range(left, right, 12):
                samples += 1
                if sum(px[x, y]) < 40:
                    row_black += 1
            if samples and row_black / samples > 0.85:
                bottom = y
            else:
                break
        bottom = max(bottom, top + 100)

        cropped = im.crop((left, top, right, bottom))
        cropped.save(out, format="JPEG", quality=92, optimize=True)
    return out


def concept_print_slices(
    path: Path,
    *,
    width_in: float = 6.85,
    page_height_in: float = 9.6,
) -> list[Path]:
    """
    Full printable-width slices cut on cream gutters between steps.
    Keeps each page's diagram fully readable (no mid-box cut-off).
    """
    from PIL import Image as PILImage

    tight = concept_print_width(path)
    cache = TMP / "balanced"
    cache.mkdir(parents=True, exist_ok=True)

    with PILImage.open(tight) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        px = im.load()

        def row_cream(y: int) -> float:
            cream = n = 0
            for x in range(int(iw * 0.15), int(iw * 0.75), 4):
                r, g, b = px[x, y]
                n += 1
                if r >= 235 and g >= 230 and b >= 220:
                    cream += 1
            return cream / max(n, 1)

        # Prefer natural breaks in cream space between cascade steps.
        bands = [y for y in range(0, ih, 20) if row_cream(y) > 0.85]
        midpoints: list[int] = []
        start = prev = None
        for y in bands:
            if start is None:
                start = prev = y
            elif y - prev <= 40:
                prev = y
            else:
                midpoints.append((start + prev) // 2)
                start = prev = y
        if start is not None:
            midpoints.append((start + prev) // 2)

        max_slice_h = max(1, int(round(iw * (page_height_in / width_in))))
        cuts: list[int] = []
        y = 0
        while y + max_slice_h < ih:
            window = [m for m in midpoints if y + int(max_slice_h * 0.45) <= m <= y + max_slice_h - 40]
            if window:
                # Cut at the cream band closest to a comfortable page fill.
                target = y + int(max_slice_h * 0.88)
                cut = min(window, key=lambda m: abs(m - target))
            else:
                cut = y + max_slice_h
            cuts.append(cut)
            y = cut
        cuts.append(ih)

        slices: list[Path] = []
        y0 = 0
        for idx, y1 in enumerate(cuts, start=1):
            box = (0, y0, iw, y1)
            out = cache / f"{path.stem}_slice_v4_{idx:02d}_w{int(width_in * 100)}.jpg"
            if not (out.exists() and out.stat().st_mtime >= tight.stat().st_mtime):
                im.crop(box).save(out, format="JPEG", quality=92, optimize=True)
            slices.append(out)
            y0 = y1
        return slices


def cover_crop(
    path: Path,
    *,
    frame_w: int = 1200,
    frame_h: int = 900,
    focus: str = "center",
) -> Path:
    """Center-crop (cover) an image into a fixed frame so paired photos match."""
    from PIL import Image as PILImage

    cache = TMP / "balanced"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{path.stem}_{frame_w}x{frame_h}_{focus}.jpg"
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out

    with PILImage.open(path) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        target_ratio = frame_w / float(frame_h)
        src_ratio = iw / float(ih)
        if src_ratio > target_ratio:
            # too wide - crop sides
            new_w = int(ih * target_ratio)
            left = (iw - new_w) // 2
            box = (left, 0, left + new_w, ih)
        else:
            # too tall - crop top/bottom (bias slightly upward for faces)
            new_h = int(iw / target_ratio)
            if focus == "top":
                top = int((ih - new_h) * 0.18)
            else:
                top = (ih - new_h) // 2
            top = max(0, min(top, ih - new_h))
            box = (0, top, iw, top + new_h)
        cropped = im.crop(box).resize((frame_w, frame_h), PILImage.Resampling.LANCZOS)
        cropped.save(out, format="JPEG", quality=88, optimize=True)
    return out


def build_docx() -> None:
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement, parse_xml
    from docx.oxml.ns import nsdecls, qn
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()
    sec = doc.sections[0]
    sec.page_width, sec.page_height = Inches(8.27), Inches(11.69)
    sec.top_margin = Inches(0.6)
    sec.bottom_margin = Inches(0.7)
    sec.left_margin = sec.right_margin = Inches(0.7)

    styles = doc.styles
    styles["Normal"].font.name = "Calibri"
    styles["Normal"].font.size = Pt(11)
    styles["Normal"].font.color.rgb = RGBColor.from_string(SLATE)
    styles["Normal"].paragraph_format.space_after = Pt(8)
    styles["Normal"].paragraph_format.line_spacing = 1.15

    for name, size, color in (
        ("Title", 26, NAVY),
        ("Heading 1", 16, NAVY),
        ("Heading 2", 13, TEAL),
    ):
        style = styles[name]
        style.font.name = "Calibri"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)

    def shade_cell(cell, hex_color: str) -> None:
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), hex_color)
        shd.set(qn("w:val"), "clear")
        tc_pr.append(shd)

    def add_para(text, *, style=None, bold=False, size=11, color=SLATE, align=None, space_after=8):
        p = doc.add_paragraph(style=style) if style else doc.add_paragraph()
        if align is not None:
            p.alignment = align
        p.paragraph_format.space_after = Pt(space_after)
        run = p.add_run(text)
        run.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.name = "Calibri"
        return p

    def add_bullets(items):
        for item in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(4)
            run = p.add_run(item)
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor.from_string(SLATE)
            run.font.name = "Calibri"

    def set_cell_border(cell, color=LINE, sz="4") -> None:
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        borders = parse_xml(
            f'<w:tcBorders {nsdecls("w")}>'
            f'<w:top w:val="nil"/>'
            f'<w:left w:val="nil"/>'
            f'<w:bottom w:val="nil"/>'
            f'<w:right w:val="nil"/>'
            f"</w:tcBorders>"
        )
        tc_pr.append(borders)

    def add_caption(text: str) -> None:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after = Pt(10)
        run = p.add_run(text)
        run.italic = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor.from_string(MUTED)
        run.font.name = "Calibri"

    def add_image(
        path: Path,
        width: float = 6.6,
        height_cap: float = 3.2,
        caption: str = "",
        *,
        balance: bool = False,
    ):
        if balance:
            # Wide banner frame so single images also feel even.
            path = cover_crop(path, frame_w=1600, frame_h=900, focus="top")
            display_w = width
        else:
            display_w = fit_width(path, width, height_cap)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after = Pt(2 if caption else 10)
        run = p.add_run()
        if balance:
            run.add_picture(str(path), width=Inches(display_w), height=Inches(display_w * 900 / 1600))
        else:
            run.add_picture(str(path), width=Inches(display_w))
        if caption:
            add_caption(caption)

    def add_photo_row(
        left_key: str,
        right_key: str,
        left_caption: str,
        right_caption: str,
        *,
        width: float = 3.25,
        height: float = 2.44,
        focus: str = "top",
    ) -> None:
        """Two related photos in equal frames, each with a caption kept in-cell."""
        # 4:3 frame keeps landscape and portrait sources visually even.
        frame_w, frame_h = 1200, 900
        left_path = cover_crop(photo(left_key), frame_w=frame_w, frame_h=frame_h, focus=focus)
        right_path = cover_crop(photo(right_key), frame_w=frame_w, frame_h=frame_h, focus=focus)

        table = doc.add_table(rows=1, cols=2)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.autofit = False
        # Keep the whole photo pair on one page (no orphan captions).
        tr = table.rows[0]._tr
        tr_pr = tr.get_or_add_trPr()
        cant = OxmlElement("w:cantSplit")
        tr_pr.append(cant)
        col_w = Inches(width + 0.12)
        for col, (path, cap) in enumerate(
            [(left_path, left_caption), (right_path, right_caption)]
        ):
            cell = table.rows[0].cells[col]
            cell.width = col_w
            set_cell_border(cell)
            cell.text = ""
            # Image
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after = Pt(2)
            p.paragraph_format.keep_with_next = True
            p.add_run().add_picture(str(path), width=Inches(width), height=Inches(height))
            # Caption stays with the image (same cell = no page-split orphan)
            cp = cell.add_paragraph()
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cp.paragraph_format.space_before = Pt(2)
            cp.paragraph_format.space_after = Pt(6)
            run = cp.add_run(cap)
            run.italic = True
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor.from_string(MUTED)
            run.font.name = "Calibri"

        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_after = Pt(4)

# ---- Cover ----
    # FairBanks brand header: logo + organisation + slogan
    brand = doc.add_table(rows=1, cols=2)
    brand.autofit = False
    brand.allow_autofit = False
    brand.columns[0].width = Inches(1.2)
    brand.columns[1].width = Inches(5.65)
    left_c, right_c = brand.rows[0].cells
    shade_cell(left_c, NAVY)
    shade_cell(right_c, NAVY)
    set_cell_border(left_c, NAVY)
    set_cell_border(right_c, NAVY)
    left_c.text = ""
    lp = left_c.paragraphs[0]
    lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    lp.paragraph_format.space_before = Pt(6)
    lp.paragraph_format.space_after = Pt(6)
    if LOGO.exists():
        run = lp.add_run()
        run.add_picture(str(LOGO), height=Inches(0.62))
    right_c.text = ""
    rp = right_c.paragraphs[0]
    rp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    rp.paragraph_format.space_before = Pt(8)
    rp.paragraph_format.space_after = Pt(0)
    r1 = rp.add_run(ORG)
    r1.bold = True
    r1.font.size = Pt(13)
    r1.font.color.rgb = RGBColor.from_string(GOLD)
    r1.font.name = "Calibri"
    rp2 = right_c.add_paragraph()
    rp2.paragraph_format.space_before = Pt(2)
    rp2.paragraph_format.space_after = Pt(0)
    r2 = rp2.add_run(SLOGAN)
    r2.italic = True
    r2.font.size = Pt(10)
    r2.font.color.rgb = RGBColor.from_string(WHITE)
    r2.font.name = "Calibri"
    rp3 = right_c.add_paragraph()
    rp3.paragraph_format.space_before = Pt(2)
    rp3.paragraph_format.space_after = Pt(6)
    r3 = rp3.add_run(PROGRAMME)
    r3.bold = True
    r3.font.size = Pt(9)
    r3.font.color.rgb = RGBColor.from_string(GOLD)
    r3.font.name = "Calibri"

    # Contact + location strip under branding
    contact_bar = doc.add_table(rows=1, cols=1)
    contact_bar.autofit = True
    cc = contact_bar.rows[0].cells[0]
    shade_cell(cc, PALE_TEAL)
    set_cell_border(cc, TEAL)
    cc.text = ""
    cp = cc.paragraphs[0]
    cp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    cp.paragraph_format.space_before = Pt(4)
    cp.paragraph_format.space_after = Pt(2)
    cr = cp.add_run(
        f"{CONTACT_NAME}  ·  {CONTACT_TITLE}\n"
        f"{LOCATION}  ·  {PHONE}  ·  {EMAIL}\n"
        f"{WEBSITE}"
    )
    cr.font.size = Pt(9)
    cr.font.color.rgb = RGBColor.from_string(NAVY)
    cr.font.name = "Calibri"

    doc.add_paragraph()
    add_para(TITLE, bold=True, size=22, color=NAVY, space_after=4)
    add_para(COVER_HEADLINE, bold=True, size=15, color=TEAL, space_after=4)
    add_para(SUBTITLE, size=12, color=MUTED, space_after=8)

    add_photo_row(
        "cover",
        "facility_branded",
        "Community Reach outreach session",
        "FairBanks Medical Centre, Kampala",
        width=3.2,
        height=2.0,
    )

    add_para(COVER_SUMMARY, size=11, color=SLATE, space_after=10)

    add_para("Why this partnership matters", bold=True, size=12, color=ORANGE, space_after=6)
    value = doc.add_table(rows=1, cols=3)
    value.alignment = WD_TABLE_ALIGNMENT.CENTER
    value.autofit = False
    value_cols = [
        (
            "For Jay Shetty's team",
            [
                "A grounded story from a working clinic and community programmes in Uganda",
                "A practical partner for wellbeing themes - care, prevention, and family health",
                "A clear route from message to community work you can point to",
            ],
            PALE_ORANGE,
        ),
        (
            "For FairBanks",
            [
                "Help telling our story to a wider audience",
                "Advice and introductions within the impact and wellbeing space",
                "Support as we develop the FCHIP MVP",
            ],
            PALE_TEAL,
        ),
        (
            "For communities we serve",
            [
                "Stronger outreach and CHW / VHT support",
                "Care nearer home, including CHIS where families need it",
                "Earlier help when health risks start to rise",
            ],
            PALE_GREEN,
        ),
    ]
    for i, (title, points, fill) in enumerate(value_cols):
        cell = value.rows[0].cells[i]
        cell.width = Inches(2.28)
        shade_cell(cell, fill)
        set_cell_border(cell, LINE)
        cell.text = ""
        tp = cell.paragraphs[0]
        tp.paragraph_format.space_after = Pt(4)
        tr = tp.add_run(title)
        tr.bold = True
        tr.font.size = Pt(10)
        tr.font.color.rgb = RGBColor.from_string(NAVY)
        tr.font.name = "Calibri"
        for point in points:
            pp = cell.add_paragraph()
            pp.paragraph_format.space_after = Pt(3)
            pr = pp.add_run(f"- {point}")
            pr.font.size = Pt(8)
            pr.font.color.rgb = RGBColor.from_string(SLATE)
            pr.font.name = "Calibri"

    doc.add_paragraph()

    ask = doc.add_table(rows=1, cols=1)
    ac = ask.rows[0].cells[0]
    shade_cell(ac, TEAL)
    ac.text = ""
    ap = ac.paragraphs[0]
    ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    ar = ap.add_run(
        "Request: a short introductory call (about 10 minutes) with your partnerships or impact team"
    )
    ar.bold = True
    ar.font.size = Pt(10)
    ar.font.color.rgb = RGBColor.from_string(WHITE)
    ar.font.name = "Calibri"

    # ---- Project summary ----
    doc.add_page_break()
    add_para("Project summary", style="Heading 1", bold=True, size=18, color=NAVY)
    add_para(
        "FairBanks is asking to explore a partnership with the Jay Shetty ecosystem. "
        "We already run a medical centre and Community Reach programmes in Kampala peri-urban "
        "communities, and we are developing an FCHIP MVP - our community health intelligence "
        "platform. We are looking for collaborators who care about wellbeing in practice - "
        "not only as a message, but as better access to care.",
        size=11,
        color=SLATE,
        space_after=10,
    )

    add_para("What FairBanks is building", style="Heading 2", bold=True, size=13, color=TEAL)
    add_bullets(
        [
            "FairBanks Medical Centre - clinical care in Kampala",
            "Community Reach - CHWs/VHTs, outreach, school health, CHIS, and livelihoods",
            "FCHIP MVP in development - tools that help teams see risk earlier using community, facility, GIS, and climate data",
            "A partnership focused on shared purpose: storytelling, field delivery, and prevention",
        ]
    )

    add_para("How we see the collaboration", style="Heading 2", bold=True, size=13, color=TEAL)
    summary = doc.add_table(rows=5, cols=2)
    summary.style = "Table Grid"
    summary_rows = [
        ("Area", "Detail"),
        ("Jay Shetty ecosystem", "Platform, storytelling reach, and connections in wellbeing and impact"),
        ("FairBanks", "Operating clinic, Community Reach work, FCHIP MVP in development, and local relationships"),
        ("Shared aim", "Stronger community health outcomes linked to compassion and practical care"),
        ("Next step", "A brief introductory conversation with partnerships or impact"),
    ]
    for i, (a, b) in enumerate(summary_rows):
        summary.rows[i].cells[0].text = a
        summary.rows[i].cells[1].text = b
        if i == 0:
            shade_cell(summary.rows[i].cells[0], TEAL)
            shade_cell(summary.rows[i].cells[1], TEAL)
        else:
            shade_cell(summary.rows[i].cells[0], PALE_TEAL)
        for cell in summary.rows[i].cells:
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.name = "Calibri"
                    if i == 0:
                        r.bold = True
                        r.font.color.rgb = RGBColor.from_string(WHITE)
                    else:
                        r.font.color.rgb = RGBColor.from_string(
                            NAVY if cell is summary.rows[i].cells[0] else SLATE
                        )
                        if cell is summary.rows[i].cells[0]:
                            r.bold = True

    doc.add_paragraph()
    add_photo_row(
        "doctor_hands",
        "audience",
        "Care at the medical centre",
        "Community members at an outreach session",
        width=3.3,
        height=2.2,
    )

    # ---- Vision ----
    add_para("1. Who FairBanks is", style="Heading 1", bold=True, size=16, color=NAVY)
    add_photo_row(
        "facility_street",
        "pharmacy",
        "Medical Centre - clinical anchor in Kampala",
        "Pharmacy - medicines beside the clinic",
    )
    add_para(
        "FairBanks is a community health organisation in Kampala. We provide primary care, "
        "run Community Reach programmes, and build FCHIP so families can get help earlier and "
        "closer to home. We believe good care should not depend on where someone lives or what "
        "they earn.",
    )
    add_para(
        "Our medical centre is where clinical care, outreach, and FCHIP meet. We test and "
        "improve ideas with the families we already serve before we ask others to use them.",
    )

    add_para("Our three parts of work", style="Heading 2", bold=True, size=13, color=TEAL)
    pillars = [
        (
            "FairBanks Medical Centre",
            "Quality family and community primary healthcare, diagnostics, pharmacy, and referrals.",
        ),
        (
            "FairBanks Community Reach",
            "Outreach, health education, screening, maternal and child health, Gericare, school health, CHIS, and livelihoods.",
        ),
        (
            "FCHIP (Community Health Intelligence Platform)",
            "An MVP in development - a platform that brings together community, facility, map, and climate signals so teams can act earlier.",
        ),
    ]
    for title, body in pillars:
        add_para(title, bold=True, size=11, color=NAVY, space_after=2)
        add_para(body, size=11, color=SLATE, space_after=8)

    add_photo_row(
        "reception",
        "gericare",
        "Reception and triage with dignity",
        "Gericare - support for older patients",
    )

    # ---- How we work ----
    add_para("2. How FairBanks works - Community Reach", style="Heading 1", bold=True, size=16, color=NAVY)
    add_para(
        "FairBanks operates a six-step cascade with a continuous feedback loop. FCHIP sits on "
        "the Data & Feedback layer - the part of the model that helps us learn from the field and improve how we serve.",
    )
    if CONCEPT.exists():
        # Directly under the section intro (no forced page break before/after).
        tight = concept_print_width(CONCEPT)
        iw, ih = _pil_size(tight)
        # Fit remaining page height after heading + intro; Word will flow if needed.
        max_h = 8.5
        fit_w = min(6.85, max_h * (iw / float(ih)))
        add_image(
            tight,
            fit_w,
            max_h + 0.05,
            caption="Community Reach - how the work connects",
        )

    add_bullets(
        [
            "Community members identify needs and own solutions",
            "CHWs / VHTs bridge homes, schools, and care",
            "Community Reach programmes deliver education, screening, and home visits",
            "FairBanks Medical Centre anchors clinical care and quality",
            "Research, partnerships, and skills strengthen evidence and capacity",
            "Economic empowerment (including CHIS and IGAs) supports resilient families",
        ]
    )
    add_photo_row(
        "audience",
        "community",
        "Listening at a community outreach session",
        "Community Reach gathering in full",
    )

    # ---- Problem / FCHIP ----
    add_para("3. The problem FCHIP helps address", style="Heading 1", bold=True, size=16, color=NAVY)
    add_photo_row(
        "outreach",
        "outreach_camp",
        "Blood-pressure screening in the community",
        "Medical camp care beyond clinic walls",
    )
    add_para(
        "Primary healthcare in underserved communities is still largely reactive. Facilities "
        "often see people only after illness appears. Data from CHWs, schools, outreach, clinics, "
        "and pharmacies stays fragmented. Climate signals rarely join the picture. Clinical records "
        "often remain locked inside facility EMR/HMS systems - so the continuous feedback loop "
        "cannot learn fast enough to prevent outbreaks, maternal complications, or stock-outs.",
    )
    add_para(
        "FCHIP is FairBanks' technology platform. It connects community members, CHWs/VHTs, "
        "schools, communities, hospitals, and clinics; safely exposes data APIs to existing EMR/HMS "
        "systems for real-time clinical ingest; and fuses signals with GIS mapping and climate APIs "
        "to surface explainable early warnings and dashboards.",
        space_after=8,
    )
    add_image(
        photo("architecture"),
        6.4,
        2.9,
        caption="FCHIP flow: community and facility capture - intelligence - action for facilities and districts",
    )
    add_para("What FCHIP includes", style="Heading 2", bold=True, size=13, color=TEAL)
    add_bullets(
        [
            "Artificial intelligence and machine learning for risk scoring and early warning",
            "GIS mapping of disease and resource geography",
            "Climate API integration for climate-sensitive health risk",
            "Secure EMR/HMS data APIs (no rip-and-replace for clinics)",
            "Offline-capable mobile capture for CHWs/VHTs",
            "Cloud sync and analytics dashboards for facilities and partners",
        ]
    )
    add_photo_row(
        "mobile",
        "dashboard",
        "Mobile tools for CHWs and VHTs",
        "Dashboards for timely decisions",
    )

    # ---- Traction / roadmap ----
    add_para("4. Where we are now", style="Heading 1", bold=True, size=16, color=NAVY)
    add_photo_row(
        "training",
        "maternal",
        "Community health education in progress",
        "Maternal and family health support",
    )
    add_para("What is already in place", style="Heading 2", bold=True, size=13, color=TEAL)
    add_bullets(
        [
            "Operational FairBanks Medical Centre in Kampala",
            "Active Community Reach across Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and surrounding communities",
            "CHW/VHT engagement, maternal and child health, Gericare, chronic-disease screening, school and corporate health",
            "CHIS and livelihood pathways linked to affordable access",
            "FCHIP MVP currently in development for structured field use and evidence building",
        ]
    )
    add_photo_row(
        "doctor_hands",
        "lab",
        "Listening before treating",
        "Point-of-care testing close to patients",
    )

    add_para("Next steps on the roadmap", style="Heading 2", bold=True, size=13, color=TEAL)
    roadmap = doc.add_table(rows=7, cols=2)
    roadmap.style = "Table Grid"
    rows = [
        ("Stage", "Status"),
        ("Community healthcare model established", "Completed"),
        ("FairBanks Medical Centre operational", "Completed"),
        ("Community Reach platform operational", "Completed"),
        ("FCHIP MVP", "In development"),
        ("Structured pilot evidence & partner onboarding", "Next phase"),
        ("District and regional expansion", "Future vision"),
    ]
    for i, (a, b) in enumerate(rows):
        roadmap.rows[i].cells[0].text = a
        roadmap.rows[i].cells[1].text = b
        if i == 0:
            shade_cell(roadmap.rows[i].cells[0], TEAL)
            shade_cell(roadmap.rows[i].cells[1], TEAL)
            for cell in roadmap.rows[i].cells:
                for p in cell.paragraphs:
                    for r in p.runs:
                        r.bold = True
                        r.font.color.rgb = RGBColor.from_string(WHITE)
                        r.font.size = Pt(10)
        else:
            for cell in roadmap.rows[i].cells:
                for p in cell.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)
                        r.font.name = "Calibri"
                        r.font.color.rgb = RGBColor.from_string(SLATE)

    doc.add_paragraph()
    add_para(
        "Longer term, we want FairBanks to be known across Africa for community health "
        "intelligence that sits beside real clinics and community programmes - not as software "
        "alone.",
    )

    # ---- Partnership ----
    add_para("5. Possible ways to work together", style="Heading 1", bold=True, size=16, color=NAVY)
    add_photo_row(
        "compassion",
        "mothers",
        "Compassion in the consulting room",
        "Mothers and families in the clinic",
    )
    add_para(
        "Jay Shetty's public work centres on wellbeing, compassion, and stronger communities. "
        "FairBanks works on the same themes through clinics, outreach, and tools that help "
        "people get care earlier and closer to home.",
    )
    add_para("Areas we would like to explore", style="Heading 2", bold=True, size=13, color=TEAL)
    add_bullets(
        [
            "Community health and wellness programmes",
            "Preventive healthcare and health promotion",
            "Mental health and holistic wellbeing initiatives",
            "Artificial intelligence for social impact",
            "Digital health innovation and research exchange",
            "Global advocacy and storytelling",
            "Impact investment and philanthropic introductions",
            "Social enterprise growth and mentorship",
        ]
    )
    add_para("What FairBanks offers", style="Heading 2", bold=True, size=13, color=TEAL)
    add_bullets(
        [
            "An operating community medical facility in Kampala",
            "Community Reach work shaped by local needs and relationships",
            "FCHIP - an MVP we are developing for careful field validation",
            "Healthcare leadership and trusted community relationships",
            "A team focused on careful, lasting community health work",
        ]
    )

    add_para("Closing", style="Heading 2", bold=True, size=13, color=TEAL)
    add_para(
        "We would welcome a conversation about how the Jay Shetty ecosystem and FairBanks can "
        "work together on community health and wellbeing. As we develop the FCHIP MVP, we are looking "
        "for partners who can help share the story, offer advice, and open useful introductions.",
    )
    add_para(
        "If this feels aligned, we would be glad to start with a short call and see where a "
        "practical collaboration might fit.",
        bold=True,
        color=NAVY,
    )
    add_para(
        "Next step: a brief 10-minute introductory conversation with your partnerships or impact team.",
        size=11,
        color=ORANGE,
        bold=True,
    )

    add_photo_row(
        "team",
        "reception_staff",
        "FairBanks team ready to partner",
        "Reception documenting care with care",
    )
    add_para(f"{CONTACT_NAME}  |  {CONTACT_TITLE}", bold=True, size=11, color=NAVY, space_after=2)
    add_para(f"{ORG}", size=10, color=MUTED, space_after=2)
    add_para(f"{LOCATION}  |  {PHONE}  |  {EMAIL}", size=10, color=MUTED, space_after=2)
    add_para(WEBSITE, size=10, color=MUTED, space_after=12)

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(DOCX)
    print(f"DOCX: {DOCX}")


def convert_pdf() -> None:
    import shutil

    import win32com.client

    TMP.mkdir(parents=True, exist_ok=True)
    tmp_pdf = TMP / "jayshetty_pdf_build.pdf"
    if tmp_pdf.exists():
        tmp_pdf.unlink()

    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    document = None
    try:
        document = word.Documents.Open(str(DOCX.resolve()), ReadOnly=True)
        document.ExportAsFixedFormat(str(tmp_pdf.resolve()), 17)
    finally:
        if document is not None:
            document.Close(False)
        word.Quit()

    if PDF.exists():
        try:
            PDF.unlink()
        except PermissionError as exc:
            raise RuntimeError(f"Close {PDF.name} in any PDF viewer, then re-run.") from exc
    shutil.move(str(tmp_pdf), str(PDF))
    print(f"PDF:  {PDF}")


def _transition(slide) -> None:
    from lxml import etree

    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sld = slide._element
    for old in sld.findall(f"{{{ns}}}transition"):
        sld.remove(old)
    transition = etree.SubElement(sld, f"{{{ns}}}transition")
    transition.set("spd", "med")
    etree.SubElement(transition, f"{{{ns}}}fade")


def _add_entrance_anims(slide) -> None:
    from lxml import etree

    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    spids = []
    for shape in slide.shapes:
        try:
            st = int(shape.shape_type) if shape.shape_type is not None else -1
        except Exception:
            st = -1
        has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
        if st == 13 or has_text:
            spids.append(str(shape.shape_id))
        if len(spids) >= 3:
            break
    if not spids:
        return
    sld = slide._element
    for old in sld.findall(f"{{{NS_P}}}timing"):
        sld.remove(old)
    children = []
    nid = 3
    for i, spid in enumerate(spids):
        delay = 0 if i == 0 else 200
        children.append(
            f'<p:par xmlns:p="{NS_P}">'
            f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
            f'fill="hold" nodeType="withEffect">'
            f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:animEffect transition="in" filter="fade">'
            f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
        )
        nid += 2
    xml = (
        f'<p:timing xmlns:p="{NS_P}">'
        f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    sld.append(etree.fromstring(xml))


def build_pptx() -> None:
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    slides_for_anim = []

    def rgb(value):
        return RGBColor.from_string(value)

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line or fill)
        return shape

    def text(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=18,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.LEFT,
        font="Calibri",
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        return box

    def bullets(slide, items, x, y, w, h, size=15, color=SLATE):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "- " + item
            p.space_after = Pt(8)
            p.font.name = "Calibri"
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
        return box

    def crop(slide, path, x, y, w, h):
        with PILImage.open(path) as im:
            iw, ih = im.size
        pic = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
        )
        frame_ratio, image_ratio = w / h, iw / ih
        if image_ratio > frame_ratio:
            amount = (1 - frame_ratio / image_ratio) / 2
            pic.crop_left = pic.crop_right = amount
        else:
            amount = (1 - image_ratio / frame_ratio) / 2
            pic.crop_top = pic.crop_bottom = amount
        return pic

    def fit(slide, path, x, y, max_w, max_h):
        with PILImage.open(path) as im:
            iw, ih = im.size
        image_ratio = iw / ih
        frame_ratio = max_w / max_h
        if image_ratio > frame_ratio:
            pw, ph = max_w, max_w / image_ratio
            px, py = x, y + (max_h - ph) / 2
        else:
            ph, pw = max_h, max_h * image_ratio
            px, py = x + (max_w - pw) / 2, y
        return slide.shapes.add_picture(
            str(path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph)
        )

    def slide():
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, 13.333, 7.5, CREAM)
        _transition(s)
        slides_for_anim.append(s)
        return s

    def band(s, kicker, title_value, subtitle=""):
        rect(s, 0, 0, 13.333, 0.12, TEAL)
        text(s, kicker.upper(), 0.55, 0.28, 6.5, 0.28, 10, ORANGE, True)
        text(s, title_value, 0.55, 0.55, 12.1, 0.55, 26, NAVY, True)
        if subtitle:
            text(s, subtitle, 0.58, 1.15, 11.8, 0.32, 13, MUTED)

    def footer(s, number):
        text(s, f"FairBanks Medical Centre  |  {LOCATION}", 0.55, 7.18, 8.5, 0.18, 9, MUTED)
        text(s, f"{number:02}", 12.2, 7.16, 0.5, 0.18, 9, MUTED, align=PP_ALIGN.RIGHT)

    def photo_pair(s, left_key, right_key, left_cap, right_cap, x=0.55, y=1.7, w=6.0, h=4.5):
        """Two related photos in equal cropped frames, with matched caption bands."""
        gap = 0.22
        pw = (w - gap) / 2
        cap_h = 0.48
        ph = h - cap_h
        crop(s, photo(left_key), x, y, pw, ph)
        crop(s, photo(right_key), x + pw + gap, y, pw, ph)
        # Equal caption bands under both frames
        text(s, left_cap, x, y + ph + 0.06, pw, cap_h - 0.06, 11, MUTED, align=PP_ALIGN.CENTER)
        text(
            s,
            right_cap,
            x + pw + gap,
            y + ph + 0.06,
            pw,
            cap_h - 0.06,
            11,
            MUTED,
            align=PP_ALIGN.CENTER,
        )

    # 1 Cover
    s = slide()
    crop(s, photo("cover"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 7.6, 7.5, NAVY)
    try:
        s.shapes.add_picture(str(photo("logo")), Inches(0.65), Inches(0.32), height=Inches(0.55))
    except Exception:
        pass
    text(s, "STRATEGIC PARTNERSHIP BRIEF", 0.65, 1.0, 6.4, 0.28, 12, GOLD, True)
    text(s, "FairBanks", 0.65, 1.4, 6.5, 0.55, 34, WHITE, True)
    text(s, "Medical Centre", 0.68, 1.95, 6.2, 0.35, 16, GOLD, True)
    text(s, COVER_HEADLINE, 0.68, 2.45, 6.2, 0.35, 16, WHITE, True)
    text(s, SUBTITLE, 0.68, 2.9, 6.2, 0.7, 14, WHITE)
    text(s, SLOGAN, 0.68, 4.55, 5.5, 0.3, 14, GOLD, True)
    text(s, f"{CONTACT_NAME}  ·  {CONTACT_TITLE}", 0.68, 5.15, 6.2, 0.28, 13, WHITE, True)
    text(s, f"{LOCATION}  ·  {PHONE}", 0.68, 5.5, 6.2, 0.28, 12, WHITE)
    text(s, f"{EMAIL}  ·  {WEBSITE}", 0.68, 5.85, 6.4, 0.28, 12, WHITE)
    text(s, "Proposed next step: a short introductory call", 0.68, 6.45, 6.2, 0.3, 12, GOLD)

    # 2 Project summary
    s = slide()
    band(
        s,
        "Project summary",
        "Why a partnership with FairBanks makes sense",
        "Shared purpose in community health and wellbeing.",
    )
    cards = [
        (
            "For Jay Shetty's team",
            ORANGE,
            PALE_ORANGE,
            [
                "A real clinic and community story",
                "Wellbeing themes linked to real care",
                "A partner grounded in daily health work",
            ],
        ),
        (
            "For FairBanks",
            TEAL,
            PALE_TEAL,
            [
                "Help sharing our work more widely",
                "Advice and useful introductions",
                "Support as we develop the FCHIP MVP",
            ],
        ),
        (
            "For communities",
            GREEN,
            PALE_GREEN,
            [
                "Stronger prevention and CHW support",
                "Care closer to home, including CHIS",
                "Earlier action when risks rise",
            ],
        ),
    ]
    for i, (title, accent, fill, items) in enumerate(cards):
        x = 0.55 + i * 4.2
        rect(s, x, 1.85, 3.95, 4.7, WHITE, LINE, True)
        rect(s, x, 1.85, 3.95, 0.55, accent)
        text(s, title, x + 0.2, 1.95, 3.5, 0.35, 15, WHITE, True)
        for j, item in enumerate(items):
            y = 2.65 + j * 0.85
            rect(s, x + 0.2, y, 3.55, 0.7, fill, fill, True)
            text(s, item, x + 0.35, y + 0.18, 3.25, 0.4, 13, NAVY, True)
    footer(s, 2)

    # 3 Who we are
    s = slide()
    band(s, "Who we are", "Clinic, community programmes, and health intelligence", "Care + Community Reach + Intelligence")
    photo_pair(
        s,
        "facility_street",
        "pharmacy",
        "Medical Centre - clinical anchor",
        "Pharmacy - care beside the clinic",
        x=0.5,
        y=1.7,
        w=7.0,
        h=5.0,
    )
    bullets(
        s,
        [
            "Ugandan social enterprise in Kampala peri-urban communities",
            "Medical centre as the clinical anchor",
            "Community Reach for prevention, education, livelihoods",
            "FCHIP closes the Data & Feedback loop",
            "Uganda first - designed for African scale",
        ],
        7.8,
        2.0,
        5.0,
        4.5,
        14,
    )
    footer(s, 3)

    # 4 Three pillars
    s = slide()
    band(s, "The FairBanks model", "Three parts of the same organisation", SLOGAN)
    cards = [
        ("Medical Centre", "Family and community primary care, diagnostics, pharmacy, referrals", photo("facility_branded")),
        ("Community Reach", "CHWs/VHTs, outreach, school health, CHIS, livelihoods", photo("canopy")),
        ("FCHIP", "AI, GIS, climate fusion, secure EMR/HMS APIs, early warning", photo("dashboard")),
    ]
    for i, (title, body, img) in enumerate(cards):
        x = 0.55 + i * 4.2
        rect(s, x, 1.75, 3.95, 5.0, WHITE, LINE, True)
        crop(s, img, x + 0.15, 1.95, 3.65, 2.2)
        text(s, title, x + 0.25, 4.35, 3.4, 0.4, 18, NAVY, True)
        text(s, body, x + 0.25, 4.85, 3.4, 1.4, 13, MUTED)
    footer(s, 4)

    # 5 Cascade - summarized Community Reach model (concept_simple)
    s = slide()
    rect(s, 0, 0, 13.333, 0.12, TEAL)
    text(s, "HOW WE CARE", 0.45, 0.22, 4.0, 0.25, 11, ORANGE, True)
    text(s, "Community Reach - summarized model", 0.45, 0.45, 9.0, 0.4, 24, NAVY, True)
    model_img = concept_print_width(CONCEPT) if CONCEPT.exists() else CONCEPT
    iw, ih = _pil_size(model_img) if model_img.exists() else (4120, 8936)
    aspect = iw / float(ih)
    # Span nearly the full slide width; height follows aspect.
    cascade_w = 12.4
    cascade_h = cascade_w / aspect
    if cascade_h > 5.55:
        cascade_h = 5.55
        cascade_w = cascade_h * aspect
    cascade_x = (13.333 - cascade_w) / 2
    cascade_y = 0.95
    rect(s, cascade_x, cascade_y, cascade_w + 0.16, cascade_h + 0.16, WHITE, LINE, True)
    if model_img.exists():
        fit(s, model_img, cascade_x + 0.08, cascade_y + 0.08, cascade_w, cascade_h)
    text(
        s,
        "FCHIP supports learning from field data  ·  It starts with the community  ·  CHWs bridge care  ·  Families grow stronger",
        0.55,
        6.85,
        12.2,
        0.35,
        12,
        MUTED,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 5)

    # 6 Problem
    s = slide()
    band(s, "The gap", "Care is still mostly reactive", "Too often, care starts only after someone is already sick.")
    photo_pair(
        s,
        "outreach",
        "outreach_camp",
        "Community BP screening",
        "Medical camp consultations",
        x=0.5,
        y=1.7,
        w=6.8,
        h=5.0,
    )
    bullets(
        s,
        [
            "Facilities react after illness appears",
            "CHW, school, clinic, pharmacy signals stay siloed",
            "Climate risk rarely joins health decisions",
            "Clinical records often lock inside EMR/HMS",
            "The feedback loop cannot learn in real time",
        ],
        7.6,
        2.0,
        5.2,
        4.5,
        14,
    )
    footer(s, 6)

    # 7 FCHIP solution
    s = slide()
    band(
        s,
        "The intelligence layer",
        "FCHIP - FairBanks Community Health Intelligence Platform",
        "Helping teams act earlier, using community and climate signals.",
    )
    photo_pair(
        s,
        "mobile",
        "dashboard",
        "Mobile capture for CHWs / VHTs",
        "Dashboards for timely action",
        x=0.5,
        y=1.7,
        w=6.5,
        h=5.0,
    )
    for i, (a, b) in enumerate(
        [
            ("Capture", "Field, school, clinic signals"),
            ("Connect", "Secure EMR/HMS APIs"),
            ("Fuse", "GIS + climate + community"),
            ("Act", "Alerts, referrals, outreach"),
        ]
    ):
        y = 1.85 + i * 1.15
        rect(s, 7.3, y, 5.4, 1.0, WHITE, LINE, True)
        rect(s, 7.3, y, 0.14, 1.0, TEAL)
        text(s, a, 7.6, y + 0.18, 1.5, 0.3, 15, NAVY, True)
        text(s, b, 9.2, y + 0.18, 3.2, 0.55, 13, MUTED)
    footer(s, 7)

    # 8 Deep tech
    s = slide()
    band(s, "Technology stack", "Built for community health work in Uganda and beyond", "AI - ML - GIS - Climate - Secure APIs - Mobile - Cloud")
    crop(s, photo("architecture"), 6.4, 1.8, 6.3, 4.55)
    text(
        s,
        "FCHIP architecture: capture to intelligence to action",
        6.4,
        6.45,
        6.3,
        0.35,
        11,
        MUTED,
        align=PP_ALIGN.CENTER,
    )
    bullets(
        s,
        [
            "AI/ML risk scoring for outbreaks, maternal risk, NCDs",
            "GIS hotspot mapping by village and catchment",
            "Climate API fusion for rainfall, heat, extremes",
            "Secure EMR/HMS data APIs for real-time clinical ingest",
            "Offline mobile tools for CHWs and VHTs",
            "Facility and partner dashboards for action",
        ],
        0.65,
        1.95,
        5.4,
        4.8,
        14,
    )
    footer(s, 8)

    # 9 Traction
    s = slide()
    band(s, "Traction", "Work already underway in Kampala communities", "Clinic, outreach, and an FCHIP MVP in development.")
    photo_pair(
        s,
        "community",
        "gericare",
        "Community Reach gatherings",
        "Gericare - dignity in ageing",
        x=0.5,
        y=1.7,
        w=8.0,
        h=5.0,
    )
    bullets(
        s,
        [
            "Live medical centre",
            "Named Kampala peri-urban communities",
            "CHW/VHT, MCH, school health",
            "CHIS and livelihood pathways",
            "FCHIP MVP currently in development",
        ],
        8.8,
        2.0,
        4.0,
        4.5,
        13,
    )
    footer(s, 9)

    # 10 Roadmap
    s = slide()
    band(s, "Roadmap", "Near-term MVP development, then careful growth", "Develop the MVP first. Grow when the evidence is clear.")
    stages = [
        ("Done", "Model + centre", "Community Reach and medical centre live"),
        ("Now", "FCHIP MVP", "MVP in development for field use"),
        ("Next", "Validate", "Structured CHW/facility cohorts + proof cases"),
        ("Next", "Partner", "Storytelling, mentorship, impact introductions"),
        ("Future", "District", "Partner clinics and district structures"),
        ("Future", "Regional", "East Africa and pan-African pathways"),
    ]
    for i, (tag, title, body) in enumerate(stages):
        col, row = i % 3, i // 3
        x, y = 0.55 + col * 4.2, 1.85 + row * 2.35
        rect(s, x, y, 3.95, 2.05, WHITE, LINE, True)
        fill = TEAL if tag == "Done" else (ORANGE if tag in ("Next", "Now") else GREEN)
        rect(s, x, y, 3.95, 0.42, fill)
        text(s, tag.upper(), x + 0.2, y + 0.08, 2.0, 0.28, 11, WHITE, True)
        text(s, title, x + 0.2, y + 0.6, 3.5, 0.35, 16, NAVY, True)
        text(s, body, x + 0.2, y + 1.05, 3.5, 0.7, 13, MUTED)
    footer(s, 10)

    # 11 Why Jay Shetty
    s = slide()
    band(
        s,
        "Why Jay Shetty",
        "From wellbeing message to community action",
        "Wellbeing, compassion, and stronger communities",
    )
    photo_pair(
        s,
        "doctor_hands",
        "maternal",
        "Compassion in the consulting room",
        "Maternal and family health support",
        x=0.5,
        y=1.7,
        w=7.0,
        h=5.0,
    )
    bullets(
        s,
        [
            "Jay Shetty champions healthier, purposeful lives",
            "FairBanks turns vision into access and prevention",
            "FCHIP helps communities see risk earlier",
            "Global storytelling + local health delivery",
            "A practical partnership, not a one-way request",
        ],
        7.8,
        2.0,
        5.0,
        4.5,
        14,
    )
    footer(s, 11)

    # 12 Partnership options
    s = slide()
    band(s, "Partnership options", "How we might work together", "Clear mutual value.")
    rect(s, 0.55, 1.65, 5.9, 3.35, WHITE, LINE, True)
    text(s, "Areas to explore", 0.8, 1.8, 5.3, 0.35, 16, TEAL, True)
    bullets(
        s,
        [
            "Community wellness and preventive campaigns",
            "Mental wellbeing and holistic health programmes",
            "AI for social impact storytelling",
            "Social enterprise growth support",
            "Introductions to aligned impact partners",
        ],
        0.8,
        2.25,
        5.3,
        2.5,
        13,
    )
    rect(s, 6.8, 1.65, 5.9, 3.35, PALE_TEAL, TEAL, True)
    text(s, "FairBanks brings", 7.05, 1.8, 5.3, 0.35, 16, NAVY, True)
    bullets(
        s,
        [
            "Operating clinic and Community Reach programmes",
            "Trusted CHW/VHT and community relationships",
            "FCHIP MVP currently in development",
            "A Uganda-first model that can grow carefully",
            "Leadership focused on practical community health results",
        ],
        7.05,
        2.25,
        5.3,
        2.5,
        13,
    )
    photo_pair(
        s,
        "compassion",
        "reception",
        "Human-centred clinical care",
        "Welcoming reception and triage",
        x=0.55,
        y=5.15,
        w=12.2,
        h=1.85,
    )
    footer(s, 12)

    # 13 Close
    s = slide()
    crop(s, photo("close"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    text(s, "NEXT STEP", 0.8, 1.1, 4.0, 0.35, 13, GOLD, True)
    text(
        s,
        "We would welcome a conversation about working together on community health.",
        0.8,
        1.7,
        10.5,
        1.5,
        28,
        WHITE,
        True,
    )
    text(
        s,
        "Next step: a brief 10-minute conversation with your partnerships or impact team.",
        0.82,
        3.5,
        10.0,
        0.6,
        16,
        WHITE,
    )
    text(s, SLOGAN, 0.82, 5.2, 4.0, 0.3, 14, GOLD, True)
    text(
        s,
        f"{CONTACT_NAME}  ·  {CONTACT_TITLE}\n"
        f"{ORG}  ·  {LOCATION}\n"
        f"{EMAIL}  ·  {PHONE}  ·  {WEBSITE}",
        0.82,
        5.7,
        10.5,
        1.1,
        13,
        WHITE,
    )

    for sl in slides_for_anim:
        _add_entrance_anims(sl)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)
    print(f"PPTX: {PPTX}")

def validate() -> None:
    from zipfile import BadZipFile, ZipFile

    from docx import Document
    from pptx import Presentation

    for path in (DOCX, PDF, PPTX):
        if not path.exists() or path.stat().st_size < 20_000:
            raise RuntimeError(f"Missing or unexpectedly small output: {path}")
    if not EMAIL_MD.exists() or EMAIL_MD.stat().st_size < 200:
        raise RuntimeError(f"Missing partnership email markdown: {EMAIL_MD}")
    for path in (DOCX, PPTX):
        try:
            with ZipFile(path) as zf:
                bad = zf.testzip()
                if bad:
                    raise RuntimeError(f"Corrupt archive member: {bad}")
        except BadZipFile as exc:
            raise RuntimeError(f"Corrupt Office file: {path}") from exc

    doc = Document(DOCX)
    content = "\n".join(p.text for p in doc.paragraphs)
    for phrase in (
        "FCHIP",
        "Jay Shetty",
        "Community Reach",
        CONTACT_NAME,
    ):
        if phrase not in content:
            raise RuntimeError(f"DOCX validation failed: missing {phrase}")
    if "Appendix - send-ready outreach email" in content:
        raise RuntimeError("DOCX still contains email appendix; move stayed incomplete")

    email_md = EMAIL_MD.read_text(encoding="utf-8")
    for phrase in ("Subject:", "FCHIP", CONTACT_NAME, "jayshetty_pdf.pdf"):
        if phrase not in email_md:
            raise RuntimeError(f"Email markdown validation failed: missing {phrase}")

    deck = Presentation(PPTX)
    if len(deck.slides) != 13:
        raise RuntimeError(f"Expected 13 slides, found {len(deck.slides)}")
    print(f"Validated: DOCX + PDF + {len(deck.slides)} PPT slides + {EMAIL_MD.name}")


def main() -> None:
    print(f"Building Jay Shetty partnership pack ({date.today().isoformat()})")
    OUT.mkdir(parents=True, exist_ok=True)
    build_docx()
    convert_pdf()
    build_pptx()
    validate()
    print("Partnership pack complete.")
    print(f"Send-ready email: {EMAIL_MD}")


if __name__ == "__main__":
    main()
