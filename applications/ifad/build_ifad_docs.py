#!/usr/bin/env python3
"""
IFAD Rural Sector Performance Assessment Grant — FairBanks / FCHIP concept note.

Generates:
  applications/ifad/documents/ifad_word.docx
  applications/ifad/documents/ifad_pdf.pdf
  applications/ifad/documents/ifad_ppt.pptx

Run: python applications/ifad/build_ifad_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "ifad"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/07/14/"
    "ifad-announces-us1-2-million-grant-call-for-proposals-to-strengthen-rural-"
    "development-policies-worldwide/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."
TAGLINE = "Health for All — Obulamu eri Bonna · Afya kwa Wote · Oburamu bwa Boona"

META = {
    "programme": "IFAD Grant — Rural Sector Performance Assessment (up to US$1.2M)",
    "doc_title": "Concept Note — Health–Livelihood Intelligence for Rural Policy",
    "subtitle": "Linking community health signals to rural development performance and investment",
    "eoi_note": "Optional non-binding Expression of Interest by 31 July 2026 · Full proposal by 4 September 2026 (12:00 CET)",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "outreach": "outreach_outdoor_clinic.jpeg",
    "architecture": "data_flow_iso_labeled.png",
    "dashboard": "dashboard_demo.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "conclusion": "outreach_audience_full_group_01.jpg",
}

SECTIONS = [
    ("1. Concept note summary", [
        "IFAD seeks proposals up to US$1.2M to leverage the Rural Sector Performance Assessment for "
        "better policy and investment — open to non-profits, private firms, universities, and think tanks.",
        "FairBanks' cascade already links health to livelihoods and resilient families. FCHIP can "
        "contribute community health and vulnerability signals that make rural performance assessment "
        "more human-centred. This pack frames a careful, partnership-ready fit — to be confirmed against "
        "IFAD's detailed thematic guidance before full proposal effort.",
    ]),
    ("2. Thematic fit verification", [
        "We treat thematic alignment as a gate, not an assumption. Before investing full proposal effort, "
        "FairBanks will submit an optional Expression of Interest (31 July 2026) and review IFAD's "
        "detailed guidance against our proposed health–livelihood indicators.",
        "If official guidance confirms fit, we proceed to a full proposal by 4 September 2026. If not, "
        "we will not stretch the narrative beyond what the call supports — protecting evaluator time and "
        "our own integrity.",
    ]),
    ("3. Health shocks and rural livelihoods", [
        "Rural policy often tracks production and infrastructure while under-weighting health shocks that "
        "wipe out household gains. A fever outbreak, a complicated pregnancy, or untreated hypertension "
        "can erase months of farm income and push families into debt.",
        "Without community health intelligence, rural performance assessments miss a core driver of poverty. "
        "FairBanks sees this daily through outreach in peri-urban and rural-facing catchments around Kampala "
        "and partner communities.",
    ]),
    ("4. FCHIP contribution to rural performance assessment", [
        "We propose a workstream where FCHIP community indicators — illness clusters, maternal risk, "
        "missed care, medicine stock-outs — inform rural performance dashboards and policy briefs, "
        "co-designed with research institutes and IFAD partners.",
        "The FairBanks cascade connects community members, CHWs/VHTs, outreach programmes, clinical care, "
        "research partnerships, and economic empowerment. FCHIP is the intelligence layer that makes "
        "health–livelihood linkages visible to policymakers.",
    ]),
    ("5. Partnership and counterpart model", [
        "FairBanks brings field operations, CHW networks, and data stewardship. We expect to partner with "
        "a university or think tank for rigorous assessment design, ethics review, and policy writing.",
        "IFAD requires 20–25% counterpart contribution. FairBanks will meet this through in-kind field "
        "operations, staff time, outreach infrastructure, and existing digital health foundations — not "
        "as a cash substitute for core research costs.",
    ]),
]

PATHWAY = [
    ("By 31 Jul 2026", "Optional non-binding Expression of Interest to test thematic fit"),
    ("Review", "Confirm alignment with IFAD Rural Sector Performance Assessment guidance"),
    ("By 4 Sep 2026", "Full proposal only if thematic fit is confirmed"),
    ("Year 1–2", "Pilot health–livelihood indicators in FairBanks / partner rural catchments"),
    ("Learning", "Policy notes for IFAD, national stakeholders, and rural investment planners"),
]

BUDGET_ROWS = [
    ("Research & assessment design", "Methods, partner contracts, ethics", "Share of grant"),
    ("Community data & FCHIP modules", "Indicators, dashboards, data quality", "Share of grant"),
    ("Policy engagement", "Briefs, stakeholder dialogues, learning events", "Share of grant"),
    ("FairBanks counterpart", "Field ops, CHW supervision, outreach, staff time", "20–25% in-kind"),
    ("Ceiling", "Per IFAD call", "Up to US$1.2M"),
]

MUTUAL_VALUE = [
    ("FairBanks gains", "Resources to connect health intelligence with rural livelihood policy; "
     "university partnerships; evidence that prevention protects income and food security."),
    ("IFAD gains", "Ground-level health and vulnerability data from community programmes; "
     "social-enterprise implementer with field access; SDG-aligned health-shock narrative; "
     "counterpart contribution via live operations."),
]

WORKSTREAMS = [
    "Health-shock markers linked to livelihood stress in rural performance dashboards",
    "Community-generated indicators from CHW outreach and clinic encounters",
    "Policy briefs translating FCHIP signals into rural investment priorities",
    "Partner-led evaluation design with academic / think-tank co-authorship",
]


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{TEAL}"/>'))
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=10, color=SLATE)
                if ri % 2:
                    c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CREAM}"/>'))
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(META["subtitle"], size=12, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption=TAGLINE)
    para(META["eoi_note"], size=10, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    for title, paras in SECTIONS:
        heading(title)
        if "FCHIP" in title:
            image("architecture", caption="Health signals on the FairBanks community cascade")
        if "Health shocks" in title:
            image("outreach", caption="Community outreach — where health meets livelihoods")
        for t in paras:
            para(t)

    heading("6. Proposed workstreams")
    bullets(WORKSTREAMS)

    heading("7. Pathway — EOI to policy learning")
    table(["Step", "Action"], PATHWAY, widths=[1.8, 4.6])

    heading("8. Indicative resources and counterpart")
    table(["Item", "Detail", "Amount"], BUDGET_ROWS, widths=[2.0, 3.2, 1.2])

    heading("9. Mutual value and honest fit")
    for label, body in MUTUAL_VALUE:
        para(f"{label}: {body}", bold=True)
    para("We will submit an EOI and proceed to full proposal only where official IFAD guidance confirms "
         "thematic fit.", italic=True)

    heading("10. Expression of interest ask")
    para("Rural development succeeds when families stay healthy enough to work and thrive. FairBanks can "
         "help IFAD see that link clearly — with rigorous partners and community-rooted data.")
    image("conclusion", caption="From community health signals to rural policy learning")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable, Table, TableStyle
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))
    st.add(ParagraphStyle("CellHead", fontName="Helvetica-Bold", fontSize=8.5, textColor=white))
    st.add(ParagraphStyle("CellBody", fontName="Helvetica", fontSize=8.5, textColor=slate))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = w * ih / iw
        if h > 2.8 * inch:
            h, w = 2.8 * inch, 2.8 * inch * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 8))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        for row in rows:
            data.append([Paragraph(str(c), st["CellBody"]) for c in row])
        widths = widths or [pw / len(headers)] * len(headers)
        t = Table(data, colWidths=widths, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal),
            ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(
        f'<font color="#{ACCENT}"><b><i>FairBanks Community Health Intelligence Platform (FCHIP)</i></b></font>',
        st["Meta"]))
    story.append(Paragraph(f'<font color="#{MUTED}"><i>{META["subtitle"]}</i></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    story.append(Spacer(1, 8))
    img("cover", cap=TAGLINE)
    story.append(Paragraph(META["eoi_note"], st["Meta"]))
    story.append(PageBreak())

    for title, paras in SECTIONS:
        story.append(Paragraph(title, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
        for t in paras:
            story.append(Paragraph(t, st["Body"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph("6. Proposed workstreams", st["H1"]))
    for w in WORKSTREAMS:
        story.append(Paragraph("• " + w, st["FBullet"]))

    story.append(Paragraph("7. Pathway — EOI to policy learning", st["H1"]))
    tbl(["Step", "Action"], PATHWAY, [pw * 0.28, pw * 0.72])

    story.append(Paragraph("8. Indicative resources and counterpart", st["H1"]))
    tbl(["Item", "Detail", "Amount"], BUDGET_ROWS, [pw * 0.28, pw * 0.52, pw * 0.2])

    story.append(Paragraph("9. Mutual value", st["H1"]))
    for label, body in MUTUAL_VALUE:
        story.append(Paragraph(f"<b>{label}:</b> {body}", st["Body"]))

    story.append(Paragraph("10. Expression of interest ask", st["H1"]))
    story.append(Paragraph(
        "We will submit an EOI and proceed to full proposal only where official IFAD guidance confirms thematic fit.",
        st["Body"]))
    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                            topMargin=0.7 * inch, bottomMargin=0.7 * inch, title="FairBanks IFAD Concept Note")
    doc.build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic_cover(sl, key):
        sl.shapes.add_picture(embed(key), 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks FCHIP | IFAD Rural Performance | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    # 1 Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.5), SW, Inches(3.5), NAVY)
    tb(s, Inches(0.6), SH - Inches(3.2), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.7), Inches(12), Inches(0.55), META["doc_title"], size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(2.15), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(1.75), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.2), Inches(12), Inches(0.35), META["eoi_note"], size=11, color="D0E8E8")

    # 2 Concept
    s = prs.slides.add_slide(blank)
    band(s, "Concept note", "Health–livelihood intelligence for rural policy")
    tb(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.5),
       "IFAD grant up to US$1.2M for Rural Sector Performance Assessment.\n\n"
       "FairBanks links health shocks to rural livelihoods through FCHIP and the community cascade.\n\n"
       "We proceed carefully — EOI first, full proposal only if thematic fit is confirmed.",
       size=20, color=MUTED)
    pic_fit(s, "dashboard", Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5))
    footer(s, 2)

    # 3 Health shocks
    s = prs.slides.add_slide(blank)
    band(s, "The gap", "Health shocks wipe out rural household gains")
    pic_fit(s, "maternal", Inches(0.45), Inches(1.15), Inches(5.8), Inches(5.6))
    tb(s, Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5),
       "• Illness erases farm income and savings\n"
       "• Performance assessments under-weight health\n"
       "• Community data could make policy human-centred\n"
       "• Prevention protects livelihoods and food security",
       size=16, color=SLATE)
    footer(s, 3)

    # 4 FCHIP layer
    s = prs.slides.add_slide(blank)
    band(s, "FCHIP layer", "Community signals → rural performance dashboards")
    pic_fit(s, "architecture", Inches(0.4), Inches(1.1), Inches(7.0), Inches(5.7))
    tb(s, Inches(7.7), Inches(1.2), Inches(5.2), Inches(5.5),
       "CHWs → outreach → clinic → research → empowerment.\n\n"
       "FCHIP adds illness clusters, maternal risk, missed care, and stock-out signals "
       "for policy briefs and investment priorities.",
       size=19, color=MUTED)
    footer(s, 4)

    # 5 Workstreams
    s = prs.slides.add_slide(blank)
    band(s, "Workstreams", "What we would implement with partners")
    for i, item in enumerate(WORKSTREAMS):
        y = Inches(1.25) + i * Inches(1.15)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.0), "FFFFFF", LINE)
        rect(s, Inches(0.5), y, Inches(0.12), Inches(1.0), TEAL)
        tb(s, Inches(0.75), y + Inches(0.15), Inches(11.5), Inches(0.7), item, size=18, color=SLATE)
    footer(s, 5)

    # 6 Counterpart
    s = prs.slides.add_slide(blank)
    band(s, "Counterpart", "20–25% in-kind from FairBanks operations")
    tb(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.5),
       "FairBanks meets IFAD counterpart through:\n\n"
       "• Field operations and CHW supervision\n"
       "• Outreach infrastructure and staff time\n"
       "• Data stewardship and existing digital foundations\n"
       "• Partner university / think tank for research design",
       size=16, color=SLATE)
    pic_fit(s, "mobile", Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5))
    footer(s, 6)

    # 7 Pathway
    s = prs.slides.add_slide(blank)
    band(s, "EOI → proposal → learning", "Gated by thematic fit verification")
    for i, (step, action) in enumerate(PATHWAY):
        y = Inches(1.15) + i * Inches(1.05)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(0.95), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.1), Inches(2.8), Inches(0.35), step, size=13, bold=True, color=TEAL)
        tb(s, Inches(3.6), y + Inches(0.1), Inches(8.5), Inches(0.75), action, size=13, color=MUTED)
    footer(s, 7)

    # 8 Budget
    s = prs.slides.add_slide(blank)
    band(s, "Resources", "Indicative allocation — up to US$1.2M ceiling")
    for i, (item, detail, amt) in enumerate(BUDGET_ROWS):
        y = Inches(1.15) + i * Inches(1.05)
        last = i == len(BUDGET_ROWS) - 1
        rect(s, Inches(0.5), y, Inches(12.2), Inches(0.95), "E8F0F0" if last else "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.1), Inches(3.5), Inches(0.35), item, size=12, bold=last, color=NAVY if last else SLATE)
        tb(s, Inches(4.3), y + Inches(0.1), Inches(6.5), Inches(0.75), detail, size=11, color=MUTED)
        tb(s, Inches(11.2), y + Inches(0.1), Inches(1.2), Inches(0.35), amt, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    footer(s, 8)

    # 9 Mutual value
    s = prs.slides.add_slide(blank)
    band(s, "Win-win", "FairBanks ↔ IFAD ↔ rural families")
    rect(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "FairBanks gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[0][1], size=18, color=MUTED)
    rect(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4), "IFAD gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[1][1], size=18, color=MUTED)
    footer(s, 9)

    # 10 Ask
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.8), SW, Inches(2.8), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.4), Inches(12), Inches(0.55),
       "EOI first — full proposal only where IFAD guidance confirms fit.", size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.5), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. ifad document set in", OUT_DIR)
