#!/usr/bin/env python3
"""Build the DPI Safeguards Accelerator 2026 FairBanks application pack.

Creates one synchronized set:
  documents/dpi-safeguards_word.docx
  documents/dpi-safeguards_pdf.pdf
  documents/dpi-safeguards_ppt.pptx

Word/PDF = Stage 1 Airtable answer bank + Stage 2 prep.
PowerPoint = elegant 10-slide visual brief for partners / internal alignment.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path


HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
ASSETS = REPO / "assets"
OUT = HERE / "documents"
DOCX = OUT / "dpi-safeguards_word.docx"
PDF = OUT / "dpi-safeguards_pdf.pdf"
PPTX = OUT / "dpi-safeguards_ppt.pptx"
PASTE = HERE / "stage1_airtable_answers.md"

NAVY = "0A1F2E"
TEAL = "0D6E6E"
GREEN = "2D7A55"
ORANGE = "C45C26"
GOLD = "D99A2B"
CREAM = "F7F5F0"
PALE_TEAL = "E8F3F2"
PALE_ORANGE = "FBEDE6"
PALE_GREEN = "E9F2EC"
SLATE = "1E2F38"
MUTED = "52636C"
LINE = "CED9D8"
WHITE = "FFFFFF"
RED = "A3312D"

PROGRAMME = "DPI Safeguards Accelerator 2026"
TITLE = "Safeguards for community health DPI in Uganda"
SUBTITLE = "Issue-specific safeguards on FairBanks FCHIP sectoral digital health infrastructure"
SLOGAN = "Your health, our mission."
DEADLINE = "30 July 2026 (Stage 1 final Wave 3)"
PATHWAY = "Non-Governmental Organization (NGO) Pathway — Uganda only"
TRACK = "Issue-specific safeguards work (with a short system-level risk map)"
AMOUNT_REQUESTED = 70_000
DURATION = "9 months"
CONFIRM = "[CONFIRM BEFORE SUBMISSION]"

ORG = "FairBanks Medical Centre"
CONTACT_NAME = "Racheal Nabukeera"
CONTACT_TITLE = "Managing Director"
WEBSITE = "https://fairbanksmedicalcentre.org"
EMAIL = "info@fairbanksmedicalcentre.org"
PHONE = "+256 772 849 258"
LOCATION = "Kampala, Uganda"
COUNTRY = "Uganda"

OFFICIAL_ACCEL = "https://www.dpi-safeguards.org/accelerator"
OFFICIAL_HOME = "https://www.dpi-safeguards.org/"
OFFICIAL_FORM = "https://airtable.com/appcltTa6tzGwskOA/pagJNFsEqox4M5WAU/form"
OFFICIAL_SUMMARY = (
    "https://www2.fundsforngos.org/construction/"
    "safeguards-accelerator-supporting-safe-and-inclusive-digital-public-infrastructure/"
)
FRAMEWORK = "https://www.dpi-safeguards.org/"  # Framework linked from home / resource hub

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "outreach": "outreach_bp_screening.jpeg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "dashboard": "dashboard_demo.png",
    "gis": "gis_hotspots.png",
    "architecture": "data_flow_iso_labeled.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "community": "outreach_audience_full_group_01.jpg",
    "facility": "facility_exterior_entrance_01.jpg",
    "team": "staff_team_reception.jpeg",
}

CONCEPT = {
    "cascade": REPO / ".cursor" / "concept_improved.jpeg",
    "simple": REPO / ".cursor" / "concept_simple.jpeg",
    "classic": REPO / ".cursor" / "concept.jpeg",
}


def photo(key: str) -> Path:
    path = ASSETS / PHOTOS[key]
    if not path.exists():
        raise FileNotFoundError(path)
    return path


def concept(key: str = "cascade") -> Path:
    path = CONCEPT[key]
    if not path.exists():
        raise FileNotFoundError(path)
    return path


CALL_FACTS = [
    ("Programme", "DPI Safeguards Accelerator (first cohort)"),
    ("Stewards", "UN ODET and UNDP; with 50-in-5, OGP, and Co-Develop"),
    ("Pathway", PATHWAY),
    ("Deadline Stage 1", DEADLINE),
    ("Stage 2 (invite only)", "Detailed proposal + government letter by 15 Aug 2026"),
    ("Funding", f"Catalytic grant up to USD {AMOUNT_REQUESTED:,} (6–9 months)"),
    ("Draft ask", f"USD {AMOUNT_REQUESTED:,} over {DURATION}"),
    ("Track", TRACK),
    ("Submit via", OFFICIAL_FORM),
]

READINESS = [
    ("NGO / CBO / non-profit legal status for NGO pathway", CONFIRM),
    ("Applicant legal name exactly as on registration documents", CONFIRM),
    ("Named government counterpart agency and contact", CONFIRM),
    ("Path to written government support letter for Stage 2", CONFIRM),
    ("Airtable access (region/VPN if form is blocked)", CONFIRM),
    ("Lead contact, email, phone, and org website", f"{CONTACT_NAME}; still verify {CONFIRM}"),
    ("Indicative budget owners and co-funding (if any)", CONFIRM),
    ("Named CSO / VHT / community champions for co-design", CONFIRM),
]

# --- Stage 1 narrative blocks (paste-ready) ---

ORG_CONTEXT = (
    "FairBanks is a community health organisation in Kampala, Uganda. We run a medical "
    "centre and FairBanks Community Reach — a six-step cascade linking community members, "
    "CHWs and VHTs, outreach programmes, clinical care, research and skills, and economic "
    "empowerment including CHIS (community health insurance) and livelihoods. Our active "
    "catchments include Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and surrounding "
    "peri-urban communities. FCHIP (FairBanks Community Health Intelligence Platform) is "
    "our working MVP digital component on the Data and Feedback loop: offline mobile "
    "capture for CHWs/VHTs, multi-source health signals, GIS, climate-aware early warning, "
    "facility dashboards, and secure authenticated data APIs so existing EMR/HMS systems "
    "can share clinical data into FCHIP in real time without replacing clinic software. "
    f"Legal entity type and registration evidence: {CONFIRM}."
)

DPI_INITIATIVE = (
    "The active DPI effort we support is sectoral digital health infrastructure for "
    "community-to-facility data exchange in Uganda — not a consumer health app. In "
    "practice this is FCHIP as a live deployment layer that connects last-mile CHW/VHT "
    "data, school and community outreach signals, FairBanks Medical Centre and partner "
    "clinic workflows, and secure EMR/HMS interoperability. This sits alongside Uganda's "
    "broader digital health and HMIS/DHIS2 landscape and is designed to strengthen trust, "
    "inclusion, and accountable data use as digital health services scale. We treat this "
    "as sectoral DPI because it is a shared digital building block for public-interest "
    "health services: identity of households for continuity of care (with consent), "
    "structured data exchange across actors, and decision support that affects real "
    f"people. National programme name / MoH digital health alignment wording: {CONFIRM}."
)

GOVERNMENT_COUNTERPART = (
    f"Intended government counterpart for Stage 1 alignment (must be confirmed): {CONFIRM} "
    "(examples to verify — do not submit unconfirmed names: Ministry of Health digital "
    "health / community health structures; Kampala Capital City Authority health "
    "directorate; or the District Health Office covering FairBanks catchments). We seek "
    "a working relationship that can progress to a formal written institutional "
    "confirmation for Stage 2 — a commitment letter supporting the safeguards "
    "intervention and willingness to participate in the cohort. Current status of "
    f"endorsement discussions: {CONFIRM}."
)

TRACK_CHOICE = (
    "Issue-specific safeguards work, with a short system-level risk map to prioritise "
    "actions. Over 6–9 months we will advance three concrete safeguards inside the live "
    "FCHIP / community health DPI deployment: (1) responsible data use, consent, and "
    "least-privilege EMR/HMS data sharing; (2) offline-first inclusion for underserved "
    "peri-urban users and frontline workers with limited connectivity or literacy; "
    "(3) a practical grievance and recourse pathway for people affected by digital health "
    "data use or alerts. The short system-level map will use the Universal DPI Safeguards "
    "Framework to show gaps, owners, and sequence — then we execute the three measures."
)

WHY_NOW = (
    "FCHIP already has a working MVP and is moving from build toward field validation "
    "inside FairBanks Community Reach. This is the moment when safeguards must be woven "
    "into live workflows — consent, API scopes, offline inclusion, and recourse — before "
    "habits harden and before EMR/HMS feeds, GIS alerts, and climate-linked early warnings "
    "scale. The Universal DPI Safeguards Framework and Accelerator emphasise that "
    "2026–2028 is a crucial window to embed operational safeguards into code, design, "
    "governance, and oversight as DPI moves from pilots to population use. Digital health "
    "is exactly the kind of sectoral DPI use case the Framework anticipates. Acting now "
    "lets Uganda produce a practical, community-rooted safeguards package others can "
    "learn from, instead of retrofitting trust after exclusion or privacy harms appear."
)

TANGIBLE_OUTPUTS = (
    "Within the 6–9 month Accelerator cycle we commit to these operational assets: "
    "(1) a safeguards risk-and-gap brief for the FCHIP community health data exchange, "
    "mapped to Universal DPI Safeguards Framework principles and practices; "
    "(2) an operational consent and data-sharing protocol for CHW/VHT capture and "
    "secure EMR/HMS APIs (authentication, purpose limitation, least privilege, audit); "
    "(3) an inclusion playbook for offline access, low-literacy UX, and reach to women, "
    "older people, persons with disabilities, and urban poor households in our catchments; "
    "(4) a community grievance / user recourse pathway for data access requests, "
    "correction, and alert-related concerns, with named owners and response times; "
    "(5) a standing multistakeholder safeguards working group (FairBanks, VHTs/CHWs, "
    "facility clinical lead, government liaison, and community champions) with meeting "
    "records; (6) a short public learning note suitable for the DPI Safeguards Resource "
    "Hub. Success means these assets are used in live workflows — not only written."
)

COORDINATION_MODEL = (
    "FairBanks team: product/clinical/data-protection leads implement protocols inside "
    "FCHIP and facility workflows; Community Reach coordinators convene field testing. "
    "Government counterpart: validates priorities, aligns with national digital health "
    "and community health guidance, and participates in the working group "
    f"({CONFIRM} named focal point). Civil society and affected users: VHTs/CHWs, "
    "community members from Bukoto, Kyebando, Kisaasi, Kamwokya, and Kikaaya, CHIS/"
    "livelihood group representatives, and local CSO champions co-design consent "
    "language, inclusion tests, and the grievance pathway — with defined roles in "
    "validation and monitoring, not one-off consultations. Private/ecosystem: EMR/HMS "
    "and facility IT contacts advise on API scopes without replacing existing systems. "
    "Rhythm: kick-off workshop; monthly working-group reviews; community listening "
    "sessions each quarter; end-of-cycle demo of live safeguards in the deployment."
)

BUDGET_NARRATIVE = (
    f"Indicative request: USD {AMOUNT_REQUESTED:,} catalytic grant over {DURATION}. "
    "This accelerates safeguards mobilisation inside an existing deployment; it does not "
    "replace FairBanks' ongoing Community Reach or clinical financing. Approximate "
    "allocation (refine in Stage 2): personnel and community stipends for co-design and "
    "safeguards facilitation (~40%); technical work to operationalise consent, API "
    "controls, offline inclusion, and grievance tooling in the MVP (~30%); workshops, "
    "travel within Kampala catchments, translation/printing (~15%); MEL, documentation, "
    f"and learning products for the Resource Hub (~15%). Co-funding / in-kind: {CONFIRM}."
)

BUDGET_ROWS = [
    ("Safeguards lead, clinical/data-protection time, VHT/CHW stipends", "28,000"),
    ("Technical implementation (consent flows, API scopes, offline UX, grievance tooling)", "21,000"),
    ("Workshops, community sessions, local travel, translation, materials", "10,500"),
    ("MEL, documentation, Resource Hub learning note, cohort participation costs", "10,500"),
    ("Total catalytic ask", "70,000"),
]

PROBLEM_SHORT = (
    "As Uganda expands digital health tools, community and facility data often move "
    "without clear consent, inclusion for offline users, or a simple way for people to "
    "challenge mistakes. Trust then lags behind technology — and exclusion grows."
)

SOLUTION_SHORT = (
    "Embed three live safeguards into FairBanks' FCHIP sectoral digital health DPI: "
    "responsible consent and EMR/HMS data sharing; offline-first inclusion; and a "
    "community grievance pathway — co-owned with government and civil society."
)

WIN_THEMES = [
    (
        "Government alignment",
        "Work with a named public health authority on an active digital health deployment; "
        "prepare for a Stage 2 commitment letter.",
    ),
    (
        "Civil society integration",
        "VHTs, households, and community champions co-design, validate, and monitor — "
        "not a one-off consultation.",
    ),
    (
        "Tangible outputs",
        "Six operational assets in 9 months, used in live FCHIP workflows and shared as "
        "learning for the global Framework.",
    ),
    (
        "Framework fit",
        "Maps to Universal DPI Safeguards priorities: inclusion, accountable data use, "
        "recourse, multistakeholder governance, and sectoral DPI (digital health).",
    ),
]

STAGE2_CHECKLIST = [
    "Granular activity plan by month with owners and milestones",
    "Finalised budget with quotations and salary scales",
    "Organisational profile + verified proof of legal status (NGO pathway)",
    "Written government support / commitment letter from collaborating authority",
    "Named civil society / affected-user roles and meeting cadence",
    "Data protection, safeguarding, and ethics approach for health data",
    "Risk register and mitigation (privacy, exclusion, alert harm, adoption)",
    "Learning and contribution plan for the DPI Safeguards Resource Hub",
]

SOURCES = [
    ("Accelerator page", OFFICIAL_ACCEL),
    ("DPI Safeguards home", OFFICIAL_HOME),
    ("NGO Stage 1 Airtable form", OFFICIAL_FORM),
    ("fundsforNGOs programme summary", OFFICIAL_SUMMARY),
]


def write_paste_sheet() -> None:
    text = f"""# DPI Safeguards Accelerator — Stage 1 paste sheet

**Deadline:** {DEADLINE}  
**Form:** {OFFICIAL_FORM}  
**Pathway:** {PATHWAY}  
**Track:** {TRACK}  
**Draft ask:** USD {AMOUNT_REQUESTED:,} / {DURATION}

> Replace every `{CONFIRM}` before submitting. Airtable may block some regions — use a permitted network if needed.

## Organisation / country / DPI rollout

{ORG_CONTEXT}

## Active DPI initiative

{DPI_INITIATIVE}

## Government counterpart

{GOVERNMENT_COUNTERPART}

## Track focus

{TRACK_CHOICE}

## Why this deployment moment matters

{WHY_NOW}

## Expected tangible outputs (6–9 months)

{TANGIBLE_OUTPUTS}

## Coordination model (team + government + civil society)

{COORDINATION_MODEL}

## Indicative budget and duration

{BUDGET_NARRATIVE}

## One-line pitch (if asked)

{SOLUTION_SHORT}

## Official links

- {OFFICIAL_ACCEL}
- {OFFICIAL_HOME}
- {OFFICIAL_FORM}
"""
    PASTE.write_text(text.strip() + "\n", encoding="utf-8")
    print(f"PASTE: {PASTE}")


def build_docx() -> None:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    def hex_rgb(h: str) -> RGBColor:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def shade(cell, color: str) -> None:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), color)
        shd.set(qn("w:val"), "clear")
        tcPr.append(shd)

    def set_run(run, size=11, bold=False, color=SLATE, font="Calibri"):
        run.font.name = font
        run._element.rPr.rFonts.set(qn("w:eastAsia"), font)
        run.font.size = Pt(size)
        run.bold = bold
        run.font.color.rgb = hex_rgb(color)

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)

    def para(text: str, size=11, bold=False, color=SLATE, space_after=8, align=None, italic=False):
        p = doc.add_paragraph()
        if align:
            p.alignment = align
        p.paragraph_format.space_after = Pt(space_after)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.line_spacing = 1.15
        run = p.add_run(text)
        set_run(run, size=size, bold=bold, color=color)
        run.italic = italic
        return p

    def heading(text: str, level=1) -> None:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14 if level == 1 else 10)
        p.paragraph_format.space_after = Pt(6)
        run = p.add_run(text)
        set_run(run, size=16 if level == 1 else 13, bold=True, color=NAVY if level == 1 else TEAL)

    def bullets(items: list[str]) -> None:
        for item in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(4)
            p.clear()
            run = p.add_run(item)
            set_run(run, size=11, color=SLATE)

    def table(headers: list[str], rows: list[tuple], widths: list[float]) -> None:
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.style = "Table Grid"
        for i, h in enumerate(headers):
            cell = t.rows[0].cells[i]
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(h)
            set_run(run, size=10, bold=True, color=WHITE)
            shade(cell, TEAL)
            cell.width = Inches(widths[i])
        for r_i, row in enumerate(rows):
            for c_i, val in enumerate(row):
                cell = t.rows[r_i + 1].cells[c_i]
                cell.text = ""
                p = cell.paragraphs[0]
                run = p.add_run(str(val))
                set_run(run, size=10, color=SLATE)
                if r_i % 2:
                    shade(cell, PALE_TEAL)

    # Cover
    if photo("logo").exists():
        doc.add_picture(str(photo("logo")), width=Inches(1.35))
    para(PROGRAMME, 12, True, TEAL, 4)
    para(TITLE, 22, True, NAVY, 6)
    para(SUBTITLE, 12, False, MUTED, 8)
    para(SLOGAN, 11, True, ORANGE, 12)
    para(
        f"Stage 1 answer bank for Airtable · Source check {date.today().isoformat()}",
        10,
        False,
        MUTED,
        10,
        italic=True,
    )

    table(
        ["Field", "Value"],
        [
            ("Organisation", ORG),
            ("Country", COUNTRY),
            ("Pathway", "NGO / civil society (single country)"),
            ("Track", "Issue-specific (+ short system-level risk map)"),
            ("Draft ask", f"USD {AMOUNT_REQUESTED:,} / {DURATION}"),
            ("Deadline", DEADLINE),
            ("Contact", f"{CONTACT_NAME}, {CONTACT_TITLE}"),
            ("Email / phone", f"{EMAIL} · {PHONE}"),
            ("Website", WEBSITE),
        ],
        [2.0, 4.6],
    )

    heading("1. Call facts")
    table(["Item", "Detail"], CALL_FACTS, [2.0, 4.6])
    para(
        "Official pages always win if this pack and the form differ. Stage 1 is a concise "
        "screening form — not a full proposal.",
        9,
        italic=True,
        color=MUTED,
    )

    heading("2. Submission readiness gates")
    para(
        "Orange items cannot be guessed. Stage 2 will require legal proof and a government letter.",
        10,
        color=ORANGE,
    )
    table(["Gate", "Status"], READINESS, [3.6, 3.0])

    heading("3. Win-win positioning")
    para(PROBLEM_SHORT)
    para(SOLUTION_SHORT, bold=True)
    bullets([f"{a}: {b}" for a, b in WIN_THEMES])
    para(
        "Visual model for partners: Community Reach cascade (concept_improved) with FCHIP "
        "on Data & Feedback — safeguards close the loop: learn, improve, and serve again.",
        10,
        italic=True,
        color=MUTED,
    )

    heading("4. Stage 1 — Airtable answer bank")
    para(
        "Paste or adapt these blocks into the NGO pathway form. Field labels may vary slightly; "
        "match meaning to the official Stage 1 topics listed on the Accelerator page.",
        10,
        italic=True,
        color=MUTED,
    )

    sections = [
        ("4.1 Organisation, country, and DPI rollout", ORG_CONTEXT),
        ("4.2 Active DPI initiative being supported", DPI_INITIATIVE),
        ("4.3 Government counterpart agency", GOVERNMENT_COUNTERPART),
        ("4.4 Track focus", TRACK_CHOICE),
        ("4.5 Why this deployment moment is critical", WHY_NOW),
        ("4.6 Expected tangible outputs (6–9 months)", TANGIBLE_OUTPUTS),
        ("4.7 Coordination model", COORDINATION_MODEL),
        ("4.8 Indicative budget and duration", BUDGET_NARRATIVE),
    ]
    for title, body in sections:
        heading(title, level=2)
        para(body)

    heading("5. Indicative budget (Stage 1 range)")
    table(["Line", "USD"], BUDGET_ROWS, [5.0, 1.6])
    para(
        "Refine line items, quotations, and any co-funding before Stage 2. Catalytic funds "
        "accelerate safeguards inside an active deployment.",
        9,
        italic=True,
        color=MUTED,
    )

    heading("6. How this maps to Universal DPI Safeguards themes")
    table(
        ["Accelerator theme", "FairBanks / FCHIP response"],
        [
            (
                "Responsible data use / consent / sharing",
                "Consent protocol + least-privilege EMR/HMS APIs + audit trail",
            ),
            (
                "Inclusion / offline access",
                "Offline CHW tools; low-literacy UX; reach to vulnerable groups",
            ),
            (
                "Grievance / recourse",
                "Named pathway for data and alert-related concerns",
            ),
            (
                "Multistakeholder governance",
                "Working group: government liaison, VHTs, community, facility",
            ),
            (
                "Sectoral DPI (digital health)",
                "Community health data exchange as public-interest infrastructure",
            ),
        ],
        [2.6, 4.0],
    )

    heading("7. Stage 2 preparation (if longlisted)")
    bullets(STAGE2_CHECKLIST)
    para(
        "Advancement to Stage 2 does not guarantee funding. Final awards follow fiduciary "
        "and programmatic clearance.",
        9,
        italic=True,
        color=MUTED,
    )

    heading("8. Submit checklist (before 30 July)")
    bullets(
        [
            "Confirm NGO-pathway legal eligibility (or eligible lead partner).",
            "Replace every CONFIRM BEFORE SUBMISSION marker with verified facts.",
            "Name the government counterpart you can truthfully claim today.",
            "Open the NGO Airtable form (use permitted network/VPN if region-blocked).",
            "Paste Stage 1 answers; keep tone practical and community-rooted.",
            "Save confirmation / screenshot of submission.",
            "Start drafting the government commitment letter request the same day.",
        ]
    )

    heading("9. Official sources")
    table(["Source", "URL"], SOURCES, [2.4, 4.2])
    para(
        f"Generated {date.today().isoformat()}. Follow {OFFICIAL_ACCEL} if anything conflicts.",
        8,
        italic=True,
        color=MUTED,
    )

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(DOCX)
    print(f"DOCX: {DOCX}")


def convert_pdf() -> None:
    import shutil
    import win32com.client

    tmp_pdf = REPO / "tmp" / "dpi-safeguards_pdf_build.pdf"
    tmp_pdf.parent.mkdir(parents=True, exist_ok=True)
    if tmp_pdf.exists():
        tmp_pdf.unlink()

    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    document = None
    try:
        document = word.Documents.Open(str(DOCX.resolve()), ReadOnly=True)
        document.ExportAsFixedFormat(str(tmp_pdf.resolve()), 17)
    finally:
        if document is not None:
            document.Close(False)
        word.Quit()

    try:
        if PDF.exists():
            PDF.unlink()
        shutil.move(str(tmp_pdf), str(PDF))
    except PermissionError as exc:
        if tmp_pdf.exists():
            tmp_pdf.unlink(missing_ok=True)
        raise PermissionError(
            f"Cannot update {PDF.name} because it is open/locked. "
            "Close it and re-run the build."
        ) from exc
    print(f"PDF:  {PDF}")


def build_pptx() -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    from lxml import etree

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def rgb(h: str) -> RGBColor:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def transition(slide) -> None:
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        tr = etree.Element(f"{{{ns}}}transition", spd="slow", advClick="1")
        etree.SubElement(tr, f"{{{ns}}}fade")
        slide._element.insert(2, tr)

    def blank():
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        transition(s)
        return s

    def rect(slide, x, y, w, h, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
        else:
            shape.line.fill.background()
        return shape

    def text(slide, content, x, y, w, h, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        run.font.name = font
        return box

    def bullets(slide, items, x, y, w, h, size=14, color=SLATE):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
            p.font.name = "Calibri"
            p.space_after = Pt(8)
        return box

    def footer(slide, n, total=10):
        text(slide, SLOGAN, 0.5, 7.1, 4.0, 0.25, 10, MUTED)
        text(slide, f"{n} / {total}", 11.5, 7.1, 1.3, 0.25, 10, MUTED, align=PP_ALIGN.RIGHT)

    def band(slide, eyebrow, title, subtitle):
        rect(slide, 0, 0, 13.333, 1.45, NAVY)
        text(slide, eyebrow.upper(), 0.55, 0.22, 8.0, 0.28, 11, GOLD, True)
        text(slide, title, 0.55, 0.5, 12.0, 0.45, 24, WHITE, True)
        text(slide, subtitle, 0.55, 1.0, 12.0, 0.3, 12, PALE_TEAL)

    def pic(slide, path: Path, x, y, w, h):
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))

    # 1 Cover
    s = blank()
    pic(s, photo("cover"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    # soft left panel feel via overlay already navy
    text(s, PROGRAMME, 0.7, 1.4, 8.0, 0.35, 14, GOLD, True)
    text(s, TITLE, 0.7, 2.0, 9.5, 1.2, 32, WHITE, True)
    text(s, SOLUTION_SHORT, 0.7, 3.5, 9.0, 1.1, 16, PALE_TEAL)
    text(s, f"NGO Pathway · Uganda · USD {AMOUNT_REQUESTED:,} · {DURATION}", 0.7, 5.0, 9.0, 0.35, 14, WHITE, True)
    text(s, f"Stage 1 deadline: {DEADLINE}", 0.7, 5.45, 9.0, 0.3, 12, GOLD)
    text(s, SLOGAN, 0.7, 6.5, 4.0, 0.3, 14, ORANGE, True)

    # 2 Why now
    s = blank()
    band(s, "Why now", "Safeguards before scale hardens bad habits", "FCHIP MVP is entering field validation — the right moment for guardrails.")
    pic(s, photo("mobile"), 0.5, 1.8, 5.2, 4.8)
    bullets(
        s,
        [
            "100+ countries advancing DPI; digital health is a sectoral use case",
            "2026–2028 window: embed safeguards in live deployments",
            "FCHIP already connects CHWs, clinics, and EMR/HMS APIs",
            "Consent, inclusion, and recourse must ship with the data — not later",
        ],
        6.1,
        2.0,
        6.5,
        4.5,
        15,
    )
    footer(s, 2)

    # 3 How FairBanks works
    s = blank()
    band(s, "Operating model", "Community Reach + FCHIP as sectoral health DPI", "Safeguards sit on Data & Feedback — not as a separate brochure.")
    cascade_h = 5.0
    cascade_w = cascade_h * 0.728
    pic(s, concept("cascade"), 0.45, 1.75, cascade_w, cascade_h)
    bullets(
        s,
        [
            "Communities → CHWs/VHTs → programmes",
            "Medical centre anchors clinical action",
            "Research, skills, partners strengthen learning",
            "CHIS and livelihoods support access",
            "FCHIP powers Data & Feedback with GIS and alerts",
        ],
        cascade_w + 0.8,
        2.1,
        6.2,
        4.2,
        15,
    )
    footer(s, 3)

    # 4 Problem / opportunity
    s = blank()
    band(s, "The gap", "Digital health without operational safeguards", PROBLEM_SHORT)
    cards = [
        ("Privacy risk", "Data moves across CHW, clinic, and API layers without clear consent scopes"),
        ("Exclusion risk", "Offline and low-literacy users drop out of digital pathways"),
        ("Accountability gap", "People lack a simple recourse when alerts or records go wrong"),
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.55 + i * 4.15
        rect(s, x, 2.1, 3.9, 4.2, WHITE, LINE)
        rect(s, x, 2.1, 3.9, 0.7, TEAL)
        text(s, t, x + 0.2, 2.25, 3.5, 0.4, 16, WHITE, True)
        text(s, b, x + 0.25, 3.1, 3.4, 2.8, 14, SLATE)
    footer(s, 4)

    # 5 Solution
    s = blank()
    band(s, "Proposal", "Three issue-specific safeguards in a live DPI deployment", TRACK)
    items = [
        ("01", "Responsible data use", "Consent + least-privilege EMR/HMS APIs + audit"),
        ("02", "Offline inclusion", "CHW tools and UX that work without always-on internet"),
        ("03", "Grievance pathway", "Named recourse for data and alert-related concerns"),
    ]
    for i, (n, t, b) in enumerate(items):
        y = 1.85 + i * 1.55
        rect(s, 0.6, y, 12.1, 1.4, PALE_TEAL, TEAL)
        text(s, n, 0.85, y + 0.35, 1.0, 0.5, 28, TEAL, True)
        text(s, t, 2.1, y + 0.25, 9.5, 0.4, 18, NAVY, True)
        text(s, b, 2.1, y + 0.75, 9.5, 0.4, 14, MUTED)
    footer(s, 5)

    # 6 Architecture
    s = blank()
    band(s, "Where safeguards land", "Capture → FCHIP intelligence → action — with guardrails", "APIs to existing EMR/HMS expand coverage without ripping out systems.")
    pic(s, photo("architecture"), 0.5, 1.75, 7.4, 5.0)
    bullets(
        s,
        [
            "CHW/VHT offline capture",
            "Secure EMR/HMS data APIs",
            "Consent and purpose limits",
            "Role-based access and audit",
            "Alerts back to field and facility",
        ],
        8.2,
        2.2,
        4.5,
        4.2,
        15,
    )
    footer(s, 6)

    # 7 Outputs
    s = blank()
    band(s, "6–9 month outputs", "Operational assets, not paper only", f"Draft catalytic ask: USD {AMOUNT_REQUESTED:,}")
    outs = [
        "Safeguards risk & gap brief (Framework-mapped)",
        "Consent & data-sharing protocol for APIs",
        "Inclusion playbook (offline + vulnerable groups)",
        "Community grievance / recourse pathway",
        "Multistakeholder working group (live cadence)",
        "Public learning note for the Resource Hub",
    ]
    for i, o in enumerate(outs):
        col, row = i % 2, i // 2
        x, y = 0.55 + col * 6.35, 1.85 + row * 1.55
        rect(s, x, y, 6.05, 1.35, WHITE, LINE)
        text(s, f"{i + 1}", x + 0.25, y + 0.4, 0.6, 0.4, 20, TEAL, True)
        text(s, o, x + 1.0, y + 0.4, 4.7, 0.6, 15, NAVY, True)
    footer(s, 7)

    # 8 Coalition
    s = blank()
    band(s, "Coordination", "Government · FairBanks · civil society", "Meaningful integration — defined roles in co-design, validation, monitoring")
    pic(s, photo("community"), 0.5, 1.8, 5.5, 4.8)
    bullets(
        s,
        [
            "Government counterpart validates priorities and joins the working group",
            "FairBanks implements protocols inside FCHIP and Community Reach",
            "VHTs and community champions co-design consent and grievance language",
            "Facility/EMR contacts advise on least-privilege API scopes",
            "Monthly reviews + quarterly listening sessions",
        ],
        6.3,
        2.0,
        6.4,
        4.5,
        14,
    )
    footer(s, 8)

    # 9 Budget & ask
    s = blank()
    band(s, "Catalytic ask", f"USD {AMOUNT_REQUESTED:,} over {DURATION}", "Accelerates safeguards mobilisation — does not replace existing Community Reach financing")
    for i, (label, amt) in enumerate(BUDGET_ROWS):
        y = 1.85 + i * 0.85
        fill = TEAL if i == len(BUDGET_ROWS) - 1 else WHITE
        tc = WHITE if i == len(BUDGET_ROWS) - 1 else NAVY
        rect(s, 0.7, y, 11.9, 0.75, fill, LINE)
        text(s, label, 0.95, y + 0.2, 8.5, 0.4, 14, tc, True)
        text(s, f"USD {amt}", 10.0, y + 0.2, 2.3, 0.4, 14, tc, True, align=PP_ALIGN.RIGHT)
    footer(s, 9)

    # 10 Call to action
    s = blank()
    pic(s, photo("facility"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    text(s, "READY FOR STAGE 1", 0.7, 1.5, 6.0, 0.35, 14, GOLD, True)
    text(s, "Submit the NGO pathway form before Wave 3 closes.", 0.7, 2.1, 10.0, 1.0, 28, WHITE, True)
    bullets(
        s,
        [
            f"Confirm legal eligibility and government counterpart ({CONFIRM})",
            "Paste Stage 1 answers from the Word pack / stage1_airtable_answers.md",
            "Start the Stage 2 government letter request the same day",
        ],
        0.85,
        3.5,
        10.0,
        2.2,
        16,
        WHITE,
    )
    text(s, OFFICIAL_FORM, 0.7, 5.9, 11.5, 0.35, 12, PALE_TEAL)
    text(s, SLOGAN, 0.7, 6.5, 5.0, 0.35, 14, ORANGE, True)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)
    print(f"PPTX: {PPTX}")


def validate() -> None:
    from zipfile import BadZipFile, ZipFile
    import fitz
    from docx import Document
    from pptx import Presentation

    for path in (DOCX, PDF, PPTX, PASTE):
        if not path.exists() or path.stat().st_size < 500:
            raise RuntimeError(f"Missing or unexpectedly small output: {path}")
    for path in (DOCX, PPTX):
        try:
            with ZipFile(path) as zf:
                bad = zf.testzip()
                if bad:
                    raise RuntimeError(f"Corrupt archive member: {bad}")
        except BadZipFile as exc:
            raise RuntimeError(f"Corrupt Office file: {path}") from exc

    doc = Document(DOCX)
    content = "\n".join(
        [p.text for p in doc.paragraphs]
        + [c.text for t in doc.tables for r in t.rows for c in r.cells]
    )
    for phrase in (
        "Stage 1 — Airtable answer bank",
        CONFIRM,
        OFFICIAL_FORM,
        "Issue-specific",
        f"{AMOUNT_REQUESTED:,}",
    ):
        if phrase not in content:
            raise RuntimeError(f"DOCX validation failed: missing {phrase}")

    pdf = fitz.open(PDF)
    if pdf.page_count < 4:
        raise RuntimeError(f"PDF has too few pages: {pdf.page_count}")
    deck = Presentation(PPTX)
    if len(deck.slides) != 10:
        raise RuntimeError(f"Expected 10 slides, found {len(deck.slides)}")
    print(f"Validated: {pdf.page_count} PDF pages | {len(deck.slides)} PPT slides")
    pdf.close()


def main() -> None:
    print(f"Building {PROGRAMME} application pack")
    print(f"Source check date: {date.today().isoformat()}")
    OUT.mkdir(parents=True, exist_ok=True)
    write_paste_sheet()
    build_docx()
    convert_pdf()
    build_pptx()
    validate()
    print("Application pack complete.")


if __name__ == "__main__":
    main()
