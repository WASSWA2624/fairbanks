#!/usr/bin/env python3
"""
U.S. Dept of State — Uganda Health System MOU (DFOP0017890) — Statement of Interest.

Win-win Phase 1 SOI under Addendum E (Uganda), Track 2, Objective 4 digital health.
Run: python applications/dos-uganda/build_dos-uganda_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT = PROJECT / "documents"
SLUG = "dos-uganda"

OUT_DOC = OUT / f"{SLUG}_word.docx"
OUT_PDF = OUT / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT / f"{SLUG}_ppt.pptx"

CALL_URL = (
    "https://opportunitiesforyouth.org/2026/06/15/"
    "u-s-department-of-state-launches-up-to-60-million-funding-opportunity-to-strengthen-"
    "ugandas-health-system-through-the-health-foreign-assistance-mou-implementation-plan/"
)
GRANTS_URL = "https://simpler.grants.gov/opportunity/44a573b8-d46a-4cec-8c84-f319cd1cc1b6"

NAVY, TEAL, ACCENT = "0A1F2E", "0D6E6E", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "dashboard": "dashboard_demo.png",
    "gis": "gis_hotspots.png",
    "facility": "facility_exterior_sign.jpeg",
    "outreach": "outreach_outdoor_clinic.jpeg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "pharmacy": "pharmacy_staff_laptop_01.jpg",
}

# ---------------------------------------------------------------------------
# Critical details (cover table — does not count toward 5-page narrative)
# ---------------------------------------------------------------------------

PROJECT_TITLE = (
    "FCHIP Last-Mile Bridge: Interoperable Community Health Data into "
    "Uganda's National Digital Health Architecture"
)
ORG_NAME = "FairBanks Medical Centre"
ADDENDUM = (
    "Addendum E: Health Foreign Assistance Memorandum of Understanding "
    "(MOU) Implementation in Uganda"
)
COUNTRY = "Uganda"
FEDERAL_SHARE = "$1,850,000"
COST_SHARE = "$0 (none proposed at SOI stage)"
PROJECT_LENGTH = "36 months"

PROGRAMME = "U.S. Department of State / GHSD — Advancing Global Health (DFOP0017890)"
DOC_TITLE = "Statement of Interest — Phase 1"
SUBTITLE = (
    "Track 2, Objective 4: Strengthen and Scale Interoperable, "
    "Fit-for-Purpose Digital Health Systems"
)
ORG_LINE = "FairBanks Medical Centre  ·  FairBanks Community Reach  ·  FCHIP"

CRITICAL_DETAILS = [
    ["Proposed Project Title", PROJECT_TITLE],
    ["Name of the Organization", ORG_NAME],
    ["Addendum to which the SOI is responding", ADDENDUM],
    ["Target Benefiting Country", COUNTRY],
    ["Total Federal Share Requested", FEDERAL_SHARE],
    ["Total Cost Share (as applicable)", COST_SHARE],
    ["Project Length", PROJECT_LENGTH],
]

# ---------------------------------------------------------------------------
# Narrative content — win-win SOI (APS §D Phase 1 structure)
# ---------------------------------------------------------------------------

ISSUE = [
    "Uganda's five-year health transition (April 2026-December 2030) asks facilities, "
    "districts, and communities to keep HIV, TB, malaria, maternal and child health, "
    "immunisation, and outbreak services strong while shifting toward government-led, "
    "performance-based programming. That shift only works if decision-makers see what is "
    "happening at the last mile - not weeks later in paper registers.",
    "Community Health Workers and Village Health Teams already walk households and "
    "refer patients. Too often their work never reaches eCHIS, facility EMRs, DHIS2/"
    "eIDSR, or the National Data Warehouse in a usable form. National digital systems "
    "then miss early fever clusters, missed antenatal visits, immunisation gaps, and "
    "stock pressure. The result is reactive care, weak accountability, and higher risk "
    "that transition gains reverse.",
    "The opportunity is clear: bridge last-mile community data into Uganda's approved "
    "National Digital Health Architecture (NDHA) so Ministry of Health (MoH) systems "
    "become more complete, more actionable, and progressively owned by Government of "
    "Uganda (GOU) institutions - without building a parallel silo.",
]

SOLUTION = [
    "FairBanks Medical Centre and FairBanks Community Reach respectfully submit this "
    "Statement of Interest under Track 2, Objective 4 of Addendum E: strengthen and "
    "scale interoperable, fit-for-purpose digital health systems with progressive GOU "
    "responsibility and sustainable operations.",
    "We propose to pilot and then transfer a community-to-facility intelligence layer - "
    "the FairBanks Community Health Intelligence Platform (FCHIP) - that captures "
    "CHW/VHT household and visit data offline, syncs securely, maps indicators to "
    "national reporting elements, and feeds facility and district action. FCHIP is not "
    "a parallel national system. It is a bridge that makes Uganda's existing digital "
    "stack more complete at community level and more useful for prevention.",
    "Primary subsystem focus (permitted under Objective 4): eCHIS-aligned household/"
    "visit capture; facility EMR encounter and referral handoff (eAfya/ClinicMaster "
    "pathways as available in the catchment); DHIS2 aggregate reporting and eIDSR "
    "surveillance signals; readiness feeds toward the National Data Warehouse; and "
    "iHRIS-aware CHW/VHT workforce directories where schemas allow. Design aligns "
    "with NDHA and OpenHIE-consistent interoperability (API-first, standards-aware "
    "exports), under MoH digital health authority.",
]

WIN_WIN_INTRO = (
    "This partnership is designed as a durable win-win: each party gains something "
    "lasting, and no party depends forever on a private tool outside GOU systems."
)

WIN_WIN = [
    [
        "Uganda / Ministry of Health",
        "Richer last-mile data into eCHIS, DHIS2/eIDSR, and related national platforms; "
        "clearer referral closure; documented handoff of staffing and recurrent costs "
        "toward GOU ownership by project close.",
    ],
    [
        "U.S. Department of State / GHSD",
        "A Ugandan implementer with a live clinic and community network that protects "
        "essential services while proving interoperable digital scale-up under the MOU - "
        "accountable, measurable, and transferable - advancing U.S. interests in health "
        "security and country self-reliance.",
    ],
    [
        "Districts & facilities",
        "Dashboards for outreach targeting, referral follow-up, and programme "
        "monitoring; GIS hotspot views that support planning before outbreaks arrive "
        "at the ward.",
    ],
    [
        "Communities & CHWs/VHTs",
        "Simple offline tools, faster support from the medical centre, and care that "
        "reaches mothers, children, older persons, and people with chronic conditions "
        "before crises.",
    ],
    [
        "FairBanks / FCHIP",
        "Room to refine and document a production-grade interoperability model inside "
        "a real catchment - then share specifications, training packages, and transfer "
        "plans with MoH-aligned partners.",
    ],
]

ACTIVITIES = [
    "A1. Deploy offline-capable CHW/VHT capture aligned with eCHIS household and "
    "visit concepts across FairBanks Community Reach catchments.",
    "A2. Establish a secure sync and role-based access pipeline linking community "
    "records to FairBanks Medical Centre workflows and authorised district viewers, "
    "consistent with Uganda data-protection and cybersecurity expectations.",
    "A3. Map community indicators to DHIS2 reporting elements and demonstrate "
    "aggregate submission pathways for MNCH, immunisation, NCD screening, and "
    "disease surveillance signals (including eIDSR-relevant alerts).",
    "A4. Deliver facility and programme dashboards plus GIS hotspot views that "
    "trigger outreach, referrals, and supply planning actions.",
    "A5. Produce interoperability specifications, data-governance protocols, "
    "training curricula, and a staged GOU-ownership transition plan (staffing and "
    "recurrent operating costs) suitable for full-proposal detailed design.",
]

APPROACH = [
    "Anchor in a live cascade: community members identify needs; CHWs/VHTs collect "
    "and educate; FairBanks Community Reach runs outreach and home visits; "
    "FairBanks Medical Centre provides clinical care, diagnostics, pharmacy, and "
    "referrals; FCHIP is the intelligence layer that closes the data-and-feedback loop.",
    "Build API-first, standards-aware exports - not a closed proprietary vault. "
    "Prioritise interoperability with eCHIS, DHIS2/eIDSR, facility EMR pathways "
    "(eAfya/ClinicMaster as applicable), and National Data Warehouse readiness when "
    "schemas are confirmed.",
    "Protect essential services during transition: structure community data around "
    "HIV/TB/malaria flags where relevant, maternal and child pathways, immunisation "
    "follow-up, and outbreak early warning so digital work reinforces - not distracts "
    "from - service continuity.",
    "Co-design forms and alerts with CHWs, facility staff, and local leaders; train "
    "Ugandan data stewards and developers; document every interface for MoH review.",
    "Start where FairBanks already operates (Bukoto, Kyebando, Kisaasi, Kamwokya, "
    "Kikaaya and surrounding communities), then design for district replication "
    "using Uganda's existing CHW/VHT architecture.",
]

OUTCOMES = [
    [
        "Save lives",
        "Earlier alerts for missed ANC, immunisation gaps, fever clusters, and "
        "referral drop-offs - linked to CHW follow-up and clinic action.",
    ],
    [
        "Strengthen systems",
        "Working mappings and demo exports for eCHIS-aligned records and DHIS2/"
        "eIDSR aggregates, with open documentation for MoH and partners.",
    ],
    [
        "Enhance efficiency",
        "Less double entry and paper lag; role-based dashboards that shorten the "
        "path from community signal to facility or district response.",
    ],
    [
        "Foster self-reliance",
        "Handoff package: source documentation, SOPs, training, hosting options, "
        "and a costed path to transfer recurrent operations toward domestic systems "
        "by project close.",
    ],
    [
        "Advance U.S. interest",
        "Measurable contribution to MOU digital transition: interoperable last-mile "
        "data under GOU authority, protecting prior U.S. health investments and "
        "strengthening epidemic early warning.",
    ],
]

CAPACITY = [
    "Ugandan social enterprise with a functioning medical centre - clinical care, "
    "diagnostics, pharmacy, and referral pathways already in daily use.",
    "Active FairBanks Community Reach programmes: maternal and child health, "
    "Gericare, chronic disease screening, school and corporate health outreach.",
    "Established CHW/VHT relationships in Kampala-area communities named above - "
    "real users, not a greenfield pilot.",
    "Digital health records foundation and pharmacy dispensing data already "
    "supporting facility operations.",
    "Research orientation and partnership practice - evidence, ethics, and skills "
    "development sit inside the FairBanks community health model.",
]

INTEROP_ROWS = [
    ["eCHIS", "Household registration, visit logs, referral flags via structured mobile capture"],
    ["Facility EMR (eAfya / ClinicMaster)", "Referral handoff and encounter linkage at FairBanks Medical Centre"],
    ["DHIS2 / eIDSR", "Aggregate exports and surveillance signals for MNCH, immunisation, outbreaks"],
    ["National Data Warehouse", "Standardised JSON/API feeds when national schema is confirmed"],
    ["iHRIS", "CHW/VHT workforce directory alignment where schemas allow"],
]

PRINCIPLES = [
    ["Sustainability & country ownership", "Design for MoH/NDHA alignment and staged transfer of operations"],
    ["Protect essential services", "Digital workflows reinforce HIV, TB, malaria, MNCH, immunisation"],
    ["Data for action & accountability", "Quality data into GOU systems - used for decisions, not parallel burden"],
    ["Integrated, people-centred care", "One community-to-facility loop; less fragmentation"],
    ["Local partnerships", "Government, CHWs/VHTs, CBOs, academic and NGO partners"],
    ["Evidence-based & context fit", "Pilot in live catchments; measure; refine before scale"],
    ["Oversight & coordination", "Align with MoH priorities and designated review platforms"],
]

PARTNERS = [
    [
        "Ministry of Health (digital health / community health)",
        "Alignment with NDHA, eCHIS, and national reporting; review of interfaces and transfer plan",
    ],
    [
        "District health teams (pilot geography)",
        "Aggregate reporting validation, dashboard use, readiness dialogue",
    ],
    [
        "CHWs, VHTs, and community leaders",
        "Primary data partners; co-design of forms and alerts",
    ],
    [
        "Academic partners (to be confirmed at Phase 2)",
        "Evaluation support and interoperability testing",
    ],
    [
        "NGOs / CBOs collaborating with FairBanks Community Reach",
        "Outreach coordination and community engagement",
    ],
]

COST_SHARE_NOTE = (
    "No cost share is proposed at the SOI stage. FairBanks will continue operating "
    "its medical centre and community outreach as complementary organisational "
    "activity. Opportunities for in-kind or matching contributions can be discussed "
    "during co-design or Phase 2 if invited."
)

CLOSING = (
    "This Statement of Interest confirms FairBanks' intent to advance Track 2, "
    "Objective 4 with a community-rooted, interoperable, government-aligned digital "
    "health proposal. We seek invitation to Phase 2 (full proposal) and/or "
    "consultative program design, where we will detail workplans, budgets, M&E, "
    "partnership letters, and a clear transition of staffing and recurrent costs "
    "toward sustainable Ugandan ownership."
)


def asset(name: str) -> Path:
    p = ASSETS / name
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def photo(key: str) -> Path:
    return asset(PHOTOS[key])


def embed(key_or_path, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key_or_path) if key_or_path in PHOTOS else Path(key_or_path)
    cache = REPO / "tmp" / f"{SLUG}_assets"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    """Submission-shaped Word: letter, 1\" margins, 12-pt Times New Roman."""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from PIL import Image as PILImage

    def rgb(c):
        return RGBColor.from_string(c)

    def font(run, size=12, bold=False, color=None, italic=False):
        run.font.name = "Times New Roman"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        if color:
            run.font.color.rgb = rgb(color)

    def shade(cell, color):
        cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))

    def border(cell, color=LINE):
        cell._tc.get_or_add_tcPr().append(parse_xml(
            f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="4" w:color="{color}"/>'
            f'<w:left w:val="single" w:sz="4" w:color="{color}"/>'
            f'<w:bottom w:val="single" w:sz="4" w:color="{color}"/>'
            f'<w:right w:val="single" w:sz="4" w:color="{color}"/></w:tcBorders>'))

    def para(text, **kw):
        defaults = dict(
            size=12, bold=False, color=None, after=6,
            align=WD_ALIGN_PARAGRAPH.LEFT, italic=False, before=0,
        )
        defaults.update(kw)
        p = doc.add_paragraph()
        p.alignment = defaults["align"]
        p.paragraph_format.space_before = Pt(defaults["before"])
        p.paragraph_format.space_after = Pt(defaults["after"])
        p.paragraph_format.line_spacing = 1.0
        font(
            p.add_run(text),
            size=defaults["size"],
            bold=defaults["bold"],
            color=defaults["color"],
            italic=defaults["italic"],
        )
        return p

    def heading(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after = Pt(6)
        p.paragraph_format.line_spacing = 1.0
        font(p.add_run(text), size=12, bold=True)

    def bullets(items, size=12):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            p.paragraph_format.line_spacing = 1.0
            p.clear()
            font(p.add_run(it), size=size)

    def table(headers, rows, widths=None, head_size=10, body_size=10):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=head_size, bold=True)
            shade(c, "E8EEEE")
            border(c, "888888")
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=body_size)
                border(c, "888888")
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    def image(key, width_in=5.5, caption=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width_in, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=10, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=8)

    doc = Document()
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.left_margin = section.right_margin = Inches(1)
    section.top_margin = section.bottom_margin = Inches(1)

    # Page numbers
    footer = section.footer
    footer.is_linked_to_previous = False
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fp.add_run("Page ")
    font(run, size=10)
    fld1 = parse_xml(
        f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>'
    )
    instr = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
    fld2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    run2 = fp.add_run()
    run2._r.append(fld1)
    run2._r.append(instr)
    run2._r.append(fld2)
    font(run2, size=10)

    # ----- Cover: Critical Details (excluded from page limit) -----
    para(PROGRAMME, size=11, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(DOC_TITLE, size=14, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(ORG_LINE, size=12, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=4)
    para(SUBTITLE, size=11, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=4)
    para(SLOGAN, size=12, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    heading("Table Listing of Critical Details")
    table(["Item", "Detail"], CRITICAL_DETAILS, widths=[2.4, 4.1], head_size=10, body_size=10)
    doc.add_page_break()

    # ----- Narrative (APS §D) -----
    heading("1. Issue / Challenge / Opportunity")
    for p in ISSUE:
        para(p)

    heading("2. Proposed Solution / Activities")
    for p in SOLUTION:
        para(p)
    para(WIN_WIN_INTRO, italic=True, after=6)
    table(["Who gains", "What they gain"], WIN_WIN, widths=[2.0, 4.5], head_size=10, body_size=10)

    para("Core activities:", bold=True, after=4)
    bullets(ACTIVITIES)

    para("Project approach:", bold=True, after=4, before=6)
    for p in APPROACH:
        para(p, after=5)

    heading("3. Anticipated Outcomes and Results")
    para(
        "Success means lives protected, systems strengthened, efficiency improved, "
        "self-reliance advanced, and U.S. investments secured through GOU-owned "
        "digital continuity:",
        after=6,
    )
    table(["Outcome lens", "What success looks like"], OUTCOMES, widths=[1.6, 4.9], head_size=10, body_size=10)

    heading("4. Organizational Capacity")
    bullets(CAPACITY)

    heading("5. Interoperability Pathway (chart)")
    table(["Platform", "Integration pathway"], INTEROP_ROWS, widths=[2.2, 4.3], head_size=10, body_size=10)
    image(
        "architecture",
        width_in=5.8,
        caption="Figure 1. Community cascade to FCHIP to national digital health systems (chart; excluded from page limit).",
    )

    heading("6. Alignment with Addendum E Guiding Principles (chart)")
    table(["Principle", "How FairBanks responds"], PRINCIPLES, widths=[2.2, 4.3], head_size=10, body_size=10)

    heading("7. List of Partner Roles and Responsibilities (preliminary)")
    para(
        "Partners below are preliminary characterisations for Phase 1. Roles, contacts, "
        "and any subrecipient federal-share splits will be refined in Phase 2 / co-design "
        "if invited. At SOI stage, FairBanks Medical Centre is proposed as prime applicant.",
        size=12, after=6,
    )
    table(
        ["Partner", "Proposed role (summary)"],
        PARTNERS,
        widths=[2.4, 4.1],
        head_size=10,
        body_size=10,
    )

    heading("8. Resource Contributions and/or Cost Share")
    para(COST_SHARE_NOTE)

    heading("9. Readiness for Phase 2")
    para(CLOSING)
    para(
        f"References: DFOP0017890 APS (Phase 1 SOI format) and Addendum E (Uganda). "
        f"Suggested upload name: FairBanks - SOI - Addendum E.pdf",
        size=10, italic=True, after=4,
    )
    para(CALL_URL, size=10, after=2)
    para(GRANTS_URL, size=10, after=8)
    para(SLOGAN, size=12, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    """Submission-shaped PDF: letter, 1\" margins, 12-pt Times-Roman."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, black
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak, KeepTogether,
    )
    from PIL import Image as PILImage

    slate = HexColor("#" + SLATE)
    line = HexColor("#888888")
    cream = HexColor("#E8EEEE")

    st = getSampleStyleSheet()
    for name, kw in [
        ("CoverProg", dict(fontName="Times-Bold", fontSize=11, leading=14, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("CoverTitle", dict(fontName="Times-Bold", fontSize=14, leading=17, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("CoverOrg", dict(fontName="Times-BoldItalic", fontSize=12, leading=15, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("CoverSub", dict(fontName="Times-Italic", fontSize=11, leading=14, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("Slogan", dict(fontName="Times-BoldItalic", fontSize=12, leading=15, textColor=black, alignment=TA_CENTER, spaceAfter=10)),
        ("H1", dict(fontName="Times-Bold", fontSize=12, leading=15, textColor=black, spaceBefore=10, spaceAfter=6, alignment=TA_LEFT)),
        ("Body", dict(fontName="Times-Roman", fontSize=12, leading=15, textColor=black, alignment=TA_JUSTIFY, spaceAfter=6)),
        ("BodyItalic", dict(fontName="Times-Italic", fontSize=12, leading=15, textColor=black, alignment=TA_JUSTIFY, spaceAfter=6)),
        ("FBullet", dict(fontName="Times-Roman", fontSize=12, leading=15, textColor=black, leftIndent=14, spaceAfter=3)),
        ("Meta", dict(fontName="Times-Italic", fontSize=10, leading=12, textColor=slate, alignment=TA_CENTER, spaceAfter=4)),
        ("CellHead", dict(fontName="Times-Bold", fontSize=10, leading=12, textColor=black)),
        ("CellBody", dict(fontName="Times-Roman", fontSize=10, leading=12, textColor=black)),
        ("Caption", dict(fontName="Times-Italic", fontSize=10, leading=12, textColor=slate, alignment=TA_CENTER, spaceAfter=8)),
        ("Link", dict(fontName="Times-Roman", fontSize=10, leading=12, textColor=slate, spaceAfter=2)),
    ]:
        st.add(ParagraphStyle(name, **kw))

    pw = letter[0] - 2 * inch
    story = []

    def img(key, w=pw * 0.92, cap=None, max_h=2.4 * inch):
        path = embed(key)
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / iw
        h = min(w * aspect, max_h)
        w = h / aspect if h == max_h else w
        block = [Image(path, width=w, height=h)]
        if cap:
            block.append(Paragraph(cap, st["Caption"]))
        story.append(KeepTogether(block))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        data += [[Paragraph(str(c), st["CellBody"]) for c in row] for row in rows]
        t = Table(data, colWidths=widths or [pw / len(headers)] * len(headers), repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), cream),
            ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    # Cover
    story.append(Paragraph(PROGRAMME, st["CoverProg"]))
    story.append(Paragraph(DOC_TITLE, st["CoverTitle"]))
    story.append(Paragraph(ORG_LINE, st["CoverOrg"]))
    story.append(Paragraph(SUBTITLE, st["CoverSub"]))
    story.append(Paragraph(SLOGAN, st["Slogan"]))
    story.append(Paragraph("Table Listing of Critical Details", st["H1"]))
    tbl(["Item", "Detail"], CRITICAL_DETAILS, [pw * 0.32, pw * 0.68])
    story.append(PageBreak())

    story.append(Paragraph("1. Issue / Challenge / Opportunity", st["H1"]))
    for p in ISSUE:
        story.append(Paragraph(p, st["Body"]))

    story.append(Paragraph("2. Proposed Solution / Activities", st["H1"]))
    for p in SOLUTION:
        story.append(Paragraph(p, st["Body"]))
    story.append(Paragraph(WIN_WIN_INTRO, st["BodyItalic"]))
    tbl(["Who gains", "What they gain"], WIN_WIN, [pw * 0.28, pw * 0.72])

    story.append(Paragraph("<b>Core activities:</b>", st["Body"]))
    for a in ACTIVITIES:
        story.append(Paragraph(f"• {a}", st["FBullet"]))

    story.append(Paragraph("<b>Project approach:</b>", st["Body"]))
    for p in APPROACH:
        story.append(Paragraph(p, st["Body"]))

    story.append(Paragraph("3. Anticipated Outcomes and Results", st["H1"]))
    story.append(Paragraph(
        "Success means lives protected, systems strengthened, efficiency improved, "
        "self-reliance advanced, and U.S. investments secured through GOU-owned "
        "digital continuity:",
        st["Body"],
    ))
    tbl(["Outcome lens", "What success looks like"], OUTCOMES, [pw * 0.24, pw * 0.76])

    story.append(Paragraph("4. Organizational Capacity", st["H1"]))
    for c in CAPACITY:
        story.append(Paragraph(f"• {c}", st["FBullet"]))

    story.append(Paragraph("5. Interoperability Pathway (chart)", st["H1"]))
    tbl(["Platform", "Integration pathway"], INTEROP_ROWS, [pw * 0.32, pw * 0.68])
    img(
        "architecture",
        cap="Figure 1. Community cascade to FCHIP to national digital health systems "
            "(chart; excluded from page limit).",
        max_h=2.2 * inch,
    )

    story.append(Paragraph("6. Alignment with Addendum E Guiding Principles (chart)", st["H1"]))
    tbl(["Principle", "How FairBanks responds"], PRINCIPLES, [pw * 0.32, pw * 0.68])

    story.append(Paragraph("7. List of Partner Roles and Responsibilities (preliminary)", st["H1"]))
    story.append(Paragraph(
        "Partners below are preliminary characterisations for Phase 1. Roles, contacts, "
        "and any subrecipient federal-share splits will be refined in Phase 2 / co-design "
        "if invited. At SOI stage, FairBanks Medical Centre is proposed as prime applicant.",
        st["Body"],
    ))
    tbl(["Partner", "Proposed role (summary)"], PARTNERS, [pw * 0.34, pw * 0.66])

    story.append(Paragraph("8. Resource Contributions and/or Cost Share", st["H1"]))
    story.append(Paragraph(COST_SHARE_NOTE, st["Body"]))

    story.append(Paragraph("9. Readiness for Phase 2", st["H1"]))
    story.append(Paragraph(CLOSING, st["Body"]))
    story.append(Paragraph(
        "References: DFOP0017890 APS (Phase 1 SOI format) and Addendum E (Uganda). "
        "Suggested upload name: FairBanks - SOI - Addendum E.pdf",
        st["Meta"],
    ))
    story.append(Paragraph(CALL_URL, st["Link"]))
    story.append(Paragraph(GRANTS_URL, st["Link"]))
    story.append(Spacer(1, 10))
    story.append(Paragraph(SLOGAN, st["Slogan"]))

    def _page(canvas, doc_):
        canvas.saveState()
        canvas.setFont("Times-Roman", 10)
        canvas.drawCentredString(letter[0] / 2, 0.55 * inch, f"Page {doc_.page}")
        canvas.restoreState()

    OUT.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(
        str(OUT_PDF),
        pagesize=letter,
        leftMargin=1 * inch,
        rightMargin=1 * inch,
        topMargin=1 * inch,
        bottomMargin=1 * inch,
    ).build(story, onFirstPage=_page, onLaterPages=_page)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    """Elegant briefing deck — same win-win story; not the official upload format."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def _add_entrance_anims(slide):
        from lxml import etree
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spids = []
        for shape in slide.shapes:
            try:
                stype = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                stype = -1
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
            if stype == 13 or has_text:
                spids.append(str(shape.shape_id))
            if len(spids) >= 3:
                break
        if not spids:
            return
        sld = slide._element
        for old in sld.findall(f"{{{NS_P}}}timing"):
            sld.remove(old)
        children = []
        nid = 3
        for i, spid in enumerate(spids):
            delay = 0 if i == 0 else 200
            children.append(
                f'<p:par xmlns:p="{NS_P}">'
                f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
                f'fill="hold" nodeType="withEffect">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
            )
            nid += 2
        xml = (
            f'<p:timing xmlns:p="{NS_P}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        )
        sld.append(etree.fromstring(xml))

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH, blank = prs.slide_width, prs.slide_height, prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        s.line.color.rgb = C(line) if line else C(fill)
        if not line:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def header(sl, roman, title):
        rect(sl, 0, 0, SW, Inches(0.08), TEAL)
        tb(sl, Inches(0.5), Inches(0.2), Inches(1.2), Inches(0.5), roman, size=28, bold=True, color=ACCENT)
        tb(sl, Inches(1.6), Inches(0.22), Inches(11), Inches(0.55), title, size=22, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.28), SW, Inches(0.28), SLATE)
        tb(sl, Inches(0.4), SH - Inches(0.26), Inches(11), Inches(0.22),
           f"FairBanks SOI  |  DFOP0017890 Objective 4  |  {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.26), Inches(0.5), Inches(0.22),
           str(n), size=9, color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Cover
    s = prs.slides.add_slide(blank)
    pic(s, "cover", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.6), SW, Inches(3.6), NAVY)
    tb(s, Inches(0.55), SH - Inches(3.3), Inches(12), Inches(0.3), PROGRAMME, size=11, bold=True, color=TEAL)
    tb(s, Inches(0.55), SH - Inches(2.85), Inches(12), Inches(0.55), DOC_TITLE, size=26, bold=True, color="FFFFFF")
    tb(s, Inches(0.55), SH - Inches(2.2), Inches(12), Inches(0.35),
       ORG_LINE, size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.55), SH - Inches(1.7), Inches(12), Inches(0.45),
       "Objective 4 — Last-mile bridge into eCHIS, EMR, DHIS2/eIDSR, and NDHA",
       size=12, italic=True, color="D0E8E8")
    tb(s, Inches(0.55), SH - Inches(1.1), Inches(12), Inches(0.3),
       f"Win-win Phase 1 SOI  |  {FEDERAL_SHARE}  |  {PROJECT_LENGTH}  |  Due 31 July 2026",
       size=12, color="FFFFFF")
    tb(s, Inches(0.55), SH - Inches(0.7), Inches(12), Inches(0.3), SLOGAN, size=13, bold=True, color="FFFFFF")

    win_body = "\n".join(
        f"• {row[0]}: {row[1][:95]}..." if len(row[1]) > 95 else f"• {row[0]}: {row[1]}"
        for row in WIN_WIN
    )
    act_body = "\n".join(f"• {a[:115]}..." if len(a) > 115 else f"• {a}" for a in ACTIVITIES[:4])
    out_body = "\n".join(f"• {r[0]} — {r[1][:85]}" for r in OUTCOMES)

    sections = [
        ("I", "The Problem", ISSUE[0] + "\n\n" + ISSUE[1], "outreach"),
        ("II", "Statement of Interest", SOLUTION[1], "architecture"),
        ("III", "Win-Win Value", win_body, None),
        ("IV", "Core Activities", act_body, "mobile"),
        ("V", "Project Approach",
         "Cascade-rooted capture → secure sync → eCHIS / EMR / DHIS2-eIDSR "
         "interoperability → dashboards & GIS → MoH/NDHA-aligned transfer package.",
         "dashboard"),
        ("VI", "Expected Outcomes", out_body, "gis"),
        ("VII", "Organizational Capacity",
         "\n".join(f"• {c}" for c in CAPACITY[:4]), "facility"),
        ("VIII", "Interoperability",
         "eCHIS household capture → facility EMR (eAfya/ClinicMaster) → DHIS2/eIDSR → "
         "National Data Warehouse readiness → iHRIS-aware workforce. Bridge layer — "
         "not a replacement silo. Aligned with NDHA / OpenHIE.",
         "pharmacy"),
        ("IX", "Sustainability & Ownership",
         "Documented transition of staffing and recurrent costs toward GOU ownership "
         "by project close.\n\n"
         "Train Ugandan developers, CHW supervisors, and data stewards.\n\n"
         "Clinic and outreach continue as FairBanks operations — the grant accelerates "
         "interoperable intelligence and transfer.",
         None),
        ("X", "Readiness for Phase 2", CLOSING, None),
    ]

    for i, (roman, title, body, img_key) in enumerate(sections, start=2):
        s = prs.slides.add_slide(blank)
        header(s, roman, title)
        if img_key:
            pic(s, img_key, Inches(0.45), Inches(1.05), Inches(5.8), Inches(5.9))
            tb(s, Inches(6.5), Inches(1.15), Inches(6.4), Inches(5.6), body, size=17, color=SLATE)
        else:
            tb(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(5.8), body, size=18, color=SLATE)
        footer(s, i)

    OUT.mkdir(parents=True, exist_ok=True)
    for _sl in prs.slides:
        _add_entrance_anims(_sl)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done.", OUT)
