#!/usr/bin/env python3
"""Build CDC-RFA-JG-26-0054 FairBanks application pack (Uganda GHS).

Creates one synchronized set:
  documents/cdc-uganda-ghs_word.docx
  documents/cdc-uganda-ghs_pdf.pdf
  documents/cdc-uganda-ghs_ppt.pptx

Curated copy-paste answers live in application_answers.md (not overwritten).
Run: python applications/cdc-uganda-ghs/build_cdc-uganda-ghs_docs.py
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
ASSETS = REPO / "assets"
OUT = HERE / "documents"
SLUG = "cdc-uganda-ghs"
OUT_DOC = OUT / f"{SLUG}_word.docx"
OUT_PDF = OUT / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT / f"{SLUG}_ppt.pptx"
PASTE = HERE / "application_answers.md"

NAVY, TEAL, ORANGE = "0A1F2E", "0D6E6E", "C45C26"
SLATE, MUTED, CREAM, LINE, WHITE = "1E2F38", "52636C", "F7F5F0", "CED9D8", "FFFFFF"
SLOGAN = "Your health, our mission."

PROGRAMME = "CDC-RFA-JG-26-0054 — Strengthening global health security through local partnerships in Uganda"
DOC_TITLE = "Win-ready project narrative pack (submission companion)"
ORG_LEGAL = "FAIRBANKS MEDICAL CENTRE LIMITED"
ORG = "FairBanks Medical Centre"
CONTACT_NAME = "Racheal Nabukeera"
CONTACT_TITLE = "Managing Director and Co-founder"
EMAIL = "info@fairbanksmedicalcentre.org"
PHONE = "+256 772 849 258"
WEBSITE = "https://www.fairbanksmedicalcentre.org/"
COMPANY_NO = "80020003843337"
TIN = "1053370026"
NSSF = "NS043295"
LOCATION = (
    "Plot 1423 and 1425 Tirupati Road, Fairbanks Medical Centre, "
    "Kampala Central Division, Kololo IV, Lugogo, Kampala, Uganda"
)
PUBLIC_LOCATION = "Kyebando–Kisalosalo, Northern Bypass, Kampala"
OFFICIAL = "https://www.grants.gov/search-results-detail/360339"
OFFICIAL_ALT = "https://simpler.grants.gov/opportunity/264249e6-fdbb-4b1c-ac90-23b7d9b07b1b"
APPLY_GUIDE = "https://grants.gov/quick-start-guide/applicants"
DEADLINE = "14 August 2026, 11:59 p.m. ET"
C1_ASK = "$2,000,000"
PERIOD = "5 years (five 12-month budget periods)"
TOTAL_ASK = "$3,300,000"
WIN_LINE = (
    "CDC and Uganda need outbreaks stopped closer to where people live. "
    "FairBanks already runs a clinic, CHW/VHT outreach, and working FCHIP tools — "
    "ready to feed MoH/NISS pathways under government leadership."
)

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "outreach": "outreach_bp_screening.jpeg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "dashboard": "dashboard_demo.png",
    "gis": "gis_hotspots.png",
    "architecture": "data_flow_iso_labeled.png",
    "facility": "facility_exterior_entrance_01.jpg",
    "community": "outreach_audience_full_group_01.jpg",
    "team": "staff_team_reception.jpeg",
}
CONCEPT = REPO / ".cursor" / "concept_improved.jpeg"


def photo(key: str) -> Path:
    path = ASSETS / PHOTOS[key]
    if not path.exists():
        raise FileNotFoundError(path)
    return path


CRITICAL = [
    ["Opportunity number", "CDC-RFA-JG-26-0054"],
    ["Applicant", ORG_LEGAL],
    ["Company No. / TIN / NSSF", f"{COMPANY_NO} · {TIN} · {NSSF}"],
    ["Country of performance", "Uganda"],
    ["Funding instrument", "Cooperative agreement"],
    ["Assistance listing", "93.318"],
    ["Period of performance", PERIOD],
    ["Year 1 total federal ask (draft)", f"{TOTAL_ASK} (CONFIRM; all components)"],
    ["Year 1 Component 1 ask (draft)", f"{C1_ASK} (CONFIRM; ceiling $5,000,000)"],
    ["Cost share", "None proposed"],
    ["Local partner preference", "Yes — Uganda-incorporated entity (+15 points; document before submit)"],
    ["Merit scoring focus", "Approach 35 · EPMP 40 · Capacity 25"],
    ["Deadline", DEADLINE],
    ["Contact", f"{CONTACT_NAME} · {EMAIL} · {PHONE}"],
]

SHARED_WIN = [
    ["CDC / United States", "Earlier containment at source — Safer, Stronger, More Prosperous"],
    ["Uganda / MoH", "Last-mile signals into NISS/NAPHS II; trained CHWs; 7-1-7 practice"],
    ["FairBanks", "Local partner with live clinic, Community Reach, and FCHIP that feeds government"],
    ["Communities", "Earlier alerts and clearer referral paths when something looks wrong"],
]

ABSTRACT = (
    f"{ORG_LEGAL} ({ORG}) is a Uganda-registered health organisation in Kampala. "
    f"We run a licensed medical centre and FairBanks Community Reach with CHWs and VHTs "
    f"in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya and nearby communities. "
    f"Slogan: {SLOGAN} "
    "Uganda’s 2023 Joint External Evaluation still flags last-mile gaps: community and "
    "facility signals that reach decision-makers too late; limited surge-ready workforce at "
    "subnational level; and incomplete links into national surveillance and emergency "
    "structures, including National Integrated Surveillance System (NISS) pathways. "
    "Under Ministry of Health leadership, Year 1 Component 1 will (1) strengthen community "
    "and facility surveillance that contributes to MoH/NISS channels using FairBanks FCHIP "
    "tools, (2) speed community-to-district detection and response using 7-1-7 style timing, "
    "(3) train CHWs/VHTs and frontline workers for One Health–aware surveillance and surge "
    "readiness, and (4) tighten community–facility links for priority diseases during routine "
    "work and outbreaks. Laboratory and border-health modules will be delivered with "
    "MoH-aligned partners. Components 2–5 are contingency surge plans that may be approved "
    "but unfunded until CDC activates emergency funding. Shared result: threats found and "
    "acted on closer to source — safer for Uganda and for the United States."
)

PROBLEM = [
    "Uganda faces repeated infectious disease threats shaped by geography, rapid urban growth, "
    "and high population movement. Outbreaks that start in peri-urban communities can grow "
    "before national systems see them. Americans and Ugandans are safer when threats are found "
    "and contained close to source.",
    "In FairBanks catchments, CHWs and VHTs already visit homes and schools, and the Medical "
    "Centre sees patients every day. Too often, community reports, clinic records, and lab "
    "referrals still sit apart. The 2023 JEE praised progress and still named gaps in "
    "information sharing, subnational surge capacity, and use of data for early action — "
    "including the need for stronger integrated surveillance such as NISS.",
    "FairBanks proposes to close the last-mile gap under MoH leadership: structured community "
    "and facility signals, faster notify–respond loops with districts, trained frontline "
    "workers, and tools that feed government systems rather than create a parallel silo.",
]

APPROACH = [
    WIN_LINE,
    "FairBanks Community Reach cascade stays the operating model: community members -> "
    "CHWs/VHTs -> Community Reach programmes -> Medical Centre -> research and skills -> "
    "economic empowerment. FCHIP is the intelligence component on the Data & Feedback loop.",
    "Award funds will support surveillance systems, training, emergency readiness, data "
    "linkages, and public health programme coordination. Routine clinical care continues on "
    "FairBanks’ own operations and is not the purpose of this cooperative agreement.",
    "We will collaborate with MoH, district/KCCA teams, other CDC partners, and community "
    "structures. Software, SOPs, and training materials developed under the award will be "
    "available to MoH and CDC for appropriate use.",
]

STRATEGIES = [
    ["1 Emergency management", "District notification SOPs, 7-1-7 drills, catchment IMS practice, AARs"],
    ["2 Workforce", "CHW/VHT and facility training; surge roster; One Health–aware sessions with districts"],
    ["3 Laboratory (support)", "Sample referral pathways; private-facility incident reporting support; biosafety basics"],
    ["4 Surveillance", "FCHIP capture -> MoH/NISS-aligned exports; data quality; simple GIS early warning"],
    ["5 Border health (support)", "RCCE and signal sharing for high-mobility corridors with MoH partners"],
    ["6 Public health programmes", "Community–facility links for HIV, TB, malaria, VHFs, mpox, immunisation campaigns"],
]

OUTCOMES = [
    ["Faster detection & response", "Shorter time from community signal to district notification (7-1-7 style)"],
    ["Stronger surveillance", "Weekly community reports and successful MoH/NISS-aligned data tests"],
    ["Ready workforce", "Trained CHWs/VHTs and surge roster exercised in drills"],
    ["Better community–facility links", "Documented referral outcomes for priority disease flags"],
    ["Government use", "Tools, SOPs, and data access handed for MoH/CDC use"],
]

CAPACITY = [
    f"{ORG_LEGAL} (No. {COMPANY_NO}; TIN {TIN}; NSSF {NSSF}) — Uganda company with principal place of business in Kampala.",
    f"Live Medical Centre and Pharmacy; Community Reach with CHWs/VHTs; working FCHIP MVP.",
    f"Public community base: {PUBLIC_LOCATION}. Catchments include Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya.",
    f"Lead: {CONTACT_NAME}, {CONTACT_TITLE} — 15+ years in Uganda private health leadership; "
    "MA Social Sector Planning (Makerere); PhD in progress; Uganda Healthcare Federation links.",
    "Company evidence pack available (incorporation, TIN, NSSF, licences, insurance) — see company-docs/.",
    "Honest note: we have not yet been prime on a multi-million CDC GHS award; financial controls "
    "and staffing for federal compliance will be strengthened as described in attachments (CONFIRM).",
]

INDICATORS = [
    "DGHP partner-level indicators from the FY26 list (surveillance, community mitigation, emergency ops, IPC, lab support) — final set agreed with CDC after award.",
    "Project measures with Year 1 targets: median days signal-to-notify (toward 7-1-7); at least 80% active CHWs with complete weekly reports by Q4; at least 70% priority referrals documented; at least 2 successful MoH export tests; at least 2 AARs.",
    "About 5–10% of funds for M&E; detailed EPMP and DMP within six months of award.",
]

BUDGET_ROWS = [
    ["1 Core GHS", C1_ASK, "$5,000,000", "Expected initial funding"],
    ["2 Small-scale response", "$350,000", "$10,000,000", "Contingency"],
    ["3 Large-scale response", "$450,000", "$15,000,000", "Contingency"],
    ["4 Emerging threats", "$250,000", "$15,000,000", "Contingency"],
    ["5 Humanitarian", "$250,000", "$20,000,000", "Contingency"],
    ["TOTAL Year 1", TOTAL_ASK, "", "All components"],
]


def build_docx() -> None:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import parse_xml

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    def set_run(run, size=12, bold=False, italic=False, color=SLATE):
        run.font.name = "Times New Roman"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def add_p(text, size=12, bold=False, italic=False, center=False, space_after=6, color=SLATE):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(space_after)
        p.paragraph_format.space_before = Pt(0)
        if center:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        set_run(run, size=size, bold=bold, italic=italic, color=color)
        return p

    def add_h(text):
        return add_p(text, size=12, bold=True, space_after=8)

    def add_bullets(items):
        for item in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            run = p.add_run(item)
            set_run(run, size=12)

    def add_table(headers, rows):
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Table Grid"
        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = ""
            run = cell.paragraphs[0].add_run(h)
            set_run(run, size=10, bold=True)
            shading = parse_xml(
                f'<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" w:fill="{CREAM}"/>'
            )
            cell._tePr = cell._tc.get_or_add_tcPr()
            cell._tc.get_or_add_tcPr().append(shading)
        for r_i, row in enumerate(rows):
            for c_i, val in enumerate(row):
                cell = table.rows[r_i + 1].cells[c_i]
                cell.text = ""
                run = cell.paragraphs[0].add_run(str(val))
                set_run(run, size=10)
        doc.add_paragraph()

    add_p(PROGRAMME, size=11, bold=True, center=True)
    add_p(DOC_TITLE, size=14, bold=True, center=True)
    add_p(f"{ORG} · FCHIP component", size=12, bold=True, italic=True, center=True)
    add_p(SLOGAN, size=12, bold=True, italic=True, center=True, space_after=12)
    try:
        doc.add_picture(str(photo("logo")), width=Inches(1.2))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    except Exception:
        pass

    add_h("Table of critical details")
    add_table(["Item", "Detail"], CRITICAL)

    add_h("Shared win")
    add_table(["Who", "What they gain"], SHARED_WIN)
    add_p(WIN_LINE)

    add_h("Project abstract summary")
    add_p(ABSTRACT)

    add_h("1. Background and approach — problem")
    for para in PROBLEM:
        add_p(para)

    add_h("2. Approach and FairBanks model")
    for para in APPROACH:
        add_p(para)
    if CONCEPT.exists():
        doc.add_picture(str(CONCEPT), width=Inches(6.0))
        cap = doc.add_paragraph()
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run("Figure: FairBanks Community Reach operating model")
        set_run(run, size=10, italic=True, color=MUTED)

    add_h("3. Year 1 strategies and activities (Component 1)")
    add_table(["Strategy", "FairBanks Year 1 focus"], STRATEGIES)

    add_h("4. Anticipated outcomes")
    add_table(["Outcome", "What success looks like"], OUTCOMES)

    add_h("5. Evaluation and performance measurement (summary)")
    add_bullets(INDICATORS)
    add_p(
        "Full EPMP and DMP text for Grants.gov paste is in application_answers.md "
        "(sections D and D7)."
    )

    add_h("6. Organizational capacity")
    add_bullets(CAPACITY)
    add_p(f"Website: {WEBSITE}")
    add_p(f"Registered address: {LOCATION}")

    add_h("7. Draft Year 1 component budgets (CONFIRM before submit)")
    add_table(["Component", "Draft ask", "Ceiling", "Notes"], BUDGET_ROWS)
    add_p(
        "Indirect: foreign organisation 8% of MTDC (CONFIRM). Separate SF-424A lines per "
        "component. Itemised budget narrative categories: salaries, fringe, travel, equipment, "
        "supplies, contractual, other, indirect."
    )

    add_h("8. Pre-submit checklist")
    add_bullets(
        [
            "SAM.gov active + UEI for physical location receiving funds (start now)",
            "Grants.gov / Login.gov active; AOR ready to Sign and Submit",
            f"Follow applicant quick start: {APPLY_GUIDE}",
            "Local partner letter + ownership/staff/board citizenship evidence from company-docs/",
            "Project narrative PDF ≤60 pages; attachments ≤50 pages; English; USD only",
            "All five components included; no research; clinical care not the award purpose",
            f"Submit by {DEADLINE} via Grants.gov",
            f"Official listing: {OFFICIAL}",
            f"Alt listing: {OFFICIAL_ALT}",
        ]
    )

    add_p(
        f"Full copy-paste bank: {PASTE.name}. Generated {date.today().isoformat()}.",
        size=10,
        italic=True,
        color=MUTED,
    )
    add_p(SLOGAN, size=12, bold=True, italic=True, center=True)

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(OUT_DOC)
    print(f"DOCX: {OUT_DOC}")


def build_pdf() -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, black
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Image,
        Table,
        TableStyle,
        PageBreak,
        KeepTogether,
        ListFlowable,
        ListItem,
    )
    from PIL import Image as PILImage

    cream = HexColor("#E8EEEE")
    line = HexColor("#888888")
    slate = HexColor("#" + SLATE)
    st = getSampleStyleSheet()
    for name, kw in [
        ("CoverProg", dict(fontName="Times-Bold", fontSize=10, leading=13, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("CoverTitle", dict(fontName="Times-Bold", fontSize=14, leading=17, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("CoverOrg", dict(fontName="Times-BoldItalic", fontSize=12, leading=15, textColor=black, alignment=TA_CENTER, spaceAfter=4)),
        ("Slogan", dict(fontName="Times-BoldItalic", fontSize=12, leading=15, textColor=black, alignment=TA_CENTER, spaceAfter=10)),
        ("H1", dict(fontName="Times-Bold", fontSize=12, leading=15, textColor=black, spaceBefore=10, spaceAfter=6, alignment=TA_LEFT)),
        ("Body", dict(fontName="Times-Roman", fontSize=12, leading=15, textColor=black, alignment=TA_JUSTIFY, spaceAfter=6)),
        ("Meta", dict(fontName="Times-Italic", fontSize=10, leading=12, textColor=slate, alignment=TA_CENTER, spaceAfter=4)),
        ("CellHead", dict(fontName="Times-Bold", fontSize=9, leading=11, textColor=black)),
        ("CellBody", dict(fontName="Times-Roman", fontSize=9, leading=11, textColor=black)),
        ("Caption", dict(fontName="Times-Italic", fontSize=10, leading=12, textColor=slate, alignment=TA_CENTER, spaceAfter=8)),
    ]:
        st.add(ParagraphStyle(name, **kw))

    pw = letter[0] - 2 * inch
    story = []

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        data += [[Paragraph(str(c), st["CellBody"]) for c in row] for row in rows]
        t = Table(data, colWidths=widths or [pw / len(headers)] * len(headers), repeatRows=1)
        t.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), cream),
                    ("GRID", (0, 0), (-1, -1), 0.4, line),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 4),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                    ("TOPPADDING", (0, 0), (-1, -1), 3),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                ]
            )
        )
        story.append(t)
        story.append(Spacer(1, 8))

    def img_path(path: Path, w=pw * 0.9, max_h=2.3 * inch, cap=None):
        if not path.exists():
            return
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / float(iw)
        h = min(w * aspect, max_h)
        w_use = h / aspect if h == max_h else w
        block = [Image(str(path), width=w_use, height=h)]
        if cap:
            block.append(Paragraph(cap, st["Caption"]))
        story.append(KeepTogether(block))

    story.append(Paragraph(PROGRAMME, st["CoverProg"]))
    story.append(Paragraph(DOC_TITLE, st["CoverTitle"]))
    story.append(Paragraph(f"{ORG} · FCHIP component", st["CoverOrg"]))
    story.append(Paragraph(SLOGAN, st["Slogan"]))
    img_path(photo("logo"), w=1.1 * inch, max_h=1.1 * inch)
    story.append(Paragraph("Table of critical details", st["H1"]))
    tbl(["Item", "Detail"], CRITICAL, [pw * 0.32, pw * 0.68])
    story.append(Paragraph("Shared win", st["H1"]))
    tbl(["Who", "What they gain"], SHARED_WIN, [pw * 0.28, pw * 0.72])
    story.append(Paragraph(WIN_LINE, st["Body"]))
    story.append(PageBreak())

    story.append(Paragraph("Project abstract summary", st["H1"]))
    story.append(Paragraph(ABSTRACT, st["Body"]))

    story.append(Paragraph("1. Background and approach — problem", st["H1"]))
    for para in PROBLEM:
        story.append(Paragraph(para, st["Body"]))
    img_path(photo("community"), cap="Community Reach in Kampala peri-urban catchments", max_h=2.0 * inch)

    story.append(Paragraph("2. Approach and FairBanks model", st["H1"]))
    for para in APPROACH:
        story.append(Paragraph(para, st["Body"]))
    if CONCEPT.exists():
        img_path(CONCEPT, cap="FairBanks Community Reach operating model", max_h=2.4 * inch)

    story.append(Paragraph("3. Year 1 strategies (Component 1)", st["H1"]))
    tbl(["Strategy", "FairBanks Year 1 focus"], STRATEGIES, [pw * 0.28, pw * 0.72])
    img_path(photo("architecture"), cap="FCHIP data flow — community and facility signals to decision support", max_h=2.1 * inch)

    story.append(Paragraph("4. Anticipated outcomes", st["H1"]))
    tbl(["Outcome", "What success looks like"], OUTCOMES, [pw * 0.28, pw * 0.72])
    img_path(photo("dashboard"), cap="Programme dashboard concept for catchment review", max_h=1.9 * inch)

    story.append(Paragraph("5. Evaluation and performance measurement (summary)", st["H1"]))
    items = [ListItem(Paragraph(x, st["Body"]), leftIndent=12, value="•") for x in INDICATORS]
    story.append(ListFlowable(items, bulletType="bullet", start="•"))
    story.append(Paragraph(
        "Paste the full EPMP and DMP from application_answers.md into the Grants.gov narrative.",
        st["Body"],
    ))

    story.append(Paragraph("6. Organizational capacity", st["H1"]))
    items = [ListItem(Paragraph(x, st["Body"]), leftIndent=12, value="•") for x in CAPACITY]
    story.append(ListFlowable(items, bulletType="bullet", start="•"))
    img_path(photo("facility"), cap="FairBanks Medical Centre — clinical anchor", max_h=1.9 * inch)

    story.append(Paragraph("7. Draft Year 1 component budgets (CONFIRM)", st["H1"]))
    tbl(["Component", "Draft ask", "Ceiling", "Notes"], BUDGET_ROWS, [pw * 0.22, pw * 0.22, pw * 0.22, pw * 0.34])

    story.append(Paragraph("8. Contacts and sources", st["H1"]))
    story.append(Paragraph(
        f"{CONTACT_NAME}, {CONTACT_TITLE}<br/>{EMAIL} · {PHONE}<br/>{WEBSITE}<br/>"
        f"Grants.gov: {OFFICIAL}<br/>"
        f"Simpler listing: {OFFICIAL_ALT}<br/>"
        f"Applicant quick start: {APPLY_GUIDE}<br/>"
        f"Grantor email: DGHPNOFOs@cdc.gov<br/>"
        f"Companion copy-paste file: application_answers.md · Generated {date.today().isoformat()}.",
        st["Body"],
    ))
    story.append(Spacer(1, 8))
    story.append(Paragraph(SLOGAN, st["Slogan"]))

    def _page(canvas, doc_):
        canvas.saveState()
        canvas.setFont("Times-Roman", 10)
        canvas.drawCentredString(letter[0] / 2, 0.55 * inch, f"Page {doc_.page}")
        canvas.restoreState()

    OUT.mkdir(parents=True, exist_ok=True)
    tmp_pdf = REPO / "tmp" / SLUG / f"{SLUG}_pdf.build.pdf"
    tmp_pdf.parent.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(
        str(tmp_pdf),
        pagesize=letter,
        leftMargin=1 * inch,
        rightMargin=1 * inch,
        topMargin=1 * inch,
        bottomMargin=1 * inch,
    ).build(story, onFirstPage=_page, onLaterPages=_page)
    try:
        tmp_pdf.replace(OUT_PDF)
    except OSError:
        import shutil

        try:
            shutil.copy2(tmp_pdf, OUT_PDF)
        except OSError as e:
            print(f"PDF built at {tmp_pdf} (could not write {OUT_PDF}: {e})")
            return
    print(f"PDF: {OUT_PDF}")


def build_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from lxml import etree

    def C(h):
        return RGBColor.from_string(h)

    def add_fade(slide):
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spids = []
        for shape in slide.shapes:
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
            try:
                stype = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                stype = -1
            if stype == 13 or has_text:
                spids.append(str(shape.shape_id))
            if len(spids) >= 3:
                break
        if not spids:
            return
        sld = slide._element
        for old in sld.findall(f"{{{NS_P}}}timing"):
            sld.remove(old)
        children = []
        nid = 3
        for i, spid in enumerate(spids):
            delay = 0 if i == 0 else 200
            children.append(
                f'<p:par xmlns:p="{NS_P}">'
                f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
                f'fill="hold" nodeType="withEffect">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                f"</p:animEffect></p:childTnLst></p:cTn></p:par>"
            )
            nid += 2
        xml = (
            f'<p:timing xmlns:p="{NS_P}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
            f"</p:childTnLst></p:cTn>"
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f"</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"
        )
        sld.append(etree.fromstring(xml))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        s.line.fill.background()
        return s

    def textbox(sl, x, y, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = C(color)
        run.font.name = "Calibri"
        return box

    def bullets(sl, x, y, w, h, lines, size=16, color=SLATE):
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(size)
            p.font.color.rgb = C(color)
            p.font.name = "Calibri"
            p.space_after = Pt(6)
        return box

    def pic(sl, path, x, y, w, h):
        sl.shapes.add_picture(str(path), x, y, width=w, height=h)

    def footer(sl, page, total=12):
        textbox(sl, Inches(0.5), SH - Inches(0.45), Inches(10), Inches(0.3), SLOGAN, size=10, color=MUTED)
        textbox(
            sl,
            SW - Inches(2.2),
            SH - Inches(0.45),
            Inches(1.8),
            Inches(0.3),
            f"{page} / {total}",
            size=10,
            color=MUTED,
            align=PP_ALIGN.RIGHT,
        )

    # 1 Cover
    sl = prs.slides.add_slide(blank)
    pic(sl, photo("cover"), Emu(0), Emu(0), SW, SH)
    rect(sl, Emu(0), Emu(0), SW, SH, "0A1F2E")
    pic(sl, photo("cover"), Emu(0), Emu(0), Inches(6.2), SH)
    rect(sl, Inches(5.8), Emu(0), SW - Inches(5.8), SH, NAVY)
    textbox(sl, Inches(6.2), Inches(1.4), Inches(6.5), Inches(0.8), "CDC · Uganda GHS", size=16, color=ORANGE)
    textbox(sl, Inches(6.2), Inches(2.1), Inches(6.5), Inches(1.8), "A local partner for earlier detection — shared win for Uganda and the U.S.", size=26, bold=True)
    textbox(sl, Inches(6.2), Inches(4.1), Inches(6.5), Inches(0.8), f"{ORG} · FCHIP component", size=16)
    slogan_box = textbox(sl, Inches(6.2), Inches(4.9), Inches(6.5), Inches(0.5), SLOGAN, size=14)
    slogan_box.text_frame.paragraphs[0].runs[0].font.italic = True
    textbox(sl, Inches(6.2), Inches(5.6), Inches(6.5), Inches(0.6), f"CDC-RFA-JG-26-0054 · Due {DEADLINE}", size=12, color="B8C4C8")
    add_fade(sl)

    # 2 Shared win
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), ORANGE)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "The shared win", size=24, bold=True)
    bullets(
        sl,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.5),
        [f"{who}: {gain}" for who, gain in SHARED_WIN],
        size=20,
    )
    footer(sl, 2)
    add_fade(sl)

    # 3 Opportunity
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "The opportunity", size=24, bold=True)
    bullets(
        sl,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5),
        [
            "Strengthen Uganda GHS capacities from the 2023 JEE (surveillance, labs, response, workforce, border health).",
            "Cooperative agreement · 5 years · apply for all five funding components.",
            f"Year 1 Component 1 draft ask: {C1_ASK} (ceiling $5M) · local partner preference +15 points.",
            "Scoring: Approach 35 · EPMP 40 · Capacity 25 — we lead with measurable last-mile surveillance.",
            "Must work in Uganda · English · USD · no research · clinical care not the award purpose.",
            f"Submit on Grants.gov by {DEADLINE}.",
        ],
        size=17,
    )
    footer(sl, 3)
    add_fade(sl)

    # 4 Problem
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "The last-mile gap", size=24, bold=True)
    pic(sl, photo("outreach"), Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.2))
    bullets(
        sl,
        Inches(6.6),
        Inches(1.4),
        Inches(6),
        Inches(5),
        [
            "Outbreaks grow when community signals stay on paper.",
            "CHWs, clinics, labs, and EOCs still share information too slowly.",
            "JEE gaps: data sharing, surge staffing, forecasting, NISS links.",
            "FairBanks already works where the gap shows up first.",
        ],
        size=17,
    )
    footer(sl, 4)
    add_fade(sl)

    # 5 Model
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), TEAL)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "How FairBanks works", size=24, bold=True)
    if CONCEPT.exists():
        pic(sl, CONCEPT, Inches(0.8), Inches(1.2), Inches(11.6), Inches(5.4))
    footer(sl, 5)
    add_fade(sl)

    # 6 Solution
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "Our offer under MoH leadership", size=24, bold=True)
    pic(sl, photo("mobile"), Inches(0.4), Inches(1.2), Inches(5.5), Inches(5.2))
    bullets(
        sl,
        Inches(6.2),
        Inches(1.3),
        Inches(6.5),
        Inches(5.5),
        [
            "Community + facility surveillance into MoH/NISS pathways via FCHIP.",
            "7-1-7 style drills with districts.",
            "CHW/VHT workforce and surge roster.",
            "Priority-disease community–facility links.",
            "Lab and border modules with specialised partners — not a parallel system.",
        ],
        size=17,
    )
    footer(sl, 6)
    add_fade(sl)

    # 7 Architecture
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), TEAL)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "FCHIP on the Data & Feedback loop", size=24, bold=True)
    pic(sl, photo("architecture"), Inches(1.2), Inches(1.2), Inches(10.8), Inches(5.4))
    footer(sl, 7)
    add_fade(sl)

    # 8 Outcomes
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "What success looks like", size=24, bold=True)
    pic(sl, photo("gis"), Inches(7.2), Inches(1.3), Inches(5.5), Inches(5.0))
    bullets(
        sl,
        Inches(0.6),
        Inches(1.4),
        Inches(6.2),
        Inches(5),
        [f"{a} — {b}" for a, b in OUTCOMES],
        size=16,
    )
    footer(sl, 8)
    add_fade(sl)

    # 9 Budget components
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), ORANGE)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "Five components — one application", size=24, bold=True)
    bullets(
        sl,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.5),
        [f"Component {c}: {ask} draft (ceiling {ceil}) — {note}" for c, ask, ceil, note in BUDGET_ROWS],
        size=17,
    )
    footer(sl, 9)
    add_fade(sl)

    # 10 Capacity
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), TEAL)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "Why FairBanks — local partner", size=24, bold=True)
    pic(sl, photo("team"), Inches(0.4), Inches(1.2), Inches(5.5), Inches(5.2))
    bullets(
        sl,
        Inches(6.2),
        Inches(1.3),
        Inches(6.5),
        Inches(5.5),
        [
            "Uganda-incorporated company with Kampala operations.",
            "Live clinic + Community Reach + working FCHIP MVP.",
            "Named catchments already in motion.",
            "Documented company standing (TIN, NSSF, licences) ready to attach.",
            "Contact: info@fairbanksmedicalcentre.org · +256 772 849 258",
        ],
        size=17,
    )
    footer(sl, 10)
    add_fade(sl)

    # 11 Grants.gov prep
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, CREAM)
    rect(sl, Emu(0), Emu(0), SW, Inches(0.9), NAVY)
    textbox(sl, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5), "Before you submit — Grants.gov", size=24, bold=True)
    bullets(
        sl,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.5),
        [
            "SAM.gov first - get UEI (can take weeks).",
            "Login.gov + Grants.gov organisation profile with the same email as SAM EBiz POC.",
            "Create a workspace for CDC-RFA-JG-26-0054; AOR must Sign and Submit.",
            "CONFIRM local partner citizenship percentages before signing the preference letter.",
            f"Deadline: {DEADLINE} · Quick start: grants.gov/quick-start-guide/applicants",
        ],
        size=18,
    )
    footer(sl, 11)
    add_fade(sl)

    # 12 Ask / close
    sl = prs.slides.add_slide(blank)
    rect(sl, Emu(0), Emu(0), SW, SH, NAVY)
    textbox(sl, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.4), "Ready to help Uganda detect and respond faster — so threats stop closer to source", size=26, bold=True, align=PP_ALIGN.CENTER)
    textbox(
        sl,
        Inches(1.5),
        Inches(3.4),
        Inches(10),
        Inches(1.2),
        f"Draft Year 1 Component 1: {C1_ASK} · Local partner · MoH-aligned · Contingency surge plans ready",
        size=16,
        align=PP_ALIGN.CENTER,
        color="B8C4C8",
    )
    textbox(
        sl,
        Inches(1.5),
        Inches(5.0),
        Inches(10),
        Inches(1.0),
        f"{CONTACT_NAME} · {EMAIL} · {PHONE}\n{SLOGAN}",
        size=16,
        align=PP_ALIGN.CENTER,
    )
    add_fade(sl)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    print(f"PPTX: {OUT_PPT}")


def main() -> None:
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. Paste-ready answers:", PASTE)


if __name__ == "__main__":
    main()
