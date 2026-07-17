"""
FID Call for Proposals 2026 — FairBanks / FCHIP proposal document set.

Generates one synchronised content set in three formats:
  - applications/fid/documents/fid_word.docx   (full narrative)
  - applications/fid/documents/fid_pdf.pdf     (full narrative, print-ready)
  - applications/fid/documents/fid_ppt.pptx    (12-slide summary deck)

Source of truth: https://opportunitiesforyouth.org/2026/06/12/fid-call-for-proposals-2026-...
Content aligns with FID evaluation criteria: (1) Evidence of Impact,
(2) Cost & Cost-Effectiveness, (3) Scale & Financial Sustainability.

Run:  python applications/fid/build_fid_docs.py
"""

from pathlib import Path

# ---------------------------------------------------------------------------
# Paths & brand
# ---------------------------------------------------------------------------
FID = Path(__file__).resolve().parent
REPO = FID.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = FID / "documents"
OUT_DOC = OUT_DIR / "fid_word.docx"
OUT_PDF = OUT_DIR / "fid_pdf.pdf"
OUT_PPT = OUT_DIR / "fid_ppt.pptx"

FID_URL = (
    "https://opportunitiesforyouth.org/2026/06/12/"
    "fid-call-for-proposals-2026-apply-for-up-to-e4-million-to-scale-innovative-"
    "solutions-tackling-poverty-and-inequality-worldwide/"
)

NAVY = "0A1F2E"
TEAL = "0D6E6E"
TEAL_L = "14A3A3"
ACCENT = "C45C26"
SLATE = "1E2F38"
MUTED = "3A4A54"
CREAM = "F7F5F0"
LINE = "D0DCDC"

SLOGAN = "Your health, our mission."
TAGLINE = "Health for All \u2014 Obulamu eri Bonna \u00b7 Afya kwa Wote \u00b7 Oburamu bwa Boona"

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "outreach_hero": "outreach_facilitator_canopy_01.jpg",
    "outreach_bp": "outreach_bp_screening.jpeg",
    "outreach_outdoor": "outreach_outdoor_clinic.jpeg",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "deep_tech": "deep_tech_collage.png",
    "gis": "gis_hotspots.png",
    "dashboard": "dashboard_demo.png",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "reactive": "reactive_clinic.png",
    "gericare": "gericare_wheelchair_assist.jpeg",
    "mothers": "waiting_room_mothers_01.jpeg",
    "facility_sign": "facility_exterior_sign.jpeg",
    "facility_entrance": "facility_exterior_entrance_01.jpg",
    "staff_field": "staff_outreach_conversation_01.jpg",
    "mvp_capture": "outreach_registration_form_01.jpg",
    "training": "indoor_training_staff_presenting_01.jpg",
    "conclusion": "outreach_audience_full_group_01.jpg",
}


def asset(name: str) -> Path:
    p = ASSETS / name
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def photo(key: str) -> Path:
    return asset(PHOTOS[key])


def embed(key_or_path, max_px: int = 1500) -> str:
    """Return a downscaled, cached copy of an image for lean document embedding."""
    from PIL import Image as PILImage

    src = photo(key_or_path) if key_or_path in PHOTOS else Path(key_or_path)
    cache = REPO / "tmp" / "fid_assets"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


# ---------------------------------------------------------------------------
# Shared content model — single source keeps all three formats in sync
# ---------------------------------------------------------------------------
META = {
    "programme": "Fund for Innovation in Development (FID) \u2014 Call for Proposals 2026",
    "title": "FairBanks Community Health Intelligence Platform (FCHIP)",
    "subtitle": "Predictive community health that keeps families out of poverty",
    "sectors": "Health \u00b7 Technology & Innovation \u00b7 Economic Inclusion",
    "geography": "Uganda (pilot) \u2192 East Africa \u2192 broader ODA-eligible markets",
    "stage": "Requested stage: Pilot Grant (up to \u20ac200,000) \u2014 pathway to Impact Evaluation & Scale-Up",
    "applicant": "FairBanks Medical Centre / FairBanks Community Reach (social enterprise, Uganda)",
}

EXEC_SUMMARY = [
    "In underserved African communities, illness is not only a health problem \u2014 it is one of the fastest routes "
    "into poverty. When an outbreak, a high-risk pregnancy, or an untreated chronic condition is caught late, "
    "families face catastrophic out-of-pocket costs, lost income, and deepening inequality. Today most primary "
    "healthcare is reactive: facilities treat patients only after they arrive sick, while the data that could "
    "warn everyone earlier sits scattered across community health workers' notebooks, outreach registers, and "
    "siloed systems.",
    "FairBanks Community Health Intelligence Platform (FCHIP) is a deep-technology innovation that turns "
    "community-generated health data into early warnings and decision support. Community Health Workers (CHWs) "
    "and Village Health Teams (VHTs) capture structured data on offline-capable mobile tools; the platform "
    "applies artificial intelligence, machine learning, and GIS mapping to predict disease outbreaks, maternal "
    "complications, chronic-disease hotspots, child-health threats, and medicine stock-outs \u2014 then routes "
    "actionable alerts to CHWs, clinics, districts, insurers, and partners before crises escalate.",
    "FairBanks is uniquely positioned to test this innovation. We already operate a functioning medical centre "
    "and active community outreach across Kampala-area communities (Bukoto, Kyebando, Kisaasi, Kamwokya, and "
    "Kikaaya), working with CHWs and VHTs on maternal and child health, geriatric care (Gericare), and "
    "chronic-disease screening. This gives us a live, real-world environment to pilot, evaluate, and refine "
    "FCHIP \u2014 exactly the evidence-first pathway FID is built to fund.",
    "We are applying for a Pilot Grant to validate FCHIP under real-world conditions and generate the early "
    "evidence needed to progress toward a rigorous impact evaluation and, ultimately, scale-up across "
    "ODA-eligible markets.",
]

PROBLEM_ROWS = [
    ["Care starts only after people fall sick", "Late detection of outbreaks and complications; higher treatment costs"],
    ["Community data sits in paper registers and silos", "Districts, NGOs, and insurers lack real-time intelligence"],
    ["High-risk pregnancies and NCDs found late", "Preventable maternal deaths and stroke/diabetes burden"],
    ["Medicine ordered on guesswork", "Stock-outs during seasonal disease surges"],
    ["Health shocks hit poor households hardest", "Catastrophic spending pushes families deeper into poverty"],
]

TOC_STEPS = [
    ("Inputs", "FID Pilot Grant; FairBanks live operations; CHW/VHT networks; mobile + cloud + AI stack."),
    ("Activities", "Build offline capture app, cloud pipeline, AI risk models, GIS layer, and dashboards; train CHWs; run pilot."),
    ("Outputs", "Structured community health data; predictive alerts; facility/district dashboards; referrals triggered."),
    ("Outcomes", "Earlier detection and referral; fewer missed high-risk cases; fewer stock-outs; better-targeted outreach."),
    ("Impact", "Reduced preventable illness and health-driven poverty; lower inequality for underserved communities."),
]

EVIDENCE_POINTS = [
    "Existing evidence: CHW/mHealth programmes and community-based early warning have documented gains in "
    "maternal, child, and communicable-disease outcomes across sub-Saharan Africa; FCHIP builds on this base by "
    "adding predictive intelligence across multiple data sources.",
    "Baseline data: FairBanks outreach screening, maternal programmes, pharmacy dispensing, and facility "
    "encounters already generate the signals needed to establish a credible pre-intervention baseline.",
    "Measurable outcomes: alert accuracy, time-to-referral, proportion of high-risk pregnancies identified "
    "early, immunisation and screening coverage, and stock-out frequency.",
    "Evaluation pathway: the pilot is designed so that a subsequent counterfactual impact evaluation "
    "(comparison communities) can rigorously test effect sizes \u2014 the qualifying step toward FID's Impact "
    "Evaluation and Scale-Up stages.",
]

COST_POINTS = [
    "Builds on assets that already exist \u2014 a live medical centre, trained CHW/VHT networks, and outreach "
    "programmes \u2014 so the pilot avoids the cost of standing up field operations from scratch.",
    "Offline-first, mobile-based capture keeps hardware and connectivity costs low and works in low-bandwidth "
    "settings.",
    "Prevention is cheaper than late treatment: early detection of outbreaks, maternal risk, and NCDs reduces "
    "costly emergency care and lost livelihoods.",
    "A shared platform serves many users (CHWs, clinics, districts, insurers), spreading cost and improving "
    "cost per beneficiary as coverage grows.",
]

SCALE_POINTS = [
    "Technical scalability: cloud architecture and standard mobile tools extend from one catchment to many "
    "districts without redesign.",
    "Model portability: CHW/VHT systems exist across Uganda and much of Africa, so the approach transfers to "
    "other ODA-eligible markets.",
    "Financial sustainability: diversified revenue \u2014 clinic/hospital subscriptions, district and Ministry of "
    "Health deployments, NGO programme monitoring, insurer analytics, and partner APIs \u2014 reduces reliance on "
    "grants over time.",
    "Policy pathway: evidence generated in the pilot supports integration into public health planning and "
    "district decision-making, aligning with FID's public-policy funding streams.",
]

TRACTION_POINTS = [
    "Functioning FairBanks Medical Centre with direct patient and community access",
    "Active Community Reach outreach in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and surrounding areas",
    "Established CHW and VHT relationships and trust in target communities",
    "Live programmes: maternal & child health, Gericare, chronic-disease screening, corporate & school health",
    "Digital health records foundation and pharmacy operations already generating data signals",
    "Research and community partnerships that support ethical data use and evaluation",
]

DEEP_TECH_ROWS = [
    ["Artificial Intelligence", "Outbreak early warning; maternal and NCD risk scoring"],
    ["Machine Learning", "Pattern learning from community data, seasonal trends, outreach outcomes"],
    ["GIS Mapping", "Disease distribution, hotspots, and resource-gap visualisation"],
    ["Mobile Data Collection", "Offline-capable CHW/VHT capture at household level"],
    ["Cloud Computing", "Secure sync across facilities, partners, administrative levels"],
    ["Analytics & NLP", "Live dashboards; local-language symptom capture where useful"],
]

USE_CASES = [
    ("Disease surveillance",
     "Rising fever reports across neighbouring villages \u2192 predicted malaria surge \u2192 targeted testing, "
     "bed-nets, and pharmacy pre-stocking."),
    ("Maternal health",
     "Home-visit BP, anaemia proxies, and ANC adherence \u2192 high-risk pregnancy flags \u2192 early CHW alerts and referral."),
    ("Non-communicable diseases",
     "Community and workplace BP/glucose screening \u2192 hypertension and diabetes hotspots \u2192 targeted campaigns."),
    ("Child health",
     "Growth, immunisation, and diarrhoea data \u2192 malnutrition and coverage-gap alerts \u2192 nutrition and immunisation drives."),
    ("Medicine demand",
     "Disease trends, rainfall, and consumption \u2192 demand forecasts \u2192 replenishment before stock-outs."),
]

BUDGET_ROWS = [
    ["Product & engineering", "Mobile capture app, cloud pipeline, AI/GIS modules, dashboards", "\u20ac70,000"],
    ["Field pilot & CHW training", "CHW/VHT onboarding, devices, supervision, community engagement", "\u20ac45,000"],
    ["Data, evaluation & M&E", "Baseline, monitoring, evaluation design toward counterfactual study", "\u20ac35,000"],
    ["Data governance & compliance", "Consent, anonymisation, Uganda Data Protection alignment", "\u20ac15,000"],
    ["Project management & operations", "Coordination, reporting, partnerships, contingency", "\u20ac35,000"],
    ["Total (indicative Pilot Grant)", "", "\u20ac200,000"],
]

ROADMAP_ROWS = [
    ["Phase 1 \u2014 Pilot", "0\u201312 mo", "MVP live with FairBanks CHWs/VHTs; validate 3 use cases; establish baseline"],
    ["Phase 2 \u2014 Impact evaluation prep", "12\u201318 mo", "Counterfactual design; comparison communities; measurement framework"],
    ["Phase 3 \u2014 District scale", "18\u201336 mo", "Partner clinics and district structures; NGO/insurer modules; APIs"],
    ["Phase 4 \u2014 Scale-up", "Year 3+", "Multi-district Uganda; East Africa entry; policy integration"],
]

SDG_ROWS = [
    ["SDG 1 \u2014 No Poverty", "Prevents health-shock-driven catastrophic spending and lost income"],
    ["SDG 3 \u2014 Good Health & Well-Being", "Predictive community health; maternal/child/NCD focus"],
    ["SDG 5 \u2014 Gender Equality", "Women-led venture; maternal health intelligence"],
    ["SDG 10 \u2014 Reduced Inequalities", "Intelligence and care for underserved communities"],
    ["SDG 17 \u2014 Partnerships", "Government, NGO, insurer, and academic collaboration"],
]

RISK_ROWS = [
    ["Data quality from field capture", "Structured forms, validation rules, CHW training and supervision"],
    ["Low CHW digital literacy", "Simple offline UX, local-language support, supervised rollout via VHTs"],
    ["Privacy and consent", "Ethical governance, anonymisation, Uganda Data Protection compliance"],
    ["Early model accuracy", "Start rule-based + limited ML; validate against FairBanks clinical outcomes"],
    ["Public-sector adoption", "Pilot evidence; district co-design; alignment with MoH community health strategy"],
]

FID_FIT_ROWS = [
    ["Evidence of impact", "Clear theory of change, baseline data, and a pathway to counterfactual evaluation"],
    ["Cost & cost-effectiveness", "Builds on existing assets; prevention reduces costly late treatment"],
    ["Scale & sustainability", "Portable model; diversified revenue; policy-integration pathway"],
    ["Poverty & inequality focus", "Targets underserved communities where health shocks drive poverty"],
]


# ===========================================================================
# DOCX + PDF (full narrative)
# ===========================================================================
def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from PIL import Image as PILImage

    def set_font(run, size=11, bold=False, color=SLATE, name="Calibri", italic=False):
        run.font.name = name
        run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def shade(cell, color):
        cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))

    def border(cell, color=LINE, sz="8"):
        cell._tc.get_or_add_tcPr().append(parse_xml(
            f'<w:tcBorders {nsdecls("w")}>'
            f'<w:top w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f'<w:left w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f'<w:bottom w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f'<w:right w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f"</w:tcBorders>"))

    def para(text, size=11, bold=False, color=SLATE, after=8, before=0,
             align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.space_before = Pt(before)
        p.paragraph_format.line_spacing = 1.22
        set_font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14, 3: 12}, {1: NAVY, 2: TEAL, 3: SLATE}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(16 if level == 1 else 10)
        p.paragraph_format.space_after = Pt(6)
        set_font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])
        return p

    def image(path, width_in=6.3, caption=None, max_h=3.6):
        path = Path(embed(str(path)))
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width_in, max_h * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(6)
        p.paragraph_format.space_after = Pt(2)
        p.add_run().add_picture(str(path), width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER,
                 italic=True, after=12, before=2)

    def bullets(items, size=11):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            p.clear()
            set_font(p.add_run(it), size=size, color=SLATE)

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            set_font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            shade(c, TEAL)
            border(c, TEAL)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                last = (ri == len(rows) - 1)
                set_font(c.paragraphs[0].add_run(str(val)), size=10,
                         bold=last and ci == 0, color=SLATE)
                if last:
                    shade(c, "E8F0F0")
                elif ri % 2:
                    shade(c, CREAM)
                border(c)
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()
        return t

    doc = Document()
    s = doc.sections[0]
    s.page_width, s.page_height = Inches(8.5), Inches(11)
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    # Cover
    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("Innovation Proposal", size=11, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    para(META["title"], size=24, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["subtitle"], size=13, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=12, italic=True)
    image(photo("cover"), width_in=6.2, max_h=3.0,
          caption="FairBanks Community Reach \u2014 " + TAGLINE)
    table(
        ["Item", "Detail"],
        [
            ["Applicant", META["applicant"]],
            ["Sectors", META["sectors"]],
            ["Geography", META["geography"]],
            ["Funding stage", META["stage"]],
        ],
        widths=[1.7, 4.7],
    )
    para("Aligned with the Fund for Innovation in Development Call for Proposals 2026.",
         size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=2, italic=True)
    para(FID_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=6)

    doc.add_page_break()

    heading("1. Executive Summary", 1)
    for i, p in enumerate(EXEC_SUMMARY):
        para(p, bold=(i == len(EXEC_SUMMARY) - 1))

    heading("2. The Problem: When Illness Deepens Poverty", 1)
    image(photo("reactive"), width_in=5.6, max_h=2.6,
          caption="Reactive care: the system waits for sickness before it responds")
    para("Across underserved African communities, a structural information gap keeps health systems reactive "
         "\u2014 and that gap has a direct poverty and inequality cost.")
    table(["Current reality", "Consequence"], PROBLEM_ROWS, widths=[3.2, 3.2])
    para("Existing digital tools \u2014 appointment booking, EMR digitisation, SMS reminders \u2014 rarely combine "
         "multi-source community data, predictive modelling, geospatial mapping, and real-time decision support. "
         "The result is fragmented data, delayed response, and missed prevention.")

    heading("3. The Innovation: FCHIP", 1)
    image(photo("architecture"), width_in=6.4,
          caption="Community signals \u2192 mobile capture \u2192 FCHIP intelligence \u2192 action")
    para("FCHIP is a deep-technology platform that turns community-generated health data into predictions and "
         "decision support. It sits between last-mile data capture and the people who can act \u2014 CHWs, clinics, "
         "districts, insurers, and partners.")
    heading("3.1 Deep technology core", 2)
    table(["Technology", "Function in FCHIP"], DEEP_TECH_ROWS, widths=[2.2, 4.2])
    heading("3.2 Predictive use cases", 2)
    for title, body in USE_CASES:
        para(title + ": " + body, after=4)

    heading("4. Theory of Change", 1)
    para("FID funds innovations with a clear, testable theory of change. Ours links a health innovation directly "
         "to poverty and inequality reduction:")
    table(["Stage", "Detail"], [[a, b] for a, b in TOC_STEPS], widths=[1.4, 5.0])

    heading("5. Evidence of Impact & Evaluation", 1)
    image(photo("gis"), width_in=5.8, max_h=2.8,
          caption="GIS hotspot mapping \u2014 measurable, testable early-warning outputs")
    bullets(EVIDENCE_POINTS)

    heading("6. Cost & Cost-Effectiveness", 1)
    bullets(COST_POINTS)

    heading("7. Scale & Financial Sustainability", 1)
    image(photo("dashboard"), width_in=6.2, max_h=2.8,
          caption="A shared intelligence layer that scales across many users and districts")
    bullets(SCALE_POINTS)

    heading("8. Why FairBanks \u2014 Traction & Foundation", 1)
    image(photo("outreach_hero"), width_in=6.4,
          caption="Live FairBanks Community Reach outreach in Kampala communities")
    bullets(TRACTION_POINTS)

    heading("9. Implementation Plan & Budget", 1)
    heading("9.1 Roadmap", 2)
    table(["Phase", "Timeline", "Milestones"], ROADMAP_ROWS, widths=[1.9, 1.2, 3.3])
    heading("9.2 Indicative pilot budget", 2)
    table(["Budget line", "Detail", "Amount"], BUDGET_ROWS, widths=[2.0, 3.3, 1.1])

    heading("10. Alignment with FID & the SDGs", 1)
    heading("10.1 FID evaluation criteria", 2)
    table(["FID criterion", "How FCHIP responds"], FID_FIT_ROWS, widths=[2.2, 4.2])
    heading("10.2 Sustainable Development Goals", 2)
    table(["SDG", "FCHIP contribution"], SDG_ROWS, widths=[2.6, 3.8])

    heading("11. Risks & Mitigation", 1)
    table(["Risk", "Mitigation"], RISK_ROWS, widths=[2.4, 4.0])

    heading("12. Conclusion", 1)
    image(photo("conclusion"), width_in=6.4,
          caption="From community signals to life-saving, poverty-reducing predictions")
    para("FCHIP is an evidence-first innovation that fits the FID model precisely: a promising solution, ready to "
         "be tested in a real-world setting, with a clear pathway from pilot to rigorous evaluation and scale. By "
         "catching health risks before they become crises, FairBanks helps underserved families avoid the "
         "catastrophic costs that trap them in poverty \u2014 a genuine win for communities and for development.")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, before=8, italic=True)
    para("Proposal aligned with the FID Call for Proposals 2026. Source: " + FID_URL,
         size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, before=10, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    try:
        doc.save(str(OUT_DOC))
        print(f"DOCX: {OUT_DOC}")
    except PermissionError:
        alt = OUT_DOC.with_name(OUT_DOC.stem + "_unlocked" + OUT_DOC.suffix)
        doc.save(str(alt))
        print(f"DOCX locked; saved as: {alt}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle,
        PageBreak, KeepTogether, HRFlowable,
    )
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = (HexColor("#" + MUTED), HexColor("#" + MUTED),
                                 HexColor("#" + CREAM), HexColor("#" + LINE))
    slate = HexColor("#" + SLATE)

    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CoverTitle", fontName="Helvetica-Bold", fontSize=22, leading=26,
                          textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=15, leading=19,
                          textColor=navy, spaceBefore=14, spaceAfter=6))
    st.add(ParagraphStyle("H2", fontName="Helvetica-Bold", fontSize=12, leading=15,
                          textColor=teal, spaceBefore=9, spaceAfter=5))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, leading=14,
                          textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("BodyBold", fontName="Helvetica-Bold", fontSize=10, leading=14,
                          textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Caption", fontName="Helvetica-Oblique", fontSize=8, leading=10,
                          textColor=muted, alignment=TA_CENTER, spaceAfter=10, spaceBefore=2))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, leading=13,
                          textColor=slate, leftIndent=12, spaceAfter=3))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, leading=12,
                          textColor=muted, alignment=TA_CENTER, spaceAfter=3))
    st.add(ParagraphStyle("CellHead", fontName="Helvetica-Bold", fontSize=8.5, leading=11, textColor=white))
    st.add(ParagraphStyle("CellBody", fontName="Helvetica", fontSize=8.5, leading=11, textColor=slate))
    st.add(ParagraphStyle("CellBodyBold", fontName="Helvetica-Bold", fontSize=8.5, leading=11, textColor=slate))

    story = []
    page_w = A4[0] - 1.6 * inch

    def img(key, w=page_w, caption=None, max_h=3.2 * inch):
        path = embed(key)
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / float(iw)
        h = w * aspect
        if h > max_h:
            h, w = max_h, max_h / aspect
        block = [Image(path, width=w, height=h)]
        block.append(Paragraph(caption, st["Caption"]) if caption else Spacer(1, 6))
        return KeepTogether(block)

    def h1(t):
        story.append(Paragraph(t, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1.2, color=teal, spaceAfter=8))

    def h2(t):
        story.append(Paragraph(t, st["H2"]))

    def body(t, bold=False):
        story.append(Paragraph(t, st["BodyBold"] if bold else st["Body"]))

    def bullets(items):
        for it in items:
            story.append(Paragraph(f"\u2022&nbsp;&nbsp;{it}", st["FBullet"]))
        story.append(Spacer(1, 4))

    def table(headers, rows, widths=None, last_bold=False):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        for ri, row in enumerate(rows):
            last = last_bold and ri == len(rows) - 1
            data.append([Paragraph(str(c), st["CellBodyBold"] if (last and ci == 0) else st["CellBody"])
                         for ci, c in enumerate(row)])
        widths = widths or [page_w / len(headers)] * len(headers)
        t = Table(data, colWidths=widths, repeatRows=1)
        cmds = [
            ("BACKGROUND", (0, 0), (-1, 0), teal),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
        ]
        if last_bold:
            cmds.append(("BACKGROUND", (0, -1), (-1, -1), HexColor("#E8F0F0")))
        t.setStyle(TableStyle(cmds))
        story.append(t)
        story.append(Spacer(1, 10))

    # Cover
    story.append(Spacer(1, 0.25 * inch))
    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph("Innovation Proposal", st["Meta"]))
    story.append(Spacer(1, 8))
    story.append(Paragraph(META["title"], st["CoverTitle"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{META["subtitle"]}</i></b></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    story.append(Spacer(1, 10))
    story.append(img("cover", w=page_w * 0.92, max_h=2.7 * inch,
                     caption="FairBanks Community Reach \u2014 " + TAGLINE))
    table(["Item", "Detail"], [
        ["Applicant", META["applicant"]],
        ["Sectors", META["sectors"]],
        ["Geography", META["geography"]],
        ["Funding stage", META["stage"]],
    ], widths=[page_w * 0.24, page_w * 0.76])
    story.append(Paragraph(
        f'Aligned with the FID Call for Proposals 2026. '
        f'<link href="{FID_URL}"><font color="#{TEAL}">Source</font></link>', st["Meta"]))
    story.append(PageBreak())

    h1("1. Executive Summary")
    for i, p in enumerate(EXEC_SUMMARY):
        body(p, bold=(i == len(EXEC_SUMMARY) - 1))

    h1("2. The Problem: When Illness Deepens Poverty")
    story.append(img("reactive", w=page_w * 0.8, max_h=2.4 * inch,
                     caption="Reactive care waits for sickness before it responds"))
    body("Across underserved African communities, a structural information gap keeps health systems reactive "
         "\u2014 with a direct poverty and inequality cost.")
    table(["Current reality", "Consequence"], PROBLEM_ROWS, widths=[page_w * 0.5, page_w * 0.5])
    body("Existing digital tools rarely combine multi-source community data, predictive modelling, geospatial "
         "mapping, and real-time decision support \u2014 leaving fragmented data and missed prevention.")

    h1("3. The Innovation: FCHIP")
    story.append(img("architecture", caption="Community signals \u2192 capture \u2192 FCHIP intelligence \u2192 action"))
    body("FCHIP turns community-generated health data into predictions and decision support, sitting between "
         "last-mile capture and the people who can act \u2014 CHWs, clinics, districts, insurers, and partners.")
    h2("3.1 Deep technology core")
    table(["Technology", "Function in FCHIP"], DEEP_TECH_ROWS, widths=[page_w * 0.32, page_w * 0.68])
    h2("3.2 Predictive use cases")
    for title, b in USE_CASES:
        body(f"<b>{title}:</b> {b}")

    h1("4. Theory of Change")
    body("A clear, testable theory of change links this health innovation directly to poverty and inequality reduction:")
    table(["Stage", "Detail"], [[a, b] for a, b in TOC_STEPS], widths=[page_w * 0.22, page_w * 0.78])

    story.append(PageBreak())
    h1("5. Evidence of Impact & Evaluation")
    story.append(img("gis", w=page_w * 0.82, max_h=2.6 * inch,
                     caption="GIS hotspot mapping \u2014 measurable, testable outputs"))
    bullets(EVIDENCE_POINTS)

    h1("6. Cost & Cost-Effectiveness")
    bullets(COST_POINTS)

    h1("7. Scale & Financial Sustainability")
    story.append(img("dashboard", caption="A shared intelligence layer that scales across users and districts"))
    bullets(SCALE_POINTS)

    h1("8. Why FairBanks \u2014 Traction & Foundation")
    story.append(img("outreach_hero", caption="Live FairBanks Community Reach outreach in Kampala communities"))
    bullets(TRACTION_POINTS)

    story.append(PageBreak())
    h1("9. Implementation Plan & Budget")
    h2("9.1 Roadmap")
    table(["Phase", "Timeline", "Milestones"], ROADMAP_ROWS,
          widths=[page_w * 0.28, page_w * 0.16, page_w * 0.56])
    h2("9.2 Indicative pilot budget")
    table(["Budget line", "Detail", "Amount"], BUDGET_ROWS,
          widths=[page_w * 0.28, page_w * 0.56, page_w * 0.16], last_bold=True)

    h1("10. Alignment with FID & the SDGs")
    h2("10.1 FID evaluation criteria")
    table(["FID criterion", "How FCHIP responds"], FID_FIT_ROWS, widths=[page_w * 0.32, page_w * 0.68])
    h2("10.2 Sustainable Development Goals")
    table(["SDG", "FCHIP contribution"], SDG_ROWS, widths=[page_w * 0.4, page_w * 0.6])

    h1("11. Risks & Mitigation")
    table(["Risk", "Mitigation"], RISK_ROWS, widths=[page_w * 0.35, page_w * 0.65])

    story.append(PageBreak())
    story.append(KeepTogether([
        Paragraph("12. Conclusion", st["H1"]),
        HRFlowable(width="100%", thickness=1.2, color=teal, spaceAfter=8),
        img("conclusion", max_h=2.7 * inch,
            caption="From community signals to life-saving, poverty-reducing predictions"),
        Paragraph(
            "FCHIP fits the FID model precisely: a promising innovation, ready to be tested in a real-world "
            "setting, with a clear pathway from pilot to rigorous evaluation and scale. By catching health risks "
            "before they become crises, FairBanks helps underserved families avoid the catastrophic costs that "
            "trap them in poverty \u2014 a genuine win for communities and for development.", st["Body"]),
        Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]),
    ]))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Proposal aligned with the FID Call for Proposals 2026.", st["Meta"]))

    def page_num(canvas, doc_):
        canvas.saveState()
        canvas.setFillColor(navy)
        canvas.rect(0, A4[1] - 10, A4[0], 10, fill=1, stroke=0)
        canvas.setFillColor(teal)
        canvas.rect(0, A4[1] - 12, A4[0], 2, fill=1, stroke=0)
        canvas.setFillColor(muted)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(0.8 * inch, 0.45 * inch, "FairBanks FCHIP  |  FID Call for Proposals 2026")
        canvas.drawRightString(A4[0] - 0.8 * inch, 0.45 * inch, f"{doc_.page}")
        canvas.restoreState()

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    def write(path):
        d = SimpleDocTemplate(str(path), pagesize=A4,
                              leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                              topMargin=0.7 * inch, bottomMargin=0.7 * inch,
                              title="FairBanks FCHIP - FID Proposal 2026",
                              author="FairBanks Community Health Intelligence Platform")
        d.build(story, onFirstPage=page_num, onLaterPages=page_num)

    try:
        write(OUT_PDF)
        print(f"PDF: {OUT_PDF}")
    except PermissionError:
        alt = OUT_PDF.with_name(OUT_PDF.stem + "_unlocked" + OUT_PDF.suffix)
        write(alt)
        print(f"PDF locked; saved as: {alt}")


# ===========================================================================
# PPTX (12-slide summary deck)
# ===========================================================================
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from PIL import Image as PILImage

    def C(hexstr):
        return RGBColor.from_string(hexstr)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, color, line_color=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(color)
        if line_color:
            shp.line.color.rgb = C(line_color)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(slide, x, y, w, h, text, size=18, bold=False, color=SLATE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri",
                line_pt=None):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        # Slight inset so glyphs don't clip the box edges
        try:
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
        except Exception:
            pass
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            # Keep multi-line titles from collapsing on top of each other
            p.space_before = Pt(0)
            p.space_after = Pt(max(4, int((line_pt or size) * 0.15))) if i < len(lines) - 1 else Pt(0)
            if line_pt is not None:
                p.line_spacing = 1.15  # multiple of font size — avoids absolute Pt overlap bugs
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            r.font.color.rgb = C(color)
        return tb

    def fit_size(box_h_in, n_items, max_size=16, min_size=11, lines_per=1, gap_pt=6):
        """Pick a font size that fills available height without overflow."""
        if n_items <= 0:
            return max_size
        usable_pt = box_h_in * 72 * 0.92  # leave a little breathing room
        # each item: lines * size * 1.25 leading + gap
        size = (usable_pt - gap_pt * max(n_items - 1, 0)) / (n_items * lines_per * 1.28)
        return max(min_size, min(max_size, int(round(size))))

    def pic_cover(slide, key, x, y, w, h):
        """Add picture cropped to fill the target box."""
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        target = w / h
        src = iw / ih
        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
        if src > target:
            crop = (1 - target / src) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (1 - src / target) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
        return pic

    def pic_fit(slide, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        px = x + (w - tw) // 2
        py = y + (h - th) // 2
        return slide.shapes.add_picture(path, px, py, width=tw, height=th)

    def band(slide, kicker, title):
        rect(slide, 0, 0, SW, Inches(1.05), CREAM)
        rect(slide, 0, 0, Inches(0.18), Inches(1.05), TEAL)
        textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.28),
                kicker.upper(), size=11, bold=True, color=ACCENT)
        textbox(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.52),
                title, size=22, bold=True, color=NAVY)

    def footer(slide, n):
        rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), NAVY)
        textbox(slide, Inches(0.4), SH - Inches(0.31), Inches(9), Inches(0.3),
                "FairBanks FCHIP  |  FID Call for Proposals 2026  |  " + SLOGAN,
                size=9, color="FFFFFF", anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, SW - Inches(1.2), SH - Inches(0.31), Inches(0.8), Inches(0.3),
                str(n), size=9, color="FFFFFF", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    def bullet_card(slide, x, y, w, h, title, body):
        rect(slide, x, y, w, h, "FFFFFF", LINE)
        rect(slide, x, y, w, Inches(0.08), TEAL)
        # Title ~18% of card; body fills the rest
        h_in = h / 914400  # EMU → inches
        title_sz = fit_size(h_in * 0.28, 1, max_size=16, min_size=13)
        body_sz = fit_size(h_in * 0.55, 1, max_size=13, min_size=11, lines_per=2)
        textbox(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(h_in * 0.28),
                title, size=title_sz, bold=True, color=TEAL)
        textbox(slide, x + Inches(0.18), y + Inches(h_in * 0.42), w - Inches(0.36), Inches(h_in * 0.5),
                body, size=body_sz, color=MUTED)

    def bullets_box(slide, x, y, w, h, items, size=None, gap=True, max_size=16, min_size=11, lines_per=2):
        h_in = h / 914400 if hasattr(h, "__int__") else float(h)
        # If caller passed Inches(...), convert; Inches returns Emu-like int
        try:
            h_in = float(h) / 914400.0
        except Exception:
            h_in = float(h)
        # Estimate wrap: long items need ~2 lines; short ~1
        avg_len = sum(len(it) for it in items) / max(len(items), 1)
        est_lines = 2 if avg_len > 70 else (1 if avg_len < 45 else lines_per)
        if size is None:
            size = fit_size(h_in, len(items), max_size=max_size, min_size=min_size,
                            lines_per=est_lines, gap_pt=8 if gap else 4)
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        try:
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
        except Exception:
            pass
        gap_pt = max(4, int(size * 0.45)) if gap else 3
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap_pt)
            p.line_spacing = Pt(int(size * 1.25))
            r = p.add_run()
            r.text = "\u2022  " + it
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            r.font.color.rgb = C(SLATE)
        return tb

    # Content area below band / above footer (shared rhythm)
    CONTENT_TOP = Inches(1.25)
    CONTENT_H = Inches(5.75)  # to just above footer

    # 1. Title — type hierarchy sized to the navy band
    # Title + acronym use SEPARATE text boxes so line-spacing never overlaps.
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover", 0, 0, SW, SH)
    band_h = Inches(3.85)
    rect(s, 0, SH - band_h, SW, band_h, NAVY)
    textbox(s, Inches(0.65), SH - Inches(3.55), Inches(12), Inches(0.32),
            META["programme"], size=13, bold=True, color=TEAL_L)
    textbox(s, Inches(0.65), SH - Inches(3.1), Inches(12), Inches(0.7),
            "FairBanks Community Health Intelligence Platform",
            size=28, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(2.35), Inches(12), Inches(0.5),
            "(FCHIP)", size=26, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(1.7), Inches(12), Inches(0.4),
            META["subtitle"], size=16, italic=True, color="F2C79B")
    textbox(s, Inches(0.65), SH - Inches(1.15), Inches(12), Inches(0.3),
            SLOGAN, size=13, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(0.75), Inches(12), Inches(0.4),
            META["stage"], size=12, color="D0E8E8")

    # 2. Problem
    s = prs.slides.add_slide(blank)
    band(s, "The problem", "When illness strikes, poverty follows")
    pic_fit(s, "reactive", Inches(0.5), CONTENT_TOP, Inches(5.7), CONTENT_H)
    probs = [
        "Care starts only after people fall sick \u2014 outbreaks caught late",
        "Community data trapped in paper registers and silos",
        "High-risk pregnancies and NCDs identified too late",
        "Medicine stock-outs during seasonal surges",
        "Health shocks push poor families into catastrophic spending",
    ]
    bullets_box(s, Inches(6.5), CONTENT_TOP + Inches(0.15), Inches(6.3), CONTENT_H - Inches(0.3),
                probs, max_size=22, min_size=15, lines_per=1)
    footer(s, 2)

    # 3. Innovation
    s = prs.slides.add_slide(blank)
    band(s, "The innovation", "FCHIP: community data becomes prediction")
    pic_fit(s, "architecture", Inches(0.4), CONTENT_TOP, Inches(7.1), CONTENT_H)
    cards = [
        ("Capture", "Offline CHW/VHT mobile tools at the household level"),
        ("Intelligence", "AI/ML risk scoring, GIS hotspots, decision support"),
        ("Action", "Alerts to CHWs, clinics, districts, insurers, partners"),
    ]
    card_h = Inches(1.7)
    card_gap = Inches(0.18)
    for i, (t, b) in enumerate(cards):
        bullet_card(s, Inches(7.85), CONTENT_TOP + Inches(0.1) + i * (card_h + card_gap),
                    Inches(5.05), card_h, t, b)
    footer(s, 3)

    # 4. Deep tech
    s = prs.slides.add_slide(blank)
    band(s, "Deep technology core", "Substantial engineering, not a light app")
    pic_fit(s, "deep_tech", Inches(0.45), CONTENT_TOP, Inches(5.5), CONTENT_H)
    n_tech = len(DEEP_TECH_ROWS)
    row_h = CONTENT_H / n_tech
    row_sz_t = fit_size(row_h / 914400 * 0.4, 1, max_size=14, min_size=12)
    row_sz_b = fit_size(row_h / 914400 * 0.45, 1, max_size=12, min_size=10)
    for i, (t, b) in enumerate(DEEP_TECH_ROWS):
        y = CONTENT_TOP + int(i * row_h)
        pad = Inches(0.06)
        rect(s, Inches(6.3), y + pad, Inches(6.55), row_h - pad * 2, "FFFFFF", LINE)
        textbox(s, Inches(6.5), y + Inches(0.12), Inches(6.2), Inches(0.32),
                t, size=row_sz_t, bold=True, color=TEAL)
        textbox(s, Inches(6.5), y + Inches(0.42), Inches(6.2), Inches(0.35),
                b, size=row_sz_b, color=MUTED)
    footer(s, 4)

    # 5. Theory of change
    s = prs.slides.add_slide(blank)
    band(s, "Theory of change", "From data to poverty reduction")
    n = len(TOC_STEPS)
    cw = Inches(2.42)
    gap = Inches(0.12)
    x0 = Inches(0.5)
    card_top = CONTENT_TOP + Inches(0.25)
    card_body_h = CONTENT_H - Inches(0.35)
    for i, (t, b) in enumerate(TOC_STEPS):
        x = x0 + i * (cw + gap)
        rect(s, x, card_top, cw, card_body_h, "FFFFFF", LINE)
        rect(s, x, card_top, cw, Inches(0.55), TEAL if i < n - 1 else ACCENT)
        textbox(s, x, card_top + Inches(0.08), cw, Inches(0.42), t,
                size=14, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
        body_sz = fit_size((card_body_h / 914400) - 0.85, 1, max_size=13, min_size=11, lines_per=5)
        textbox(s, x + Inches(0.12), card_top + Inches(0.7), cw - Inches(0.24),
                card_body_h - Inches(0.85), b, size=body_sz, color=MUTED)
    footer(s, 5)

    # 6. Evidence & evaluation
    s = prs.slides.add_slide(blank)
    band(s, "Evidence of impact", "Built to be measured and evaluated")
    pic_fit(s, "gis", Inches(0.45), CONTENT_TOP, Inches(5.5), CONTENT_H)
    bullets_box(s, Inches(6.3), CONTENT_TOP + Inches(0.1), Inches(6.5), CONTENT_H - Inches(0.2),
                EVIDENCE_POINTS, max_size=18, min_size=13, lines_per=3)
    footer(s, 6)

    # 7. Cost-effectiveness
    s = prs.slides.add_slide(blank)
    band(s, "Cost & cost-effectiveness", "Prevention costs less than late treatment")
    bullets_box(s, Inches(0.55), CONTENT_TOP + Inches(0.1), Inches(6.2), CONTENT_H - Inches(0.2),
                COST_POINTS, max_size=18, min_size=13, lines_per=3)
    pic_fit(s, "outreach_bp", Inches(7.0), CONTENT_TOP, Inches(5.9), CONTENT_H)
    footer(s, 7)

    # 8. Scale & sustainability
    s = prs.slides.add_slide(blank)
    band(s, "Scale & sustainability", "A portable model with diversified revenue")
    pic_fit(s, "dashboard", Inches(0.45), CONTENT_TOP, Inches(5.9), CONTENT_H)
    bullets_box(s, Inches(6.7), CONTENT_TOP + Inches(0.1), Inches(6.2), CONTENT_H - Inches(0.2),
                SCALE_POINTS, max_size=18, min_size=13, lines_per=3)
    footer(s, 8)

    # 9. Traction
    s = prs.slides.add_slide(blank)
    band(s, "Why FairBanks", "A live venture, ready to pilot")
    pic_fit(s, "outreach_hero", Inches(0.45), CONTENT_TOP, Inches(6.1), CONTENT_H)
    bullets_box(s, Inches(6.9), CONTENT_TOP + Inches(0.1), Inches(6.0), CONTENT_H - Inches(0.2),
                TRACTION_POINTS, max_size=18, min_size=13, lines_per=2)
    footer(s, 9)

    # 10. Roadmap & budget
    s = prs.slides.add_slide(blank)
    band(s, "Plan & budget", "Pilot first, then evaluate and scale")
    textbox(s, Inches(0.5), CONTENT_TOP, Inches(6), Inches(0.35),
            "Roadmap", size=15, bold=True, color=TEAL)
    n_rm = len(ROADMAP_ROWS)
    rm_h = (CONTENT_H - Inches(0.45)) / n_rm
    for i, (ph, tl, ms) in enumerate(ROADMAP_ROWS):
        y = CONTENT_TOP + Inches(0.4) + int(i * rm_h)
        pad = Inches(0.05)
        rect(s, Inches(0.5), y + pad, Inches(6.1), rm_h - pad * 2, "FFFFFF", LINE)
        rect(s, Inches(0.5), y + pad, Inches(0.1), rm_h - pad * 2, TEAL if i < 3 else ACCENT)
        textbox(s, Inches(0.75), y + Inches(0.1), Inches(5.7), Inches(0.3),
                f"{ph}  ({tl})", size=13, bold=True, color=NAVY)
        textbox(s, Inches(0.75), y + Inches(0.4), Inches(5.7), Inches(0.4),
                ms, size=11, color=MUTED)
    textbox(s, Inches(7.0), CONTENT_TOP, Inches(6), Inches(0.35),
            "Indicative pilot budget", size=15, bold=True, color=TEAL)
    n_bg = len(BUDGET_ROWS)
    bg_h = (CONTENT_H - Inches(0.45)) / n_bg
    for i, (line_item, _detail, amt) in enumerate(BUDGET_ROWS):
        y = CONTENT_TOP + Inches(0.4) + int(i * bg_h)
        pad = Inches(0.04)
        last = i == n_bg - 1
        rect(s, Inches(7.0), y + pad, Inches(5.9), bg_h - pad * 2,
             "E8F0F0" if last else "FFFFFF", LINE)
        textbox(s, Inches(7.2), y + Inches(0.12), Inches(4.3), Inches(0.4),
                line_item, size=12, bold=last, color=NAVY if last else SLATE,
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(11.4), y + Inches(0.12), Inches(1.35), Inches(0.4),
                amt, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 10)

    # 11. FID fit & SDGs
    s = prs.slides.add_slide(blank)
    band(s, "Why it fits FID", "Evidence, cost-effectiveness, scale")
    n_fit = len(FID_FIT_ROWS)
    fit_h = CONTENT_H / n_fit
    for i, (t, b) in enumerate(FID_FIT_ROWS):
        y = CONTENT_TOP + int(i * fit_h)
        pad = Inches(0.05)
        rect(s, Inches(0.45), y + pad, Inches(7.2), fit_h - pad * 2, "FFFFFF", LINE)
        textbox(s, Inches(0.65), y + Inches(0.1), Inches(2.5), fit_h - Inches(0.2),
                t, size=13, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(3.2), y + Inches(0.1), Inches(4.25), fit_h - Inches(0.2),
                b, size=12, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(8.0), CONTENT_TOP, Inches(4.9), Inches(0.35),
            "SDG alignment", size=14, bold=True, color=TEAL)
    bullets_box(s, Inches(8.0), CONTENT_TOP + Inches(0.4), Inches(4.9), CONTENT_H - Inches(0.5),
                [f"{a}: {b}" for a, b in SDG_ROWS], max_size=13, min_size=11, lines_per=2)
    footer(s, 11)

    # 12. Ask / vision — type sized to the lower band
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.7), SW, Inches(3.7), NAVY)
    textbox(s, Inches(0.65), SH - Inches(3.4), Inches(12), Inches(0.32),
            "The ask", size=13, bold=True, color=TEAL_L)
    textbox(s, Inches(0.65), SH - Inches(2.95), Inches(12), Inches(0.65),
            "A Pilot Grant to validate FCHIP in FairBanks' live catchment —",
            size=22, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(2.25), Inches(12), Inches(0.55),
            "and prove the pathway to evaluation and scale.",
            size=22, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(1.5), Inches(12), Inches(0.4),
            "Predictive community health that keeps families out of poverty.",
            size=15, italic=True, color="F2C79B")
    textbox(s, Inches(0.65), SH - Inches(0.9), Inches(12), Inches(0.35),
            SLOGAN, size=14, bold=True, color="FFFFFF")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUT_PPT))
        print(f"PPTX: {OUT_PPT}")
    except PermissionError:
        alt = OUT_PPT.with_name(OUT_PPT.stem + "_unlocked" + OUT_PPT.suffix)
        prs.save(str(alt))
        print(f"PPTX locked; saved as: {alt}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. FID document set generated in", OUT_DIR)
