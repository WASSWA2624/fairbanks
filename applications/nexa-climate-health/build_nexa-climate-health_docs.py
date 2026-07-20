#!/usr/bin/env python3
"""Build the Nexa Climate and Health 2026 FairBanks application pack.

Creates one synchronized set:
  documents/nexa-climate-health_word.docx
  documents/nexa-climate-health_pdf.pdf
  documents/nexa-climate-health_ppt.pptx

The Word/PDF files are a Proof-of-Concept answer bank and submission workbook.
The PowerPoint is a 12-slide visual proposal summary.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import Iterable


HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
ASSETS = REPO / "assets"
OUT = HERE / "documents"
DOCX = OUT / "nexa-climate-health_word.docx"
PDF = OUT / "nexa-climate-health_pdf.pdf"
PPTX = OUT / "nexa-climate-health_ppt.pptx"

OFFICIAL_PAGE = "https://www.grandchallenges.ca/rfp-nexa/"
OFFICIAL_RFP = "https://www.grandchallenges.ca/wp-content/uploads/2026/06/EN-Nexa-RFP-Jun-18-2026.pdf"
OFFICIAL_POC = "https://www.grandchallenges.ca/wp-content/uploads/2026/06/Nexa-POC-application_for-webpage.pdf"
OFFICIAL_TTS = "https://www.grandchallenges.ca/wp-content/uploads/2026/06/Nexa-TTS-Stage-1-application_for-webpage.pdf"
FLUXX = "https://gcc.fluxx.io"

NAVY = "0A1F2E"
TEAL = "0D6E6E"
GREEN = "2D7A55"
ORANGE = "C45C26"
GOLD = "D99A2B"
CREAM = "F7F5F0"
PALE_TEAL = "E8F3F2"
PALE_ORANGE = "FBEDE6"
PALE_GREEN = "E9F2EC"
SLATE = "1E2F38"
MUTED = "52636C"
LINE = "CED9D8"
WHITE = "FFFFFF"
RED = "A3312D"

PROGRAMME = "Nexa Climate and Health Initiative 2026"
TITLE = "FCHIP Climate Health Early Action"
SUBTITLE = "Climate-informed early warning that drives timely health action in Uganda"
SLOGAN = "Your health, our mission."
DEADLINE = "22 July 2026, 2:00 p.m. ET / 6:00 p.m. UTC"
TRACK = "Proof of Concept - climate-informed early warning and monitoring systems"
POC_MAX = 200_000
AMOUNT_REQUESTED = 180_000  # 90% of PoC maximum
REQUEST = f"USD {AMOUNT_REQUESTED:,} (90% of PoC maximum USD {POC_MAX:,}) for 18, 21, or 24 months"
CONFIRM = "[CONFIRM BEFORE SUBMISSION]"

# Direct USD 163,640 + indirect USD 16,360 (10% of direct, rounded under the cap) = 180,000
BUDGET_ROWS = [
    (
        "Remuneration",
        "74,000",
        "Project Lead, clinical lead, product/data roles, MEL, community engagement, "
        "and CHW/VHT stipends across the investment period",
    ),
    (
        "Subcontractor fees",
        "30,000",
        "Climate/epidemiology support, independent MEL or statistician, and contracted "
        "engineering for offline capture, secure EMR/HMS APIs, and GIS action board",
    ),
    (
        "Travel costs",
        "11,000",
        "Community co-design and feedback visits, site supervision, partner and district "
        "meetings within Uganda",
    ),
    (
        "Goods and supplies",
        "16,500",
        "Pilot response supplies (e.g. malaria testing materials for alert follow-up), "
        "outreach materials, airtime/data, printing, and consumables",
    ),
    (
        "Equipment costs",
        "13,500",
        "Offline-capable phones/tablets for CHW/VHT capture, protective cases, and basic "
        "field accessories",
    ),
    (
        "Project administration costs",
        "10,640",
        "Finance coordination, banking fees, insurance allocation, and day-to-day project "
        "administration in Uganda",
    ),
    (
        "Sub-grants",
        "8,000",
        "Small community or VHT-linked implementation support for co-design, outreach, "
        "and feedback activities",
    ),
    (
        "Indirect costs",
        "16,360",
        "10% of direct costs (USD 163,640); within the 10% Fluxx cap",
    ),
    (
        "Total",
        "180,000",
        "90% of PoC maximum (USD 200,000); majority spend in Uganda",
    ),
]

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "outreach": "outreach_bp_screening.jpeg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "dashboard": "dashboard_demo.png",
    "gis": "gis_hotspots.png",
    "architecture": "data_flow_iso_labeled.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "elderly": "gericare_wheelchair_assist.jpeg",
    "community": "outreach_audience_full_group_01.jpg",
    "facility": "facility_exterior_entrance_01.jpg",
    "team": "staff_team_reception.jpeg",
}

CONCEPT = {
    "cascade": REPO / ".cursor" / "concept_improved.jpeg",
    "simple": REPO / ".cursor" / "concept_simple.jpeg",
    "classic": REPO / ".cursor" / "concept.jpeg",
}


def photo(key: str) -> Path:
    path = ASSETS / PHOTOS[key]
    if not path.exists():
        raise FileNotFoundError(path)
    return path


def concept(key: str = "cascade") -> Path:
    path = CONCEPT[key]
    if not path.exists():
        raise FileNotFoundError(path)
    return path


CALL_FACTS = [
    ("Call", "Nexa - Climate and Health Innovation: Africa, Latin America and the Caribbean"),
    ("Co-leads", "Grand Challenges Canada and the Science for Africa Foundation"),
    ("Deadline", DEADLINE),
    ("Submission", "Fluxx only; complete application in English or French"),
    ("Recommended stage", "Proof of Concept; current evidence does not support Transition to Scale"),
    ("Amount", f"Draft ask USD {AMOUNT_REQUESTED:,} (90% of PoC maximum USD {POC_MAX:,})"),
    ("Duration", "18, 21, or 24 months"),
    ("Applicant", "Incorporated or equivalent in Africa or Brazil; Uganda is eligible"),
    ("Implementation", "Uganda is an eligible Proof-of-Concept implementation country"),
    ("Community threshold", "At least community linked"),
    ("Climate bar", "Must be climate-integrated: climate shapes design, function, and deployment"),
    ("Innovation Screen", "Only seven Overview answers reviewed first; typically over 80% declined there"),
]

READINESS = [
    ("Applicant legal name and entity type", CONFIRM),
    ("Incorporated, active, and in good standing in Uganda", CONFIRM),
    ("Project Lead identity, role, email, and one-application confirmation", CONFIRM),
    (
        "Requested amount and 18/21/24-month duration",
        f"Amount draft USD {AMOUNT_REQUESTED:,}; duration still " + CONFIRM,
    ),
    ("Collaborators and written roles", CONFIRM),
    ("Climate/weather data source and permission or licence", CONFIRM),
    ("District, facility, community, ethics, and data approvals", CONFIRM),
    (
        "Baseline values, sample, targets, and budget",
        f"Budget draft total USD {AMOUNT_REQUESTED:,}; baselines/targets still " + CONFIRM,
    ),
    ("Prior GCC applications/awards and third-party IP", CONFIRM),
]

PROBLEM = (
    "In Kampala peri-urban communities served by FairBanks Community Reach — including "
    "Bukoto, Kyebando, Kisaasi, Kamwokya, and Kikaaya — climate change is shifting when "
    "and where health risk rises. Changing rainfall, temperature, and humidity alter "
    "mosquito breeding and malaria transmission, while extreme heat raises risk for "
    "pregnant women and people living with cardiovascular disease or diabetes. Local "
    "actors already collect pieces of the picture: CHW and VHT household reports, "
    "outreach screening, facility visits, medicine use, and weather information. Clinical "
    "records often stay locked inside existing EMR or hospital management systems, so "
    "community, climate, and facility signals rarely fuse in real time. The gap is not "
    "another stand-alone forecast. It is the missing end-to-end bridge from climate-driven "
    "risk signals to timely health service action — named owners, protocols, referral, "
    "follow-up, and supply readiness. Without that bridge, households face delayed "
    "prevention, testing, and continuity of care, and local systems cannot adapt services "
    "during rainfall or heat stress periods. FairBanks needs a practical climate-informed "
    "early warning and monitoring workflow that builds local resilience through earlier, "
    "targeted action — not a dashboard that stops at an alert."
)

INNOVATION = (
    "FCHIP Climate Health Early Action is a climate-integrated FairBanks platform "
    "component on the Community Reach Data and Feedback loop. It directly addresses "
    "Nexa's climate-informed early warning and monitoring focus by integrating approved "
    "rainfall, temperature, humidity, and seasonal data with structured CHW/VHT reports, "
    "outreach screening, facility trends, medicine-use signals, and secure authenticated "
    "data APIs to existing EMR/HMS systems so clinical feeds can enter in real time "
    "without replacing software facilities already use. An offline-capable mobile "
    "workflow supports frontline capture and review. Explainable rules first score "
    "time-and-place malaria risk under changing mosquito ecology and heat-sensitive risk "
    "for pregnancy hypertension/heat stress, cardiovascular disease, and diabetes. A GIS "
    "action board shows risk, data quality, and agreed triggers. Each trigger links to a "
    "named service protocol: targeted messages, household follow-up, malaria testing and "
    "referral, outreach scheduling, clinician review, or medicine and supply preparation. "
    "This Proof of Concept will test whether the full loop — sense, interpret, act, learn "
    "— is feasible, accepted, timely, and able to improve access to climate-responsive "
    "care. Models will not diagnose or replace clinical judgement."
)

POPULATION = (
    "The pilot focuses on climate-vulnerable households in FairBanks Community Reach "
    "catchments in and around Kampala, including underserved peri-urban and informal "
    "settlement settings. Priority groups match Nexa's focus: pregnant women; children; "
    "older people; and people living with cardiovascular disease or diabetes — groups more "
    "exposed to malaria risk, heat stress, disrupted care, and financial barriers. CHWs "
    "and VHTs will help identify needs, test language and workflows, explain consent, "
    "collect only necessary data, and connect households to prevention, testing, referral, "
    "and follow-up. Community members will join design sessions, usability tests, feedback "
    "meetings, and review of what actions followed alerts, with gender-responsive and "
    "disability-aware engagement. Access will not depend on owning a smartphone: frontline "
    "workers will use offline-capable tools and existing outreach channels so care remains "
    "reachable during heat, heavy rain, or poor connectivity. Exact participation targets, "
    "sites, safeguarding steps, and compensation: " + CONFIRM + "."
)

INNOVATIVE = (
    "Current approaches often keep weather information, community reports, and clinical "
    "records in separate systems — including EMR/HMS platforms that do not easily share "
    "data with community tools. Many early warning products stop at forecasting and leave "
    "local teams without a named action pathway. FCHIP is designed as an end-to-end "
    "climate-informed early warning and monitoring system in Nexa's sense: monitoring and "
    "prediction linked to communication, preparedness, and timely health service delivery "
    "action. The core innovation is the climate-triggered action loop — climate plus "
    "community and facility signals, explainable trigger, named actor, protocol, referral, "
    "and follow-up result — not a weather layer bolted onto a health register. It combines "
    "technological innovation (offline capture, secure EMR/HMS data APIs, interpretable "
    "analytics, GIS), social innovation (CHW/VHT co-design and community feedback), and a "
    "delivery model rooted in a functioning medical centre and Community Reach programmes. "
    "It fits lower-resource and climate-stressed settings because it starts with simple "
    "validated rules, works offline, keeps clinicians responsible for care decisions, and "
    "does not require replacing existing EMR/HMS software. The pilot will compare "
    "timeliness, completion, acceptability, and cost with current manual planning before "
    "any machine learning is added."
)

OUTCOMES = (
    "This Proof of Concept aims to improve health access — especially timeliness — for "
    "Nexa priority outcomes, not only alert accuracy. FCHIP will combine climate signals "
    "with multi-source health data from CHWs/VHTs, FairBanks and partner clinics and "
    "hospitals, and secure EMR/HMS feeds so community and clinical evidence inform the "
    "same triggers. Primary pathways: (1) malaria — rainfall plus community and facility "
    "symptom or case signals trigger earlier prevention messages, testing, outreach, "
    "referral, and supply readiness; (2) pregnancy-related outcomes — heat thresholds "
    "trigger follow-up for hypertension, gestational diabetes risk where captured, and "
    "heat stress counselling; (3) climate-sensitive NCDs — heat triggers continuity "
    "checks and clinician review for cardiovascular disease and diabetes, using outreach "
    "and EMR/HMS care-continuity signals where available. We will measure whether more "
    "at-risk people receive a documented health service action after a climate-informed "
    "alert; whether median time from trigger to action falls; and whether referral and "
    "follow-up completion improve. Intermediary outcomes include local actors' capacity "
    "to interpret climate-health signals, ability to adapt outreach and surge during "
    "hazard periods, and continuity of essential services under rainfall or heat stress. "
    "Climate adaptiveness will be assessed by comparing performance across hazard and "
    "non-hazard periods. False alerts, missed events, equity of reach, and safety events "
    "will be reported. Baseline values and numeric targets: " + CONFIRM + ". "
    "Forecast accuracy alone will not count as success."
)

CLIMATE_DESIGN = (
    "Under Nexa's climate-innovation framework, FCHIP Climate Health Early Action is a "
    "climate-integrated innovation: it began as community health intelligence and is being "
    "fundamentally re-engineered so climate factors shape the core design, functioning, and "
    "deployment. Design: rainfall, temperature, humidity, seasonality, GIS risk geography, "
    "and locally agreed hazard thresholds become first-class inputs, not optional add-ons. "
    "Function: explainable malaria and heat-risk rules, alert timing, and service protocols "
    "respond to climate conditions and data quality. Deployment: outreach, staffing, "
    "follow-up intensity, and supply preparation are scheduled by risk period and place, "
    "including during heat, heavy rainfall, and intermittent connectivity. The pilot will "
    "test climate integration in a real FairBanks catchment and document whether "
    "performance and service response change across climate-risk periods. Exact climate "
    "data provider, resolution, history, update frequency, licences, and thresholds: "
    + CONFIRM + "."
)

TEAM_CONNECTION = (
    "FairBanks has an active local presence through its medical centre and Community Reach "
    "cascade in Kampala-area communities including Bukoto, Kyebando, Kisaasi, Kamwokya, "
    "and Kikaaya. The cascade links community members, CHWs/VHTs, outreach programmes, "
    "clinical care, research and skills, and economic empowerment including CHIS where "
    "relevant — giving a practical route for co-design, referral, and follow-up. Key local "
    "roles include frontline CHW/VHT engagement, clinical review, outreach coordination, "
    "and community feedback. The application must still name the Project Lead, clinical "
    "and technical leads, CHW/VHT representatives, MEL and data-protection leads, and "
    "confirmed collaborators, with training, lived experience, location, role, and prior "
    "partnership outcomes. Select the community connection category only from verified "
    "evidence; do not claim community-owned or community-led status until leadership and "
    "residence evidence confirms it. Names and evidence: " + CONFIRM + "."
)

OVERVIEW_QUESTIONS = [
    ("1. Specific climate-driven health gap [1,500]", PROBLEM, 1500),
    ("2. Proposed innovation [2,000]", INNOVATION, 2000),
    ("3. Population design and access [1,500]", POPULATION, 1500),
    ("4. Innovation in the target setting [2,000]", INNOVATIVE, 2000),
    ("5. Potential to improve priority outcomes [2,000]", OUTCOMES, 2000),
    ("6. Climate influence on design and deployment [1,500]", CLIMATE_DESIGN, 1500),
    ("7. Team and community connection [1,500]", TEAM_CONNECTION, 1500),
]

EXECUTION = [
    (
        "1. Execution plan [3,000]",
        "Months 1-3 — Govern and co-design: confirm governance, communities, users, climate "
        "data source and licence, baseline indicators, malaria and heat response protocols, "
        "ethics and safeguarding; run design sessions with CHWs/VHTs, clinicians, and "
        "community representatives. Months 4-7 — Build: offline forms, secure EMR/HMS data "
        "APIs, data pipeline, explainable climate-health rules, GIS action board, audit "
        "trail, and role-based access. Months 8-10 — Test: train users; run a small "
        "usability and data-quality pilot; fix workflow friction before scale-up. Months "
        "11-18 — Field pilot: operate malaria and heat-sensitive care workflows across "
        "selected sites and seasons, with monthly safety, equity, and learning reviews; "
        "refine thresholds using observed signals and service capacity. Months 19-21 — "
        "Analyse: quantitative and qualitative results, cost and sustainability scenarios, "
        "and climate-adaptiveness comparisons across hazard periods. Months 22-24 — Decide: "
        "validate findings with communities and partners, publish a responsible evidence "
        "package, and make a clear scale/no-scale decision. Learning and refinement are "
        "built into every stage. Exact sites, sample, milestones, owners, and duration "
        "(18/21/24 months): " + CONFIRM + ".",
    ),
    (
        "2. Community engagement [2,000]",
        "Use a human-rights and gender-responsive approach so marginalized voices shape "
        "design, testing, evaluation, and iteration. Fairly support CHW/VHT and community "
        "representation in co-design, usability tests, interpretation workshops, and "
        "governance reviews. Hold separate listening sessions where needed for pregnant "
        "women, older people, people with disabilities, and low-income or informal-settlement "
        "residents. Test language, consent, alert burden, referral barriers, preferred "
        "communication channels, and disability access. Publish a feedback route and show "
        "what changed after community input. Track participation by sex, age, disability "
        "where appropriate, and location without collecting unnecessary personal data. "
        "Safeguarding, complaints handling, and informed consent will be active throughout. "
        "Community representatives, compensation, safeguarding route, and engagement "
        "targets: " + CONFIRM + ".",
    ),
    (
        "3. Risks and mitigation [2,000]",
        "Key risks: weak or biased field data; false or missed alerts; alert fatigue; "
        "privacy harm; poor connectivity; low adoption; exclusion of women, older people, "
        "or people with disabilities; unsafe or incomplete referrals; partner delay; "
        "climate-data licensing limits; corruption or misuse of funds; and supply "
        "constraints after an alert. Mitigation: data minimisation and consent; "
        "role-based access, encryption, and audit logs; clinician oversight and "
        "conservative thresholds; offline workflows; named escalation and pause criteria; "
        "safeguarding and complaints handling; anti-fraud and segregated duties; partner "
        "agreements; and pre-agreed service capacity before triggering outreach. EMR/HMS "
        "APIs will use authentication, consent, and least-privilege scopes. Applicable "
        "Uganda approvals, ethics review, incident response, insurance, and organisational "
        "policies: " + CONFIRM + ".",
    ),
    (
        "4. Expected impact [2,000]",
        "Over the investment period we expect: improved local capacity to interpret "
        "climate-health signals; faster documented action after malaria and heat triggers; "
        "more at-risk people reached with prevention, testing, referral, or continuity-of-"
        "care support; improved referral completion; and better preparation of outreach and "
        "supplies during rainfall and heat periods. Reach and outcomes will be reported by "
        "priority group and location. Early signals of health-access impact are the PoC "
        "focus; we will not promise reductions in illness or hospitalisation without a "
        "powered evaluation design. Key assumptions: climate and health data arrive on "
        "time; service owners can act after alerts; communities accept the workflow; and "
        "ethics approvals are secured. Verified baseline, sample size, effect assumptions, "
        "targets, and attribution approach: " + CONFIRM + ".",
    ),
    (
        "5. Proof of Concept and objectives [2,000]",
        "Proof of Concept to establish: an end-to-end climate-informed early warning and "
        "monitoring workflow is technically feasible, acceptable, safe, affordable enough "
        "to continue, and able to improve timely access to named health service actions for "
        "malaria and heat-sensitive pregnancy/NCD care in the FairBanks catchment. Proposed "
        "measurable objectives: (1) achieve agreed completeness and timeliness for core "
        "climate and health data; (2) validate malaria and heat-risk triggers against "
        "observed local data with transparent false-positive/negative reporting; (3) reduce "
        "median time from trigger to documented service action versus baseline; (4) improve "
        "referral and follow-up completion for alerted cases; (5) reach priority groups "
        "equitably without serious safety or privacy events; and (6) produce a costed "
        "continuation decision from FairBanks and confirmed partners. Numeric targets and "
        "success thresholds: " + CONFIRM + ".",
    ),
    (
        "6. Monitoring and evaluation [2,500]",
        "Use a prospective mixed-methods pilot with a pre-defined theory of change linking "
        "climate signals to intermediary system response and priority health-access "
        "outcomes. Compare baseline and implementation periods and, if feasible, matched "
        "workflows or sites. Measure data quality; alert performance; response time; "
        "documented service actions; referrals; continuity of care; reach and equity; "
        "acceptability; adoption; cost per person reached; and safety events. Assess "
        "climate adaptiveness by comparing results across rainfall, heat, and service-"
        "pressure periods. Sources: system logs, CHW forms, facility records including "
        "approved EMR/HMS feeds, referral registers, climate data, surveys, interviews, "
        "and focus groups. Use confidence intervals and transparent missing-data analysis; "
        "report false positives and negatives. Limitations may include small sample, "
        "seasonal confounding, and incomplete comparator data — these will be stated "
        "openly. Independent statistician or research partner, final design, power/sample "
        "calculation, tools, approvals, and baseline: " + CONFIRM + ".",
    ),
    (
        "7. Sustainability [2,000]",
        "At this early stage, sustainability means designing for real operating cost and "
        "local ownership, not assuming grant-free scale. Build on existing FairBanks "
        "outreach and clinical workflows; use low-cost offline tools; train local users; "
        "document protocols; and test willingness and ability to continue after the grant. "
        "Future routes may include clinic subscriptions, programme implementation "
        "contracts, district or NGO deployments, and approved integrations including "
        "secure EMR/HMS data APIs. Barriers include climate-data cost, maintenance, "
        "workforce time, public-system fit, and reliable service capacity after alerts. "
        "The pilot will produce total-cost-of-ownership and financing scenarios and a "
        "clear scale/no-scale decision. Confirmed buyers, government pathway, partner "
        "commitments, and pricing evidence: " + CONFIRM + ".",
    ),
    (
        "8. Project team [1,500]",
        "FairBanks contributes local primary-care delivery, Community Reach outreach, "
        "CHW/VHT links, and community follow-up — the validation environment for this "
        "Proof of Concept. The project requires named leadership across clinical care, "
        "climate or epidemiology, product engineering, data protection, MEL, community "
        "engagement, finance, safeguarding, and partnerships, with time commitments "
        "matched to the workplan. Collaborators should close any skill gaps (for example "
        "climate data science, independent MEL, or district liaison) under written roles. "
        "Names, qualifications, employment links, CV evidence, and collaborator "
        "commitments: " + CONFIRM + ".",
    ),
    (
        "9. Prior partnerships [1,500]",
        "Describe only verified work with community groups, health facilities, government, "
        "academic, private-sector, and development partners. For each partner state dates, "
        "role, activity, and result so reviewers can judge stakeholder engagement capacity. "
        "Repository material supports an active community-health ecosystem but does not "
        "verify every formal partnership or outcome. Evidence and approved wording: "
        + CONFIRM + ".",
    ),
    (
        "10. Budget",
        f"We request USD {AMOUNT_REQUESTED:,} — 90% of the Proof-of-Concept maximum "
        f"(USD {POC_MAX:,}). Direct costs total USD 163,640 and indirect costs USD 16,360 "
        "(10% of direct costs, within the Fluxx 10% cap). The draft allocation funds "
        "remuneration for the core team and CHW/VHT stipends; subcontracted climate, MEL, "
        "and engineering support; Uganda travel for co-design and supervision; goods for "
        "alert-linked response and outreach; offline field equipment; project administration; "
        "and small community/VHT sub-grants. Most activities and expenses will occur in "
        "Uganda. Confirm salary scales, quotations, partner sub-grant agreements, currency "
        "assumptions, and final category totals before Fluxx certification: " + CONFIRM + ".",
    ),
]

TECH = [
    ("Climate inputs", "Rainfall, temperature, humidity, seasonality, and approved thresholds"),
    ("Health inputs", "CHW/VHT reports, outreach screening, facility trends, medicine use, and secure EMR/HMS clinical feeds"),
    ("Secure EMR/HMS APIs", "Authenticated, consent-aware APIs so existing facility systems share clinical data in real time"),
    ("Responsible analytics", "Explainable rules first; validation, monitoring, and clinician review"),
    ("Action layer", "GIS risk board, named response owner, protocol, referral, and follow-up"),
    ("Learning loop", "Response outcome, data quality, equity, safety, cost, and threshold refinement"),
]

MEL = [
    ("Data readiness", "Completeness, timeliness, missingness, geographic coverage"),
    ("Alert performance", "Sensitivity, specificity, false alerts, missed events, lead time"),
    ("Intermediary response", "Local actor capacity to interpret signals; surge/outreach adaptation; service continuity"),
    ("Priority outcomes", "Timely prevention, testing, care access, continuity; malaria and heat-sensitive care"),
    ("Climate adaptiveness", "Performance compared across rainfall, heat, and service-pressure periods"),
    ("Equity and safety", "Reach by priority group, complaints, adverse events, privacy incidents"),
    ("Feasibility", "Adoption, acceptability, uptime, cost per person reached, staff burden"),
]

CASCADE = [
    ("1. Community members", "Needs, participation, and household climate-health signals"),
    ("2. CHWs / VHTs", "Bridge: outreach, education, referral, offline data capture"),
    ("3. Community Reach programmes", "Screening, MCH, NCD, school/community education"),
    ("4. FairBanks Medical Centre", "Clinical review, diagnostics, pharmacy, referral QA; secure EMR/HMS APIs for real-time clinical ingest"),
    ("5. Research · partnerships · skills", "Evidence, climate expertise, training, ethics"),
    ("6. Economic empowerment / CHIS", "Affordable access and resilient households where data exists"),
]

RISKS = [
    ("Overclaiming climate fit", "Make climate data and hazard-responsive action central, not decorative."),
    ("Weak evidence", "Set baseline, comparison logic, sample, tools, and analysis before launch."),
    ("Unsafe automation", "No diagnosis; clinician oversight, conservative thresholds, and pause rules."),
    ("Privacy and ethics", "Minimise data, obtain consent/approval, control access, and log activity; EMR/HMS APIs use authentication and least-privilege scopes."),
    ("EMR/HMS friction", "Expose stable documented APIs; start with FairBanks systems; do not require replacing existing EMR/HMS."),
    ("Exclusion", "Offline CHW-mediated access; disability, gender, language, and poverty checks."),
    ("No capacity after alert", "Pre-agree service protocols, owners, supplies, referral routes, and escalation."),
    ("Financial misuse", "Segregated duties, procurement controls, declarations, audit trail, and reporting."),
]

SCREEN_WIN = [
    "Area of focus named: climate-informed early warning and monitoring systems.",
    "Hazards in scope: changing mosquito ecology (malaria) and extreme heat; air quality out unless data and workflow are confirmed.",
    "Priority outcomes named: malaria; pregnancy hypertension / heat stress / GDM where captured; CVD and diabetes access.",
    "Climate-integrated framing: design, function, and deployment all shaped by climate factors.",
    "End-to-end loop: climate + health signals -> explainable trigger -> named actor -> service action -> follow-up.",
    "Health access/timeliness is the success bar — not alert accuracy alone; include intermediary system-response outcomes.",
    "Climate adaptiveness: compare performance across rainfall, heat, and service-pressure periods.",
    "Community connection at least linked/partnered from verified evidence; CHW/VHT co-design visible.",
    "PoC stage honest: working FairBanks field base; climate-health impact still to be proven — no TTS claims.",
    "Every CONFIRM BEFORE SUBMISSION field replaced before Fluxx submit.",
]

SOURCES = [
    ("Official Nexa call page", OFFICIAL_PAGE),
    ("English Funding Opportunity", OFFICIAL_RFP),
    ("Proof-of-Concept questions", OFFICIAL_POC),
    ("Transition-to-Scale questions", OFFICIAL_TTS),
    ("Fluxx application portal", FLUXX),
    ("Nexa initiative site", "https://www.nexaclimate.org"),
]


def _shade(cell, color: str) -> None:
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls
    cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))


def build_docx() -> None:
    from docx import Document
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()
    sec = doc.sections[0]
    sec.page_width, sec.page_height = Inches(8.27), Inches(11.69)
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.85)
    sec.left_margin = sec.right_margin = Inches(0.72)
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10)
    styles["Normal"].font.color.rgb = RGBColor.from_string(SLATE)
    for name, size, color in (("Title", 29, NAVY), ("Heading 1", 20, NAVY), ("Heading 2", 14, TEAL)):
        style = styles[name]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)

    def run(r, size=10, bold=False, color=SLATE, italic=False):
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = RGBColor.from_string(color)

    def para(value="", size=10, bold=False, color=SLATE, italic=False, align=None, after=5):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(after)
        if align is not None:
            p.alignment = align
        run(p.add_run(value), size, bold, color, italic)
        return p

    def heading(value, level=1):
        p = doc.add_heading(value, level)
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(5)

    def bullets(items: Iterable[str]):
        for item in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(2)
            run(p.add_run(item))

    def table(headers, rows, widths=None, compact=False):
        t = doc.add_table(rows=1, cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        t.autofit = False
        for i, value in enumerate(headers):
            cell = t.rows[0].cells[i]
            _shade(cell, NAVY)
            cell.text = ""
            run(cell.paragraphs[0].add_run(str(value)), 8.5, True, WHITE)
        for idx, values in enumerate(rows):
            cells = t.add_row().cells
            for i, value in enumerate(values):
                cell = cells[i]
                cell.text = ""
                if idx % 2 == 0:
                    _shade(cell, PALE_TEAL)
                if "CONFIRM BEFORE SUBMISSION" in str(value):
                    _shade(cell, PALE_ORANGE)
                run(
                    cell.paragraphs[0].add_run(str(value)),
                    8 if compact else 9,
                    i == 0,
                    RED if "CONFIRM BEFORE SUBMISSION" in str(value) else SLATE,
                )
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
        if widths:
            for row in t.rows:
                for i, width in enumerate(widths):
                    row.cells[i].width = Inches(width)
        doc.add_paragraph().paragraph_format.space_after = Pt(1)

    def add_photo(key, width=6.6, caption=None):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after = Pt(4)
        p.add_run().add_picture(str(photo(key)), width=Inches(width))
        if caption:
            para(caption, 8, color=MUTED, italic=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=8)

    hp = sec.header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run(hp.add_run("FairBanks | Nexa Climate and Health 2026"), 8, color=MUTED)
    fp = sec.footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run(fp.add_run(f"{SLOGAN}  |  Proof-of-Concept working pack  |  Page "), 8, color=MUTED)
    field = OxmlElement("w:fldSimple")
    field.set(qn("w:instr"), "PAGE")
    fp._p.append(field)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(photo("logo")), width=Inches(1.8))
    para(PROGRAMME, 13, True, TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=8)
    title = doc.add_paragraph(style="Title")
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.add_run(TITLE)
    para(SUBTITLE, 14, True, ORANGE, align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
    add_photo("cover", 6.7)
    para(SLOGAN, 12, True, TEAL, align=WD_ALIGN_PARAGRAPH.CENTER)
    para(
        "Proof-of-Concept Fluxx answer bank | Uganda | Confirm orange fields before submit",
        9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER,
    )
    doc.add_page_break()

    heading("Submission readiness gate")
    para(
        "The strongest fit is Proof of Concept. Do not apply as Transition to Scale without "
        "credible, quantitative climate-health results. Do not submit until every orange field "
        "has been verified by the applicant.",
        10, True, RED,
    )
    table(["Critical item", "Current position"], READINESS, [2.8, 3.8])

    heading("Innovation Screen win checklist")
    para(
        "Typically over 80% of applications are declined at the Innovation Screen, which "
        "reviews only the seven Overview answers. Paste Overview text only after this "
        "checklist is true.",
        9, italic=True, color=MUTED,
    )
    bullets(SCREEN_WIN)

    heading("1. Official call brief")
    table(["Item", "Verified information"], CALL_FACTS, [1.8, 4.8])
    heading("1.1 Why this write-up can win", 2)
    bullets([
        "Direct fit with climate-informed early warning and monitoring systems.",
        "Climate-integrated under Nexa's framework: climate shapes design, function, and deployment.",
        "End-to-end loop from climate and health signals to named service-delivery actions.",
        "Priority outcomes mapped: malaria; heat-sensitive pregnancy; CVD and diabetes access.",
        "Live FairBanks medical centre, Community Reach cascade, and CHW/VHT links for real-world PoC testing.",
        "Uganda is eligible; target groups match Nexa's climate-vulnerable priorities.",
        "Honest PoC ask: prove climate-responsive access and adaptiveness — not decorative weather data.",
    ])
    heading("1.2 Relevant hazards and outcomes", 2)
    table(
        ["Nexa scope", "Proposed focus"],
        [
            ("Changing mosquito ecology", "Malaria risk, prevention, testing, referral, and supply readiness"),
            ("Extreme heat", "Pregnancy hypertension/heat stress, cardiovascular disease, and diabetes care"),
            ("Poor air quality", "Out of initial scope unless reliable local data and respiratory workflow are confirmed"),
            ("Priority groups", "Pregnant women, children, older people, underserved urban communities, climate-sensitive NCDs"),
        ],
        [2.0, 4.6],
    )

    heading("2. Proposed investment details")
    table(
        ["Portal field", "Draft / action"],
        [
            ("Applicant organisation", CONFIRM),
            ("Project Lead", CONFIRM + " - formally affiliated, key team member, one application only"),
            ("Project title (100 characters)", TITLE),
            ("Amount requested", f"USD {AMOUNT_REQUESTED:,} — 90% of PoC maximum USD {POC_MAX:,}"),
            ("Duration", CONFIRM + " - select 18, 21, or 24 months"),
            ("Collaborators", CONFIRM + " - list role and legal name"),
            ("Country of incorporation", "Uganda - confirm against incorporation documents"),
            ("Legal entity type", CONFIRM),
            ("Country of implementation", "Uganda"),
            ("Community connection", CONFIRM + " - likely community partnered/linked; choose from evidence"),
            ("Area of focus", "Climate-informed early warning and monitoring systems"),
            ("Language", "English"),
        ],
        [2.2, 4.4],
        compact=True,
    )

    heading("3. Innovation Overview - exact character-limited answer bank")
    para(
        "Only paste verified text. Character counts below include spaces and exclude the question label. "
        "Keep a small buffer because Fluxx may count formatting differently.",
        9, italic=True, color=MUTED,
    )
    for question, answer, limit in OVERVIEW_QUESTIONS:
        heading(question, 2)
        para(answer)
        count = len(answer)
        para(
            f"Draft count: {count:,} / {limit:,} characters | Buffer: {limit - count:,}",
            8, True, GREEN if count <= limit else RED,
        )

    heading("4. Full Proof-of-Concept application")
    para(
        "These fields are reviewed after the Innovation Screen. They are application-ready structures, "
        "but facts and numbers marked for confirmation must be supplied before submission.",
        9, italic=True,
    )
    for question, answer in EXECUTION:
        heading(question, 2)
        para(answer)

    heading("4.1 Budget workbook", 2)
    para(
        f"Draft Fluxx allocation totaling USD {AMOUNT_REQUESTED:,} (90% of the USD {POC_MAX:,} "
        "PoC ceiling). Confirm salary scales, quotations, and sub-grant partners before certification.",
        9, italic=True, color=MUTED,
    )
    table(
        ["Fluxx category", "Amount (USD)", "Description / basis"],
        BUDGET_ROWS,
        [2.0, 1.25, 3.35],
        compact=True,
    )

    heading("5. Product and action pathway")
    add_photo("architecture", 6.0, "Illustrative FCHIP architecture; not evidence of a deployed climate-health product.")
    table(["Layer", "Proof-of-Concept role"], TECH, [1.8, 4.8])
    heading("5.1 Responsible AI and clinical safety", 2)
    bullets([
        "Start with explainable thresholds and rules; add machine learning only after sufficient data and validation.",
        "Do not diagnose, prescribe, or automatically deny or delay care.",
        "Show uncertainty, data quality, and the reason for each alert.",
        "Require clinical review for care decisions and define escalation and pause criteria.",
        "Monitor bias, false alerts, missed events, workload, privacy, and unintended harm.",
    ])

    heading("6. Monitoring, evaluation, and learning")
    add_photo("dashboard", 6.0, "Illustrative dashboard concept; measures and thresholds remain to be validated.")
    table(["Domain", "Measures"], MEL, [1.8, 4.8])
    para(
        "Nexa requires meaningful health access or health outcome evidence. Forecast accuracy alone is not enough.",
        10, True, ORANGE,
    )

    heading("7. Community Reach cascade and delivery")
    add_photo_path = concept("cascade")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(4)
    # Portrait cascade diagram — keep width modest so caption/table never collide with footer.
    p.add_run().add_picture(str(add_photo_path), width=Inches(4.55))
    para(
        "Canonical FairBanks Community Reach model. FCHIP sits on the Data and Feedback loop.",
        8, color=MUTED, italic=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=8,
    )
    table(["Cascade layer", "Role in this Proof of Concept"], CASCADE, [2.4, 4.2], compact=True)
    doc.add_page_break()
    bullets([
        "Community members identify needs and test usability, language, and fairness of workflows.",
        "CHWs and VHTs bridge households to prevention, screening, referral, and follow-up.",
        "Community Reach programmes coordinate outreach and health education.",
        "FairBanks Medical Centre provides clinical review, diagnostics, pharmacy, referral, and quality assurance.",
        "Research and partners support evidence, climate expertise, training, and ethical evaluation.",
        "Economic empowerment and CHIS support affordable access where enrolment data exists.",
        "Learning returns to communities and strengthens resilience, not only project reporting.",
    ])

    heading("8. Risks and safeguards")
    table(["Risk", "Mitigation"], RISKS, [1.9, 4.7], compact=True)

    heading("9. Additional portal information")
    table(
        ["Field", "Action"],
        [
            ("Previous GCC application", CONFIRM + " - Yes/No and details"),
            ("Previous GCC award", CONFIRM + " - Yes/No, grant IDs, and difference from prior work"),
            ("Citations", "Add recent, relevant climate-malaria, heat-health, CHW, and early-action evidence"),
            ("Third-party IP", CONFIRM + " - include software, models, maps, climate and health datasets"),
            ("Licences/access", CONFIRM + " - verify rights allow implementation and future scale"),
            ("Certification", "Authorised signatory must accept call terms and privacy notice"),
            ("Confidentiality", "Do not submit trade secrets or sensitive data that are unnecessary for review"),
            ("AI disclosure", "Applicant remains responsible for original, true, accurate, complete, and cited content"),
        ],
        [2.0, 4.6],
        compact=True,
    )

    heading("10. Final submission checklist")
    bullets([
        "Recheck the live call, deadline, and any revised documents.",
        "Register in Fluxx early and confirm the legal organisation record.",
        "Confirm Proof-of-Concept stage, amount, duration, and Uganda implementation.",
        "Replace every CONFIRM BEFORE SUBMISSION marker.",
        "Keep all seven Innovation Overview answers within character limits.",
        "Verify that climate factors shape design, function, and deployment.",
        "Link every risk signal to a named health service action and outcome.",
        "Complete the 24-month plan, MEL design, objectives, safeguards, and budget.",
        "Verify Project Lead, team, collaborators, community role, approvals, IP, and prior GCC history.",
        "Review for confidential information, plagiarism, unsupported claims, and citation quality.",
        "Submit in Fluxx before the deadline; save the confirmation.",
    ])

    heading("11. Official sources")
    table(["Source", "URL"], SOURCES, [2.2, 4.4], compact=True)
    para(
        f"Source check date: {date.today().isoformat()}. Official Nexa documents and Fluxx always win.",
        8, italic=True, color=MUTED,
    )

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(DOCX)
    print(f"DOCX: {DOCX}")


def convert_pdf() -> None:
    import shutil
    import win32com.client

    tmp_pdf = REPO / "tmp" / "nexa-climate-health_pdf_build.pdf"
    tmp_pdf.parent.mkdir(parents=True, exist_ok=True)
    if tmp_pdf.exists():
        tmp_pdf.unlink()

    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    document = None
    try:
        document = word.Documents.Open(str(DOCX.resolve()), ReadOnly=True)
        document.ExportAsFixedFormat(str(tmp_pdf.resolve()), 17)
    finally:
        if document is not None:
            document.Close(False)
        word.Quit()

    try:
        if PDF.exists():
            PDF.unlink()
        shutil.move(str(tmp_pdf), str(PDF))
    except PermissionError as exc:
        if tmp_pdf.exists():
            tmp_pdf.unlink(missing_ok=True)
        raise PermissionError(
            f"Cannot update {PDF.name} because it is open/locked. "
            "Close it in Foxit/Cursor and re-run the build."
        ) from exc
    print(f"PDF:  {PDF}")


def active_pdf() -> Path:
    return PDF


def _transition(slide) -> None:
    from lxml import etree
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    transition = etree.Element(f"{{{ns}}}transition", spd="slow", advClick="1")
    etree.SubElement(transition, f"{{{ns}}}fade")
    slide._element.insert(2, transition)


def build_pptx() -> None:
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(value):
        return RGBColor.from_string(value)

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line or fill)
        return shape

    def text(slide, value, x, y, w, h, size=18, color=SLATE, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        return box

    def bullets(slide, items, x, y, w, h, size=16, color=SLATE):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.space_after = Pt(8)
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
        return box

    def crop(slide, path, x, y, w, h):
        """Fill the frame (cover). Fine for photos; do not use for labelled diagrams."""
        with PILImage.open(path) as im:
            iw, ih = im.size
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        frame_ratio, image_ratio = w / h, iw / ih
        if image_ratio > frame_ratio:
            amount = (1 - frame_ratio / image_ratio) / 2
            pic.crop_left = pic.crop_right = amount
        else:
            amount = (1 - image_ratio / frame_ratio) / 2
            pic.crop_top = pic.crop_bottom = amount
        return pic

    def fit(slide, path, x, y, max_w, max_h):
        """Place the full image inside the box (contain). No content cropped."""
        with PILImage.open(path) as im:
            iw, ih = im.size
        image_ratio = iw / ih
        frame_ratio = max_w / max_h
        if image_ratio > frame_ratio:
            pw, ph = max_w, max_w / image_ratio
            px, py = x, y + (max_h - ph) / 2
        else:
            ph, pw = max_h, max_h * image_ratio
            px, py = x + (max_w - pw) / 2, y
        return slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))

    def slide():
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, 13.333, 7.5, CREAM)
        _transition(s)
        return s

    def band(s, kicker, title_value, subtitle=""):
        rect(s, 0, 0, 13.333, 0.12, TEAL)
        text(s, kicker.upper(), 0.55, 0.28, 5.2, 0.28, 10, ORANGE, True)
        text(s, title_value, 0.55, 0.58, 12.1, 0.55, 26, NAVY, True, font="Aptos Display")
        if subtitle:
            text(s, subtitle, 0.58, 1.18, 11.8, 0.32, 12, MUTED)

    def footer(s, number):
        text(s, "FairBanks | FCHIP Climate Health", 0.55, 7.18, 4.2, 0.18, 8, MUTED)
        text(s, f"{number:02}", 12.2, 7.16, 0.5, 0.18, 8, MUTED, align=PP_ALIGN.RIGHT)

    s = slide()
    crop(s, photo("cover"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 7.7, 7.5, NAVY)
    text(s, PROGRAMME.upper(), 0.7, 0.62, 6.1, 0.4, 11, GOLD, True)
    text(s, TITLE, 0.7, 1.2, 6.15, 1.65, 31, WHITE, True, font="Aptos Display")
    text(s, SUBTITLE, 0.72, 3.25, 5.9, 1.1, 18, WHITE, True)
    text(s, SLOGAN, 0.72, 5.9, 4.0, 0.35, 13, GOLD, True)
    text(s, "Proof of Concept | Uganda | July 2026", 0.72, 6.38, 4.5, 0.28, 10, WHITE)

    s = slide()
    band(s, "Climate-health gap", "Local actors see the risk too late", "Climate and health signals remain separate from service action.")
    crop(s, photo("outreach"), 0.55, 1.95, 5.3, 4.85)
    bullets(s, [
        "Changing mosquito ecology and extreme heat raise malaria and heat-health risk.",
        "CHW, outreach, facility, medicine, and weather signals stay fragmented.",
        "Clinical data often stays locked inside existing EMR/HMS systems.",
        "Pregnant women, children, older people, and people with CVD/diabetes face delayed care.",
        "Forecasts rarely name the local action, owner, referral, or follow-up.",
    ], 6.25, 1.85, 6.25, 3.9, 15)
    # Keep callout clear of bullets and footer.
    rect(s, 6.35, 6.15, 5.75, 0.48, PALE_ORANGE, ORANGE, True)
    text(s, "Nexa gap: climate signal -> timely health service action.", 6.55, 6.26, 5.4, 0.28, 12, ORANGE, True)
    footer(s, 2)

    s = slide()
    band(s, "The innovation", "FCHIP Climate Health Early Action", "An end-to-end climate-informed workflow built for low-connectivity care.")
    crop(s, photo("dashboard"), 6.55, 1.9, 6.15, 4.9)
    for i, (a, b) in enumerate([
        ("Sense", "Climate + community + secure EMR/HMS feeds"),
        ("Interpret", "Explainable malaria and heat-risk triggers"),
        ("Act", "Outreach, testing, referral, follow-up, supplies"),
        ("Learn", "Outcome, safety, equity, cost, and refinement"),
    ]):
        y = 2.0 + i * 1.08
        rect(s, 0.65, y, 5.35, 0.82, WHITE, LINE, True)
        rect(s, 0.65, y, 0.12, 0.82, TEAL)
        text(s, a, 0.95, y + 0.12, 1.2, 0.25, 14, NAVY, True)
        text(s, b, 2.15, y + 0.12, 3.45, 0.45, 11, MUTED)
    footer(s, 3)

    s = slide()
    band(s, "Climate integration", "Climate shapes design, function, and deployment", "Climate-integrated under Nexa: re-engineered, not weather bolted on.")
    crop(s, photo("gis"), 0.55, 1.9, 5.4, 4.9)
    bullets(s, [
        "Design: rainfall, temperature, humidity, seasonality, GIS, and local thresholds.",
        "Function: time-and-place malaria and heat-sensitive risk rules.",
        "Deployment: risk-based outreach, follow-up, staffing, and supply preparation.",
        "Adaptiveness: compare performance across rainfall, heat, and service pressure.",
    ], 6.3, 2.0, 6.1, 4.2, 17)
    footer(s, 4)

    s = slide()
    band(s, "Priority outcomes", "Alerts matter only when care improves", "Technical accuracy is necessary, but it is not the final outcome.")
    cards = [
        ("Malaria", "Earlier prevention, testing, referral, and supply readiness"),
        ("Pregnancy", "Heat stress, hypertension, and GDM follow-up where captured"),
        ("NCD care", "Continuity for cardiovascular disease and diabetes during heat"),
        ("System response", "Faster action, completed referrals, climate-period adaptiveness"),
    ]
    for i, (a, b) in enumerate(cards):
        col, row = i % 2, i // 2
        x, y = 0.7 + col * 6.2, 2.0 + row * 2.15
        rect(s, x, y, 5.75, 1.65, WHITE, LINE, True)
        rect(s, x, y, 5.75, 0.15, GREEN)
        text(s, a, x + 0.3, y + 0.38, 1.65, 0.35, 17, NAVY, True)
        text(s, b, x + 2.0, y + 0.35, 3.35, 0.85, 13, MUTED)
    footer(s, 5)

    s = slide()
    band(s, "People", "Designed with climate-vulnerable communities", "Access does not depend on owning a smartphone.")
    crop(s, photo("maternal"), 0.55, 1.9, 4.0, 4.85)
    crop(s, photo("elderly"), 4.75, 1.9, 3.2, 4.85)
    bullets(s, [
        "Pregnant women and children",
        "Older people",
        "People with cardiovascular disease or diabetes",
        "Underserved urban households",
        "CHWs and VHTs as frontline co-designers",
    ], 8.3, 2.05, 4.2, 4.2, 15)
    footer(s, 6)

    s = slide()
    band(s, "Proof of Concept", "What we must establish", "Feasible, accepted, safe, useful, and affordable enough to continue.")
    for i, item in enumerate([
        "Reliable core climate and health data",
        "Validated malaria and heat-risk triggers",
        "Faster documented service action",
        "Better referral and follow-up completion",
        "Equitable reach without serious harm",
        "A costed continuation and scale decision",
    ]):
        col, row = i % 3, i // 3
        x, y = 0.65 + col * 4.2, 2.0 + row * 2.0
        rect(s, x, y, 3.75, 1.5, WHITE, LINE, True)
        text(s, str(i + 1), x + 0.22, y + 0.28, 0.5, 0.45, 20, ORANGE, True)
        text(s, item, x + 0.85, y + 0.28, 2.55, 0.85, 13, NAVY, True)
    footer(s, 7)

    s = slide()
    band(s, "Execution", "A staged 24-month learning plan", "Final duration, sites, sample, owners, and targets require confirmation.")
    stages = [
        ("1-3", "Govern", "Approvals, users, baseline, climate data, protocols"),
        ("4-7", "Build", "Offline capture, EMR/HMS APIs, rules, GIS, audit"),
        ("8-10", "Test", "Training, usability, data quality, safety checks"),
        ("11-18", "Pilot", "Malaria and heat workflows across risk periods"),
        ("19-21", "Learn", "Outcomes, cost, equity, threshold refinement"),
        ("22-24", "Decide", "Evidence package and scale/no-scale decision"),
    ]
    for i, (m, a, b) in enumerate(stages):
        col, row = i % 3, i // 3
        x, y = 0.65 + col * 4.2, 1.95 + row * 2.1
        rect(s, x, y, 3.75, 1.62, WHITE, LINE, True)
        rect(s, x, y, 3.75, 0.44, TEAL)
        text(s, f"MONTHS {m}", x + 0.12, y + 0.11, 1.15, 0.18, 9, WHITE, True)
        text(s, a, x + 0.25, y + 0.7, 1.0, 0.3, 15, NAVY, True)
        text(s, b, x + 1.3, y + 0.66, 2.1, 0.65, 11, MUTED)
    footer(s, 8)

    s = slide()
    band(s, "MEL", "Measure the whole path to health action", "Mixed methods, climate adaptiveness, and no unsupported impact claims.")
    crop(s, photo("mobile"), 0.55, 1.9, 4.1, 4.9)
    for i, (a, b) in enumerate(MEL):
        y = 1.85 + i * 0.7
        text(s, a, 5.0, y, 1.85, 0.22, 11, TEAL, True)
        text(s, b, 6.9, y, 5.5, 0.4, 10, SLATE)
        rect(s, 5.0, y + 0.48, 7.0, 0.015, LINE)
    footer(s, 9)

    s = slide()
    band(s, "How FairBanks works", "Community Reach cascade + FCHIP intelligence", "FCHIP is the Data and Feedback layer — not a clinic-only app.")
    # Portrait cascade diagram must stay fully visible (no cover-crop clipping).
    cascade_h = 5.2
    cascade_w = cascade_h * 0.728
    cascade_x = 0.45
    cascade_y = 1.65
    rect(s, cascade_x, cascade_y, cascade_w + 0.16, cascade_h + 0.16, WHITE, LINE, True)
    fit(s, concept("cascade"), cascade_x + 0.08, cascade_y + 0.08, cascade_w, cascade_h)
    bullets(s, [
        "Communities -> CHWs/VHTs -> programmes",
        "Medical centre anchors clinical action",
        "Research, skills, and partners strengthen evidence",
        "CHIS and livelihoods support access",
        "Climate + GIS alerts close the loop",
    ], cascade_x + cascade_w + 0.45, 2.0, 13.0 - (cascade_x + cascade_w + 0.7), 4.4, 15)
    footer(s, 10)

    s = slide()
    band(s, "Readiness", "Critical facts before Fluxx submission", "Orange items cannot be guessed or inferred.")
    rect(s, 0.7, 1.95, 11.9, 4.85, PALE_ORANGE, ORANGE, True)
    text(s, "CONFIRM BEFORE SUBMISSION", 1.05, 2.25, 4.0, 0.35, 16, RED, True)
    bullets(s, [
        "Legal applicant, good standing, Project Lead, team, and collaborators",
        "Amount, duration, sites, sample, baseline, numeric targets, and full budget",
        "Climate-data source, licence, resolution, history, and thresholds",
        "Community category, approvals, safeguarding, ethics, and data protection",
        "Partner commitments, prior GCC history, third-party IP, and citations",
        "Every Innovation Overview answer within its exact character limit",
    ], 1.05, 2.9, 10.9, 3.4, 16)
    footer(s, 11)

    s = slide()
    crop(s, photo("cover"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    text(s, "THE PROOF-OF-CONCEPT ASK", 0.75, 0.75, 4.2, 0.35, 12, GOLD, True)
    text(s, "Help local health actors act before climate-driven risk becomes harm.", 0.75, 1.35, 8.6, 1.45, 30, WHITE, True, font="Aptos Display")
    bullets(s, [
        f"USD {AMOUNT_REQUESTED:,} requested (90% of PoC maximum)",
        "18, 21, or 24 months",
        "Climate-informed early warning linked to timely care",
    ], 0.78, 3.35, 6.3, 2.0, 18, WHITE)
    text(s, SLOGAN, 0.78, 6.35, 3.6, 0.35, 13, GOLD, True)
    text(s, FLUXX, 8.2, 6.4, 4.3, 0.25, 9, WHITE, align=PP_ALIGN.RIGHT)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)
    print(f"PPTX: {PPTX}")


def validate() -> None:
    from zipfile import BadZipFile, ZipFile
    import fitz
    from docx import Document
    from pptx import Presentation

    for path in (DOCX, PDF, PPTX):
        if not path.exists() or path.stat().st_size < 20_000:
            raise RuntimeError(f"Missing or unexpectedly small output: {path}")
    for path in (DOCX, PPTX):
        try:
            with ZipFile(path) as zf:
                bad = zf.testzip()
                if bad:
                    raise RuntimeError(f"Corrupt archive member: {bad}")
        except BadZipFile as exc:
            raise RuntimeError(f"Corrupt Office file: {path}") from exc

    doc = Document(DOCX)
    content = "\n".join(
        [p.text for p in doc.paragraphs]
        + [c.text for t in doc.tables for r in t.rows for c in r.cells]
    )
    for phrase in ("Submission readiness gate", "Innovation Overview", "CONFIRM BEFORE SUBMISSION", OFFICIAL_POC):
        if phrase not in content:
            raise RuntimeError(f"DOCX validation failed: {phrase}")
    for question, answer, limit in OVERVIEW_QUESTIONS:
        if len(answer) > limit:
            raise RuntimeError(f"Character limit exceeded: {question}")

    direct = sum(int(row[1].replace(",", "")) for row in BUDGET_ROWS[:-2])
    indirect = int(BUDGET_ROWS[-2][1].replace(",", ""))
    total = int(BUDGET_ROWS[-1][1].replace(",", ""))
    if total != AMOUNT_REQUESTED:
        raise RuntimeError(f"Budget total {total} != amount requested {AMOUNT_REQUESTED}")
    if direct + indirect != total:
        raise RuntimeError(f"Budget parts do not sum: {direct}+{indirect}!={total}")
    if indirect > direct * 0.10 + 0.01:
        raise RuntimeError(f"Indirect {indirect} exceeds 10% of direct {direct}")

    pdf = fitz.open(active_pdf())
    if pdf.page_count < 10:
        raise RuntimeError(f"PDF has too few pages: {pdf.page_count}")
    pdf_text = "\n".join(page.get_text() for page in pdf)
    if f"{AMOUNT_REQUESTED:,}" not in pdf_text and str(AMOUNT_REQUESTED) not in pdf_text:
        raise RuntimeError("PDF missing amount requested")
    for page in pdf:
        blocks = page.get_text("blocks")
        for block in blocks:
            x0, y0, x1, y1, text_value = block[:5]
            if x0 < -2 or y0 < -2 or x1 > page.rect.width + 2 or y1 > page.rect.height + 2:
                raise RuntimeError(f"PDF text outside page bounds: {text_value[:60]}")
    deck = Presentation(PPTX)
    if len(deck.slides) != 12:
        raise RuntimeError(f"Expected 12 slides, found {len(deck.slides)}")
    print(f"Validated: {pdf.page_count} PDF pages | {len(deck.slides)} PPT slides")
    pdf.close()


def main() -> None:
    print(f"Building {PROGRAMME} application pack")
    print(f"Source check date: {date.today().isoformat()}")
    OUT.mkdir(parents=True, exist_ok=True)
    build_docx()
    convert_pdf()
    build_pptx()
    validate()
    print("Application pack complete.")


if __name__ == "__main__":
    main()
