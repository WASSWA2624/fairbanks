#!/usr/bin/env python3
"""Build the Africa Health-Tech Accelerator 2026 application pack.

Creates one synchronized application set in documents/:
  documents/africa-health-tech-accelerator_word.docx
  documents/africa-health-tech-accelerator_pdf.pdf
  documents/africa-health-tech-accelerator_ppt.pptx

Also creates a separate pitch-deck folder (does not replace documents/):
  pitch-deck/FairBanks_FCHIP_Africa_Health_Tech_Accelerator_2026_Pitch.pptx
  pitch-deck/FairBanks_FCHIP_Africa_Health_Tech_Accelerator_2026_Pitch.pdf

The pitch-deck pair is copied to a dedicated Google Drive folder for the
portal Pitch Deck (upload link) field.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable


HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
OUT = HERE / "documents"
OUT_DOCX = OUT / "africa-health-tech-accelerator_word.docx"
OUT_PDF = OUT / "africa-health-tech-accelerator_pdf.pdf"
OUT_PPTX = OUT / "africa-health-tech-accelerator_ppt.pptx"
PITCH_DIR = HERE / "pitch-deck"
PITCH_PPTX = (
    PITCH_DIR / "FairBanks_FCHIP_Africa_Health_Tech_Accelerator_2026_Pitch.pptx"
)
PITCH_PDF = (
    PITCH_DIR / "FairBanks_FCHIP_Africa_Health_Tech_Accelerator_2026_Pitch.pdf"
)
ASSETS = REPO / "assets"
TMP = REPO / "tmp"
# Previous Drive folder kept as-is; pitch uploads go to this new folder only.
DRIVE_PREVIOUS_DIR = Path(
    r"G:\My Drive\Application Docs\FairBanks\Africa Health-Tech Accelerator 2026"
)
DRIVE_DIR = Path(
    r"G:\My Drive\Application Docs\FairBanks"
    r"\Africa Health-Tech Accelerator 2026 - Pitch Deck"
)

OFFICIAL_HOME = "https://accelerator.africahealthexcon.com/"
OFFICIAL_APPLY = "https://accelerator.africahealthexcon.com/apply"
OFFICIAL_PUBLICATIONS = "https://accelerator.africahealthexcon.com/publications"
OFFICIAL_LINKEDIN = (
    "https://www.linkedin.com/posts/africa-health-tech-accelerator_"
    "healthtechstartup-healthtech-digitalhealth-activity-"
    "7477401395688574976-TAof"
)

NAVY = "0A1F2E"
TEAL = "0D6E6E"
TEAL_LIGHT = "14A3A3"
ORANGE = "C45C26"
GOLD = "D99A2B"
CREAM = "F7F5F0"
PALE_TEAL = "E8F3F2"
PALE_ORANGE = "FBEDE6"
SLATE = "1E2F38"
MUTED = "52636C"
LINE = "CED9D8"
WHITE = "FFFFFF"
RED = "A3312D"

SLOGAN = "Your health, our mission."
TITLE = "FairBanks Community Health Intelligence Platform (FCHIP)"
PROGRAMME = "Africa Health-Tech Accelerator 2026"
DEADLINE = "20 July 2026 (extended deadline - recheck before submission)"
STATUS = (
    "Application draft - FCHIP MVP confirmed; founder identity, commercial "
    "fields, and public Drive link still require confirmation"
)
PRODUCT_STAGE = "MVP (working product; validating in FairBanks Community Reach)"
PITCH_DRIVE_NOTE = (
    "Use files from the local pitch-deck/ folder (or the Drive folder "
    f"{DRIVE_DIR.name}). Enable Anyone-with-the-link access, test while "
    "signed out, then paste the share URL here. Leave the previous Drive "
    "folder untouched."
)

ASSET_MAP = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "outreach": "outreach_bp_screening.jpeg",
    "dashboard": "dashboard_demo.png",
    "architecture": "data_flow_iso_labeled.png",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "gis": "gis_hotspots.png",
    "team": "staff_team_reception.jpeg",
    "facility": "facility_exterior_entrance_01.jpg",
    "community": "outreach_audience_full_group_01.jpg",
    "maternal": "waiting_room_mothers_01.jpeg",
    "concept": str(REPO / ".cursor" / "concept_improved.jpeg"),
}


def asset(key: str) -> Path:
    value = ASSET_MAP[key]
    path = Path(value)
    return path if path.is_absolute() else ASSETS / path


CALL_FACTS = [
    ("Programme", "Six-month pan-African health-tech accelerator"),
    ("Delivery", "Primarily online, with one in-person milestone event"),
    ("Applicant", "Africa-based early-stage health-tech startup or SME"),
    ("Eligibility", "At least two founders; minimum MVP; 80% session commitment"),
    (
        "Benefits",
        "Training, one-to-one mentorship, seed-funding opportunity, investor "
        "access, Demo Day exposure, and continental visibility",
    ),
    ("Deadline", DEADLINE),
]

READINESS_GATES = [
    "Confirm at least two founders and list every founder with a role.",
    "FCHIP MVP is confirmed - keep a short demo or screenshots ready for "
    "screening calls.",
    "Confirm the presenting founder can attend at least 80% of sessions.",
    "Replace every CONFIRM BEFORE SUBMISSION field with verified information.",
    "Upload the pitch deck (PPT + PDF) to Google Drive with public access; "
    "test the link while signed out.",
]

PROBLEM_ANSWER = (
    "Primary healthcare in underserved African communities is still reactive. "
    "Facilities often see people only after illness has worsened, while useful "
    "signals from CHW and VHT visits, outreach screenings, maternal care, "
    "pharmacy activity, and facility records remain fragmented. Clinical data "
    "often stays locked inside existing EMR or hospital management systems, so "
    "community and facility intelligence cannot fuse in real time. This delays "
    "outbreak detection, high-risk pregnancy referral, NCD follow-up, and "
    "medicine planning. Clinics and district teams therefore lack a timely, "
    "community-level view of who is at risk, where health threats are growing, "
    "and what action should happen next."
)

SOLUTION_ANSWER = (
    "FCHIP is an offline-capable community health intelligence platform with a "
    "working MVP. CHWs and VHTs capture structured household and outreach data "
    "through mobile tools; approved facility and programme data are synchronized "
    "securely; and FCHIP safely exposes authenticated data APIs to existing EMR "
    "and HMS systems so clinics can share clinical records into the platform in "
    "real time without replacing the software they already use. An analytics "
    "layer produces risk flags, GIS hotspot maps, referrals, and dashboards. "
    "The first release focuses on three practical use cases: maternal risk and "
    "missed-care alerts, hypertension and diabetes hotspot detection, and "
    "community disease surveillance. FairBanks is validating the MVP in its "
    "existing medical centre and Community Reach network before expanding to "
    "partner clinics and districts."
)

TARGET_CUSTOMERS = (
    "Private and faith-based clinics, district health offices, the Ministry of "
    "Health, Uganda health federations, NGOs and development programmes, "
    "universities and other teaching and research institutions, and health "
    "insurers that need better community-level data, follow-up, early warning, "
    "and programme monitoring. CHWs and VHTs are the core frontline users."
)

UNIQUE_ANSWER = (
    "FairBanks combines what many health-tech startups and facilities lack: a "
    "live medical centre, active outreach programmes, CHW and VHT links, "
    "digital health workflows, and a working FCHIP MVP designed from real "
    "community care. This gives the product a direct test environment and a "
    "trusted route to users. Secure EMR/HMS data APIs let facilities plug "
    "clinical feeds into FCHIP without ripping out existing systems. The "
    "platform is not only an electronic register; it closes the loop from "
    "community and facility signal to risk insight, referral, targeted "
    "outreach, and follow-up."
)

BOTTLENECK_ANSWER = (
    "Our biggest bottleneck is converting a working MVP and a strong real-world "
    "health delivery foundation into validated, investment-ready traction. We "
    "need disciplined product-market fit evidence with CHWs and facility teams, "
    "clearer buying pathways beyond one FairBanks catchment, and partnerships "
    "that support deployment at district and programme scale. We also need to "
    "strengthen our evidence, data-governance, and fundraising package without "
    "overbuilding before the highest-value workflows are proven."
)

MOTIVATION_ANSWER = (
    "The accelerator directly matches FCHIP's next milestone: validate an "
    "Africa-built digital health MVP and prepare it for cross-border scale. "
    "Mentorship can sharpen our product choices and business model; healthcare "
    "and public-sector networks can open credible pilot and integration "
    "pathways; and investor-readiness support can turn field evidence into a "
    "financeable growth plan. We will contribute a practical Ugandan "
    "perspective grounded in community outreach, primary care, CHW and VHT "
    "workflows, and the daily realities of underserved urban communities."
)

APPLICATION_SECTIONS = [
    (
        "Section 1 of 6 - Startup and founding team",
        [
            ("Startup Name*", TITLE),
            ("Founder Name*", "[CONFIRM BEFORE SUBMISSION]"),
            ("Founder Gender*", "[CONFIRM BEFORE SUBMISSION: Male or Female]"),
            ("Founder Mobile*", "[CONFIRM BEFORE SUBMISSION]"),
            ("Founder Email*", "[CONFIRM BEFORE SUBMISSION]"),
            ("Founder LinkedIn", "[CONFIRM BEFORE SUBMISSION]"),
            ("Startup Website*", "[CONFIRM BEFORE SUBMISSION]"),
            (
                "Startup Stage*",
                "MVP - working FCHIP product ready for structured validation "
                "in FairBanks Community Reach (applicant-confirmed).",
            ),
            ("Country Headquarter*", "Uganda"),
            ("Primary African office if headquarters is outside Africa", "Not applicable"),
        ],
    ),
    (
        "Section 2 of 6 - Startup overview",
        [
            ("What Problem Are You Solving?*", PROBLEM_ANSWER),
            ("What is Your Solution?*", SOLUTION_ANSWER),
            ("Who Are Your Target Customers?*", TARGET_CUSTOMERS),
            (
                "Startup Innovation Area*",
                "Select: Digital Health Platforms; Health Data & Analytics; "
                "Artificial Intelligence in Healthcare; Electronic Health "
                "Records; Remote Patient Monitoring; Preventive Health & "
                "Wellness Technologies. Add Pharmaceutical Supply Chain "
                "Solutions only if medicine forecasting is in the MVP.",
            ),
            ("Other areas", "Community health intelligence and CHW/VHT decision support"),
            ("What Makes Your Solution Unique?*", UNIQUE_ANSWER),
            ("Which Countries Are You Operating In?*", "Uganda"),
            (
                "Pitch Deck (upload link)*",
                f"[CONFIRM AFTER UPLOAD] {PITCH_DRIVE_NOTE}",
            ),
        ],
    ),
    (
        "Section 3 of 6 - Eligibility screening",
        [
            ("Is your startup based in Africa?*", "Yes - Uganda"),
            ("Does your startup provide a Health-Tech product/service?*", "Yes"),
            ("Co-Founder Names and Roles*", "[CONFIRM BEFORE SUBMISSION]"),
            ("Total Founders (including yourself)*", "[CONFIRM BEFORE SUBMISSION: minimum 2]"),
            (
                "Can you commit to 80%+ of programme sessions?*",
                "[CONFIRM BEFORE SUBMISSION: select Yes only if true]",
            ),
        ],
    ),
    (
        "Section 4 of 6 - Traction and readiness",
        [
            (
                "Monthly Revenue Range*",
                "[CONFIRM BEFORE SUBMISSION: Less than $5,000 / $5,000-$20,000 "
                "/ $20,000-$50,000 / $50,000+]",
            ),
            (
                "Number of Paying Customers*",
                "[CONFIRM BEFORE SUBMISSION: 0 / 1-10 / 11-50 / 51-200 / 200+]",
            ),
            (
                "Fundraising in the next 6 months?*",
                "[CONFIRM BEFORE SUBMISSION: Yes or No]",
            ),
            ("Raised funding before?*", "[CONFIRM BEFORE SUBMISSION: Yes or No]"),
            ("If yes, amount raised (USD)", "[CONFIRM BEFORE SUBMISSION]"),
            (
                "Current Partners*",
                "[CONFIRM BEFORE SUBMISSION. Do not count a prospective partner. "
                "Portal options: Hospitals/Clinics; Pharma; Insurers; NGOs; "
                "Government; Corporates; Universities; None Yet.]",
            ),
            ("Other Partners", "[CONFIRM BEFORE SUBMISSION]"),
            (
                "Next 6-Month Milestone*",
                "Validate product-market fit for the working FCHIP MVP with "
                "CHWs, facility teams, and partner pathways.",
            ),
        ],
    ),
    (
        "Section 5 of 6 - Expected support",
        [
            (
                "Top 3 Priorities*",
                "1. Validation & Market Research; 2. Product Development "
                "(UX iteration on the MVP); 3. Partnerships & Ecosystem Access.",
            ),
            (
                "Alternative priority",
                "Replace Product Development with Investment Readiness & "
                "Fundraising if the MVP is already stable and validated.",
            ),
        ],
    ),
    (
        "Section 6 of 6 - Final notes",
        [
            ("What is your biggest bottleneck right now?*", BOTTLENECK_ANSWER),
            (
                "Why do you want to join Africa Health-Tech Accelerator?*",
                MOTIVATION_ANSWER,
            ),
            (
                "How did you hear about the programme?*",
                "[CONFIRM BEFORE SUBMISSION: choose the truthful portal option. "
                "Use Other if discovered through an opportunity listing.]",
            ),
        ],
    ),
]

TECH_STACK = [
    ("Offline mobile capture", "Structured CHW/VHT visits, vitals, referrals, and follow-up"),
    ("Secure EMR/HMS data APIs", "Authenticated, consent-aware APIs so existing EMR/HMS systems can share clinical data in real time"),
    ("Secure data pipeline", "Validation, consent, role-based access, and approved integrations"),
    ("Analytics and AI", "Start with interpretable rules; add validated ML as data quality grows"),
    ("GIS risk layer", "Community hotspots, coverage gaps, and outreach priorities"),
    ("Action dashboard", "Alerts, referrals, trends, data quality, and medicine planning"),
]

EVIDENCE_BASE = [
    "Functioning FairBanks Medical Centre with direct access to care workflows.",
    "Community Reach activity in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and nearby areas.",
    "CHW and VHT links for outreach, referral, follow-up, and community data collection.",
    "Internal application materials cite 40+ CHWs and VHTs collecting field data weekly "
    "(confirm the current number before using it in the portal).",
    "Existing maternal and child health, Gericare, NCD screening, school, and corporate health work.",
    "Digital health records and pharmacy workflows that inform product integration.",
    "Secure EMR/HMS data APIs so existing facility systems can feed clinical "
    "data into FCHIP in real time without replacement.",
    "Working FCHIP MVP (applicant-confirmed) ready for structured validation "
    "inside the live FairBanks cascade.",
]

SIX_MONTH_PLAN = [
    ("Month 1", "Lock users, workflows, governance, demo evidence, and baseline indicators."),
    ("Months 2-3", "Run CHW/facility usability testing; improve offline capture and dashboards."),
    ("Months 3-4", "Validate maternal, NCD, and surveillance workflows in FairBanks catchment."),
    ("Months 4-5", "Test willingness to pay and partnership model with clinics, NGOs, and districts."),
    ("Month 6", "Publish pilot evidence, investor data room, partnership pipeline, and scale plan."),
]

BUSINESS_MODEL = [
    ("Clinics and hospitals", "Subscription plus onboarding and support"),
    ("Districts and ministries", "Deployment, implementation, and analytics contracts"),
    ("NGOs and programmes", "Monitoring, evaluation, workflow, and reporting contracts"),
    ("Research partners", "Ethical study tools and approved anonymised analytics"),
    ("Digital health partners", "Secure EMR/HMS data APIs and other integration services"),
    ("CHW organisations", "Training and implementation support"),
]

RISKS = [
    ("Eligibility", "Confirm two founders before submission; MVP is already confirmed."),
    ("Data quality", "Use simple forms, validation rules, supervision, and data-quality dashboards."),
    ("Privacy", "Consent, minimisation, role-based access, audit logs, and Uganda compliance; EMR/HMS APIs use authentication and least-privilege scopes."),
    ("EMR/HMS friction", "Expose stable documented APIs; start with FairBanks systems; do not require replacing existing EMR/HMS."),
    ("AI accuracy", "Begin with interpretable rules and clinician review; validate before automation."),
    ("Adoption", "Co-design with CHWs, facilities, districts, and programme partners."),
    ("Commercial proof", "Test willingness to pay and buying pathways during the pilot."),
]


def _set_cell_shading(cell, color: str) -> None:
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    cell._tc.get_or_add_tcPr().append(
        parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>')
    )


def _set_cell_margins(cell, top=100, start=120, bottom=100, end=120) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for margin, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{margin}"))
        if node is None:
            node = OxmlElement(f"w:{margin}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def build_docx() -> None:
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()
    section = doc.sections[0]
    section.page_width = Inches(8.27)
    section.page_height = Inches(11.69)
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.72)
    section.right_margin = Inches(0.72)

    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10)
    styles["Normal"].font.color.rgb = RGBColor.from_string(SLATE)
    styles["Title"].font.name = "Aptos Display"
    styles["Title"].font.size = Pt(30)
    styles["Title"].font.bold = True
    styles["Title"].font.color.rgb = RGBColor.from_string(NAVY)
    for style_name, size, color in (
        ("Heading 1", 20, NAVY),
        ("Heading 2", 14, TEAL),
        ("Heading 3", 11, ORANGE),
    ):
        style = styles[style_name]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)

    def set_run(run, size=10, bold=False, color=SLATE, italic=False):
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text="", size=10, bold=False, color=SLATE, align=None, italic=False, after=5):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(after)
        if align is not None:
            p.alignment = align
        run = p.add_run(text)
        set_run(run, size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text: str, level=1):
        p = doc.add_heading(text, level=level)
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(5)
        return p

    def bullets(items: Iterable[str]):
        for item in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(2)
            set_run(p.add_run(item), size=10)

    def table(headers, rows, widths=None, compact=False):
        t = doc.add_table(rows=1, cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        t.autofit = False
        for i, header in enumerate(headers):
            cell = t.rows[0].cells[i]
            cell.text = ""
            _set_cell_shading(cell, NAVY)
            _set_cell_margins(cell, top=80, bottom=80)
            p = cell.paragraphs[0]
            set_run(p.add_run(str(header)), size=9, bold=True, color=WHITE)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for row_index, values in enumerate(rows):
            cells = t.add_row().cells
            for i, value in enumerate(values):
                cell = cells[i]
                cell.text = ""
                _set_cell_margins(cell, top=70 if compact else 100, bottom=70 if compact else 100)
                if row_index % 2 == 0:
                    _set_cell_shading(cell, PALE_TEAL)
                if "CONFIRM BEFORE SUBMISSION" in str(value):
                    _set_cell_shading(cell, PALE_ORANGE)
                p = cell.paragraphs[0]
                set_run(
                    p.add_run(str(value)),
                    size=8.5 if compact else 9,
                    bold=i == 0,
                    color=RED if "CONFIRM BEFORE SUBMISSION" in str(value) else SLATE,
                )
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
        if widths:
            for row in t.rows:
                for i, width in enumerate(widths):
                    row.cells[i].width = Inches(width)
        doc.add_paragraph().paragraph_format.space_after = Pt(1)
        return t

    def add_photo(key, width=6.8, caption=None):
        path = asset(key)
        if not path.exists():
            return
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(path), width=Inches(width))
        if caption:
            cp = doc.add_paragraph()
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            set_run(cp.add_run(caption), size=8, color=MUTED, italic=True)

    # Header/footer
    header = section.header
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    set_run(hp.add_run("FCHIP | Africa Health-Tech Accelerator 2026"), size=8, color=MUTED)
    footer = section.footer
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_run(
        fp.add_run(f"{SLOGAN}  |  Application working pack  |  Page "),
        size=8,
        color=MUTED,
    )
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), "PAGE")
    fp._p.append(fld)

    # Cover
    if asset("logo").exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(asset("logo")), width=Inches(2.0))
    para(PROGRAMME, size=13, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=8)
    title_p = doc.add_paragraph(style="Title")
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.add_run(TITLE)
    para(
        "From community health signals to timely action across Africa",
        size=15,
        bold=True,
        color=ORANGE,
        align=WD_ALIGN_PARAGRAPH.CENTER,
        after=10,
    )
    add_photo("cover", width=6.8)
    para(SLOGAN, size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=8)
    table(
        ["Application detail", "Working position"],
        [
            ("Deadline", DEADLINE),
            ("Headquarters", "Uganda"),
            ("Product stage", PRODUCT_STAGE),
            ("Document status", STATUS),
        ],
        widths=[1.8, 4.8],
        compact=True,
    )
    para(
        "Prepared from the live official programme site and all six sections of "
        f"the application portal on {datetime.now().strftime('%d %B %Y')}.",
        size=8,
        italic=True,
        color=MUTED,
        align=WD_ALIGN_PARAGRAPH.CENTER,
    )
    doc.add_page_break()

    # Readiness gate
    heading("Submission readiness gate", 1)
    para(
        "The portal requires at least two founders and a minimum viable product. "
        "FCHIP MVP status is confirmed. Do not invent paying FCHIP customers, "
        "AI performance scores, or unsigned partners. Complete founder and "
        "commercial CONFIRM fields before submission.",
        bold=True,
        color=ORANGE,
    )
    bullets(READINESS_GATES)
    table(
        ["Critical field", "Current status"],
        [
            ("Founder identity and contacts", "CONFIRM BEFORE SUBMISSION"),
            ("Co-founders and roles", "CONFIRM BEFORE SUBMISSION"),
            ("Working FCHIP MVP", "Confirmed - keep demo/screenshots ready"),
            ("Revenue and paying customers", "CONFIRM BEFORE SUBMISSION"),
            ("Funding raised and partner list", "CONFIRM BEFORE SUBMISSION"),
            ("Website", "CONFIRM BEFORE SUBMISSION"),
            ("Public pitch-deck Drive link", "CONFIRM AFTER UPLOAD (PPT + PDF)"),
        ],
        widths=[2.6, 4.0],
    )

    heading("1. Official programme brief", 1)
    table(["Item", "Verified call information"], CALL_FACTS, widths=[1.6, 5.0])
    heading("1.1 Strong fit with FCHIP", 2)
    bullets(
        [
            "Digital health platform with health data and analytics at its core.",
            "AI-assisted risk insight, used only after validation and clinical oversight.",
            "Preventive health, remote monitoring, secure EMR/HMS APIs, and supply planning.",
            "Africa-based product with a clear Uganda-to-East-Africa scale path.",
            "A live care and outreach environment for real-world product validation.",
        ]
    )
    heading("1.2 Programme journey", 2)
    table(
        ["Stage", "Official programme description"],
        [
            ("Open call", "Application and screening; deadline extended to 20 July 2026"),
            ("Bootcamp and selection", "Intensive virtual bootcamp; up to 50 startups targeted"),
            ("Accelerator", "Weekly sessions, office hours, and Demo Day preparation"),
            ("Closing event", "Pitching, winner announcement, and funding pathway"),
        ],
        widths=[1.8, 4.8],
    )

    doc.add_page_break()
    heading("2. Portal answer bank", 1)
    para(
        "Copy only verified answers into the live portal. The wording below is "
        "submission-ready where facts are known; orange/red fields require the applicant.",
        italic=True,
    )
    for section_title, rows in APPLICATION_SECTIONS:
        heading(section_title, 2)
        table(["Portal question", "Draft answer / action"], rows, widths=[2.15, 4.45], compact=True)

    doc.add_page_break()
    heading("3. Venture narrative", 1)
    add_photo(
        "outreach",
        width=6.7,
        caption="FairBanks outreach provides a real environment for FCHIP product validation.",
    )
    heading("3.1 Problem", 2)
    para(PROBLEM_ANSWER)
    heading("3.2 Solution", 2)
    para(SOLUTION_ANSWER)
    heading("3.3 Competitive advantage", 2)
    para(UNIQUE_ANSWER)

    heading("4. Product and technology", 1)
    add_photo(
        "dashboard",
        width=6.7,
        caption="FCHIP dashboard concept - illustrative, not evidence of a deployed product.",
    )
    table(["Product layer", "Role"], TECH_STACK, widths=[2.0, 4.6])
    para(
        "Responsible AI position: begin with simple, explainable rules and "
        "clinician review. Add machine-learning models only when the training "
        "data, bias checks, monitoring, and clinical validation are adequate.",
        bold=True,
    )

    heading("5. Target customers and business model", 1)
    para(TARGET_CUSTOMERS)
    table(["Customer", "Revenue path"], BUSINESS_MODEL, widths=[2.1, 4.5])
    heading("5.1 Market-entry logic", 2)
    bullets(
        [
            "Beachhead: FairBanks catchment and Kampala-area outreach workflows.",
            "First external buyers: clinics and health programmes needing follow-up and reporting.",
            "Institutional scale: district, NGO, research, insurer, and ministry deployments.",
            "Regional scale: adapt to countries with established CHW systems and similar care gaps.",
        ]
    )

    doc.add_page_break()
    heading("6. Traction and honest evidence", 1)
    add_photo(
        "concept",
        width=6.0,
        caption="FairBanks Community Reach operating model - FCHIP is the intelligence layer.",
    )
    bullets(EVIDENCE_BASE)
    para(
        "These points prove delivery access and validation readiness. They do "
        "not by themselves prove FCHIP product traction, paying customers, or "
        "model performance. Those claims must be supported separately.",
        bold=True,
        color=RED,
    )

    heading("7. Six-month milestone plan", 1)
    para(
        "Recommended portal milestone: Validate product-market fit, provided a "
        "working MVP already exists.",
        bold=True,
    )
    table(["Timing", "Milestone"], SIX_MONTH_PLAN, widths=[1.25, 5.35])
    heading("7.1 Suggested pilot indicators", 2)
    bullets(
        [
            "CHW/VHT weekly active use and form-completion rate.",
            "Referral completion and time from risk flag to clinical review.",
            "Maternal, NCD, and surveillance alerts reviewed and acted upon.",
            "Data completeness, duplicate rate, and correction time.",
            "Facility/partner willingness to pay and signed pilot commitments.",
        ]
    )

    heading("8. Accelerator support requested", 1)
    table(
        ["Priority", "Why it matters now"],
        [
            ("Validation & Market Research", "Prove user need, buying pathway, and measurable outcomes."),
            ("Product Development", "Improve offline UX, workflow fit, dashboards, and secure EMR/HMS APIs."),
            ("Partnerships & Ecosystem Access", "Secure clinics, districts, programmes, and scale partners."),
        ],
        widths=[2.3, 4.3],
    )
    para(MOTIVATION_ANSWER)

    heading("9. Risk and governance", 1)
    table(["Risk", "Mitigation"], RISKS, widths=[1.8, 4.8], compact=True)

    doc.add_page_break()
    heading("10. Submission checklist", 1)
    checklist = [
        "Recheck the live deadline and portal availability.",
        "Verify the legal/startup name used in the form.",
        "Complete founder identity, contacts, gender, LinkedIn, and website.",
        "List at least one co-founder and verify total founders is two or more.",
        "Select Startup Stage = MVP; keep a short demo or screenshots ready.",
        "Choose only innovation areas included in the current product.",
        "Enter exact revenue, customer, fundraising, and funding figures.",
        "Name only confirmed current partners.",
        "Confirm 80% session attendance and in-person milestone availability.",
        "Upload PPT + pitch PDF to Google Drive; enable public link; test signed out.",
        "Paste the Drive share URL into Pitch Deck (upload link).",
        "Proofread all answers against this pack and the official portal.",
        "Submit before the deadline and save the confirmation.",
    ]
    for item in checklist:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        set_run(p.add_run("[  ]  " + item), size=10)

    heading("11. Official sources", 1)
    sources = [
        ("Programme home", OFFICIAL_HOME),
        ("Application portal", OFFICIAL_APPLY),
        ("Official programme article", OFFICIAL_PUBLICATIONS),
        ("Official deadline-extension post", OFFICIAL_LINKEDIN),
    ]
    table(["Source", "URL"], sources, widths=[2.1, 4.5], compact=True)
    para(
        "Source check date: "
        f"{datetime.now().strftime('%d %B %Y')}. The official portal wins if any detail changes.",
        size=9,
        italic=True,
        color=MUTED,
    )

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(OUT_DOCX)
    print(f"DOCX: {OUT_DOCX}")


def convert_docx_to_pdf() -> None:
    """Use Microsoft Word so the PDF exactly mirrors the Word application pack."""
    import shutil
    import win32com.client

    # Export to tmp first so a locked official PDF does not block Word COM.
    tmp_pdf = REPO / "tmp" / "africa-health-tech-accelerator_export.pdf"
    tmp_pdf.parent.mkdir(parents=True, exist_ok=True)
    if tmp_pdf.exists():
        tmp_pdf.unlink()

    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    document = None
    try:
        document = word.Documents.Open(str(OUT_DOCX.resolve()), ReadOnly=True)
        document.ExportAsFixedFormat(
            OutputFileName=str(tmp_pdf.resolve()),
            ExportFormat=17,  # wdExportFormatPDF
            OpenAfterExport=False,
            OptimizeFor=0,
            Range=0,
            Item=0,
            IncludeDocProps=True,
            KeepIRM=True,
            CreateBookmarks=1,
            DocStructureTags=True,
            BitmapMissingFonts=True,
            UseISO19005_1=False,
        )
    except Exception as exc:
        raise RuntimeError(
            "PDF export failed. Close the official Word/PDF files if they are "
            "open in another app, then rebuild."
        ) from exc
    finally:
        if document is not None:
            document.Close(False)
        word.Quit()

    try:
        shutil.copy2(tmp_pdf, OUT_PDF)
    except PermissionError as exc:
        raise RuntimeError(
            "PDF was exported, but the official PDF is locked. Close "
            f"{OUT_PDF.name} in your viewer, then rebuild."
        ) from exc
    print(f"PDF:  {OUT_PDF}")


def _ppt_add_transition(slide) -> None:
    from lxml import etree

    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    transition = etree.Element(f"{{{ns}}}transition", spd="slow", advClick="1")
    etree.SubElement(transition, f"{{{ns}}}fade")
    slide._element.insert(2, transition)


def build_pptx() -> None:
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(hex_value):
        return RGBColor.from_string(hex_value)

    def rect(slide, x, y, w, h, fill, line=None, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line or fill)
        return shape

    def text(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=20,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.LEFT,
        font="Aptos",
        valign=MSO_ANCHOR.TOP,
        margin=0.05,
    ):
        box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = valign
        p = frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = value
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def bullets(slide, items, x, y, w, h, size=18, color=SLATE):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.05)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.level = 0
            p.space_after = Pt(10)
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
        return box

    def image_crop(slide, path: Path, x, y, w, h):
        if not path.exists():
            rect(slide, x, y, w, h, PALE_TEAL)
            text(slide, "Visual unavailable", x, y + h / 2 - 0.2, w, 0.4, 12, MUTED, align=PP_ALIGN.CENTER)
            return None
        with PILImage.open(path) as im:
            iw, ih = im.size
        frame_ratio = w / h
        image_ratio = iw / ih
        pic = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
        )
        if image_ratio > frame_ratio:
            visible = frame_ratio / image_ratio
            crop = (1 - visible) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            visible = image_ratio / frame_ratio
            crop = (1 - visible) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
        return pic

    def footer(slide, number):
        text(slide, "FairBanks | FCHIP", 0.55, 7.12, 3.2, 0.2, 8, MUTED)
        text(slide, str(number).zfill(2), 12.25, 7.1, 0.5, 0.2, 8, MUTED, align=PP_ALIGN.RIGHT)

    def top_band(slide, kicker, title_value, subtitle=None):
        rect(slide, 0, 0, 13.333, 0.12, TEAL)
        text(slide, kicker.upper(), 0.55, 0.35, 4.5, 0.3, 10, ORANGE, True)
        text(slide, title_value, 0.55, 0.72, 12.0, 0.72, 27, NAVY, True, font="Aptos Display")
        if subtitle:
            text(slide, subtitle, 0.58, 1.45, 11.8, 0.45, 13, MUTED)

    def add_slide():
        slide = prs.slides.add_slide(blank)
        rect(slide, 0, 0, 13.333, 7.5, CREAM)
        _ppt_add_transition(slide)
        return slide

    # 1 Cover
    s = add_slide()
    image_crop(s, asset("cover"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 7.6, 7.5, NAVY)
    text(s, PROGRAMME.upper(), 0.7, 0.65, 6.1, 0.4, 12, GOLD, True)
    text(s, TITLE, 0.7, 1.25, 6.1, 1.9, 31, WHITE, True, font="Aptos Display")
    text(
        s,
        "Working MVP. Community signals. Predictive insight. Timely action.",
        0.72,
        3.45,
        5.8,
        0.85,
        18,
        WHITE,
        True,
    )
    text(s, SLOGAN, 0.72, 5.85, 4.5, 0.35, 13, GOLD, True)
    text(s, "Uganda | MVP | July 2026", 0.72, 6.35, 4.0, 0.3, 10, WHITE)

    # 2 Problem
    s = add_slide()
    top_band(s, "The problem", "Care reacts too late", "Community data exists, but it rarely becomes timely action.")
    image_crop(s, asset("outreach"), 0.55, 2.05, 5.4, 4.45)
    bullets(
        s,
        [
            "Outreach, maternal, pharmacy, and facility signals remain fragmented.",
            "Clinical records often stay locked inside existing EMR/HMS systems.",
            "High-risk pregnancies and chronic conditions are found late.",
            "Districts and clinics lack a live view of community risk.",
            "Medicine and outreach plans rely on delayed or incomplete data.",
        ],
        6.35,
        2.05,
        6.25,
        4.0,
        16,
    )
    rect(s, 6.4, 6.02, 5.7, 0.52, PALE_ORANGE, ORANGE, True)
    text(s, "The gap is not data collection alone. It is intelligence and follow-through.", 6.65, 6.12, 5.25, 0.25, 12, ORANGE, True)
    footer(s, 2)

    # 3 Solution
    s = add_slide()
    top_band(s, "The solution", "FCHIP turns signals into action", "An offline-capable intelligence layer for community primary care.")
    image_crop(s, asset("dashboard"), 6.4, 1.9, 6.35, 4.75)
    cards = [
        ("1", "Capture", "CHW/VHT plus secure EMR/HMS APIs"),
        ("2", "Understand", "Risk rules, analytics, and GIS"),
        ("3", "Act", "Alerts, referrals, and outreach"),
    ]
    for i, (n, t, b) in enumerate(cards):
        y = 2.0 + i * 1.45
        rect(s, 0.65, y, 5.2, 1.05, WHITE, LINE, True)
        rect(s, 0.82, y + 0.18, 0.55, 0.55, TEAL, TEAL, True)
        text(s, n, 0.82, y + 0.28, 0.55, 0.24, 13, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, t, 1.55, y + 0.16, 1.5, 0.3, 14, NAVY, True)
        text(s, b, 1.55, y + 0.52, 3.85, 0.28, 11, MUTED)
    footer(s, 3)

    # 4 Product architecture
    s = add_slide()
    top_band(s, "Product", "Designed for low-connectivity workflows", "Simple at the frontline; disciplined and secure behind the scenes.")
    image_crop(s, asset("architecture"), 0.6, 1.95, 7.45, 4.75)
    stack = [
        ("Offline mobile", "Structured visits and referrals"),
        ("EMR/HMS APIs", "Real-time clinical data, safely"),
        ("Secure pipeline", "Consent, validation, role access"),
        ("AI + GIS dashboards", "Explainable risk and action"),
    ]
    for i, (t, b) in enumerate(stack):
        y = 2.0 + i * 1.08
        rect(s, 8.45, y, 4.25, 0.82, WHITE, LINE, True)
        rect(s, 8.45, y, 0.12, 0.82, TEAL)
        text(s, t, 8.75, y + 0.11, 3.6, 0.25, 13, NAVY, True)
        text(s, b, 8.75, y + 0.43, 3.6, 0.2, 10, MUTED)
    footer(s, 4)

    # 5 Use cases
    s = add_slide()
    top_band(s, "Use cases", "Start focused. Prove value.", "Three live MVP workflows with clear users and actions.")
    image_crop(s, asset("mobile"), 0.55, 1.95, 4.45, 4.8)
    use_cases = [
        ("Maternal risk", "Flag missed ANC, high BP, and referral urgency."),
        ("NCD hotspots", "Map BP/glucose patterns and target screening."),
        ("Disease surveillance", "Detect rising symptom clusters and coordinate response."),
    ]
    for i, (t, b) in enumerate(use_cases):
        y = 2.0 + i * 1.4
        rect(s, 5.4, y, 7.15, 1.05, WHITE, LINE, True)
        rect(s, 5.65, y + 0.2, 0.55, 0.55, ORANGE, ORANGE, True)
        text(s, str(i + 1), 5.65, y + 0.3, 0.55, 0.2, 13, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, t, 6.45, y + 0.15, 2.1, 0.3, 15, NAVY, True)
        text(s, b, 8.55, y + 0.16, 3.6, 0.55, 12, MUTED)
    footer(s, 5)

    # 6 Market
    s = add_slide()
    top_band(s, "Market", "A shared platform for the primary-care stack", "Frontline users act; institutions buy.")
    image_crop(s, asset("gis"), 7.55, 1.85, 5.2, 4.95)
    segments = [
        ("Clinics", "Follow-up, outreach, population visibility"),
        ("Districts & MoH", "Early warning, planning, and national oversight"),
        ("Health federations", "Network coordination and shared standards"),
        ("NGOs", "Programme monitoring and evidence"),
        ("Universities", "Teaching, research, and ethical study analytics"),
        ("Insurers", "Prevention-focused population insight"),
    ]
    for i, (t, b) in enumerate(segments):
        y = 1.88 + i * 0.78
        text(s, t, 0.7, y, 2.0, 0.28, 13, TEAL, True)
        text(s, b, 2.85, y, 4.15, 0.4, 12, SLATE)
        rect(s, 0.7, y + 0.48, 6.2, 0.02, LINE)
    footer(s, 6)

    # 7 Business model
    s = add_slide()
    top_band(s, "Business model", "Diversified revenue, one intelligence core", "Recurring software plus implementation and partnership services.")
    model = [
        ("Subscriptions", "Clinics and hospitals"),
        ("Deployments", "Districts and ministries"),
        ("Contracts", "NGO monitoring and reporting"),
        ("Integrations", "Secure EMR/HMS APIs and partners"),
        ("Training", "CHW onboarding and implementation"),
        ("Research", "Approved study and analytics support"),
    ]
    for i, (t, b) in enumerate(model):
        col = i % 3
        row = i // 3
        x = 0.65 + col * 4.18
        y = 2.0 + row * 2.05
        rect(s, x, y, 3.75, 1.55, WHITE, LINE, True)
        rect(s, x, y, 3.75, 0.15, TEAL)
        text(s, t, x + 0.25, y + 0.35, 3.1, 0.35, 16, NAVY, True)
        text(s, b, x + 0.25, y + 0.86, 3.1, 0.42, 12, MUTED)
    text(s, "Beachhead: FairBanks catchment -> Kampala partners -> Uganda districts -> East Africa", 1.2, 6.35, 10.9, 0.35, 14, ORANGE, True, align=PP_ALIGN.CENTER)
    footer(s, 7)

    # 8 Validation advantage
    s = add_slide()
    top_band(s, "Why FairBanks", "A real environment to validate health-tech", "Field access, care delivery, and community trust reduce pilot friction.")
    image_crop(s, asset("facility"), 0.55, 1.9, 4.0, 4.9)
    image_crop(s, asset("team"), 4.75, 1.9, 3.55, 2.3)
    image_crop(s, asset("community"), 4.75, 4.4, 3.55, 2.4)
    bullets(
        s,
        [
            "Functioning medical centre",
            "Community Reach programmes",
            "CHW and VHT links",
            "Maternal, NCD, geriatric, school, and corporate health workflows",
            "Digital records and pharmacy signals",
        ],
        8.65,
        2.05,
        4.0,
        4.3,
        15,
    )
    footer(s, 8)

    # 9 Six-month plan
    s = add_slide()
    top_band(s, "Roadmap", "Six months to product-market evidence", "Working MVP now. Next: validate fit, buyers, and partners.")
    for i, (timing, milestone) in enumerate(SIX_MONTH_PLAN):
        x = 0.75 + i * 2.48
        rect(s, x, 2.05, 2.08, 3.65, WHITE, LINE, True)
        rect(s, x, 2.05, 2.08, 0.52, TEAL)
        text(s, timing, x + 0.08, 2.18, 1.92, 0.2, 11, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, str(i + 1), x + 0.65, 2.9, 0.78, 0.65, 28, ORANGE, True, align=PP_ALIGN.CENTER)
        text(s, milestone, x + 0.2, 3.8, 1.68, 1.45, 11, SLATE, align=PP_ALIGN.CENTER)
    text(
        s,
        "Outcome: usable evidence for buyers, partners, and Demo Day.",
        0.9,
        6.25,
        11.5,
        0.45,
        14,
        ORANGE,
        True,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 9)

    # 10 Scale
    s = add_slide()
    top_band(s, "Scale", "Built for pan-African adaptation", "Standard core; configurable workflows, language, and integrations.")
    image_crop(s, asset("concept"), 0.55, 1.85, 5.65, 4.95)
    scale_steps = [
        ("1. Validate", "FairBanks catchment with working MVP"),
        ("2. Replicate", "Partner clinics and programmes"),
        ("3. Integrate", "District and national systems"),
        ("4. Expand", "East African CHW markets"),
    ]
    for i, (t, b) in enumerate(scale_steps):
        y = 1.95 + i * 1.12
        rect(s, 6.65, y, 5.85, 0.85, WHITE, LINE, True)
        text(s, t, 6.95, y + 0.13, 1.45, 0.28, 14, TEAL, True)
        text(s, b, 8.55, y + 0.14, 3.45, 0.28, 13, SLATE)
    footer(s, 10)

    # 11 Why now
    s = add_slide()
    top_band(
        s,
        "Why now",
        "Operating care. Working MVP. Ready to validate at scale.",
        "FairBanks already runs the cascade FCHIP is built to digitise.",
    )
    image_crop(s, asset("team"), 0.55, 1.9, 4.4, 4.75)
    strengths = [
        "Working FCHIP MVP with demo-ready workflows",
        "Secure EMR/HMS data APIs for real-time clinical ingest",
        "Clear six-month product-market validation plan",
        "Accelerator fit: mentorship, networks, and Demo Day readiness",
        "Buyer evidence path for clinics, districts, and NGO partners",
        "Uganda beachhead ready to expand across East Africa",
    ]
    rect(s, 5.35, 1.95, 7.35, 4.75, WHITE, LINE, True)
    rect(s, 5.35, 1.95, 0.14, 4.75, TEAL)
    text(s, "VALIDATION ADVANTAGE", 5.75, 2.2, 6.55, 0.35, 14, TEAL, True)
    bullets(s, strengths, 5.75, 2.75, 6.2, 3.5, 15, SLATE)
    footer(s, 11)

    # 12 Ask
    s = add_slide()
    image_crop(s, asset("cover"), 0, 0, 13.333, 7.5)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    text(s, "THE ASK", 0.75, 0.75, 2.0, 0.35, 12, GOLD, True)
    text(s, "Help a working MVP become validated African scale.", 0.75, 1.35, 7.8, 1.35, 28, WHITE, True, font="Aptos Display")
    asks = [
        "Validation and market research",
        "Product and UX iteration on the MVP",
        "Partnerships and ecosystem access",
    ]
    for i, item in enumerate(asks):
        y = 3.25 + i * 0.75
        rect(s, 0.78, y, 0.38, 0.38, ORANGE, ORANGE, True)
        text(s, str(i + 1), 0.78, y + 0.08, 0.38, 0.14, 10, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, item, 1.4, y - 0.02, 5.6, 0.42, 16, WHITE, True)
    text(s, SLOGAN, 0.78, 6.35, 3.6, 0.35, 13, GOLD, True)
    text(s, OFFICIAL_APPLY, 8.15, 6.4, 4.4, 0.25, 8, WHITE, align=PP_ALIGN.RIGHT)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"PPTX: {OUT_PPTX}")


def publish_pitch_deck() -> tuple[Path, Path]:
    """Write pitch PPT+PDF into pitch-deck/ and copy to the new Drive folder.

    Leaves documents/ and the previous Drive folder unchanged for this step
    (documents PPT is still produced by build_pptx as part of the sync set).
    """
    import shutil
    import win32com.client

    if not OUT_PPTX.exists():
        raise RuntimeError(f"Missing source deck: {OUT_PPTX}")

    PITCH_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(OUT_PPTX, PITCH_PPTX)

    # Export PDF from the pitch-deck PPT copy so documents/ stays untouched.
    if PITCH_PDF.exists():
        try:
            PITCH_PDF.unlink()
        except PermissionError as exc:
            raise RuntimeError(
                f"Close {PITCH_PDF.name} if it is open, then rebuild."
            ) from exc

    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = None
    try:
        presentation = powerpoint.Presentations.Open(
            str(PITCH_PPTX.resolve()), WithWindow=False
        )
        presentation.SaveAs(str(PITCH_PDF.resolve()), 32)  # ppSaveAsPDF
    except Exception as exc:
        raise RuntimeError(
            "Pitch PDF export failed. Close the PowerPoint file if it is open, then rebuild."
        ) from exc
    finally:
        if presentation is not None:
            presentation.Close()
        powerpoint.Quit()

    if not PITCH_PDF.exists() or PITCH_PDF.stat().st_size < 20_000:
        raise RuntimeError("Pitch PDF export produced an empty or missing file.")

    print(f"PITCH PPT: {PITCH_PPTX}")
    print(f"PITCH PDF: {PITCH_PDF}")

    if not DRIVE_DIR.parent.exists():
        raise RuntimeError(
            f"Google Drive path not found: {DRIVE_DIR.parent}. "
            "Start Google Drive for desktop and rebuild."
        )
    DRIVE_DIR.mkdir(parents=True, exist_ok=True)
    ppt_dest = DRIVE_DIR / PITCH_PPTX.name
    pdf_dest = DRIVE_DIR / PITCH_PDF.name
    shutil.copy2(PITCH_PPTX, ppt_dest)
    shutil.copy2(PITCH_PDF, pdf_dest)
    print(f"DRIVE PPT: {ppt_dest}")
    print(f"DRIVE PDF: {pdf_dest}")
    print(f"Previous Drive folder retained: {DRIVE_PREVIOUS_DIR}")
    return ppt_dest, pdf_dest


def validate_outputs() -> None:
    from zipfile import ZipFile, BadZipFile

    import fitz
    from docx import Document
    from pptx import Presentation

    required = [OUT_DOCX, OUT_PDF, OUT_PPTX]
    for path in required:
        if not path.exists() or path.stat().st_size < 20_000:
            raise RuntimeError(f"Missing or unexpectedly small output: {path}")

    for archive in (OUT_DOCX, OUT_PPTX, PITCH_PPTX):
        if not archive.exists():
            raise RuntimeError(f"Missing pitch or documents archive: {archive}")
        try:
            with ZipFile(archive) as zf:
                bad = zf.testzip()
                if bad is not None:
                    raise RuntimeError(f"Corrupt archive member in {archive}: {bad}")
        except BadZipFile as exc:
            raise RuntimeError(f"Corrupt Office archive: {archive}") from exc

    if not PITCH_PDF.exists() or PITCH_PDF.stat().st_size < 20_000:
        raise RuntimeError(f"Missing or small pitch PDF: {PITCH_PDF}")

    # Fail if any official deliverable is older than the generator source.
    builder_mtime = Path(__file__).stat().st_mtime
    for path in required:
        if path.stat().st_mtime + 1 < builder_mtime:
            raise RuntimeError(f"Stale official output: {path.name}")

    doc = Document(OUT_DOCX)
    paragraph_text = [p.text for p in doc.paragraphs]
    table_text = [
        cell.text
        for table in doc.tables
        for row in table.rows
        for cell in row.cells
    ]
    text_content = "\n".join(paragraph_text + table_text)
    for phrase in (
        "Submission readiness gate",
        "Section 6 of 6 - Final notes",
        "CONFIRM BEFORE SUBMISSION",
        "Working FCHIP MVP",
        PRODUCT_STAGE,
        OFFICIAL_APPLY,
    ):
        if phrase not in text_content:
            raise RuntimeError(f"DOCX validation failed: {phrase}")

    pdf = fitz.open(OUT_PDF)
    if pdf.page_count < 8:
        raise RuntimeError(f"PDF has too few pages: {pdf.page_count}")

    ppt = Presentation(OUT_PPTX)
    if len(ppt.slides) != 12:
        raise RuntimeError(f"Pitch deck should have 12 slides, found {len(ppt.slides)}")

    pitch_ppt = Presentation(PITCH_PPTX)
    if len(pitch_ppt.slides) != 12:
        raise RuntimeError(
            f"pitch-deck PPT should have 12 slides, found {len(pitch_ppt.slides)}"
        )

    pitch_pdf = fitz.open(PITCH_PDF)
    if pitch_pdf.page_count != 12:
        raise RuntimeError(
            f"pitch-deck PDF should have 12 pages, found {pitch_pdf.page_count}"
        )

    print(
        f"Validated: {len(doc.paragraphs)} DOCX paragraphs | "
        f"{pdf.page_count} PDF pages | {len(ppt.slides)} PPT slides | "
        f"pitch-deck {len(pitch_ppt.slides)} slides / {pitch_pdf.page_count} pages"
    )
    pdf.close()
    pitch_pdf.close()


def main() -> None:
    print(f"Building {PROGRAMME} application pack")
    print(f"Source check date: {datetime.now().strftime('%Y-%m-%d')}")
    OUT.mkdir(parents=True, exist_ok=True)
    try:
        build_docx()
        convert_docx_to_pdf()
        build_pptx()
        publish_pitch_deck()
    except PermissionError as exc:
        raise SystemExit(
            "Build failed because an official deliverable is locked. "
            "Close the Word, PDF, or PowerPoint files in documents/, then rebuild."
        ) from exc
    validate_outputs()
    print("Application pack complete.")
    print(f"Pitch deck folder: {PITCH_DIR}")
    print(f"New Drive folder:  {DRIVE_DIR}")
    print(f"Previous Drive folder retained: {DRIVE_PREVIOUS_DIR}")
    print(
        "Next: share the new Drive pitch PPT/PDF (Anyone with the link), "
        "paste the URL into the portal Pitch Deck field, and test signed out."
    )


if __name__ == "__main__":
    main()
