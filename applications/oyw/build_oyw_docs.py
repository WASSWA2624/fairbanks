#!/usr/bin/env python3
"""
Brandtech Group AI for Good Scholarship 2026 (One Young World) — nominee profile.

Generates synced Word, PDF, and PowerPoint for FairBanks / FCHIP nominee.
Run: python applications/oyw/build_oyw_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT = PROJECT / "documents"
SLUG = "oyw"

OUT_DOC = OUT / f"{SLUG}_word.docx"
OUT_PDF = OUT / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT / f"{SLUG}_ppt.pptx"

CALL_URL = (
    "https://opportunitiesforyouth.org/2026/07/15/"
    "ai-for-good-scholarship-2025-empowering-young-innovators-to-address-global-challenges/"
)

NAVY, TEAL, ACCENT = "0A1F2E", "0D6E6E", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "deep_tech": "deep_tech_collage.png",
    "mobile": "hero_chw_mobile.png",
    "outreach": "outreach_facilitator_canopy_01.jpg",
    "dashboard": "dashboard_demo.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "audience": "outreach_hands_raised_01.jpg",
}

PROGRAMME = "Brandtech Group AI for Good Scholarship 2026 — One Young World"
DOC_TITLE = "Nominee Profile — Responsible AI for Community Health"
SUBTITLE = "From Kampala outreach to Cape Town — AI that serves communities first"

NOMINEE_STORY = [
    "I grew up seeing neighbours wait until illness became an emergency. At FairBanks Community Reach "
    "I work where care actually starts — in streets, schools, and home visits with Community Health "
    "Workers — not only inside a clinic building. That experience shaped why I build: technology should "
    "help CHWs see risk early, not add another app nobody uses.",
    "FairBanks Community Health Intelligence Platform (FCHIP) is my answer — responsible AI on top of "
    "our community health cascade. We combine offline mobile capture, machine learning, and GIS mapping "
    "so fever clusters, high-risk pregnancies, and NCD hotspots surface before they become crises.",
]

RESPONSIBLE_AI = [
    "Human-in-the-loop: CHWs and nurses validate alerts — AI supports, never replaces, clinical judgment",
    "Offline-first design for low-bandwidth communities; no dependency on constant connectivity",
    "Consent and anonymisation aligned with Uganda Data Protection principles",
    "Local-language symptom capture where it improves trust and accuracy",
    "Transparent limits: we publish what models can and cannot do in community settings",
    "Bias awareness: training data includes women, children, elderly, and urban-poor outreach cohorts",
]

CASCADE = [
    ("Community members", "Identify needs; participate in solutions"),
    ("CHWs / VHTs", "Capture data, educate, refer, follow up"),
    ("FCHIP intelligence", "Predictive alerts, GIS hotspots, decision support"),
    ("FairBanks Medical Centre", "Clinical care, pharmacy, specialist referral"),
    ("Partners & research", "Evidence, scale, policy dialogue"),
]

PROOF = [
    "Live outreach in Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya",
    "Maternal, child, Gericare, and NCD screening programmes generating real field data",
    "Functioning medical centre anchoring referrals and quality assurance",
    "CHW/VHT relationships built on trust — essential for ethical AI adoption",
]

CAPE_TOWN = [
    "Share a practitioner story: how AI can strengthen — not bypass — community health workers",
    "Join health and technology sessions with examples from Ugandan primary care",
    "Connect with African innovators using AI for social impact — explore partnerships for FCHIP",
    "Bring Cape Town conversations back to FairBanks staff, CHWs, and youth volunteers",
    "Represent responsible AI in health at a summit where marketing hype often drowns out field reality",
]

ETHICS = [
    ["Privacy", "Household consent; minimum necessary data; secure storage"],
    ["Equity", "Design for users with low digital literacy and basic phones"],
    ["Accountability", "FairBanks clinical team reviews flagged cases"],
    ["Sustainability", "Open documentation so partners can audit and extend"],
]


def asset(name: str) -> Path:
    p = ASSETS / name
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def photo(key: str) -> Path:
    return asset(PHOTOS[key])


def embed(key_or_path, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key_or_path) if key_or_path in PHOTOS else Path(key_or_path)
    cache = REPO / "tmp" / f"{SLUG}_assets"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from PIL import Image as PILImage

    def rgb(c):
        return RGBColor.from_string(c)

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = rgb(color)

    def shade(cell, color):
        cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))

    def border(cell, color=LINE):
        cell._tc.get_or_add_tcPr().append(parse_xml(
            f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:left w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:bottom w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:right w:val="single" w:sz="8" w:color="{color}"/></w:tcBorders>'))

    def para(text, **kw):
        d = dict(size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False)
        d.update(kw)
        p = doc.add_paragraph()
        p.alignment = d["align"]
        p.paragraph_format.space_after = Pt(d["after"])
        font(p.add_run(text), size=d["size"], bold=d["bold"], color=d["color"], italic=d["italic"])

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14 if level == 1 else 8)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            p.clear()
            font(p.add_run(it))

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            shade(c, TEAL)
            border(c, TEAL)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=10)
                if ri % 2:
                    shade(c, CREAM)
                border(c)
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    def image(key, width_in=6.0, caption=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width_in, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(PROGRAMME, size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(DOC_TITLE, size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=4)
    para(SUBTITLE, size=12, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=6)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=12)
    image("cover", caption="AI for good means nothing without communities in the loop")
    doc.add_page_break()

    heading("My story")
    for p in NOMINEE_STORY:
        para(p)
    image("outreach", width_in=6.0, caption="Community outreach — where this work began")

    heading("The problem I see every week")
    para("Health facilities wait for patients. CHWs see warning signs days or weeks earlier — in fevers, "
         "missed ANC visits, blood pressure readings, and malnutrition. Without intelligence on that data, "
         "we lose the chance to prevent.")
    image("maternal", width_in=5.4, caption="Maternal outreach — early signals matter")

    heading("FCHIP — responsible AI in practice")
    bullets(RESPONSIBLE_AI)
    image("deep_tech", width_in=5.6, caption="Deep technology in service of community health, not hype")

    heading("Where AI sits in the FairBanks cascade")
    table(["Layer", "Role"], [[a, b] for a, b in CASCADE], widths=[2.2, 4.2])
    image("architecture", caption="Community signals flow up; care and action flow back")

    heading("Proof of work")
    bullets(PROOF)
    image("dashboard", width_in=5.4, caption="Dashboards turn field data into decisions")

    heading("Cape Town contribution plan")
    bullets(CAPE_TOWN)
    image("audience", width_in=5.6, caption="Communities participate — summit lessons must return to them")

    heading("Ethical commitments")
    table(["Principle", "Practice"], ETHICS, widths=[1.6, 4.8])

    heading("Mutual value")
    para("For One Young World and Brandtech: a nominee who builds AI in real African primary care — "
         "with ethics, offline constraints, and CHW partnership at the centre.")
    para("For FairBanks: global visibility, partnerships, and skills that strengthen FCHIP and "
         "Community Reach after Cape Town.")
    para("Scholarship details: see the official AI for Good programme page.", size=9, color=MUTED, italic=True)
    para(CALL_URL, size=8, color=TEAL)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak, KeepTogether
    from PIL import Image as PILImage

    teal, accent = HexColor("#" + TEAL), HexColor("#" + ACCENT)
    navy, slate, muted, cream, line = HexColor("#" + NAVY), HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    for name, kw in [
        ("CoverTitle", dict(fontName="Helvetica-Bold", fontSize=20, leading=24, textColor=navy, alignment=TA_CENTER, spaceAfter=6)),
        ("H1", dict(fontName="Helvetica-Bold", fontSize=14, leading=18, textColor=navy, spaceBefore=12, spaceAfter=6)),
        ("Body", dict(fontName="Helvetica", fontSize=10, leading=14, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=6)),
        ("Meta", dict(fontName="Helvetica", fontSize=9, leading=12, textColor=muted, alignment=TA_CENTER)),
        ("FBullet", dict(fontName="Helvetica", fontSize=10, leading=13, textColor=slate, leftIndent=12, spaceAfter=3)),
        ("CellHead", dict(fontName="Helvetica-Bold", fontSize=8.5, leading=11, textColor=white)),
        ("CellBody", dict(fontName="Helvetica", fontSize=8.5, leading=11, textColor=slate)),
    ]:
        st.add(ParagraphStyle(name, **kw))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None, max_h=2.7 * inch):
        path = embed(key)
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / iw
        h = min(w * aspect, max_h)
        w = h / aspect if h == max_h else w
        block = [Image(path, width=w, height=h)]
        if cap:
            block.append(Paragraph(cap, st["Meta"]))
        story.append(KeepTogether(block))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        data += [[Paragraph(str(c), st["CellBody"]) for c in row] for row in rows]
        t = Table(data, colWidths=widths or [pw / len(headers)] * len(headers), repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal), ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Paragraph(PROGRAMME, st["Meta"]))
    story.append(Paragraph(DOC_TITLE, st["CoverTitle"]))
    story.append(Paragraph(
        f'<b><i><font color="#{ACCENT}">FairBanks Community Health Intelligence Platform (FCHIP)</font></i></b>',
        st["Meta"]))
    story.append(Paragraph(f'<i><font color="#{MUTED}">{SUBTITLE}</font></i>', st["Meta"]))
    story.append(Paragraph(f'<b><i><font color="#{ACCENT}">{SLOGAN}</font></i></b>', st["Meta"]))
    story.append(Spacer(1, 10))
    story.append(img("cover", cap="AI for good means nothing without communities in the loop"))
    story.append(PageBreak())

    story.append(Paragraph("My story", st["H1"]))
    for p in NOMINEE_STORY:
        story.append(Paragraph(p, st["Body"]))

    story.append(Paragraph("FCHIP — responsible AI in practice", st["H1"]))
    for r in RESPONSIBLE_AI:
        story.append(Paragraph(f"• {r}", st["FBullet"]))

    story.append(Paragraph("Where AI sits in the FairBanks cascade", st["H1"]))
    tbl(["Layer", "Role"], CASCADE, [pw * 0.32, pw * 0.68])

    story.append(PageBreak())
    story.append(Paragraph("Proof of work", st["H1"]))
    for p in PROOF:
        story.append(Paragraph(f"• {p}", st["FBullet"]))

    story.append(Paragraph("Cape Town contribution plan", st["H1"]))
    for c in CAPE_TOWN:
        story.append(Paragraph(f"• {c}", st["FBullet"]))

    story.append(Paragraph("Ethical commitments", st["H1"]))
    tbl(["Principle", "Practice"], ETHICS, [pw * 0.24, pw * 0.76])

    story.append(Paragraph("Mutual value", st["H1"]))
    story.append(Paragraph(
        "A nominee building AI in real African primary care — with ethics and CHW partnership at the centre.",
        st["Body"]))
    story.append(Spacer(1, 10))
    story.append(Paragraph(f'<b><i><font color="#{ACCENT}">{SLOGAN}</font></i></b>', st["Meta"]))

    OUT.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                      topMargin=0.7 * inch, bottomMargin=0.7 * inch).build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH, blank = prs.slide_width, prs.slide_height, prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.28), Inches(11), Inches(0.25),
           f"FairBanks FCHIP  |  OYW Cape Town 2026  |  {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.9), SH - Inches(0.28), Inches(0.6), Inches(0.25), str(n), size=9, color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Cover
    s = prs.slides.add_slide(blank)
    pic(s, "mobile", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.3), SW, Inches(3.3), NAVY)
    tb(s, Inches(0.55), SH - Inches(3.05), Inches(12), Inches(0.35), PROGRAMME, size=12, bold=True, color=TEAL)
    tb(s, Inches(0.55), SH - Inches(2.5), Inches(12), Inches(0.6), DOC_TITLE, size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.55), SH - Inches(1.9), Inches(12), Inches(0.35),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.55), SH - Inches(1.45), Inches(12), Inches(0.35), SUBTITLE, size=12, italic=True, color="D0E8E8")
    tb(s, Inches(0.55), SH - Inches(0.95), Inches(12), Inches(0.35), SLOGAN, size=13, bold=True, color="FFFFFF")

    deck = [
        ("The journey", NOMINEE_STORY[0][:280] + "…", "outreach"),
        ("Problem on the ground", "CHWs see warning signs weeks before emergencies. FCHIP closes that gap.", "maternal"),
        ("Responsible AI", "\n".join(f"• {r[:70]}…" if len(r) > 70 else f"• {r}" for r in RESPONSIBLE_AI[:4]), "deep_tech"),
        ("Community cascade", "AI sits between CHW capture and clinic action — not above communities.", "architecture"),
        ("Proof", "\n".join(f"• {p}" for p in PROOF), "dashboard"),
        ("Cape Town", CAPE_TOWN[0], "audience"),
        ("Ethics", f"{ETHICS[0][0]}: {ETHICS[0][1]}\n{ETHICS[1][0]}: {ETHICS[1][1]}", None),
    ]

    for i, (title, body, img_key) in enumerate(deck, start=2):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, SW, Inches(1.05), CREAM)
        rect(s, Inches(0.5), Inches(0.18), Inches(0.08), Inches(0.7), ACCENT)
        tb(s, Inches(0.75), Inches(0.22), Inches(12), Inches(0.65), title, size=26, bold=True, color=NAVY)
        if img_key:
            pic(s, img_key, Inches(0.4), Inches(1.15), Inches(5.8), Inches(5.85))
            tb(s, Inches(6.4), Inches(1.2), Inches(6.6), Inches(5.7), body, size=19, color=SLATE)
        else:
            tb(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(5.6), body, size=18, color=SLATE)
        footer(s, i)

    # Close
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, SW, SH, NAVY)
    tb(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2),
       "Communities first.\nAI second.\nImpact always.",
       size=32, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    tb(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.5), SLOGAN, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
    footer(s, 9)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done.", OUT)
