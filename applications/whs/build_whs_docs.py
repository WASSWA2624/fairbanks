#!/usr/bin/env python3
"""
World Health Summit 2026 Youth Group Opportunity — delegation application pack.

Generates synced Word, PDF, and PowerPoint for FairBanks youth delegation.
Run: python applications/whs/build_whs_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT = PROJECT / "documents"
SLUG = "whs"

OUT_DOC = OUT / f"{SLUG}_word.docx"
OUT_PDF = OUT / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT / f"{SLUG}_ppt.pptx"

CALL_URL = (
    "https://opportunitiesforyouth.org/2026/06/26/"
    "apply-now-world-health-summit-2025-stipend-program-for-global-health-changemakers/"
)

NAVY, TEAL, ACCENT = "0A1F2E", "0D6E6E", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "outreach": "outreach_audience_full_group_01.jpg",
    "training": "indoor_training_audience_01.jpg",
    "bp": "outreach_bp_screening.jpeg",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "team": "outreach_cofacilitators_01.jpg",
    "community": "outreach_community_speaker_01.jpg",
}

PROGRAMME = "World Health Summit 2026 — Youth Group Opportunity (Berlin)"
DOC_TITLE = "FairBanks Youth Delegation Application"
SUBTITLE = "Learning global prevention practice to strengthen community health at home"

WHO_WE_ARE = [
    "FairBanks Community Reach is a youth-led community health initiative rooted in Kampala, "
    "Uganda. We work alongside Community Health Workers and Village Health Teams in "
    "neighbourhoods including Bukoto, Kyebando, Kisaasi, Kamwokya, and Kikaaya — running "
    "outreach, screening, maternal support, and health education where clinics alone cannot reach.",
    "We are building FairBanks Community Health Intelligence Platform (FCHIP) — the intelligence "
    "layer on our community cascade — so field data becomes early warnings instead of paper "
    "that arrives too late. Our delegation would bring that lived, last-mile perspective to Berlin.",
]

PURPOSE = [
    "Represent Ugandan community health practice at a global table — not as observers, but as "
    "practitioners who screen, refer, and follow up every week.",
    "Learn how leading health systems, researchers, and youth networks approach prevention, "
    "digital health governance, and pandemic preparedness — then translate lessons for our CHWs.",
    "Connect FairBanks outreach and FCHIP pilot work to partners who care about primary care "
    "in underserved regions.",
]

LEARNING = [
    ["Prevention & primary care", "Shift from reactive treatment to community-level early action"],
    ["Digital health & ethics", "Responsible data use in low-resource, offline-first settings"],
    ["Youth leadership in health", "How young changemakers influence policy and programme design"],
    ["Maternal & child health", "Best practice we can adapt for home visits and ANC follow-up"],
    ["NCDs & ageing", "Screening models relevant to our chronic-care and Gericare programmes"],
    ["Global health partnerships", "NGO, academic, and government collaboration models"],
]

BRING = [
    "Stories from real outreach — what works and what fails in last-mile Uganda",
    "A prevention-first narrative grounded in CHW/VHT practice, not hospital-centric theory",
    "Experience bridging community dialogues, clinic care, and emerging digital tools",
    "Commitment to share the summit with peers who cannot travel — via briefing and open notes",
]

POST_SUMMIT = [
    "Within 30 days of return: host a community briefing for CHWs, VHTs, and FairBanks staff",
    "Publish a simple learning note (English + local-language summary where possible) for partners",
    "Integrate at least three summit insights into FairBanks outreach plans for Q4 2026",
    "Present findings to FairBanks Medical Centre leadership and interested district contacts",
    "Maintain connections made in Berlin through FCHIP and Community Reach partnership channels",
]

TRAVEL_NOTE = (
    "We understand that complimentary WHS tickets cover summit registration only. FairBanks "
    "will self-fund travel, visa, and accommodation for nominated delegates aged 28 or under, "
    "as required by the programme. This commitment is documented here — not on the cover page."
)


def asset(name: str) -> Path:
    p = ASSETS / name
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def photo(key: str) -> Path:
    return asset(PHOTOS[key])


def embed(key_or_path, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key_or_path) if key_or_path in PHOTOS else Path(key_or_path)
    cache = REPO / "tmp" / f"{SLUG}_assets"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from PIL import Image as PILImage

    def rgb(c):
        return RGBColor.from_string(c)

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = rgb(color)

    def shade(cell, color):
        cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))

    def border(cell, color=LINE):
        cell._tc.get_or_add_tcPr().append(parse_xml(
            f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:left w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:bottom w:val="single" w:sz="8" w:color="{color}"/>'
            f'<w:right w:val="single" w:sz="8" w:color="{color}"/></w:tcBorders>'))

    def para(text, **kw):
        d = dict(size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False)
        d.update(kw)
        p = doc.add_paragraph()
        p.alignment = d["align"]
        p.paragraph_format.space_after = Pt(d["after"])
        font(p.add_run(text), size=d["size"], bold=d["bold"], color=d["color"], italic=d["italic"])

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14 if level == 1 else 8)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            p.clear()
            font(p.add_run(it))

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            shade(c, TEAL)
            border(c, TEAL)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=10)
                if ri % 2:
                    shade(c, CREAM)
                border(c)
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    def image(key, width_in=6.0, caption=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width_in, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(PROGRAMME, size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(DOC_TITLE, size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=4)
    para(SUBTITLE, size=12, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=6)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=12)
    image("cover", caption="Youth-led community health — from Kampala neighbourhoods to global dialogue")
    doc.add_page_break()

    heading("Our organisation")
    for p in WHO_WE_ARE:
        para(p)
    image("outreach", width_in=6.0, caption="Community outreach — where our delegation's stories begin")

    heading("Why we seek WHS 2026 tickets")
    bullets(PURPOSE)
    image("team", width_in=5.6, caption="Young facilitators and CHW partners in the field")

    heading("Learning agenda in Berlin")
    table(["Theme", "What we will absorb and bring home"], LEARNING, widths=[2.2, 4.2])
    image("training", width_in=5.4, caption="We train and learn continuously — WHS extends that habit globally")

    heading("What we contribute to the summit")
    bullets(BRING)
    image("bp", width_in=5.4, caption="Prevention in practice — screening at community level")

    heading("Post-summit briefing commitment")
    para("FairBanks treats summit access as a responsibility to our community, not a private benefit.")
    bullets(POST_SUMMIT)
    image("community", width_in=5.6, caption="Knowledge returns to community dialogues and CHW networks")

    heading("Travel & logistics")
    para(TRAVEL_NOTE, italic=True)
    para("Programme details and deadlines: see the official WHS youth group opportunity page.",
         size=9, color=MUTED, italic=True)
    para(CALL_URL, size=8, color=TEAL)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak, KeepTogether
    from PIL import Image as PILImage

    teal, accent = HexColor("#" + TEAL), HexColor("#" + ACCENT)
    navy, slate, muted, cream, line = HexColor("#" + NAVY), HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    for name, kw in [
        ("CoverTitle", dict(fontName="Helvetica-Bold", fontSize=20, leading=24, textColor=navy, alignment=TA_CENTER, spaceAfter=6)),
        ("H1", dict(fontName="Helvetica-Bold", fontSize=14, leading=18, textColor=navy, spaceBefore=12, spaceAfter=6)),
        ("Body", dict(fontName="Helvetica", fontSize=10, leading=14, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=6)),
        ("Meta", dict(fontName="Helvetica", fontSize=9, leading=12, textColor=muted, alignment=TA_CENTER)),
        ("FBullet", dict(fontName="Helvetica", fontSize=10, leading=13, textColor=slate, leftIndent=12, spaceAfter=3)),
        ("CellHead", dict(fontName="Helvetica-Bold", fontSize=8.5, leading=11, textColor=white)),
        ("CellBody", dict(fontName="Helvetica", fontSize=8.5, leading=11, textColor=slate)),
    ]:
        st.add(ParagraphStyle(name, **kw))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None, max_h=2.7 * inch):
        path = embed(key)
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / iw
        h = min(w * aspect, max_h)
        w = h / aspect if h == max_h else w
        block = [Image(path, width=w, height=h)]
        if cap:
            block.append(Paragraph(cap, st["Meta"]))
        story.append(KeepTogether(block))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        data += [[Paragraph(str(c), st["CellBody"]) for c in row] for row in rows]
        t = Table(data, colWidths=widths or [pw / len(headers)] * len(headers), repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal), ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Paragraph(PROGRAMME, st["Meta"]))
    story.append(Paragraph(DOC_TITLE, st["CoverTitle"]))
    story.append(Paragraph(
        f'<b><i><font color="#{ACCENT}">FairBanks Community Health Intelligence Platform (FCHIP)</font></i></b>',
        st["Meta"]))
    story.append(Paragraph(f'<i><font color="#{MUTED}">{SUBTITLE}</font></i>', st["Meta"]))
    story.append(Paragraph(f'<b><i><font color="#{ACCENT}">{SLOGAN}</font></i></b>', st["Meta"]))
    story.append(Spacer(1, 10))
    story.append(img("cover", cap="Youth-led community health — Kampala to global dialogue"))
    story.append(PageBreak())

    story.append(Paragraph("Our organisation", st["H1"]))
    for p in WHO_WE_ARE:
        story.append(Paragraph(p, st["Body"]))

    story.append(Paragraph("Why we seek WHS 2026 tickets", st["H1"]))
    for p in PURPOSE:
        story.append(Paragraph(f"• {p}", st["FBullet"]))

    story.append(Paragraph("Learning agenda in Berlin", st["H1"]))
    tbl(["Theme", "Take-home"], LEARNING, [pw * 0.32, pw * 0.68])

    story.append(Paragraph("What we contribute to the summit", st["H1"]))
    for b in BRING:
        story.append(Paragraph(f"• {b}", st["FBullet"]))

    story.append(PageBreak())
    story.append(Paragraph("Post-summit briefing commitment", st["H1"]))
    for ps in POST_SUMMIT:
        story.append(Paragraph(f"• {ps}", st["FBullet"]))

    story.append(Paragraph("Travel & logistics", st["H1"]))
    story.append(Paragraph(TRAVEL_NOTE, st["Body"]))
    story.append(Spacer(1, 10))
    story.append(Paragraph(f'<b><i><font color="#{ACCENT}">{SLOGAN}</font></i></b>', st["Meta"]))

    OUT.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                      topMargin=0.7 * inch, bottomMargin=0.7 * inch).build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH, blank = prs.slide_width, prs.slide_height, prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        s.line.fill.background()

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.28), Inches(11), Inches(0.25),
           f"FairBanks Youth Delegation  |  WHS 2026 Berlin  |  {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.9), SH - Inches(0.28), Inches(0.6), Inches(0.25), str(n), size=9, color="FFFFFF", align=PP_ALIGN.RIGHT)

    # 1 Cover
    s = prs.slides.add_slide(blank)
    pic(s, "cover", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.0), SW, Inches(3.0), NAVY)
    tb(s, Inches(0.55), SH - Inches(2.75), Inches(12), Inches(0.35), PROGRAMME, size=12, bold=True, color=TEAL)
    tb(s, Inches(0.55), SH - Inches(2.35), Inches(12), Inches(0.55), DOC_TITLE, size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.55), SH - Inches(1.8), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.55), SH - Inches(1.4), Inches(12), Inches(0.35), SUBTITLE, size=12, italic=True, color="D0E8E8")
    tb(s, Inches(0.55), SH - Inches(0.85), Inches(12), Inches(0.35), SLOGAN, size=13, bold=True, color="FFFFFF")

    slides = [
        ("Youth-led health in Uganda", WHO_WE_ARE[0], "outreach"),
        ("Why Berlin", PURPOSE[1], "maternal"),
        ("Learning agenda", "\n".join(f"• {t}" for t, _ in LEARNING[:4]), "training"),
        ("We bring lived practice", BRING[0], "mobile"),
        ("After the summit", POST_SUMMIT[0] + "\n\n" + POST_SUMMIT[2], "community"),
        ("Travel note", TRAVEL_NOTE, "team"),
    ]

    for i, (title, body, img) in enumerate(slides, start=2):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, SW, Inches(0.95), CREAM)
        rect(s, 0, 0, Inches(0.12), Inches(0.95), ACCENT)
        tb(s, Inches(0.45), Inches(0.2), Inches(12), Inches(0.6), title, size=24, bold=True, color=NAVY)
        pic(s, img, Inches(0.4), Inches(1.05), Inches(5.6), Inches(5.9))
        tb(s, Inches(6.2), Inches(1.1), Inches(6.7), Inches(5.8), body, size=19, color=SLATE)
        footer(s, i)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done.", OUT)
