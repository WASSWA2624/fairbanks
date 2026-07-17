#!/usr/bin/env python3
"""
Young Feminist AI School 2026 (UN Women) — FairBanks learning plan & project set.

Generates:
  applications/feminist-ai/documents/feminist-ai_word.docx
  applications/feminist-ai/documents/feminist-ai_pdf.pdf
  applications/feminist-ai/documents/feminist-ai_ppt.pptx

Run: python applications/feminist-ai/build_feminist-ai_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "feminist-ai"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/07/16/"
    "ai-for-gender-equality-un-women-ai-school-opens-for-changemakers/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

META = {
    "programme": "Young Feminist AI School 2026 — UN Women",
    "doc_title": "Learning Plan & Gender-Responsive AI Project",
    "subtitle": "Maternal and adolescent health intelligence — built with communities, not on them",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "mothers": "waiting_room_mothers_01.jpeg",
    "architecture": "data_flow_iso_labeled.png",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "training": "indoor_training_facilitator_01.jpg",
    "deep_tech": "deep_tech_collage.png",
    "outreach": "outreach_audience_discussion_01.jpg",
    "conclusion": "outreach_hands_raised_01.jpg",
}


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


LEARNING_MODULES = [
    ("Foundations of feminist AI", "Bias in health datasets; who is missing from training data; participatory design with women and girls."),
    ("Responsible ML for community health", "Explainable alerts, consent, and offline-first capture for CHWs in low-bandwidth settings."),
    ("Gender-responsive product design", "Interfaces and workflows that work for mothers, adolescents, and CHWs — not only engineers."),
    ("Policy & advocacy", "Translating model outputs into action without stigmatising vulnerable groups."),
]

PROJECT_PHASES = [
    ("Discover", "Co-design with adolescent girls and mothers in FairBanks outreach — what questions are safe and useful?"),
    ("Prototype", "Lightweight risk flags for missed ANC and adolescent SRH education gaps using anonymised community signals."),
    ("Validate", "Compare alerts against clinical outcomes and CHW feedback; document bias checks and corrections."),
    ("Share back", "Open learning brief for UN Women cohort — methods, failures, and ethical guardrails."),
]

SCHOOL_GAINS = [
    "A real HealthTech field site with maternal and adolescent programmes",
    "Documented gender-responsive AI design process from Uganda",
    "Case material linking feminist AI principles to community health",
]

FAIRBANKS_GAINS = [
    "Structured feminist AI curriculum and peer changemakers",
    "Mentorship on bias, safety, and inclusive model design",
    "Continental network to stress-test FCHIP ethics and narrative",
]


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption="Community health ecosystem — intelligence with dignity")
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    heading("1. Why this school fits my work")
    para("AI in health fails women and girls when models are trained on incomplete data and deployed without "
         "community voice. FairBanks runs maternal and adolescent programmes through outreach and CHWs — "
         "exactly where gender-responsive AI must be tested honestly.")
    para("The Young Feminist AI School offers structured learning I cannot replicate alone: feminist ethics, "
         "inclusive design, and a cohort of changemakers holding each other accountable.")

    heading("2. My learning plan")
    image("training", caption="Learning that connects classroom principles to field practice")
    for title, body in LEARNING_MODULES:
        para(title + ": " + body, bold=True)

    heading("3. Gender-responsive AI project — maternal & adolescent health")
    image("maternal", caption="Maternal health outreach — co-design starts here")
    para("Project title: Safe Signals — early support flags for mothers and adolescents without blame or exposure.")
    para("Goal: Use FCHIP’s community data pipeline to prototype explainable alerts for (a) missed antenatal care "
         "and anaemia risk patterns, and (b) adolescent SRH education gaps — always with consent, anonymisation, "
         "and CHW-led follow-up.")
    for phase, detail in PROJECT_PHASES:
        para(f"{phase}: {detail}")

    heading("4. FCHIP as the field laboratory")
    image("architecture", caption="Community signals → FCHIP → gender-aware decision support")
    bullets([
        "Live medical centre and outreach generating structured maternal and adolescent touchpoints",
        "CHW/VHT networks for ethical follow-up — technology never replaces human care",
        "Dashboards that can be audited for bias across age, gender, and neighbourhood",
    ])

    heading("5. Ethics, bias, and inclusion commitments")
    image("deep_tech", width=5.8, caption="Deep tech with feminist guardrails")
    bullets([
        "No individual public labelling of adolescents; aggregate and CHW-private alerts only",
        "Community advisory input before any model influences outreach messaging",
        "Document false positives/negatives transparently for the UN Women cohort",
        "Align with Uganda Data Protection principles and FairBanks consent practice",
    ])

    heading("6. Learning together — what each side gains")
    para("The school strengthens changemakers; FairBanks offers a honest field lab. That exchange is the point.")
    para("UN Women cohort and faculty gain:", bold=True)
    bullets(SCHOOL_GAINS)
    para("FairBanks gains:", bold=True)
    bullets(FAIRBANKS_GAINS)

    image("conclusion", caption="Participatory health AI — led with communities")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted = HexColor("#" + SLATE), HexColor("#" + MUTED)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.88, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = min(w * ih / iw, 2.6 * inch)
        w = h * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(
        f'<font color="#{ACCENT}"><b><i>FairBanks Community Health Intelligence Platform (FCHIP)</i></b></font>',
        st["Meta"]))
    story.append(Paragraph(f'<font color="#{MUTED}"><i>{META["subtitle"]}</i></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    img("cover")
    story.append(PageBreak())

    blocks = [
        ("1. Why this school fits my work", [
            "AI in health fails women and girls when models skip community voice. FairBanks maternal and adolescent "
            "programmes are the right place to test gender-responsive AI honestly.",
            "The Young Feminist AI School offers feminist ethics, inclusive design, and peer accountability.",
        ], None),
        ("2. My learning plan", [f"<b>{t}:</b> {b}" for t, b in LEARNING_MODULES], "training"),
        ("3. Gender-responsive AI project", [
            "<b>Safe Signals</b> — explainable alerts for missed ANC/anaemia patterns and adolescent SRH education gaps.",
            "Co-design with mothers and adolescents; CHW-led follow-up; no public individual labelling.",
        ] + [f"<b>{p}:</b> {d}" for p, d in PROJECT_PHASES], "maternal"),
        ("4. FCHIP as the field laboratory", [
            "Live outreach and CHW networks · Auditable dashboards · Human care always central",
        ], "architecture"),
        ("5. Ethics commitments", [
            "Community advisory before outreach messaging changes",
            "Transparent documentation of model limits for the cohort",
            "Uganda Data Protection alignment",
        ], None),
        ("6. Learning together", [
            "<b>UN Women cohort gains:</b> " + "; ".join(SCHOOL_GAINS),
            "<b>FairBanks gains:</b> " + "; ".join(FAIRBANKS_GAINS),
        ], None),
    ]

    for title, paras, img_key in blocks:
        story.append(Paragraph(title, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
        if img_key:
            img(img_key)
        for t in paras:
            story.append(Paragraph(t, st["Body"]))
        story.append(Spacer(1, 4))

    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                      topMargin=0.7 * inch, bottomMargin=0.7 * inch).build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    def _add_entrance_anims(slide):
        """Light fade-in on pictures and key text (documents.mdc motion)."""
        from lxml import etree
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spids = []
        for shape in slide.shapes:
            try:
                st = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                st = -1
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
            if st == 13 or has_text:
                spids.append(str(shape.shape_id))
            if len(spids) >= 3:
                break
        if not spids:
            return
        sld = slide._element
        for old in sld.findall(f"{{{NS_P}}}timing"):
            sld.remove(old)
        children = []
        nid = 3
        for i, spid in enumerate(spids):
            delay = 0 if i == 0 else 200
            children.append(
                f'<p:par xmlns:p="{NS_P}">'
                f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
                f'fill="hold" nodeType="withEffect">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
            )
            nid += 2
        xml = (
            f'<p:timing xmlns:p="{NS_P}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        )
        sld.append(etree.fromstring(xml))

    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic_cover(sl, key):
        sl.shapes.add_picture(embed(key), 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks FCHIP | Young Feminist AI School | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.2), SW, Inches(3.2), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.9), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.4), Inches(12), Inches(0.5), META["doc_title"], size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.9), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(1.45), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")

    # Learning objectives — card grid
    s = prs.slides.add_slide(blank)
    band(s, "Curriculum map", "Four learning blocks I will complete")
    for i, (t, b) in enumerate(LEARNING_MODULES):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = Inches(1.25) + row * Inches(2.85)
        rect(s, x, y, Inches(6.1), Inches(2.6), "FFFFFF", LINE)
        rect(s, x, y, Inches(6.1), Inches(0.12), TEAL)
        tb(s, x + Inches(0.2), y + Inches(0.25), Inches(5.7), Inches(0.45), t, size=15, bold=True, color=TEAL)
        tb(s, x + Inches(0.2), y + Inches(0.75), Inches(5.7), Inches(1.6), b, size=13, color=MUTED)
    footer(s, 2)

    # Gender gap
    s = prs.slides.add_slide(blank)
    band(s, "Problem", "Health AI often ignores women and girls")
    pic_fit(s, "mothers", Inches(0.45), Inches(1.15), Inches(5.8), Inches(5.6))
    tb(s, Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5),
       "Missing data → biased models → harmful or useless alerts\n\n"
       "Maternal & adolescent programmes need participatory design\n\n"
       "FairBanks outreach is where ethics meets engineering",
       size=17, color=SLATE)
    footer(s, 3)

    # Project
    s = prs.slides.add_slide(blank)
    band(s, "Project", "Safe Signals — maternal & adolescent support flags")
    for i, (p, d) in enumerate(PROJECT_PHASES):
        y = Inches(1.25) + i * Inches(1.15)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.0), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.12), Inches(2.0), Inches(0.75), p, size=14, bold=True, color=ACCENT)
        tb(s, Inches(2.8), y + Inches(0.12), Inches(9.5), Inches(0.75), d, size=18, color=SLATE)
    footer(s, 4)

    # Field lab
    s = prs.slides.add_slide(blank)
    band(s, "Field lab", "FCHIP — community data with feminist guardrails")
    pic_fit(s, "architecture", Inches(0.4), Inches(1.1), Inches(7.0), Inches(5.7))
    tb(s, Inches(7.6), Inches(1.2), Inches(5.3), Inches(5.5),
       "CHW-led follow-up\nConsent & anonymisation\nAuditable dashboards\nNo stigma by design",
       size=17, color=MUTED)
    footer(s, 5)

    # Exchange
    s = prs.slides.add_slide(blank)
    band(s, "Exchange", "School ↔ FairBanks — mutual learning")
    rect(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.0), CREAM, LINE)
    tb(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "Cohort & faculty receive", size=16, bold=True, color=TEAL)
    tb(s, Inches(0.7), Inches(2.0), Inches(5.5), Inches(3.8), "\n".join(f"• {g}" for g in SCHOOL_GAINS), size=18, color=SLATE)
    rect(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.0), CREAM, LINE)
    tb(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4), "I receive", size=16, bold=True, color=TEAL)
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.8), "\n".join(f"• {g}" for g in FAIRBANKS_GAINS), size=18, color=SLATE)
    footer(s, 6)

    # Timeline
    s = prs.slides.add_slide(blank)
    band(s, "Timeline", "From classroom to community pilot")
    pic_fit(s, "mobile", Inches(0.45), Inches(1.15), Inches(5.5), Inches(5.6))
    tb(s, Inches(6.2), Inches(1.2), Inches(6.5), Inches(5.5),
       "Weeks 1–4: Feminist AI foundations + dataset audit\n"
       "Weeks 5–8: Prototype Safe Signals with CHWs\n"
       "Weeks 9–12: Validate, document bias checks, share open brief",
       size=16, color=SLATE)
    footer(s, 7)

    # Close
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.5), SW, Inches(2.5), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.1), Inches(12), Inches(0.55),
       "Gender-responsive AI — learned with UN Women, tested with communities.", size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.2), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for _sl in prs.slides:
        _add_entrance_anims(_sl)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. feminist-ai document set in", OUT_DIR)
