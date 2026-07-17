#!/usr/bin/env python3
"""
GirlCode Hackathon 2026 — Kampala challenge brief & mentorship offer.

Generates:
  applications/girlcode/documents/girlcode_word.docx
  applications/girlcode/documents/girlcode_pdf.pdf
  applications/girlcode/documents/girlcode_ppt.pptx

Run: python applications/girlcode/build_girlcode_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "girlcode"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/07/15/"
    "building-africas-next-generation-of-ai-powered-female-innovators-apply-for-the-"
    "2026-girlcode-hackathon-across-six-african-countries/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

META = {
    "programme": "GirlCode Hackathon 2026 — Kampala",
    "doc_title": "Challenge Brief & Mentorship Offer",
    "subtitle": "30-hour build sprint · UICT Nakawa · 5–6 September 2026",
    "event": "UICT Nakawa Campus, New Portbell Road, Kampala",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "dashboard": "dashboard_demo.png",
    "gis": "gis_hotspots.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "outreach": "outreach_mobile_phone_demo_01.jpg",
    "training": "indoor_training_audience_01.jpg",
    "conclusion": "outreach_facilitator_group_01.jpg",
}

CHALLENGES = [
    {
        "id": "A",
        "title": "Maternal risk early warning",
        "problem": "High-risk pregnancies are often found late. Can your team turn outreach-style signals "
                   "(missed ANC, BP patterns, anaemia proxies) into timely CHW alerts?",
        "stretch": "Explainable scores · offline-first · no stigma",
    },
    {
        "id": "B",
        "title": "Adolescent SRH companion",
        "problem": "Adolescents need private, trusted SRH information and referral paths. Build a lightweight "
                   "companion that works with school and community outreach — not against parents and CHWs.",
        "stretch": "Local languages · consent · safe escalation to FairBanks outreach",
    },
    {
        "id": "C",
        "title": "Outbreak early signal",
        "problem": "When fever reports rise across neighbouring villages, clinics react too late. Can you "
                   "cluster community reports and flag a possible surge before the pharmacy runs dry?",
        "stretch": "GIS view · simple rules + optional ML · stock-out prevention",
    },
]

MENTORSHIP = [
    "FairBanks clinicians and outreach leads on realistic workflows",
    "Sample (anonymised) data schemas from CHW capture — not raw patient records",
    "Office hours during the hackathon on UX for low-literacy and offline settings",
    "Path to continue building with FairBanks if your MVP wins locally or regionally",
]

TEAM_WINS = [
    "Real community health problem statements grounded in Kampala outreach",
    "Mentorship from an operating medical centre — not a fictional case study",
    "Portfolio project with social impact narrative for jobs and investors",
    "Visibility with FairBanks partners in maternal, adolescent, and NCD work",
]


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption=META["subtitle"])
    para(META["event"], size=10, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    heading("1. To Kampala builders — why FairBanks is sponsoring problem statements")
    para("GirlCode’s 30-hour hackathon is where Africa’s next women builders meet real problems. FairBanks is a "
         "community health ecosystem — medical centre, outreach, CHWs, and FCHIP, our intelligence layer. "
         "We are offering three challenge tracks and hands-on mentorship because the best MVPs come from "
         "teams who understand users on the ground.")

    heading("2. The platform you will build on — FCHIP in one page")
    image("architecture", caption="Community → CHW capture → FCHIP → action")
    para("FCHIP connects last-mile health data to predictions: outbreaks, maternal risk, chronic hotspots, "
         "child health gaps, and medicine demand. Your hackathon prototype can plug into this vision — "
         "especially offline capture, alert UX, and maps.")

    for ch in CHALLENGES:
        heading(f"3.{ch['id']} Challenge {ch['id']}: {ch['title']}", 2)
        para("Problem: " + ch["problem"])
        para("Stretch goals: " + ch["stretch"], italic=True, color=MUTED)

    heading("4. Data, APIs, and mentorship from FairBanks")
    image("training", caption="Mentorship during the hackathon — clinicians + outreach leads")
    bullets(MENTORSHIP)

    heading("5. What strong teams walk away with")
    bullets(TEAM_WINS)
    para("GirlCode gains fresh AI-powered prototypes on maternal and adolescent health — FairBanks gains "
         "talented women builders stress-testing FCHIP ideas in a 30-hour sprint. That is how ecosystems grow.")

    image("conclusion", caption="Build with us in Kampala — Your health, our mission.")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted = HexColor("#" + SLATE), HexColor("#" + MUTED)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("H2", fontName="Helvetica-Bold", fontSize=12, textColor=teal, spaceBefore=8, spaceAfter=4))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.88, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = min(w * ih / iw, 2.5 * inch)
        w = h * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>FCHIP — {META["subtitle"]}</i></b></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    img("cover")
    story.append(Paragraph(META["event"], st["Meta"]))
    story.append(PageBreak())

    story.append(Paragraph("1. To Kampala builders", st["H1"]))
    story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
    story.append(Paragraph(
        "FairBanks offers three FCHIP challenge tracks and mentorship because the best MVPs come from "
        "teams who understand users on the ground.", st["Body"]))
    img("architecture", cap="FCHIP — community intelligence stack")

    for ch in CHALLENGES:
        story.append(Paragraph(f"Challenge {ch['id']}: {ch['title']}", st["H2"]))
        story.append(Paragraph(f"<b>Problem:</b> {ch['problem']}", st["Body"]))
        story.append(Paragraph(f"<i>Stretch: {ch['stretch']}</i>", st["Body"]))

    story.append(Paragraph("Mentorship offer", st["H1"]))
    for m in MENTORSHIP:
        story.append(Paragraph("• " + m, st["FBullet"]))
    story.append(Paragraph("What teams gain", st["H1"]))
    for t in TEAM_WINS:
        story.append(Paragraph("• " + t, st["FBullet"]))
    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                      topMargin=0.7 * inch, bottomMargin=0.7 * inch).build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    def _add_entrance_anims(slide):
        """Light fade-in on pictures and key text (documents.mdc motion)."""
        from lxml import etree
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spids = []
        for shape in slide.shapes:
            try:
                st = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                st = -1
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
            if st == 13 or has_text:
                spids.append(str(shape.shape_id))
            if len(spids) >= 3:
                break
        if not spids:
            return
        sld = slide._element
        for old in sld.findall(f"{{{NS_P}}}timing"):
            sld.remove(old)
        children = []
        nid = 3
        for i, spid in enumerate(spids):
            delay = 0 if i == 0 else 200
            children.append(
                f'<p:par xmlns:p="{NS_P}">'
                f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
                f'fill="hold" nodeType="withEffect">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
            )
            nid += 2
        xml = (
            f'<p:timing xmlns:p="{NS_P}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        )
        sld.append(etree.fromstring(xml))

    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic_cover(sl, key):
        sl.shapes.add_picture(embed(key), 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks FCHIP | GirlCode Kampala 2026 | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.4), SW, Inches(3.4), NAVY)
    tb(s, Inches(0.6), SH - Inches(3.1), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.65), Inches(12), Inches(0.5), META["doc_title"], size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(2.1), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(1.7), Inches(12), Inches(0.3), META["subtitle"], size=12, color="D0E8E8")
    tb(s, Inches(0.6), SH - Inches(1.15), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")

    # Invite
    s = prs.slides.add_slide(blank)
    band(s, "Hackathon", "Build on real community health problems")
    tb(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
       "FairBanks = medical centre + outreach + CHWs + FCHIP intelligence\n\n"
       "30 hours · women builders · Kampala · UICT Nakawa\n\n"
       "We provide problem statements, sample schemas, and mentor hours.",
       size=18, color=SLATE)
    footer(s, 2)

    # FCHIP
    s = prs.slides.add_slide(blank)
    band(s, "Platform", "FCHIP — what you are extending")
    pic_fit(s, "architecture", Inches(0.4), Inches(1.1), Inches(7.2), Inches(5.7))
    tb(s, Inches(7.9), Inches(1.2), Inches(5.0), Inches(5.5),
       "Offline capture → predictions → alerts\n\nOutbreak · maternal · NCD · child health · stock-outs",
       size=20, color=MUTED)
    footer(s, 3)

    # Three challenges — one slide each
    imgs = ["maternal", "mobile", "gis"]
    for i, (ch, img_key) in enumerate(zip(CHALLENGES, imgs)):
        s = prs.slides.add_slide(blank)
        band(s, f"Track {ch['id']}", ch["title"])
        pic_fit(s, img_key, Inches(0.45), Inches(1.15), Inches(5.5), Inches(5.6))
        tb(s, Inches(6.2), Inches(1.2), Inches(6.5), Inches(3.5), ch["problem"], size=19, color=SLATE)
        tb(s, Inches(6.2), Inches(4.8), Inches(6.5), Inches(1.5), "Stretch: " + ch["stretch"],
           size=13, italic=True, color=ACCENT)
        footer(s, 4 + i)

    # Mentorship
    s = prs.slides.add_slide(blank)
    band(s, "Mentorship", "FairBanks office hours during the sprint")
    for i, m in enumerate(MENTORSHIP):
        y = Inches(1.25) + i * Inches(1.15)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.0), "FFFFFF", LINE)
        tb(s, Inches(0.75), y + Inches(0.2), Inches(11.5), Inches(0.6), m, size=19, color=SLATE)
    footer(s, 7)

    # Team wins
    s = prs.slides.add_slide(blank)
    band(s, "Outcomes", "Why teams should pick a FairBanks track")
    pic_fit(s, "dashboard", Inches(0.45), Inches(1.15), Inches(5.5), Inches(5.6))
    tb(s, Inches(6.2), Inches(1.2), Inches(6.5), Inches(5.5),
       "\n".join(f"• {t}" for t in TEAM_WINS), size=20, color=SLATE)
    footer(s, 8)

    # Close
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.5), SW, Inches(2.5), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.1), Inches(12), Inches(0.55),
       "See you at UICT Nakawa — 5–6 September 2026.", size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.2), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for _sl in prs.slides:
        _add_entrance_anims(_sl)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. girlcode document set in", OUT_DIR)
