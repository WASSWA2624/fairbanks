#!/usr/bin/env python3
"""
World Bank GovTech Innovation Challenge 2026 — FairBanks / FCHIP readiness brief.

Generates:
  applications/govtech/documents/govtech_word.docx
  applications/govtech/documents/govtech_pdf.pdf
  applications/govtech/documents/govtech_ppt.pptx

Run: python applications/govtech/build_govtech_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "govtech"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/07/10/"
    "world-bank-govtech-innovation-challenge-2026-global-call-for-proposals-to-build-"
    "digital-solutions-for-governments/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."
TAGLINE = "Health for All — Obulamu eri Bonna · Afya kwa Wote · Oburamu bwa Boona"

META = {
    "programme": "World Bank GovTech Innovation Challenge 2026",
    "doc_title": "GovTech Readiness Brief — District Health Intelligence",
    "subtitle": "Monitoring for matching challenges · ready PoC partner when health calls open",
    "monitor_note": "Challenge deadlines vary — we submit only for health / Africa-aligned use cases",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "dashboard": "dashboard_demo.png",
    "gis": "gis_hotspots.png",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "pharmacy": "pharmacy_staff_laptop_01.jpg",
    "outreach": "outreach_facilitator_canopy_01.jpg",
    "conclusion": "outreach_audience_full_group_01.jpg",
}

SECTIONS = [
    ("1. Readiness brief summary", [
        "The World Bank GovTech Innovation Challenge connects technology companies with government "
        "digitalisation needs — offering capacity building, a Switzerland bootcamp, and proof-of-concept "
        "work with government partners.",
        "FairBanks maintains this readiness brief because current open challenges may focus elsewhere "
        "(for example public-finance audit use cases in Morocco). We do not pretend those are our target. "
        "When a health or Africa-aligned challenge appears, we can move fast with a district health "
        "intelligence proof-of-concept rooted in live operations.",
    ]),
    ("2. How we monitor GovTech challenges", [
        "We review World Bank GovTech challenge pages weekly. We track new calls for district health, "
        "community data, CHW systems, or African government digitalisation partners.",
        "We submit tailored proposals only when the use case matches — protecting evaluator time and "
        "our own credibility. This pack is a baseline, not a blind submission to unrelated audit challenges.",
    ]),
    ("3. District health intelligence PoC pattern", [
        "FCHIP offers a GovTech-ready pattern: secure community capture on offline-capable mobile tools; "
        "cloud sync and validation; analytics and GIS hotspot views; role-based dashboards co-designed "
        "with district health managers; configurable early-warning alerts aligned to local protocols.",
        "The PoC measures decision-cycle gains — how fast managers see coverage gaps, respond to signals, "
        "and improve CHW data quality for government reporting.",
    ]),
    ("4. What we are not applying for", [
        "Open Moroccan public-finance audit challenges are not FairBanks' domain. We will not repurpose "
        "a community health narrative to fit unrelated audit software requirements.",
        "Our value is district health intelligence from an operator that already runs clinic and community "
        "programmes — not a generic vendor stretching into every challenge category.",
    ]),
    ("5. FCHIP architecture for government partners", [
        "FairBanks Community Health Intelligence Platform sits on the community cascade: members identify "
        "needs; CHWs/VHTs bridge homes to care; outreach programmes deliver services; the medical centre "
        "provides clinical care; research and partnerships strengthen the system; livelihoods make gains last.",
        "For government partners, FCHIP complements national systems — feeding cleaner CHW data, early "
        "warnings, and facility decision support without replacing existing EMR or DHIS2 investments where "
        "interoperability is preferred.",
    ]),
]

MONITOR_PLAN = [
    ("Weekly", "Scan GovTech challenge pages for health / Africa digital-government calls"),
    ("Match gate", "Submit only when use case aligns with district health intelligence PoC"),
    ("If selected", "Attend capacity building / Switzerland bootcamp as offered"),
    ("Co-design", "Work with government partner on scoped PoC and success metrics"),
    ("Measure", "Track decision-cycle time, data quality, and manager usability"),
]

POC_COMPONENTS = [
    ("Community capture", "Offline CHW/VHT forms with validation rules"),
    ("Data pipeline", "Secure sync, anonymisation, and quality dashboards"),
    ("GIS layer", "Hotspot maps for managers and partners"),
    ("Alerts", "Configurable early warnings aligned to local protocols"),
    ("Handover", "Documentation and training for sustained use"),
]

MUTUAL_VALUE = [
    ("FairBanks gains", "Government PoC pathways; World Bank / SECO ecosystem visibility; "
     "learning that hardens FCHIP for public-sector deployment; bootcamp access if selected."),
    ("GovTech programme gains", "Vendor that runs clinic and community operations — not only software; "
     "health intelligence use cases governments need; African implementer ready for co-design; "
     "reusable PoC patterns for district health offices."),
]


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{TEAL}"/>'))
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                font(c.paragraphs[0].add_run(str(val)), size=10, color=SLATE)
                if ri % 2:
                    c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CREAM}"/>'))
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(META["subtitle"], size=12, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption=TAGLINE)
    para(META["monitor_note"], size=10, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    for title, paras in SECTIONS:
        heading(title)
        if "PoC" in title or "architecture" in title.lower():
            image("architecture", caption="District health intelligence — community to government dashboards")
        if "not applying" in title.lower():
            para("We respect the programme by staying in our lane.", italic=True, color=MUTED)
        for t in paras:
            para(t)

    heading("6. PoC components")
    table(["Component", "Function"], POC_COMPONENTS, widths=[2.2, 4.2])

    heading("7. Monitor → match → co-design plan")
    table(["Step", "Action"], MONITOR_PLAN, widths=[1.6, 4.8])

    heading("8. Field foundation")
    image("outreach", caption="Live FairBanks Community Reach — not a slide-deck-only vendor")
    bullets([
        "FairBanks Medical Centre with clinic, pharmacy, and diagnostics",
        "CHW/VHT outreach in Kampala-area communities",
        "Maternal, child, chronic, and geriatric programmes generating data signals",
        "Team prepared for capacity building and bootcamp if a matching challenge selects us",
    ])

    heading("9. Mutual value with GovTech partners")
    for label, body in MUTUAL_VALUE:
        para(f"{label}: {body}", bold=True)

    heading("10. Readiness statement")
    para("GovTech works when government problems meet operators who already serve people. FairBanks is "
         "ready when the right health or Africa digital-government challenge opens.")
    para("We ask to be considered for matching challenges — not to force-fit unrelated audit use cases.")
    image("conclusion", caption="District intelligence ready for the right GovTech call")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable, Table, TableStyle
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))
    st.add(ParagraphStyle("CellHead", fontName="Helvetica-Bold", fontSize=8.5, textColor=white))
    st.add(ParagraphStyle("CellBody", fontName="Helvetica", fontSize=8.5, textColor=slate))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.9, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = w * ih / iw
        if h > 2.8 * inch:
            h, w = 2.8 * inch, 2.8 * inch * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 8))

    def tbl(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        for row in rows:
            data.append([Paragraph(str(c), st["CellBody"]) for c in row])
        widths = widths or [pw / len(headers)] * len(headers)
        t = Table(data, colWidths=widths, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal),
            ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(
        f'<font color="#{ACCENT}"><b><i>FairBanks Community Health Intelligence Platform (FCHIP)</i></b></font>',
        st["Meta"]))
    story.append(Paragraph(f'<font color="#{MUTED}"><i>{META["subtitle"]}</i></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    story.append(Spacer(1, 8))
    img("cover", cap=TAGLINE)
    story.append(Paragraph(META["monitor_note"], st["Meta"]))
    story.append(PageBreak())

    for title, paras in SECTIONS:
        story.append(Paragraph(title, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
        for t in paras:
            story.append(Paragraph(t, st["Body"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph("6. PoC components", st["H1"]))
    tbl(["Component", "Function"], POC_COMPONENTS, [pw * 0.32, pw * 0.68])

    story.append(Paragraph("7. Monitor → match → co-design plan", st["H1"]))
    tbl(["Step", "Action"], MONITOR_PLAN, [pw * 0.22, pw * 0.78])

    story.append(Paragraph("9. Mutual value", st["H1"]))
    for label, body in MUTUAL_VALUE:
        story.append(Paragraph(f"<b>{label}:</b> {body}", st["Body"]))

    story.append(Paragraph("10. Readiness statement", st["H1"]))
    story.append(Paragraph(
        "We ask to be considered for matching health and Africa digital-government challenges as they appear.",
        st["Body"]))
    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                            topMargin=0.7 * inch, bottomMargin=0.7 * inch, title="FairBanks GovTech Readiness")
    doc.build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    def _add_entrance_anims(slide):
        """Light fade-in on pictures and key text (documents.mdc motion)."""
        from lxml import etree
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spids = []
        for shape in slide.shapes:
            try:
                st = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                st = -1
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
            if st == 13 or has_text:
                spids.append(str(shape.shape_id))
            if len(spids) >= 3:
                break
        if not spids:
            return
        sld = slide._element
        for old in sld.findall(f"{{{NS_P}}}timing"):
            sld.remove(old)
        children = []
        nid = 3
        for i, spid in enumerate(spids):
            delay = 0 if i == 0 else 200
            children.append(
                f'<p:par xmlns:p="{NS_P}">'
                f'<p:cTn id="{nid}" presetID="10" presetClass="entr" presetSubtype="0" '
                f'fill="hold" nodeType="withEffect">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{nid + 1}" dur="450"/>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
            )
            nid += 2
        xml = (
            f'<p:timing xmlns:p="{NS_P}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="true" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{nid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(children)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        )
        sld.append(etree.fromstring(xml))

    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic_cover(sl, key):
        sl.shapes.add_picture(embed(key), 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks FCHIP | GovTech Readiness | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    # 1 Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.5), SW, Inches(3.5), NAVY)
    tb(s, Inches(0.6), SH - Inches(3.2), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.7), Inches(12), Inches(0.55), META["doc_title"], size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(2.15), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(1.75), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.2), Inches(12), Inches(0.35), META["monitor_note"], size=11, color="D0E8E8")

    # 2 Readiness
    s = prs.slides.add_slide(blank)
    band(s, "Readiness brief", "Ready PoC partner — not a blind submitter")
    tb(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.5),
       "GovTech connects tech firms with government digitalisation needs.\n\n"
       "Capacity building · Switzerland bootcamp · government PoC\n\n"
       "We maintain this pack and submit only when health / Africa challenges match.",
       size=20, color=MUTED)
    pic_fit(s, "dashboard", Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5))
    footer(s, 2)

    # 3 Monitor
    s = prs.slides.add_slide(blank)
    band(s, "Challenge monitoring", "Weekly scan · match gate · tailored submit")
    for i, (step, action) in enumerate(MONITOR_PLAN[:3]):
        y = Inches(1.25) + i * Inches(1.35)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.2), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.15), Inches(2.5), Inches(0.35), step, size=14, bold=True, color=TEAL)
        tb(s, Inches(3.3), y + Inches(0.15), Inches(9.0), Inches(0.9), action, size=18, color=MUTED)
    footer(s, 3)

    # 4 Not audit
    s = prs.slides.add_slide(blank)
    band(s, "Honest scope", "Moroccan audit challenges are not our target")
    tb(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
       "Current open public-finance audit use cases are outside FairBanks' domain.\n\n"
       "We will not stretch a community health story to fit unrelated software requirements.\n\n"
       "Our lane: district health intelligence from a live clinic + CHW operator.",
       size=18, color=SLATE)
    footer(s, 4)

    # 5 PoC pattern
    s = prs.slides.add_slide(blank)
    band(s, "PoC pattern", "District health intelligence for government partners")
    pic_fit(s, "architecture", Inches(0.4), Inches(1.1), Inches(7.0), Inches(5.7))
    tb(s, Inches(7.7), Inches(1.2), Inches(5.2), Inches(5.5),
       "Capture → sync → GIS → dashboards → alerts\n\n"
       "Measure decision-cycle gains and CHW data quality for reporting.",
       size=19, color=MUTED)
    footer(s, 5)

    # 6 Components
    s = prs.slides.add_slide(blank)
    band(s, "PoC stack", "Five components co-designed with district users")
    for i, (comp, func) in enumerate(POC_COMPONENTS):
        y = Inches(1.15) + i * Inches(1.05)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(0.95), "FFFFFF", LINE)
        tb(s, Inches(0.7), y + Inches(0.1), Inches(3.0), Inches(0.35), comp, size=13, bold=True, color=TEAL)
        tb(s, Inches(3.8), y + Inches(0.1), Inches(8.5), Inches(0.75), func, size=13, color=MUTED)
    footer(s, 6)

    # 7 GIS
    s = prs.slides.add_slide(blank)
    band(s, "Decision support", "Managers see hotspots and coverage gaps")
    pic_fit(s, "gis", Inches(0.45), Inches(1.15), Inches(6.0), Inches(5.6))
    tb(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
       "• Population risk views\n• Early-warning alerts\n• CHW validation rules\n"
       "• Export hooks for national systems where permitted",
       size=16, color=SLATE)
    footer(s, 7)

    # 8 Field ops
    s = prs.slides.add_slide(blank)
    band(s, "Operator advantage", "Clinic + community — not software alone")
    pic_fit(s, "outreach", Inches(0.45), Inches(1.15), Inches(5.8), Inches(5.6))
    tb(s, Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5),
       "• FairBanks Medical Centre + Community Reach\n"
       "• CHW/VHT networks and outreach programmes\n"
       "• Team ready for bootcamp if selected\n"
       "• PoC measured in real manager workflows",
       size=16, color=SLATE)
    footer(s, 8)

    # 9 Mutual value
    s = prs.slides.add_slide(blank)
    band(s, "Win-win", "GovTech ↔ FairBanks ↔ districts")
    rect(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "FairBanks gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[0][1], size=18, color=MUTED)
    rect(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.2), "FFFFFF", LINE)
    tb(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4), "GovTech gains", size=16, bold=True, color=TEAL)
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0), MUTUAL_VALUE[1][1], size=18, color=MUTED)
    footer(s, 9)

    # 10 Ask
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.8), SW, Inches(2.8), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.4), Inches(12), Inches(0.55),
       "Ready when the right health or Africa GovTech challenge opens.", size=24, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.5), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for _sl in prs.slides:
        _add_entrance_anims(_sl)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. govtech document set in", OUT_DIR)
