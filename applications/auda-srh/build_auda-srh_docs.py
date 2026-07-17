#!/usr/bin/env python3
"""
AUDA-NEPAD Youth SRH Ambassadors Initiative 2026 — nomination & advocacy set.

Generates:
  applications/auda-srh/documents/auda-srh_word.docx
  applications/auda-srh/documents/auda-srh_pdf.pdf
  applications/auda-srh/documents/auda-srh_ppt.pptx

Run: python applications/auda-srh/build_auda-srh_docs.py
"""

from pathlib import Path

PROJECT = Path(__file__).resolve().parent
REPO = PROJECT.parents[1]
ASSETS = REPO / "assets"
OUT_DIR = PROJECT / "documents"
SLUG = "auda-srh"
CACHE = REPO / "tmp" / f"{SLUG}_assets"

OUT_DOC = OUT_DIR / f"{SLUG}_word.docx"
OUT_PDF = OUT_DIR / f"{SLUG}_pdf.pdf"
OUT_PPT = OUT_DIR / f"{SLUG}_ppt.pptx"

PROGRAMME_URL = (
    "https://opportunitiesforyouth.org/2026/06/13/"
    "african-union-youth-srh-ambassadors-initiative-2026-apply-to-become-an-auda-nepad-youth-srh-ambassador/"
)

NAVY, TEAL, TEAL_L, ACCENT = "0A1F2E", "0D6E6E", "14A3A3", "C45C26"
SLATE, MUTED, CREAM, LINE = "1E2F38", "3A4A54", "F7F5F0", "D0DCDC"
SLOGAN = "Your health, our mission."

META = {
    "programme": "AUDA-NEPAD Youth SRH Ambassadors Initiative 2026",
    "doc_title": "Ambassador Nomination — Advocacy & Community Mobilisation",
    "subtitle": "Bridging grassroots outreach with continental youth SRH leadership",
    "deadline": "Application deadline: 17 July 2026",
}

PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "outreach": "outreach_facilitator_canopy_01.jpg",
    "discussion": "outreach_audience_discussion_01.jpg",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "group": "outreach_audience_full_group_01.jpg",
    "training": "outreach_training_presenter_01.jpg",
    "conclusion": "outreach_hands_raised_01.jpg",
}

ADVOCACY_TRACK = [
    "School and community dialogues on adolescent SRH through FairBanks outreach",
    "Maternal health education linking ANC uptake, safe pregnancy, and partner involvement",
    "Digital content that meets young people where they are — without shame or jargon",
    "Peer mobilisation through CHWs/VHTs who already hold community trust",
]

ENERGIZE_ALIGN = [
    "Youth as co-creators — not passive beneficiaries — in SRH programming",
    "Grassroots realities from Kampala communities informing continental advocacy",
    "Linking SRH to livelihoods, education, and the demographic dividend",
    "Amplifying AUDA-NEPAD frameworks through lived FairBanks field experience",
]

COMMITMENTS = [
    "Monthly advocacy content tied to outreach calendars and adolescent/maternal programmes",
    "Report anonymised community insights upward — what youth actually ask in the field",
    "Support A2DSRH and Energize Africa messaging with honest, local-language examples",
    "Mentor younger peer advocates through FairBanks Community Reach volunteer pathways",
]

CONTINENTAL_ADDS = [
    "Official ambassador recognition and continental SRH platforms",
    "Training, mentorship, and networking with policy and youth leaders",
    "Visibility to translate Kampala outreach lessons into Africa-wide dialogue",
    "Sponsored participation in selected AUDA-NEPAD engagements where appropriate",
]


def photo(key: str) -> Path:
    p = ASSETS / PHOTOS[key]
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def embed(key: str, max_px: int = 1500) -> str:
    from PIL import Image as PILImage

    src = photo(key)
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def build_docx():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    def font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def para(text, size=11, bold=False, color=SLATE, after=8, align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.22
        font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)

    def heading(text, level=1):
        sizes, colors = {1: 20, 2: 14}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(6)
        font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width=6.2, caption=None):
        from PIL import Image as PILImage
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width, 3.4 * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(path, width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True, after=10)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            font(p.add_run(it), size=11, color=SLATE)

    doc = Document()
    s = doc.sections[0]
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    para(META["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(META["doc_title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para("FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True,
         color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption=META["subtitle"])
    para(META["deadline"], size=10, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=2, bold=True)
    para(PROGRAMME_URL, size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=14)
    doc.add_page_break()

    heading("1. Nomination statement")
    para("I am applying to serve as an AUDA-NEPAD Youth SRH Ambassador because sexual and reproductive health "
         "is central to Africa’s demographic dividend — and because I already mobilise young people through "
         "FairBanks Community Reach in Kampala. Ambassadors should bridge grassroots truth and continental policy; "
         "my daily work lives in that bridge.")

    heading("2. SRH advocacy through FairBanks outreach")
    image("outreach", caption="Community mobilisation under the canopy — where SRH conversations start")
    bullets(ADVOCACY_TRACK)
    para("FCHIP, our community health intelligence layer, helps us see where adolescent and maternal gaps cluster — "
         "so advocacy is evidence-informed, not guesswork.")

    heading("3. Community mobilisation in Kampala")
    image("discussion", caption="Dialogue-first mobilisation — peers, CHWs, and facilitators together")
    para("In Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya, and surrounding communities we run outreach on "
         "maternal health, adolescent wellbeing, chronic screening, and home visits. CHWs and VHTs translate "
         "SRH messages into trusted, local action.")

    heading("4. Energize Africa and Agenda 2063 alignment")
    bullets(ENERGIZE_ALIGN)
    para("Goal 18 — engaged and empowered youth — is not a poster on our wall. It is the design principle for "
         "how we run programmes and how I would represent youth SRH continentally.")

    heading("5. Ambassador commitments if selected")
    image("training", caption="Training the next peer advocates through FairBanks pathways")
    bullets(COMMITMENTS)

    heading("6. What the continental platform adds to this work")
    para("FairBanks will keep delivering care and outreach regardless. Ambassador status amplifies what already "
         "works — connecting Kampala field lessons to Africa-wide SRH dialogue.")
    bullets(CONTINENTAL_ADDS)

    image("conclusion", caption="Youth SRH leadership rooted in community — Your health, our mission.")
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)
    para(f"Source: {PROGRAMME_URL}", size=8, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, italic=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_DOC))
    print(f"DOCX: {OUT_DOC}")


def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, HRFlowable
    from PIL import Image as PILImage

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted = HexColor("#" + SLATE), HexColor("#" + MUTED)
    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CT", fontName="Helvetica-Bold", fontSize=20, textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, textColor=navy, spaceBefore=12, spaceAfter=6))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, textColor=slate, alignment=TA_JUSTIFY, spaceAfter=7))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, textColor=muted, alignment=TA_CENTER, spaceAfter=4))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, textColor=slate, leftIndent=14, spaceAfter=3))

    pw = A4[0] - 1.6 * inch
    story = []

    def img(key, w=pw * 0.88, cap=None):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = min(w * ih / iw, 2.5 * inch)
        w = h * iw / ih
        story.append(Image(path, width=w, height=h))
        if cap:
            story.append(Paragraph(cap, st["Meta"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph(META["programme"], st["Meta"]))
    story.append(Paragraph(META["doc_title"], st["CT"]))
    story.append(Paragraph(
        f'<font color="#{ACCENT}"><b><i>FairBanks Community Health Intelligence Platform (FCHIP)</i></b></font>',
        st["Meta"]))
    story.append(Paragraph(f'<font color="#{MUTED}"><i>{META["subtitle"]}</i></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    img("cover")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b>{META["deadline"]}</b></font>', st["Meta"]))
    story.append(PageBreak())

    sections = [
        ("1. Nomination statement", [
            "I mobilise young people through FairBanks Community Reach in Kampala. Ambassadors should bridge "
            "grassroots truth and continental policy — my daily work lives in that bridge.",
        ], None),
        ("2. SRH advocacy through FairBanks outreach", ADVOCACY_TRACK, "outreach"),
        ("3. Community mobilisation in Kampala", [
            "Outreach across Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya with CHWs/VHTs.",
            "FCHIP helps us see where adolescent and maternal gaps cluster.",
        ], "discussion"),
        ("4. Energize Africa alignment", ENERGIZE_ALIGN, None),
        ("5. Ambassador commitments", COMMITMENTS, "training"),
        ("6. Continental platform value", CONTINENTAL_ADDS, None),
    ]

    for title, items, img_key in sections:
        story.append(Paragraph(title, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1, color=teal, spaceAfter=6))
        if img_key:
            img(img_key)
        for it in items:
            if it.startswith("Outreach") or it.startswith("I mobilise") or "Goal 18" in it or "FairBanks will":
                story.append(Paragraph(it, st["Body"]))
            else:
                story.append(Paragraph("• " + it, st["FBullet"]))
        story.append(Spacer(1, 4))

    img("conclusion")
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(str(OUT_PDF), pagesize=A4, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                      topMargin=0.7 * inch, bottomMargin=0.7 * inch).build(story)
    print(f"PDF: {OUT_PDF}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = C(fill)
        if line:
            s.line.color.rgb = C(line)
        else:
            s.line.fill.background()
        return s

    def tb(sl, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=False):
        """Text box with balanced sizing for tall side-by-side photo panels."""
        box = sl.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        h_in = float(h) / 914400.0
        w_in = float(w) / 914400.0
        # Tall panels beside photos: grow type and vertically centre (not titles/footers)
        body = h_in >= 4.0 and 13 <= size <= 22
        if body:
            try:
                tf.anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        raw = str(text).replace("\r\n", "\n")
        if "\n\n" in raw:
            parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
        else:
            parts = list(raw.split("\n"))
        if not parts:
            parts = [""]
        cpl = max(18, int(w_in * 6.5))
        est = 0.0
        nonempty = 0
        for part in parts:
            if not str(part).strip():
                est += 0.35
                continue
            nonempty += 1
            est += max(1, (len(part) + cpl - 1) // cpl)
        est = max(est, float(nonempty or 1))
        use_size = size
        if body and est > 0:
            fitted = int((h_in * 0.82 * 72) / (est * 1.32))
            # Grow into empty space; shrink slightly if denser than the requested size
            use_size = max(15, min(26, fitted))
            if fitted >= size:
                use_size = max(size, use_size)
        gap = max(10, int(use_size * 0.55)) if body else 3
        for i, ln in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(gap if body and i < len(parts) - 1 else 2)
            try:
                p.line_spacing = 1.28 if body else 1.15
            except Exception:
                pass
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(use_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = C(color)
            r.font.name = "Calibri"
        return box

    def pic_cover(sl, key):
        sl.shapes.add_picture(embed(key), 0, 0, width=SW, height=SH)

    def pic_fit(sl, key, x, y, w, h):
        path = embed(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        sl.shapes.add_picture(path, x + (w - tw) // 2, y + (h - th) // 2, width=tw, height=th)

    def band(sl, kicker, title):
        rect(sl, 0, 0, SW, Inches(1.0), CREAM)
        rect(sl, 0, 0, Inches(0.15), Inches(1.0), TEAL)
        tb(sl, Inches(0.45), Inches(0.1), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=ACCENT)
        tb(sl, Inches(0.45), Inches(0.42), Inches(12), Inches(0.5), title, size=24, bold=True, color=NAVY)

    def footer(sl, n):
        rect(sl, 0, SH - Inches(0.3), SW, Inches(0.3), NAVY)
        tb(sl, Inches(0.4), SH - Inches(0.29), Inches(10), Inches(0.28),
           f"FairBanks | AUDA-NEPAD Youth SRH Ambassador | {SLOGAN}", size=9, color="FFFFFF")
        tb(sl, SW - Inches(0.8), SH - Inches(0.29), Inches(0.5), Inches(0.28), str(n), size=9,
           color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover")
    rect(s, 0, SH - Inches(3.4), SW, Inches(3.4), NAVY)
    tb(s, Inches(0.6), SH - Inches(3.1), Inches(12), Inches(0.35), META["programme"], size=13, bold=True, color=TEAL_L)
    tb(s, Inches(0.6), SH - Inches(2.65), Inches(12), Inches(0.5), META["doc_title"], size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(2.1), Inches(12), Inches(0.3),
       "FairBanks Community Health Intelligence Platform (FCHIP)", size=13, bold=True, italic=True, color="F2C79B")
    tb(s, Inches(0.6), SH - Inches(1.65), Inches(12), Inches(0.3), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.15), Inches(12), Inches(0.35), META["deadline"], size=12, bold=True, color="F2C79B")

    # Nomination
    s = prs.slides.add_slide(blank)
    band(s, "Nomination", "Grassroots mobiliser → continental SRH voice")
    tb(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
       "FairBanks Community Reach runs maternal & adolescent programmes with CHWs/VHTs.\n\n"
       "Ambassadors must connect field truth to policy — that is already my work.\n\n"
       "FCHIP adds evidence on where SRH gaps cluster.",
       size=18, color=SLATE)
    footer(s, 2)

    # Advocacy track
    s = prs.slides.add_slide(blank)
    band(s, "Advocacy", "SRH through FairBanks outreach")
    pic_fit(s, "outreach", Inches(0.45), Inches(1.15), Inches(5.8), Inches(5.6))
    tb(s, Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5),
       "\n".join(f"• {a}" for a in ADVOCACY_TRACK), size=19, color=SLATE)
    footer(s, 3)

    # Mobilisation
    s = prs.slides.add_slide(blank)
    band(s, "Mobilisation", "Kampala communities — dialogue first")
    pic_fit(s, "discussion", Inches(0.45), Inches(1.15), Inches(6.0), Inches(5.6))
    tb(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.5),
       "Bukoto · Kyebando · Kisaasi · Kamwokya · Kikaaya\n\n"
       "CHWs translate SRH into trusted local action\n\n"
       "Digital + in-person — no shame, no jargon",
       size=20, color=MUTED)
    footer(s, 4)

    # Energize Africa
    s = prs.slides.add_slide(blank)
    band(s, "Alignment", "Energize Africa · Agenda 2063 Goal 18")
    for i, e in enumerate(ENERGIZE_ALIGN):
        y = Inches(1.25) + i * Inches(1.15)
        rect(s, Inches(0.5), y, Inches(12.2), Inches(1.0), "FFFFFF", LINE)
        tb(s, Inches(0.75), y + Inches(0.2), Inches(11.5), Inches(0.6), e, size=19, color=SLATE)
    footer(s, 5)

    # Commitments
    s = prs.slides.add_slide(blank)
    band(s, "If selected", "Ambassador commitments")
    pic_fit(s, "mobile", Inches(0.45), Inches(1.15), Inches(5.2), Inches(5.6))
    tb(s, Inches(6.0), Inches(1.2), Inches(6.8), Inches(5.5),
       "\n".join(f"• {c}" for c in COMMITMENTS), size=18, color=SLATE)
    footer(s, 6)

    # Continental value
    s = prs.slides.add_slide(blank)
    band(s, "Amplification", "What the ambassador platform adds")
    tb(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
       "\n".join(f"• {c}" for c in CONTINENTAL_ADDS), size=20, color=SLATE)
    pic_fit(s, "group", Inches(6.6), Inches(1.15), Inches(6.2), Inches(5.6))
    footer(s, 7)

    # Close
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion")
    rect(s, 0, SH - Inches(2.6), SW, Inches(2.6), NAVY)
    tb(s, Inches(0.6), SH - Inches(2.2), Inches(12), Inches(0.55),
       "Youth SRH advocacy rooted in community — ready to serve continentally.", size=22, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(1.25), Inches(12), Inches(0.35), SLOGAN, size=14, bold=True, color="FFFFFF")
    tb(s, Inches(0.6), SH - Inches(0.75), Inches(12), Inches(0.3), META["deadline"], size=12, color="F2C79B")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"PPTX: {OUT_PPT}")


if __name__ == "__main__":
    build_docx()
    build_pdf()
    build_pptx()
    print("Done. auda-srh document set in", OUT_DIR)
