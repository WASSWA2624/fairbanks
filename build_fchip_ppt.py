#!/usr/bin/env python3
"""
FCHIP overview deck — 10 slides at repo root.
FairBanks Community Health Intelligence Platform.
"""

from pathlib import Path

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent
ASSETS = REPO / "assets"
CURSOR = REPO / ".cursor"
LOGO = ASSETS / "fairbanks_logo.jpeg"
OUT_PPT = REPO / "FCHIP_ppt.pptx"

NAVY, TEAL, GREEN = "0A1F2E", "0D6E6E", "2D7A55"
ORANGE, GOLD, CREAM = "C45C26", "D99A2B", "F7F5F0"
WHITE = "FFFFFF"
SLATE, MUTED, LINE, LIGHT = "1E2F38", "52636C", "CED9D8", "D4E8DC"
PALE_TEAL = "E8F3F2"

ORG = "FairBanks Medical Centre"
SLOGAN = "Your health, our mission."
FOOTER = f"{ORG}  ·  {SLOGAN}"
FCHIP = "FairBanks Community Health Intelligence Platform"

PHOTOS = {
    "cover": ASSETS / "cover_hero_cinematic.jpg",
    "facility": ASSETS / "facility_exterior_branded_entrance_01.jpeg",
    "facility_wide": ASSETS / "facility_exterior_branded_entrance_02.jpeg",
    "outreach": ASSETS / "outreach_bp_screening.jpeg",
    "dashboard": ASSETS / "dashboard_demo.png",
    "architecture": ASSETS / "data_flow_iso_labeled.png",
    "gis": ASSETS / "gis_hotspots.png",
    "mobile": ASSETS / "outreach_mobile_phone_demo_01.jpg",
    "concept": CURSOR / "concept_simple.jpeg",
    "community": ASSETS / "outreach_audience_full_group_01.jpg",
}

ML, MR = 0.55, 0.55
SW, SH = 13.333, 7.5
CW = SW - ML - MR
TOP_BAR = 0.06
LOGO_H = 0.40
LOGO_W = LOGO_H * (269 / 101)
LOGO_Y = 0.24
FOOTER_Y = 7.02
CONTENT_BOTTOM = 6.62
HEADER_Y = 1.05
HEADER_Y_SUB = 1.42


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run(run, text, size, color, bold=False, font="Calibri"):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_fade(slide):
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    transition = etree.Element(qn("p:transition"), {"spd": "med"})
    etree.SubElement(transition, qn("p:fade"))
    cSld = sld.find(qn("p:cSld"))
    if cSld is not None:
        cSld.addnext(transition)
    else:
        sld.insert(0, transition)


def build():
    missing = [p for p in [LOGO, *PHOTOS.values()] if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing assets:\n" + "\n".join(map(str, missing)))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
        else:
            shape.line.fill.background()
        if rounded:
            try:
                shape.adjustments[0] = 0.09
            except Exception:
                pass
        try:
            spPr = shape._element.spPr
            for child in list(spPr):
                if "effectLst" in child.tag or "effectDag" in child.tag:
                    spPr.remove(child)
            etree.SubElement(spPr, qn("a:effectLst"))
        except Exception:
            pass
        return shape

    def textbox(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=22,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        line_spacing=1.5,
        para_gap=12,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(22000)
        tf.margin_top = tf.margin_bottom = Emu(8000)
        tf.vertical_anchor = valign
        lines = value.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(0)
            if len(lines) > 1:
                p.space_after = Pt(22 if line.strip() == "" else para_gap)
            else:
                p.space_after = Pt(0)
            p.line_spacing = line_spacing
            r = p.add_run()
            set_run(r, line, size, color, bold)
        return box

    def bullets(slide, items, x, y, w, h, size=22, color=SLATE, space=18):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(22000)
        tf.margin_top = tf.margin_bottom = Emu(6000)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4)
            p.space_after = Pt(space)
            p.line_spacing = 1.5
            r = p.add_run()
            set_run(r, "•  " + item, size, color)
        return box

    def add_logo(slide, x, y, height=LOGO_H):
        pad = 0.08
        py = max(y - pad, TOP_BAR + 0.05)
        rect(slide, x - pad, py, LOGO_W + 2 * pad, height + (y - py) + pad, WHITE, LINE, True)
        return slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), height=Inches(height))

    def crop_photo(slide, path, x, y, w, h):
        with PILImage.open(path) as im:
            iw, ih = im.size
        pic = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
        )
        fr, ir = w / h, iw / ih
        if ir > fr:
            a = (1 - fr / ir) / 2
            pic.crop_left = pic.crop_right = a
        else:
            a = (1 - ir / fr) / 2
            pic.crop_top = pic.crop_bottom = a
        return pic

    def new_slide(bg=CREAM):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, SW, SH, bg)
        add_fade(s)
        return s

    def header(s, kicker, title, subtitle=""):
        rect(s, 0, 0, SW, TOP_BAR, TEAL)
        rect(s, 0, TOP_BAR, 0.08, SH - TOP_BAR, ORANGE)
        logo_x = SW - MR - LOGO_W
        add_logo(s, logo_x, LOGO_Y)
        tw = logo_x - ML - 0.40
        if kicker.strip():
            box = s.shapes.add_textbox(Inches(ML), Inches(0.28), Inches(tw), Inches(0.62))
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(18000)
            tf.margin_top = tf.margin_bottom = Emu(4000)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.15
            r1 = p.add_run()
            set_run(r1, kicker.strip() + ": ", 28, ORANGE, True)
            r2 = p.add_run()
            set_run(r2, title, 28, NAVY, True)
            y = 1.00
        else:
            textbox(s, title, ML, 0.24, tw, 0.62, 34, NAVY, True)
            y = 0.98
        if subtitle:
            textbox(s, subtitle, ML, y, CW, 0.42, 20, MUTED)
            return HEADER_Y_SUB
        return HEADER_Y

    def footer(s, number, total):
        rect(s, ML, FOOTER_Y, CW, 0.012, LINE)
        textbox(s, FOOTER, ML, FOOTER_Y + 0.08, 9.8, 0.28, 14, MUTED)
        textbox(
            s,
            f"{number}  /  {total}",
            SW - MR - 1.3,
            FOOTER_Y + 0.08,
            1.3,
            0.28,
            14,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )

    slides = []

    # ------------------------------------------------------------------
    # 1 Cover
    # ------------------------------------------------------------------
    s = new_slide(NAVY)
    panel = 7.85
    crop_photo(s, PHOTOS["cover"], 0, 0, SW, SH)
    rect(s, 0, 0, panel, SH, NAVY)
    rect(s, panel, 0, 0.12, SH, GOLD)
    add_logo(s, 0.55, 0.28, 0.50)

    textbox(s, "COMMUNITY HEALTH INTELLIGENCE", 0.55, 1.05, 7.0, 0.36, 18, GOLD, True)
    textbox(s, ORG, 0.55, 1.50, 7.0, 0.40, 22, TEAL, True)
    textbox(s, "FCHIP", 0.55, 2.10, 7.0, 0.70, 52, WHITE, True)
    textbox(s, FCHIP, 0.55, 2.85, 7.0, 0.55, 20, GOLD, True)
    rect(s, 0.55, 3.55, 2.80, 0.09, GOLD)

    textbox(
        s,
        "From reactive treatment to proactive prevention —\n"
        "powered by community data, AI, GIS, and climate signals.",
        0.55,
        3.85,
        7.0,
        0.90,
        20,
        LIGHT,
        para_gap=8,
    )
    textbox(s, SLOGAN, 0.55, 5.00, 7.0, 0.40, 24, GOLD, True)

    rect(s, 0, 6.25, panel, 1.25, TEAL)
    textbox(s, "A FairBanks platform component", 0.55, 6.42, 7.0, 0.38, 19, WHITE, True)
    textbox(
        s,
        "Kampala, Uganda  ·  fairbanksmedicalcentre.org/fchip",
        0.55,
        6.88,
        7.0,
        0.35,
        17,
        LIGHT,
    )
    slides.append(s)

    # ------------------------------------------------------------------
    # 2 Problem
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "The gap",
        "Healthcare still waits for people to get sick",
        "Community signals stay fragmented — so outbreaks, maternal risk, and stock-outs arrive late.",
    )
    left_w = CW * 0.52
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["outreach"], ML, y0, left_w, h)
    rect(s, ML + left_w + gap, y0, right_w, h, WHITE, LINE, True)
    problems = [
        ("Late detection", "Facilities treat illness after it appears"),
        ("Siloed data", "CHW, school, clinic, and EMR records stay separate"),
        ("Climate blind", "Weather and place signals rarely join health data"),
        ("Guesswork stock", "Medicine orders follow history, not early risk"),
    ]
    accents = [ORANGE, TEAL, GREEN, GOLD]
    row_h = (h - 0.40) / 4
    for i, ((title, body), accent) in enumerate(zip(problems, accents)):
        y = y0 + 0.20 + i * row_h
        rx = ML + left_w + gap
        rect(s, rx + 0.22, y + 0.12, 0.10, row_h - 0.28, accent)
        textbox(s, title, rx + 0.50, y + 0.10, right_w - 0.75, 0.36, 22, NAVY, True)
        textbox(s, body, rx + 0.50, y + 0.48, right_w - 0.75, 0.40, 18, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 3 Community Reach
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Operating model",
        "How FairBanks Community Reach works",
        "FCHIP sits on the Data & Feedback loop — it does not replace the cascade.",
    )
    left_w = CW * 0.46
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["concept"], ML, y0, left_w, h)
    steps = [
        ("1", "Community members"),
        ("2", "CHWs / VHTs — the bridge"),
        ("3", "Community Reach programmes"),
        ("4", "FairBanks Medical Centre"),
        ("5", "Research · partnerships · skills"),
        ("6", "Empowerment · CHIS · livelihoods"),
    ]
    accents = [TEAL, ORANGE, GREEN, TEAL, ORANGE, GREEN]
    cell_h = (h - 5 * 0.10) / 6
    for i, ((num, label), accent) in enumerate(zip(steps, accents)):
        y = y0 + i * (cell_h + 0.10)
        x = ML + left_w + gap
        rect(s, x, y, right_w, cell_h, WHITE, LINE, True)
        rect(s, x, y, 0.70, cell_h, accent)
        textbox(
            s,
            num,
            x,
            y,
            0.70,
            cell_h,
            22,
            WHITE,
            True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        textbox(
            s,
            label,
            x + 0.90,
            y,
            right_w - 1.10,
            cell_h,
            20,
            NAVY,
            True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slides.append(s)

    # ------------------------------------------------------------------
    # 4 Solution — FCHIP
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Solution", "FCHIP — the intelligence layer")
    left_w = CW * 0.48
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    rect(s, ML, y0, left_w, h, TEAL, rounded=True)
    textbox(s, "What FCHIP does", ML + 0.35, y0 + 0.28, left_w - 0.70, 0.40, 20, GOLD, True)
    textbox(
        s,
        "Connects communities, CHWs/VHTs, schools, clinics, and hospitals into one intelligent network.\n\n"
        "Fuses community signals with GIS maps and climate APIs.\n\n"
        "Safely exposes data APIs to existing EMR/HMS — without replacing systems clinics already use.",
        ML + 0.35,
        y0 + 0.85,
        left_w - 0.70,
        h - 1.20,
        20,
        WHITE,
        para_gap=10,
    )
    crop_photo(s, PHOTOS["dashboard"], ML + left_w + gap, y0, right_w, h * 0.58)
    chips = [
        ("Predict", "Risk before crisis"),
        ("Map", "Hotspots by place"),
        ("Alert", "CHWs & facilities"),
        ("Learn", "Close the loop"),
    ]
    fx = ML + left_w + gap
    fy = y0 + h * 0.58 + 0.14
    fh = CONTENT_BOTTOM - fy
    gap_x = 0.12
    cell_w = (right_w - gap_x) / 2
    cell_h = (fh - gap_x) / 2
    accents = [TEAL, ORANGE, GREEN, GOLD]
    for i, ((t, b), accent) in enumerate(zip(chips, accents)):
        col, row = i % 2, i // 2
        x = fx + col * (cell_w + gap_x)
        y = fy + row * (cell_h + gap_x)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x, y, 0.10, cell_h, accent)
        textbox(s, t, x + 0.22, y + 0.12, cell_w - 0.35, 0.32, 18, NAVY, True)
        textbox(s, b, x + 0.22, y + 0.44, cell_w - 0.35, 0.36, 16, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 5 Deep technology
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Deep tech",
        "More than an app — an intelligence stack",
        "AI · ML · GIS · climate APIs · secure EMR/HMS APIs · mobile edge · cloud",
    )
    tech = [
        ("AI / ML", "Disease risk, maternal flags, NCD hotspots", TEAL),
        ("GIS mapping", "Village and parish risk geography", ORANGE),
        ("Climate APIs", "Rainfall, heat, and flood-linked warning", GREEN),
        ("EMR / HMS APIs", "Real-time clinical ingest, consent-aware", GOLD),
        ("Mobile capture", "Offline CHW/VHT forms and worklists", TEAL),
        ("Dashboards", "Facility, programme, and partner views", ORANGE),
    ]
    gap_x, gap_y = 0.20, 0.16
    cell_w = (CW - 2 * gap_x) / 3
    cell_h = (CONTENT_BOTTOM - y0 - gap_y) / 2
    for i, (title, body, accent) in enumerate(tech):
        col, row = i % 3, i // 3
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x, y, cell_w, 0.12, accent)
        textbox(s, title, x + 0.28, y + 0.40, cell_w - 0.56, 0.45, 24, NAVY, True)
        textbox(s, body, x + 0.28, y + 1.00, cell_w - 0.56, cell_h - 1.25, 19, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 6 How it works
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Architecture",
        "From last-mile capture to action",
        "Data in → FCHIP intelligence → alerts, maps, and referrals back to the field.",
    )
    left_w = CW * 0.58
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["architecture"], ML, y0, left_w, h)
    flow = [
        ("1", "Capture", "CHWs, clinics, EMR/HMS, climate feeds"),
        ("2", "Analyse", "AI, GIS, climate fusion, risk scores"),
        ("3", "Act", "Alerts, dashboards, outreach, stock"),
    ]
    accents = [TEAL, ORANGE, GREEN]
    cell_h = (h - 2 * 0.16) / 3
    for i, ((num, title, body), accent) in enumerate(zip(flow, accents)):
        y = y0 + i * (cell_h + 0.16)
        x = ML + left_w + gap
        rect(s, x, y, right_w, cell_h, WHITE, LINE, True)
        rect(s, x, y, 0.85, cell_h, accent)
        textbox(
            s,
            num,
            x,
            y,
            0.85,
            cell_h,
            28,
            WHITE,
            True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        textbox(s, title, x + 1.05, y + 0.28, right_w - 1.25, 0.40, 22, NAVY, True)
        textbox(s, body, x + 1.05, y + 0.72, right_w - 1.25, cell_h - 0.90, 18, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 7 Use cases
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(s, "Proof cases", "What FCHIP helps predict and prevent")
    cases = [
        ("Disease surveillance", "Fever clusters + rainfall → early malaria / outbreak flags", TEAL),
        ("Maternal health", "Home-visit signals → high-risk pregnancy alerts to CHWs", ORANGE),
        ("NCDs", "BP and glucose trends → hypertension and diabetes hotspots", GREEN),
        ("Child health", "Growth, immunisation, diarrhoea → community flags", GOLD),
        ("Medicine demand", "Disease + climate + history → pharmacy pre-stocking", TEAL),
        ("Cascade metrics", "Referrals, CHIS, school sessions → close the feedback loop", ORANGE),
    ]
    gap_x, gap_y = 0.18, 0.14
    cell_w = (CW - gap_x) / 2
    cell_h = (CONTENT_BOTTOM - y0 - 2 * gap_y) / 3
    for i, (title, body, accent) in enumerate(cases):
        col, row = i % 2, i // 2
        x = ML + col * (cell_w + gap_x)
        y = y0 + row * (cell_h + gap_y)
        rect(s, x, y, cell_w, cell_h, WHITE, LINE, True)
        rect(s, x + 0.18, y + 0.18, 0.10, cell_h - 0.36, accent)
        textbox(s, title, x + 0.45, y + 0.22, cell_w - 0.70, 0.40, 20, NAVY, True)
        textbox(s, body, x + 0.45, y + 0.68, cell_w - 0.70, cell_h - 0.90, 17, MUTED)
    slides.append(s)

    # ------------------------------------------------------------------
    # 8 Traction & advantage
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Why FairBanks",
        "Live field base + working FCHIP MVP",
        "Design, pilot, validate, and refine inside a real Community Reach ecosystem.",
    )
    left_w = CW * 0.42
    gap = 0.24
    right_w = CW - left_w - gap
    h = CONTENT_BOTTOM - y0
    crop_photo(s, PHOTOS["community"], ML, y0, left_w, h * 0.62)
    rect(s, ML, y0 + h * 0.62 + 0.14, left_w, h * 0.38 - 0.14, TEAL, rounded=True)
    textbox(
        s,
        "Bukoto · Kyebando · Kisaasi\nKamwokya · Kikaaya & nearby",
        ML + 0.22,
        y0 + h * 0.62 + 0.28,
        left_w - 0.44,
        h * 0.38 - 0.45,
        17,
        WHITE,
        para_gap=6,
    )
    assets = [
        "Licensed medical centre in Kampala",
        "Active CHW / VHT Community Reach",
        "Maternal, child, Gericare, NCD programmes",
        "CHIS and livelihood pathways",
        "Working FCHIP MVP — validating now",
        "Secure EMR/HMS API path planned",
    ]
    bullets(s, assets, ML + left_w + gap, y0, right_w, h, 20, SLATE, 16)
    slides.append(s)

    # ------------------------------------------------------------------
    # 9 MVP & roadmap
    # ------------------------------------------------------------------
    s = new_slide()
    y0 = header(
        s,
        "Roadmap",
        "Validate → district → regional scale",
        "Working MVP exists. Phase 1 proves value in the live FairBanks catchment.",
    )
    phases = [
        ("1", "Validate", "0–12 mo", "CHW tools, GIS + climate, EMR APIs, 3 use cases", TEAL),
        ("2", "District", "12–24 mo", "Partner clinics, Kampala structures, NGO M&E", ORANGE),
        ("3", "Regional", "24–36 mo", "Multi-district Uganda · East Africa entry", GREEN),
        ("4", "Platform", "Year 3+", "Clinical support, NLP, research modules", GOLD),
    ]
    gap = 0.18
    card_w = (CW - 3 * gap) / 4
    card_h = CONTENT_BOTTOM - y0
    for i, (num, label, when, detail, accent) in enumerate(phases):
        x = ML + i * (card_w + gap)
        rect(s, x, y0, card_w, card_h, WHITE, LINE, True)
        rect(s, x, y0, card_w, 1.35, accent)
        textbox(s, f"Phase {num}", x + 0.18, y0 + 0.22, card_w - 0.36, 0.32, 16, WHITE, True)
        textbox(s, label, x + 0.18, y0 + 0.55, card_w - 0.36, 0.40, 24, WHITE, True)
        textbox(s, when, x + 0.18, y0 + 0.95, card_w - 0.36, 0.28, 16, LIGHT)
        textbox(s, detail, x + 0.18, y0 + 1.60, card_w - 0.36, card_h - 1.85, 17, SLATE)
    slides.append(s)

    # ------------------------------------------------------------------
    # 10 Closing
    # ------------------------------------------------------------------
    s = new_slide(NAVY)
    panel = 7.85
    crop_photo(s, PHOTOS["facility_wide"], 0, 0, SW, SH)
    rect(s, 0, 0, panel, SH, NAVY)
    rect(s, panel, 0, 0.12, SH, GOLD)
    add_logo(s, 0.55, 0.28, 0.50)

    textbox(s, "LET'S BUILD TOGETHER", 0.55, 1.05, 7.0, 0.36, 18, GOLD, True)
    textbox(s, "FCHIP", 0.55, 1.55, 7.0, 0.55, 40, WHITE, True)
    textbox(
        s,
        "Africa's community health intelligence —\nrooted in FairBanks Community Reach.",
        0.55,
        2.25,
        7.0,
        0.90,
        22,
        LIGHT,
        para_gap=8,
    )
    rect(s, 0.55, 3.35, 2.80, 0.09, GOLD)
    textbox(s, SLOGAN, 0.55, 3.65, 7.0, 0.40, 24, GOLD, True)

    textbox(
        s,
        "Health for All — Obulamu eri Bonna\nAfya kwa Wote · Oburamu bwa Boona",
        0.55,
        4.25,
        7.0,
        0.70,
        18,
        LIGHT,
        para_gap=6,
    )

    rect(s, 0.55, 5.20, 6.90, 1.80, TEAL, rounded=True)
    textbox(s, "Contact", 0.80, 5.35, 6.4, 0.32, 17, GOLD, True)
    textbox(
        s,
        "Racheal Nabukeera · Managing Director & Co-founder\n"
        "info@fairbanksmedicalcentre.org  ·  +256 772 849 258\n"
        "fairbanksmedicalcentre.org  ·  /fchip",
        0.80,
        5.75,
        6.4,
        1.05,
        18,
        WHITE,
        para_gap=8,
    )
    slides.append(s)

    total = len(slides)
    assert total == 10, f"Expected 10 slides, got {total}"
    for i, s in enumerate(slides):
        if i not in (0, total - 1):
            footer(s, i + 1, total)

    prs.save(OUT_PPT)
    print(f"Wrote {OUT_PPT} ({total} slides)")


if __name__ == "__main__":
    build()
