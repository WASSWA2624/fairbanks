"""
AWIEF Pitch n Grow 2026 — FairBanks / FCIN / CHIP investor pitch deck (12 slides).
Image-rich, lively widescreen deck. Strictly capped at 10–12 slides.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Brand palette (HealthTech: deep teal + navy) ---
NAVY = RGBColor(0x0A, 0x1F, 0x2E)
TEAL = RGBColor(0x0D, 0x6E, 0x6E)
TEAL_LIGHT = RGBColor(0x14, 0xA3, 0xA3)
CREAM = RGBColor(0xF7, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x2C, 0x3E, 0x4A)
MUTED = RGBColor(0x5A, 0x6B, 0x75)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)
LINE = RGBColor(0xD0, 0xDC, 0xDC)
SOFT_TEAL = RGBColor(0xE4, 0xF2, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ASSETS = Path(__file__).resolve().parent / "assets"


def img(name: str) -> str:
    path = ASSETS / name
    if not path.exists():
        raise FileNotFoundError(f"Missing asset: {path}")
    return str(path)


def set_run(run, size=18, bold=False, color=SLATE, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(shape, color)
    return shape


def add_round_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(shape, color)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def textbox(slide, left, top, width, height, text, size=18, bold=False, color=SLATE,
            align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def bullet_box(slide, left, top, width, height, items, size=14, color=SLATE, bullet="•"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5 if i else 0)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        set_run(run, size=size, bold=False, color=color)
    return box


def footer(slide, page, total=12, dark=False):
    bar = NAVY if not dark else RGBColor(0x06, 0x16, 0x22)
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), bar)
    textbox(
        slide, Inches(0.5), Inches(7.18), Inches(9), Inches(0.28),
        "AWIEF Pitch n Grow 2026  |  Deep Roots. Digital Futures.  |  HealthTech · AI/ML",
        size=10, color=RGBColor(0xA8, 0xC4, 0xC4),
    )
    textbox(
        slide, Inches(11.2), Inches(7.18), Inches(1.6), Inches(0.28),
        f"{page} / {total}", size=10, color=RGBColor(0xA8, 0xC4, 0xC4), align=PP_ALIGN.RIGHT,
    )


def add_picture_cover(slide, path, left, top, width, height):
    """Place image cropped to fill a frame (cover-style via centered crop approx)."""
    pic = slide.shapes.add_picture(path, left, top, width=width)
    # If taller than frame, crop; if wider, crop sides — keep aspect via crop
    nat_w, nat_h = pic.image.size  # pixels
    target_ratio = width / height
    img_ratio = nat_w / nat_h
    if img_ratio > target_ratio:
        # too wide — crop left/right
        new_w = height * img_ratio
        pic.height = height
        pic.width = int(new_w)
        crop_frac = (pic.width - width) / (2 * pic.width)
        pic.crop_left = crop_frac
        pic.crop_right = crop_frac
        pic.left = left
        pic.top = top
        pic.width = width
        pic.height = height
    else:
        # too tall — crop top/bottom
        new_h = width / img_ratio
        pic.width = width
        pic.height = int(new_h)
        crop_frac = (pic.height - height) / (2 * pic.height)
        pic.crop_top = crop_frac
        pic.crop_bottom = crop_frac
        pic.left = left
        pic.top = top
        pic.width = width
        pic.height = height
    return pic


def add_picture_fit(slide, path, left, top, width, height):
    """Fit image inside box preserving aspect (letterbox not drawn)."""
    pic = slide.shapes.add_picture(path, left, top, width=width)
    nat_w, nat_h = pic.image.size
    img_ratio = nat_w / nat_h
    box_ratio = width / height
    if img_ratio > box_ratio:
        pic.width = width
        pic.height = int(width / img_ratio)
        pic.top = int(top + (height - pic.height) / 2)
        pic.left = left
    else:
        pic.height = height
        pic.width = int(height * img_ratio)
        pic.left = int(left + (width - pic.width) / 2)
        pic.top = top
    return pic


def image_card(slide, path, left, top, width, height, caption=None):
    # Shadow-like border frame
    frame = add_round_rect(slide, left, top, width, height, WHITE)
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1)
    # Image inset
    inset = Inches(0.06)
    add_picture_cover(
        slide, path,
        left + inset, top + inset,
        width - 2 * inset,
        height - (Inches(0.38) if caption else 2 * inset),
    )
    if caption:
        textbox(
            slide, left + Inches(0.12), top + height - Inches(0.34),
            width - Inches(0.2), Inches(0.28),
            caption, size=10, color=MUTED, align=PP_ALIGN.CENTER,
        )


def section_band(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), TEAL)
    add_rect(slide, 0, Inches(0.12), SLIDE_W, Inches(7.03), CREAM)
    textbox(slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.4),
            title, size=26, bold=True, color=NAVY)
    if subtitle:
        textbox(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.32),
                subtitle, size=13, color=MUTED)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ========== 1. TITLE — full-bleed hero ==========
    s = prs.slides.add_slide(blank)
    add_picture_cover(s, img("hero_chw_mobile.png"), 0, 0, SLIDE_W, SLIDE_H)
    # Dark gradient overlays via solid panels
    add_rect(s, 0, 0, Inches(7.2), SLIDE_H, NAVY)
    # Soften right edge with teal accent bar
    add_rect(s, Inches(7.05), 0, Inches(0.12), SLIDE_H, TEAL_LIGHT)
    textbox(s, Inches(0.55), Inches(1.0), Inches(6.2), Inches(0.35),
            "AWIEF Pitch n Grow 2026  ·  Startup Track  ·  HealthTech",
            size=12, color=TEAL_LIGHT)
    textbox(s, Inches(0.55), Inches(1.55), Inches(6.3), Inches(1.0),
            "FairBanks Community\nIntelligence Network",
            size=30, bold=True, color=WHITE)
    textbox(s, Inches(0.55), Inches(2.75), Inches(6.2), Inches(0.4),
            "Community Health Intelligence Platform (CHIP)",
            size=16, color=TEAL_LIGHT)
    textbox(s, Inches(0.55), Inches(3.4), Inches(6.2), Inches(1.2),
            "AI-powered community health intelligence that turns last-mile data into "
            "predictions — so African primary care prevents crises instead of waiting for them.",
            size=15, color=RGBColor(0xD0, 0xDC, 0xE0))
    textbox(s, Inches(0.55), Inches(4.85), Inches(6.2), Inches(0.35),
            "Theme: Deep Roots. Digital Futures.",
            size=13, bold=True, color=ACCENT)
    textbox(s, Inches(0.55), Inches(5.35), Inches(6.2), Inches(0.35),
            "Uganda pilot  →  East Africa  →  Pan-African scale",
            size=13, color=RGBColor(0xA8, 0xC4, 0xC4))
    textbox(s, Inches(0.55), Inches(6.3), Inches(6.2), Inches(0.4),
            "Woman-led deep-tech  |  Live medical centre + CHW/VHT foundation",
            size=11, color=MUTED)

    # ========== 2. PROBLEM — image + pain cards ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "The problem", "Primary healthcare waits for sickness — intelligence arrives too late.")
    image_card(s, img("reactive_clinic.png"), Inches(0.45), Inches(1.15), Inches(5.4), Inches(5.7),
               "Reactive clinic flow — patients arrive after the crisis starts")
    problems = [
        ("Facilities wait for sickness", "Outbreaks detected too late"),
        ("Care ends at the gate", "No village-level visibility"),
        ("Outreach data siloed", "Districts lack live intel"),
        ("Guesswork procurement", "Seasonal stock-outs"),
        ("Late high-risk detection", "Preventable maternal & NCD harm"),
    ]
    for i, (title, desc) in enumerate(problems):
        top = Inches(1.2) + Inches(i * 1.05)
        card = add_round_rect(s, Inches(6.1), top, Inches(6.7), Inches(0.95), WHITE)
        card.line.color.rgb = LINE
        add_rect(s, Inches(6.1), top, Inches(0.12), Inches(0.95), ACCENT if i == 0 else TEAL)
        textbox(s, Inches(6.45), top + Inches(0.15), Inches(6.1), Inches(0.35),
                title, size=14, bold=True, color=NAVY)
        textbox(s, Inches(6.45), top + Inches(0.5), Inches(6.1), Inches(0.3),
                desc, size=12, color=MUTED)
    footer(s, 2)

    # ========== 3. SOLUTION — split with dashboard ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "The solution: FCIN / CHIP",
                 "From community medical practice to Africa's Community Health Intelligence Company.")
    textbox(s, Inches(0.5), Inches(1.15), Inches(6.3), Inches(0.9),
            "CHIP ingests community-generated health data, runs AI/ML and GIS analytics, "
            "and delivers predictive alerts to CHWs, facilities, districts, and partners — before crises escalate.",
            size=14, color=SLATE)
    pillars = [
        ("FCIN", "Venture brand connecting communities, providers, governments & partners"),
        ("CHIP", "Deep-tech platform: capture → intelligence → action"),
        ("FairBanks edge", "Live centre + CHW/VHT networks to pilot & validate in the field"),
    ]
    for i, (t, b) in enumerate(pillars):
        top = Inches(2.2) + Inches(i * 1.4)
        card = add_round_rect(s, Inches(0.5), top, Inches(6.3), Inches(1.25), WHITE)
        card.line.color.rgb = LINE
        add_rect(s, Inches(0.5), top, Inches(0.12), Inches(1.25), TEAL)
        textbox(s, Inches(0.85), top + Inches(0.2), Inches(5.7), Inches(0.35),
                t, size=16, bold=True, color=TEAL)
        textbox(s, Inches(0.85), top + Inches(0.6), Inches(5.7), Inches(0.5),
                b, size=12, color=SLATE)
    image_card(s, img("dashboard_demo.png"), Inches(7.05), Inches(1.15), Inches(5.8), Inches(5.7),
               "Facility & district intelligence dashboard (concept)")
    footer(s, 3)

    # ========== 4. DEEP TECH — collage + stack ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "Deep technology core",
                 "Substantial engineering — not a lightweight booking app.")
    image_card(s, img("deep_tech_collage.png"), Inches(0.45), Inches(1.15), Inches(6.0), Inches(5.7),
               "AI · ML · GIS · mobile edge · cloud · NLP")
    techs = [
        ("Artificial Intelligence", "Outbreak early warning; maternal & NCD risk scoring"),
        ("Machine Learning", "Seasonal patterns from community & outreach data"),
        ("GIS Mapping", "Hotspots, disease distribution, resource gaps"),
        ("Mobile Data Collection", "Offline CHW/VHT apps at household level"),
        ("Cloud Computing", "Secure sync across facilities & partners"),
        ("Analytics + NLP", "Live dashboards; local-language symptom capture"),
    ]
    for i, (t, b) in enumerate(techs):
        top = Inches(1.2) + Inches(i * 0.9)
        card = add_round_rect(s, Inches(6.7), top, Inches(6.1), Inches(0.82), WHITE)
        card.line.color.rgb = LINE
        textbox(s, Inches(6.95), top + Inches(0.1), Inches(5.6), Inches(0.3),
                t, size=13, bold=True, color=TEAL)
        textbox(s, Inches(6.95), top + Inches(0.42), Inches(5.6), Inches(0.3),
                b, size=11, color=MUTED)
    footer(s, 4)

    # ========== 5. HOW IT WORKS — data flow visual ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "How it works", "Last-mile signals become predictions — then action.")
    image_card(s, img("data_flow_iso.png"), Inches(0.4), Inches(1.1), Inches(7.5), Inches(5.75),
               "Community → Capture → CHIP intelligence → Decision-makers")
    layers = [
        ("1. Sources", "CHWs, clinics, pharmacies, schools, ANC/PNC, labs, HMIS"),
        ("2. Capture", "Offline mobile & facility apps · cloud sync"),
        ("3. CHIP core", "AI/ML · Predictive analytics · GIS · CDS"),
        ("4. Action", "Alerts · dashboards · referrals · stock · outreach"),
    ]
    for i, (t, b) in enumerate(layers):
        top = Inches(1.2) + Inches(i * 1.35)
        card = add_round_rect(s, Inches(8.15), top, Inches(4.7), Inches(1.2), WHITE)
        card.line.color.rgb = LINE
        add_rect(s, Inches(8.15), top, Inches(4.7), Inches(0.1), TEAL if i < 3 else ACCENT)
        textbox(s, Inches(8.4), top + Inches(0.25), Inches(4.2), Inches(0.3),
                t, size=14, bold=True, color=NAVY)
        textbox(s, Inches(8.4), top + Inches(0.6), Inches(4.2), Inches(0.45),
                b, size=12, color=MUTED)
    footer(s, 5)

    # ========== 6. USE CASES — photo mosaic ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "Predictive use cases", "Proof-of-concept scenarios in FairBanks' live catchment.")
    # Three demo images across top
    image_card(s, img("maternal_visit.png"), Inches(0.4), Inches(1.1), Inches(4.1), Inches(3.0),
               "Maternal risk — home-visit signals")
    image_card(s, img("gis_hotspots.png"), Inches(4.65), Inches(1.1), Inches(4.1), Inches(3.0),
               "Surveillance — GIS hotspot early warning")
    image_card(s, img("pharmacy_stock.png"), Inches(8.9), Inches(1.1), Inches(4.0), Inches(3.0),
               "Medicine demand — forecast before stock-outs")

    cases = [
        ("Disease surveillance", "Fever clusters → malaria risk ~14 days → testing & bed-nets"),
        ("Maternal health", "BP/Hb/ANC gaps → high-risk flags → CHW referral"),
        ("NCDs", "Community BP/glucose → hotspots → targeted screening"),
        ("Child health", "Growth & immunisation → coverage gaps → drives"),
        ("Medicine demand", "Trends + seasonality → procurement before stock-outs"),
    ]
    for i, (t, b) in enumerate(cases):
        left = Inches(0.4) + Inches(i * 2.55)
        card = add_round_rect(s, left, Inches(4.3), Inches(2.45), Inches(2.5), WHITE)
        card.line.color.rgb = LINE
        add_rect(s, left, Inches(4.3), Inches(2.45), Inches(0.1), TEAL if i % 2 == 0 else ACCENT)
        textbox(s, left + Inches(0.12), Inches(4.55), Inches(2.2), Inches(0.55),
                t, size=12, bold=True, color=NAVY)
        textbox(s, left + Inches(0.12), Inches(5.2), Inches(2.2), Inches(1.35),
                b, size=11, color=MUTED)
    footer(s, 6)

    # ========== 7. MARKET — visual + segments ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "Market & customers", "B2B / B2G buyers across Africa's primary-care stack.")
    image_card(s, img("mobile_capture.png"), Inches(0.4), Inches(1.15), Inches(5.5), Inches(5.7),
               "CHW/VHT mobile capture — the frontline product surface")
    customers = [
        ("Medical centres & clinics", "Outreach planning & population visibility"),
        ("District health offices", "Disease intelligence & early warning"),
        ("NGOs & partners", "Real-time M&E and impact evidence"),
        ("CHWs / VHTs", "Tools, workflows, decision support"),
        ("Ministries of health", "Planning & outbreak preparedness"),
        ("Insurers & research", "Prevention insights; ethical datasets"),
    ]
    for i, (t, b) in enumerate(customers):
        col = i % 2
        row = i // 2
        left = Inches(6.15) + Inches(col * 3.4)
        top = Inches(1.2) + Inches(row * 1.7)
        card = add_round_rect(s, left, top, Inches(3.25), Inches(1.5), WHITE)
        card.line.color.rgb = LINE
        textbox(s, left + Inches(0.18), top + Inches(0.25), Inches(2.9), Inches(0.5),
                t, size=13, bold=True, color=TEAL)
        textbox(s, left + Inches(0.18), top + Inches(0.8), Inches(2.9), Inches(0.5),
                b, size=11, color=MUTED)
    footer(s, 7)

    # ========== 8. BUSINESS MODEL ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "Business model", "Diversified HealthTech revenue — B2B, B2G, and partner channels.")
    # Decorative side image strip
    image_card(s, img("dashboard_demo.png"), Inches(9.0), Inches(1.15), Inches(3.9), Inches(5.7),
               "SaaS dashboards & partner analytics")
    streams = [
        "Subscription licences for clinics and hospitals",
        "District health office deployments (SaaS + implementation)",
        "NGO programme monitoring contracts",
        "Ministry of Health national / sub-national implementations",
        "Custom analytics and reporting for partners",
        "Research collaborations with universities",
        "API integrations for digital health partners",
        "CHW training and certification on platform use",
    ]
    for i, item in enumerate(streams):
        top = Inches(1.2) + Inches(i * 0.68)
        card = add_round_rect(s, Inches(0.45), top, Inches(8.3), Inches(0.6), WHITE)
        card.line.color.rgb = LINE
        add_rect(s, Inches(0.45), top, Inches(0.1), Inches(0.6), TEAL)
        textbox(s, Inches(0.8), top + Inches(0.12), Inches(7.7), Inches(0.4),
                f"{i + 1}.  {item}", size=13, color=SLATE)
    footer(s, 8)

    # ========== 9. TRACTION ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "Traction & unfair advantage",
                 "Not a concept detached from the field — a live operating foundation.")
    image_card(s, img("hero_chw_mobile.png"), Inches(0.4), Inches(1.15), Inches(5.6), Inches(3.35),
               "Active CHW/VHT engagement across Kampala communities")
    image_card(s, img("maternal_visit.png"), Inches(0.4), Inches(4.6), Inches(2.7), Inches(2.25),
               "Maternal & child")
    image_card(s, img("pharmacy_stock.png"), Inches(3.25), Inches(4.6), Inches(2.75), Inches(2.25),
               "Pharmacy & NCD")

    assets = [
        "Functioning FairBanks medical centre",
        "Community Reach Programme (live)",
        "CHW/VHT: Bukoto, Kyebando, Kisaasi, Kamwokya, Kikaaya+",
        "Maternal & child health initiatives",
        "Gericare + chronic disease screening",
        "Corporate & school health programmes",
        "Digital health records foundation",
        "Research & community partnerships",
    ]
    right = add_round_rect(s, Inches(6.25), Inches(1.15), Inches(6.6), Inches(3.5), WHITE)
    right.line.color.rgb = LINE
    textbox(s, Inches(6.5), Inches(1.3), Inches(6.1), Inches(0.35),
            "Existing FairBanks ecosystem", size=15, bold=True, color=TEAL)
    bullet_box(s, Inches(6.5), Inches(1.8), Inches(6.1), Inches(2.6), assets, size=12)

    win = add_round_rect(s, Inches(6.25), Inches(4.8), Inches(6.6), Inches(2.05), NAVY)
    textbox(s, Inches(6.5), Inches(4.95), Inches(6.1), Inches(0.3),
            "Why FairBanks wins", size=14, bold=True, color=TEAL_LIGHT)
    textbox(s, Inches(6.5), Inches(5.35), Inches(6.1), Inches(1.3),
            "Technology + field access in one venture.\n"
            "Context-rooted AI from African community realities.\n"
            "Design → pilot → validate → refine on a live site.\n"
            "Recommended track: Startup.",
            size=12, color=WHITE)
    footer(s, 9)

    # ========== 10. MVP & ROADMAP ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "MVP & roadmap", "Prove predictive intelligence in FairBanks' catchment first.")
    image_card(s, img("mobile_capture.png"), Inches(0.4), Inches(1.1), Inches(4.0), Inches(2.7),
               "MVP: CHW/VHT offline app")
    image_card(s, img("gis_hotspots.png"), Inches(4.55), Inches(1.1), Inches(4.0), Inches(2.7),
               "MVP: GIS risk layer v1")
    image_card(s, img("dashboard_demo.png"), Inches(8.7), Inches(1.1), Inches(4.2), Inches(2.7),
               "MVP: Facility analytics dashboard")

    mvp_bits = [
        "Cloud sync & validated data pipeline",
        "AI Module v1 — fever clusters, maternal flags, BP trends",
        "Pilot integration with FairBanks records & outreach",
    ]
    for i, t in enumerate(mvp_bits):
        left = Inches(0.4) + Inches(i * 4.25)
        card = add_round_rect(s, left, Inches(4.0), Inches(4.1), Inches(0.7), SOFT_TEAL)
        textbox(s, left + Inches(0.15), Inches(4.15), Inches(3.8), Inches(0.45),
                t, size=11, bold=True, color=NAVY)

    phases = [
        ("Phase 1 · 0–12 mo", "Pilot MVP; validate 3 use cases"),
        ("Phase 2 · 12–24 mo", "District scale; NGO M&E; APIs"),
        ("Phase 3 · 24–36 mo", "Multi-district UG + East Africa"),
        ("Phase 4 · Year 3+", "CDS, local NLP, research modules"),
    ]
    for i, (t, b) in enumerate(phases):
        left = Inches(0.4) + Inches(i * 3.2)
        add_round_rect(s, left, Inches(4.9), Inches(3.05), Inches(1.85), NAVY)
        textbox(s, left + Inches(0.15), Inches(5.1), Inches(2.75), Inches(0.4),
                t, size=12, bold=True, color=TEAL_LIGHT)
        textbox(s, left + Inches(0.15), Inches(5.55), Inches(2.75), Inches(0.9),
                b, size=12, color=WHITE)
    footer(s, 10)

    # ========== 11. IMPACT ==========
    s = prs.slides.add_slide(blank)
    section_band(s, "Impact & SDG alignment",
                 "From sick-care to predictive, community-centred prevention.")
    image_card(s, img("maternal_visit.png"), Inches(0.4), Inches(1.15), Inches(5.5), Inches(5.7),
               "Earlier intervention at the household and village level")
    outcomes = [
        "Earlier outbreak detection; lower communicable morbidity",
        "Fewer maternal & neonatal complications",
        "Reduced NCD burden via hotspot screening",
        "Improved child nutrition & immunisation coverage",
        "Fewer medicine stock-outs; smarter allocation",
        "Evidence base for NGOs, donors & governments",
    ]
    card = add_round_rect(s, Inches(6.15), Inches(1.15), Inches(6.7), Inches(3.4), WHITE)
    card.line.color.rgb = LINE
    textbox(s, Inches(6.4), Inches(1.3), Inches(6.2), Inches(0.35),
            "Expected outcomes", size=15, bold=True, color=TEAL)
    bullet_box(s, Inches(6.4), Inches(1.8), Inches(6.2), Inches(2.5), outcomes, size=13)

    sdgs = [("SDG 3", "Health"), ("SDG 5", "Gender"), ("SDG 9", "Innovation"),
            ("SDG 10", "Equality"), ("SDG 17", "Partnerships")]
    for i, (code, name) in enumerate(sdgs):
        left = Inches(6.15) + Inches(i * 1.35)
        add_round_rect(s, left, Inches(4.75), Inches(1.25), Inches(2.1), NAVY)
        textbox(s, left + Inches(0.05), Inches(5.15), Inches(1.15), Inches(0.35),
                code, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        textbox(s, left + Inches(0.05), Inches(5.55), Inches(1.15), Inches(0.5),
                name, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    footer(s, 11)

    # ========== 12. ASK — hero close ==========
    s = prs.slides.add_slide(blank)
    add_picture_cover(s, img("gis_hotspots.png"), Inches(5.8), 0, Inches(7.533), SLIDE_H)
    add_rect(s, 0, 0, Inches(6.2), SLIDE_H, NAVY)
    add_rect(s, Inches(6.05), 0, Inches(0.15), SLIDE_H, TEAL_LIGHT)

    textbox(s, Inches(0.5), Inches(0.55), Inches(5.4), Inches(0.45),
            "The ask & next steps", size=26, bold=True, color=WHITE)
    textbox(s, Inches(0.5), Inches(1.15), Inches(5.4), Inches(1.2),
            "Support to pilot the CHIP MVP — mobile CHW/VHT capture, cloud sync, "
            "predictive analytics, and facility/district dashboards — validated in "
            "FairBanks' live ecosystem before district and regional scale-up.",
            size=13, color=RGBColor(0xD0, 0xDC, 0xE0))

    steps = [
        "Finalise MVP scope & AI priorities",
        "Submit AWIEF application + PDF deck",
        "Record 3-minute technology demo",
        "Launch CHW/VHT pilot in catchment communities",
    ]
    for i, step in enumerate(steps):
        top = Inches(2.55) + Inches(i * 0.65)
        add_round_rect(s, Inches(0.5), top, Inches(5.4), Inches(0.55), RGBColor(0x12, 0x32, 0x42))
        textbox(s, Inches(0.7), top + Inches(0.1), Inches(5.0), Inches(0.35),
                f"{i + 1}.  {step}", size=12, color=WHITE)

    textbox(s, Inches(0.5), Inches(5.35), Inches(5.4), Inches(0.4),
            "Vision: Africa's leading Community Health Intelligence Platform.",
            size=13, bold=True, color=TEAL_LIGHT)
    textbox(s, Inches(0.5), Inches(5.9), Inches(5.4), Inches(0.35),
            "FairBanks Community Intelligence Network (FCIN) · CHIP",
            size=11, color=RGBColor(0xA8, 0xC4, 0xC4))
    textbox(s, Inches(0.5), Inches(6.35), Inches(5.4), Inches(0.35),
            "Contact: [founder · email · phone · web]",
            size=11, color=MUTED)
    textbox(s, Inches(0.5), Inches(6.75), Inches(5.4), Inches(0.3),
            "awief.untap.us/pitch-n-grow2026  ·  Cape Town 10–11 Nov 2026",
            size=10, color=TEAL_LIGHT)

    out = Path(__file__).resolve().parent / "FairBanks_FCIN_CHIP_AWIEF_Pitch_n_Grow_2026.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        out = Path(__file__).resolve().parent / "FairBanks_FCIN_CHIP_AWIEF_Pitch_n_Grow_2026_visual.pptx"
        prs.save(str(out))
        print(f"Main file locked; saved as: {out}")
    else:
        print(out)
    return str(out)


if __name__ == "__main__":
    build()
