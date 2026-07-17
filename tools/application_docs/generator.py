"""Generate synced Word, PDF, and PowerPoint application packs.

Each opportunity project passes a SPEC dict. Output lives in
applications/{slug}/documents/ as {slug}_{word,pdf,ppt}.*
"""

from __future__ import annotations

from pathlib import Path

NAVY = "0A1F2E"
TEAL = "0D6E6E"
TEAL_L = "14A3A3"
ACCENT = "C45C26"
SLATE = "1E2F38"
MUTED = "3A4A54"
CREAM = "F7F5F0"
LINE = "D0DCDC"

SLOGAN = "Your health, our mission."
TAGLINE = "Health for All — Obulamu eri Bonna · Afya kwa Wote · Oburamu bwa Boona"

DEFAULT_PHOTOS = {
    "cover": "cover_hero_cinematic.jpg",
    "logo": "fairbanks_logo.jpeg",
    "outreach": "outreach_facilitator_canopy_01.jpg",
    "architecture": "data_flow_iso_labeled.png",
    "deep_tech": "deep_tech_collage.png",
    "gis": "gis_hotspots.png",
    "dashboard": "dashboard_demo.png",
    "mobile": "outreach_mobile_phone_demo_01.jpg",
    "reactive": "reactive_clinic.png",
    "maternal": "bloom_maternal_health_participant_01.jpg",
    "training": "indoor_training_staff_presenting_01.jpg",
    "conclusion": "outreach_audience_full_group_01.jpg",
    "bp": "outreach_bp_screening.jpeg",
}


def _repo_root(project_dir: Path) -> Path:
    # applications/{slug} -> repo root
    return project_dir.resolve().parents[1]


def _embed(repo: Path, name: str, cache_tag: str, max_px: int = 1400) -> str:
    from PIL import Image as PILImage

    src = repo / "assets" / name
    if not src.exists():
        raise FileNotFoundError(src)
    cache = repo / "tmp" / f"{cache_tag}_assets"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{src.stem}_opt.jpg"
    if out.exists() and out.stat().st_mtime >= src.stat().st_mtime:
        return str(out)
    with PILImage.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        scale = min(1.0, max_px / float(max(iw, ih)))
        if scale < 1.0:
            im = im.resize((int(iw * scale), int(ih * scale)), PILImage.Resampling.LANCZOS)
        im.save(out, format="JPEG", quality=82, optimize=True)
    return str(out)


def _photo(spec: dict, key: str) -> str:
    photos = {**DEFAULT_PHOTOS, **spec.get("photos", {})}
    return photos[key]


def build_all(spec: dict, project_dir: Path) -> None:
    """Write Word, PDF, and PPTX for one opportunity application."""
    project_dir = project_dir.resolve()
    out_dir = project_dir / "documents"
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = spec["slug"]
    repo = _repo_root(project_dir)
    paths = {
        "word": out_dir / f"{slug}_word.docx",
        "pdf": out_dir / f"{slug}_pdf.pdf",
        "ppt": out_dir / f"{slug}_ppt.pptx",
    }
    build_docx(spec, repo, paths["word"])
    build_pdf(spec, repo, paths["pdf"])
    build_pptx(spec, repo, paths["ppt"])
    print(f"Done. {slug} document set in {out_dir}")


def build_docx(spec: dict, repo: Path, out: Path) -> None:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from PIL import Image as PILImage

    tag = spec["slug"]

    def emb(key: str) -> str:
        return _embed(repo, _photo(spec, key), tag)

    def set_font(run, size=11, bold=False, color=SLATE, italic=False):
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)

    def shade(cell, color):
        cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color}"/>'))

    def border(cell, color=LINE, sz="8"):
        cell._tc.get_or_add_tcPr().append(parse_xml(
            f'<w:tcBorders {nsdecls("w")}>'
            f'<w:top w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f'<w:left w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f'<w:bottom w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f'<w:right w:val="single" w:sz="{sz}" w:color="{color}"/>'
            f"</w:tcBorders>"))

    def para(text, size=11, bold=False, color=SLATE, after=8, before=0,
             align=WD_ALIGN_PARAGRAPH.LEFT, italic=False):
        p = doc.add_paragraph()
        p.alignment = align
        p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.space_before = Pt(before)
        p.paragraph_format.line_spacing = 1.2
        set_font(p.add_run(text), size=size, bold=bold, color=color, italic=italic)
        return p

    def heading(text, level=1):
        sizes, colors = {1: 18, 2: 13}, {1: NAVY, 2: TEAL}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14 if level == 1 else 9)
        p.paragraph_format.space_after = Pt(5)
        set_font(p.add_run(text), size=sizes[level], bold=True, color=colors[level])

    def image(key, width_in=6.2, caption=None, max_h=3.2):
        path = Path(emb(key))
        with PILImage.open(path) as im:
            iw, ih = im.size
        w = min(width_in, max_h * iw / ih)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after = Pt(2)
        p.add_run().add_picture(str(path), width=Inches(w))
        if caption:
            para(caption, size=9, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER,
                 italic=True, after=10, before=2)

    def bullets(items):
        for it in items:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            p.clear()
            set_font(p.add_run(it), size=11, color=SLATE)

    def table(headers, rows, widths=None):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i]
            c.text = ""
            set_font(c.paragraphs[0].add_run(h), size=10, bold=True, color="FFFFFF")
            shade(c, TEAL)
            border(c, TEAL)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = t.rows[ri + 1].cells[ci]
                c.text = ""
                set_font(c.paragraphs[0].add_run(str(val)), size=10, color=SLATE)
                if ri % 2:
                    shade(c, CREAM)
                border(c)
        if widths:
            for row in t.rows:
                for i, w in enumerate(widths):
                    row.cells[i].width = Inches(w)
        doc.add_paragraph()

    doc = Document()
    s = doc.sections[0]
    s.page_width, s.page_height = Inches(8.5), Inches(11)
    s.left_margin = s.right_margin = Inches(0.85)
    s.top_margin = s.bottom_margin = Inches(0.75)

    meta = spec["meta"]
    para(meta["programme"], size=12, bold=True, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(spec.get("doc_label", "Application Pack"), size=11, color=MUTED,
         align=WD_ALIGN_PARAGRAPH.CENTER, after=12)
    para(meta["title"], size=22, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    para(meta["subtitle"], size=12, bold=True, color=ACCENT,
         align=WD_ALIGN_PARAGRAPH.CENTER, after=4, italic=True)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, after=10, italic=True)
    image("cover", caption="FairBanks Community Reach — " + TAGLINE)
    table(
        ["Item", "Detail"],
        [
            ["Applicant", meta["applicant"]],
            ["Programme", meta["programme"]],
            ["Deadline", meta["deadline"]],
            ["Geography", meta["geography"]],
            ["The ask", meta["ask"]],
        ],
        widths=[1.6, 4.8],
    )
    para(spec["url"], size=8, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER, after=6)

    doc.add_page_break()
    heading("1. Executive Summary")
    for i, p in enumerate(spec["exec_summary"]):
        para(p, bold=(i == len(spec["exec_summary"]) - 1))

    heading("2. Win-Win Value")
    para("This application is designed as a true partnership: FairBanks grows capacity and reach; "
         "the programme advances its own mission with a grounded African health-tech partner.")
    heading("2.1 What FairBanks / FCHIP gains", 2)
    bullets(spec["win_for_us"])
    heading("2.2 What the programme / partners gain", 2)
    bullets(spec["win_for_them"])

    heading("3. The Problem We Address")
    image("reactive", width_in=5.5, max_h=2.5,
          caption="Reactive care waits for sickness before it responds")
    for p in spec["problem"]:
        para(p)

    heading("4. Our Solution: FCHIP & Community Reach")
    image("architecture", caption="Community signals → capture → FCHIP intelligence → action")
    for p in spec["solution"]:
        para(p)
    if spec.get("deep_tech"):
        heading("4.1 Deep technology core", 2)
        table(["Technology", "Role"], spec["deep_tech"], widths=[2.2, 4.2])
    if spec.get("use_cases"):
        heading("4.2 Priority use cases for this call", 2)
        for title, body in spec["use_cases"]:
            para(f"{title}: {body}", after=4)

    heading("5. Fit to Programme Criteria")
    table(["Criterion / theme", "How we respond"], spec["fit_rows"], widths=[2.4, 4.0])

    heading("6. Why FairBanks Is Ready")
    image("outreach", caption="Live outreach in Kampala-area communities")
    bullets(spec["traction"])

    heading("7. Proposed Plan")
    if spec.get("plan_rows"):
        table(["Phase / step", "Detail"], spec["plan_rows"], widths=[2.0, 4.4])
    if spec.get("plan_bullets"):
        bullets(spec["plan_bullets"])

    if spec.get("budget_rows"):
        heading("8. Indicative Resources")
        table(["Item", "Detail", "Amount"], spec["budget_rows"], widths=[2.0, 3.2, 1.2])
        next_n = 9
    else:
        next_n = 8

    heading(f"{next_n}. Risks & Safeguards")
    table(["Risk", "Mitigation"], spec["risks"], widths=[2.4, 4.0])

    heading(f"{next_n + 1}. Closing Ask")
    image("conclusion", caption="From community needs to lasting health impact")
    for p in spec["closing"]:
        para(p)
    para(SLOGAN, size=12, bold=True, color=ACCENT, align=WD_ALIGN_PARAGRAPH.CENTER, before=8, italic=True)
    para(f"Aligned with: {spec['url']}", size=8, color=MUTED,
         align=WD_ALIGN_PARAGRAPH.CENTER, before=8, italic=True)

    try:
        doc.save(str(out))
        print(f"DOCX: {out}")
    except PermissionError:
        alt = out.with_name(out.stem + "_unlocked" + out.suffix)
        doc.save(str(alt))
        print(f"DOCX locked; saved as: {alt}")


def build_pdf(spec: dict, repo: Path, out: Path) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle,
        PageBreak, KeepTogether, HRFlowable,
    )
    from PIL import Image as PILImage

    tag = spec["slug"]
    meta = spec["meta"]

    def emb(key: str) -> str:
        return _embed(repo, _photo(spec, key), tag)

    navy, teal, accent = HexColor("#" + NAVY), HexColor("#" + TEAL), HexColor("#" + ACCENT)
    slate, muted, cream, line = HexColor("#" + SLATE), HexColor("#" + MUTED), HexColor("#" + CREAM), HexColor("#" + LINE)

    st = getSampleStyleSheet()
    st.add(ParagraphStyle("CoverTitle", fontName="Helvetica-Bold", fontSize=20, leading=24,
                          textColor=navy, alignment=TA_CENTER, spaceAfter=6))
    st.add(ParagraphStyle("H1", fontName="Helvetica-Bold", fontSize=14, leading=18,
                          textColor=navy, spaceBefore=12, spaceAfter=5))
    st.add(ParagraphStyle("H2", fontName="Helvetica-Bold", fontSize=11, leading=14,
                          textColor=teal, spaceBefore=8, spaceAfter=4))
    st.add(ParagraphStyle("Body", fontName="Helvetica", fontSize=10, leading=13,
                          textColor=slate, alignment=TA_JUSTIFY, spaceAfter=6))
    st.add(ParagraphStyle("BodyBold", fontName="Helvetica-Bold", fontSize=10, leading=13,
                          textColor=slate, alignment=TA_JUSTIFY, spaceAfter=6))
    st.add(ParagraphStyle("Caption", fontName="Helvetica-Oblique", fontSize=8, leading=10,
                          textColor=muted, alignment=TA_CENTER, spaceAfter=8, spaceBefore=2))
    st.add(ParagraphStyle("FBullet", fontName="Helvetica", fontSize=10, leading=12,
                          textColor=slate, leftIndent=12, spaceAfter=2))
    st.add(ParagraphStyle("Meta", fontName="Helvetica", fontSize=9, leading=11,
                          textColor=muted, alignment=TA_CENTER, spaceAfter=3))
    st.add(ParagraphStyle("CellHead", fontName="Helvetica-Bold", fontSize=8.5, leading=11, textColor=white))
    st.add(ParagraphStyle("CellBody", fontName="Helvetica", fontSize=8.5, leading=11, textColor=slate))

    story = []
    page_w = A4[0] - 1.6 * inch

    def img(key, w=page_w, caption=None, max_h=2.8 * inch):
        path = emb(key)
        with PILImage.open(path) as pi:
            iw, ih = pi.size
        aspect = ih / float(iw)
        h = w * aspect
        if h > max_h:
            h, w = max_h, max_h / aspect
        block = [Image(path, width=w, height=h)]
        block.append(Paragraph(caption, st["Caption"]) if caption else Spacer(1, 4))
        return KeepTogether(block)

    def h1(t):
        story.append(Paragraph(t, st["H1"]))
        story.append(HRFlowable(width="100%", thickness=1.1, color=teal, spaceAfter=6))

    def h2(t):
        story.append(Paragraph(t, st["H2"]))

    def body(t, bold=False):
        story.append(Paragraph(t, st["BodyBold"] if bold else st["Body"]))

    def bullets(items):
        for it in items:
            story.append(Paragraph(f"•&nbsp;&nbsp;{it}", st["FBullet"]))
        story.append(Spacer(1, 3))

    def table(headers, rows, widths=None):
        data = [[Paragraph(h, st["CellHead"]) for h in headers]]
        for row in rows:
            data.append([Paragraph(str(c), st["CellBody"]) for c in row])
        widths = widths or [page_w / len(headers)] * len(headers)
        t = Table(data, colWidths=widths, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), teal),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4), ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("GRID", (0, 0), (-1, -1), 0.4, line),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, cream]),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph(meta["programme"], st["Meta"]))
    story.append(Paragraph(spec.get("doc_label", "Application Pack"), st["Meta"]))
    story.append(Spacer(1, 6))
    story.append(Paragraph(meta["title"], st["CoverTitle"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{meta["subtitle"]}</i></b></font>', st["Meta"]))
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))
    story.append(Spacer(1, 8))
    story.append(img("cover", w=page_w * 0.9, max_h=2.5 * inch,
                     caption="FairBanks Community Reach — " + TAGLINE))
    table(["Item", "Detail"], [
        ["Applicant", meta["applicant"]],
        ["Programme", meta["programme"]],
        ["Deadline", meta["deadline"]],
        ["Geography", meta["geography"]],
        ["The ask", meta["ask"]],
    ], widths=[page_w * 0.24, page_w * 0.76])
    story.append(Paragraph(f'<link href="{spec["url"]}"><font color="#{TEAL}">Official call</font></link>', st["Meta"]))
    story.append(PageBreak())

    h1("1. Executive Summary")
    for i, p in enumerate(spec["exec_summary"]):
        body(p, bold=(i == len(spec["exec_summary"]) - 1))

    h1("2. Win-Win Value")
    body("FairBanks grows capacity and reach; the programme advances its mission with a grounded African health-tech partner.")
    h2("2.1 What FairBanks / FCHIP gains")
    bullets(spec["win_for_us"])
    h2("2.2 What the programme / partners gain")
    bullets(spec["win_for_them"])

    h1("3. The Problem We Address")
    story.append(img("reactive", w=page_w * 0.75, max_h=2.2 * inch,
                     caption="Reactive care waits for sickness before it responds"))
    for p in spec["problem"]:
        body(p)

    h1("4. Our Solution: FCHIP & Community Reach")
    story.append(img("architecture", caption="Community signals → capture → intelligence → action"))
    for p in spec["solution"]:
        body(p)
    if spec.get("deep_tech"):
        h2("4.1 Deep technology core")
        table(["Technology", "Role"], spec["deep_tech"], widths=[page_w * 0.32, page_w * 0.68])
    if spec.get("use_cases"):
        h2("4.2 Priority use cases for this call")
        for title, b in spec["use_cases"]:
            body(f"<b>{title}:</b> {b}")

    story.append(PageBreak())
    h1("5. Fit to Programme Criteria")
    table(["Criterion / theme", "How we respond"], spec["fit_rows"],
          widths=[page_w * 0.35, page_w * 0.65])

    h1("6. Why FairBanks Is Ready")
    story.append(img("outreach", caption="Live outreach in Kampala-area communities"))
    bullets(spec["traction"])

    h1("7. Proposed Plan")
    if spec.get("plan_rows"):
        table(["Phase / step", "Detail"], spec["plan_rows"],
              widths=[page_w * 0.3, page_w * 0.7])
    if spec.get("plan_bullets"):
        bullets(spec["plan_bullets"])

    if spec.get("budget_rows"):
        h1("8. Indicative Resources")
        table(["Item", "Detail", "Amount"], spec["budget_rows"],
              widths=[page_w * 0.28, page_w * 0.52, page_w * 0.2])
        risk_h, close_h = "9. Risks & Safeguards", "10. Closing Ask"
    else:
        risk_h, close_h = "8. Risks & Safeguards", "9. Closing Ask"

    h1(risk_h)
    table(["Risk", "Mitigation"], spec["risks"], widths=[page_w * 0.35, page_w * 0.65])

    story.append(PageBreak())
    story.append(KeepTogether([
        Paragraph(close_h, st["H1"]),
        HRFlowable(width="100%", thickness=1.1, color=teal, spaceAfter=6),
        img("conclusion", max_h=2.4 * inch,
            caption="From community needs to lasting health impact"),
    ]))
    for p in spec["closing"]:
        body(p)
    story.append(Paragraph(f'<font color="#{ACCENT}"><b><i>{SLOGAN}</i></b></font>', st["Meta"]))

    def page_num(canvas, doc_):
        canvas.saveState()
        canvas.setFillColor(navy)
        canvas.rect(0, A4[1] - 8, A4[0], 8, fill=1, stroke=0)
        canvas.setFillColor(muted)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(0.8 * inch, 0.4 * inch, f"FairBanks FCHIP  |  {meta['programme'][:48]}")
        canvas.drawRightString(A4[0] - 0.8 * inch, 0.4 * inch, str(doc_.page))
        canvas.restoreState()

    def write(path):
        d = SimpleDocTemplate(
            str(path), pagesize=A4,
            leftMargin=0.8 * inch, rightMargin=0.8 * inch,
            topMargin=0.65 * inch, bottomMargin=0.65 * inch,
            title=f"FairBanks — {meta['programme']}",
            author="FairBanks Community Health Intelligence Platform",
        )
        d.build(story, onFirstPage=page_num, onLaterPages=page_num)

    try:
        write(out)
        print(f"PDF: {out}")
    except PermissionError:
        alt = out.with_name(out.stem + "_unlocked" + out.suffix)
        write(alt)
        print(f"PDF locked; saved as: {alt}")


def build_pptx(spec: dict, repo: Path, out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    tag = spec["slug"]
    meta = spec["meta"]

    def emb(key: str) -> str:
        return _embed(repo, _photo(spec, key), tag)

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, color, line_color=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(color)
        if line_color:
            shp.line.color.rgb = C(line_color)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(slide, x, y, w, h, text, size=18, bold=False, color=SLATE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(4) if i < len(lines) - 1 else Pt(0)
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = "Calibri"
            r.font.color.rgb = C(color)
        return tb

    def pic_cover(slide, key, x, y, w, h):
        path = emb(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        target = w / h
        src = iw / ih
        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
        if src > target:
            crop = (1 - target / src) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (1 - src / target) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop

    def pic_fit(slide, key, x, y, w, h):
        path = emb(key)
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = ih / iw
        tw, th = w, int(w * aspect)
        if th > h:
            th, tw = h, int(h / aspect)
        px = x + (w - tw) // 2
        py = y + (h - th) // 2
        return slide.shapes.add_picture(path, px, py, width=tw, height=th)

    def band(slide, kicker, title):
        rect(slide, 0, 0, SW, Inches(1.05), CREAM)
        rect(slide, 0, 0, Inches(0.18), Inches(1.05), TEAL)
        textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.28),
                kicker.upper(), size=11, bold=True, color=ACCENT)
        textbox(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.52),
                title, size=22, bold=True, color=NAVY)

    def footer(slide, n):
        rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), NAVY)
        textbox(slide, Inches(0.4), SH - Inches(0.31), Inches(10), Inches(0.3),
                f"FairBanks FCHIP  |  {meta['programme'][:40]}  |  {SLOGAN}",
                size=9, color="FFFFFF", anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, SW - Inches(1.2), SH - Inches(0.31), Inches(0.8), Inches(0.3),
                str(n), size=9, color="FFFFFF", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    def bullets_box(slide, x, y, w, h, items, size=14):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = "•  " + it
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            r.font.color.rgb = C(SLATE)

    CONTENT_TOP = Inches(1.25)
    CONTENT_H = Inches(5.75)

    # 1 Title
    s = prs.slides.add_slide(blank)
    pic_cover(s, "cover", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.7), SW, Inches(3.7), NAVY)
    textbox(s, Inches(0.65), SH - Inches(3.4), Inches(12), Inches(0.35),
            meta["programme"], size=13, bold=True, color=TEAL_L)
    textbox(s, Inches(0.65), SH - Inches(2.9), Inches(12), Inches(0.7),
            meta["title"], size=26, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(2.05), Inches(12), Inches(0.4),
            meta["subtitle"], size=15, italic=True, color="F2C79B")
    textbox(s, Inches(0.65), SH - Inches(1.45), Inches(12), Inches(0.3),
            SLOGAN, size=13, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(0.95), Inches(12), Inches(0.4),
            meta["ask"], size=12, color="D0E8E8")

    # 2 Win-win
    s = prs.slides.add_slide(blank)
    band(s, "Win-win", "Mutual value for FairBanks and the programme")
    textbox(s, Inches(0.5), CONTENT_TOP, Inches(6), Inches(0.35),
            "What we gain", size=16, bold=True, color=TEAL)
    bullets_box(s, Inches(0.5), CONTENT_TOP + Inches(0.4), Inches(6), CONTENT_H - Inches(0.5),
                spec["win_for_us"][:5], size=13)
    textbox(s, Inches(6.9), CONTENT_TOP, Inches(6), Inches(0.35),
            "What they gain", size=16, bold=True, color=TEAL)
    bullets_box(s, Inches(6.9), CONTENT_TOP + Inches(0.4), Inches(5.9), CONTENT_H - Inches(0.5),
                spec["win_for_them"][:5], size=13)
    footer(s, 2)

    # 3 Problem
    s = prs.slides.add_slide(blank)
    band(s, "The problem", spec.get("problem_slide_title", "Health systems react too late"))
    pic_fit(s, "reactive", Inches(0.45), CONTENT_TOP, Inches(5.8), CONTENT_H)
    bullets_box(s, Inches(6.6), CONTENT_TOP + Inches(0.1), Inches(6.2), CONTENT_H - Inches(0.2),
                spec.get("problem_bullets", spec["problem"][:5]), size=14)
    footer(s, 3)

    # 4 Solution
    s = prs.slides.add_slide(blank)
    band(s, "The solution", "FCHIP — community health intelligence")
    pic_fit(s, "architecture", Inches(0.4), CONTENT_TOP, Inches(7.0), CONTENT_H)
    bullets_box(s, Inches(7.7), CONTENT_TOP + Inches(0.1), Inches(5.2), CONTENT_H - Inches(0.2),
                spec.get("solution_bullets", [
                    "CHWs and VHTs capture data at household level",
                    "AI and GIS turn signals into early warnings",
                    "Clinics, districts, and partners act before crises",
                    "Built on live FairBanks Medical Centre outreach",
                ]), size=14)
    footer(s, 4)

    # 5 Fit
    s = prs.slides.add_slide(blank)
    band(s, "Programme fit", "Why this call matches FCHIP")
    n = len(spec["fit_rows"])
    row_h = CONTENT_H / max(n, 1)
    for i, (t, b) in enumerate(spec["fit_rows"]):
        y = CONTENT_TOP + int(i * row_h)
        pad = Inches(0.05)
        rect(s, Inches(0.45), y + pad, Inches(12.4), row_h - pad * 2, "FFFFFF", LINE)
        textbox(s, Inches(0.65), y + Inches(0.12), Inches(3.2), row_h - Inches(0.25),
                t, size=13, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(4.0), y + Inches(0.12), Inches(8.5), row_h - Inches(0.25),
                b, size=12, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 5)

    # 6 Traction
    s = prs.slides.add_slide(blank)
    band(s, "Why FairBanks", "A live community health ecosystem")
    pic_fit(s, "outreach", Inches(0.45), CONTENT_TOP, Inches(6.0), CONTENT_H)
    bullets_box(s, Inches(6.8), CONTENT_TOP + Inches(0.1), Inches(6.0), CONTENT_H - Inches(0.2),
                spec["traction"][:6], size=14)
    footer(s, 6)

    # 7 Plan
    s = prs.slides.add_slide(blank)
    band(s, "The plan", spec.get("plan_slide_title", "What we will do if selected"))
    items = spec.get("plan_bullets") or [f"{a}: {b}" for a, b in spec.get("plan_rows", [])]
    bullets_box(s, Inches(0.55), CONTENT_TOP + Inches(0.1), Inches(6.3), CONTENT_H - Inches(0.2),
                items[:6], size=15)
    pic_fit(s, "training", Inches(7.1), CONTENT_TOP, Inches(5.7), CONTENT_H)
    footer(s, 7)

    # 8 Deep tech / use cases
    s = prs.slides.add_slide(blank)
    band(s, "Deep technology", "Intelligence rooted in African community care")
    pic_fit(s, "deep_tech", Inches(0.45), CONTENT_TOP, Inches(5.6), CONTENT_H)
    tech = spec.get("deep_tech") or [
        ["AI", "Risk prediction and early warning"],
        ["Mobile capture", "Offline CHW/VHT tools"],
        ["GIS", "Hotspot maps for action"],
        ["Cloud analytics", "Shared dashboards for partners"],
    ]
    n = len(tech)
    row_h = CONTENT_H / n
    for i, (t, b) in enumerate(tech):
        y = CONTENT_TOP + int(i * row_h)
        pad = Inches(0.06)
        rect(s, Inches(6.4), y + pad, Inches(6.4), row_h - pad * 2, "FFFFFF", LINE)
        textbox(s, Inches(6.6), y + Inches(0.12), Inches(6.0), Inches(0.3),
                t, size=14, bold=True, color=TEAL)
        textbox(s, Inches(6.6), y + Inches(0.42), Inches(6.0), Inches(0.4),
                b, size=12, color=MUTED)
    footer(s, 8)

    # 9 Cascade / community model
    s = prs.slides.add_slide(blank)
    band(s, "Our model", "Communities first — intelligence that serves care")
    cascade = [
        "Community members identify needs and own solutions",
        "CHWs / VHTs bridge homes to clinics",
        "FairBanks Community Reach programmes deliver outreach",
        "FairBanks Medical Centre provides clinical care",
        "Research, partnerships, and skills strengthen the system",
        "Livelihoods and empowerment make health gains last",
    ]
    bullets_box(s, Inches(0.55), CONTENT_TOP + Inches(0.1), Inches(6.3), CONTENT_H - Inches(0.2),
                cascade, size=14)
    pic_fit(s, "gis", Inches(7.1), CONTENT_TOP, Inches(5.7), CONTENT_H)
    footer(s, 9)

    # 10 Ask
    s = prs.slides.add_slide(blank)
    pic_cover(s, "conclusion", 0, 0, SW, SH)
    rect(s, 0, SH - Inches(3.5), SW, Inches(3.5), NAVY)
    textbox(s, Inches(0.65), SH - Inches(3.2), Inches(12), Inches(0.3),
            "The ask", size=13, bold=True, color=TEAL_L)
    textbox(s, Inches(0.65), SH - Inches(2.7), Inches(12), Inches(0.9),
            meta["ask"], size=22, bold=True, color="FFFFFF")
    textbox(s, Inches(0.65), SH - Inches(1.6), Inches(12), Inches(0.4),
            meta["subtitle"], size=14, italic=True, color="F2C79B")
    textbox(s, Inches(0.65), SH - Inches(1.0), Inches(12), Inches(0.35),
            SLOGAN, size=14, bold=True, color="FFFFFF")

    try:
        prs.save(str(out))
        print(f"PPTX: {out}")
    except PermissionError:
        alt = out.with_name(out.stem + "_unlocked" + out.suffix)
        prs.save(str(alt))
        print(f"PPTX locked; saved as: {alt}")
